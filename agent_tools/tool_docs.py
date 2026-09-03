# -*- coding: utf-8 -*-
"""tool_docs —— P0-1 批量拆分（工具域模块）：📊 数据与文档.

共享符号策略：permissions / security / shared / toolkit 为独立模块直接 import；
阈值常量/锁统一从 shared 导入（P1-3 下沉：见 shared.py「工具域阈值与锁」节）；
仅剩余辅助函数仍依赖主文件加载顺序契约（在 `from agent_tools import *` 前已定义）。
"""

import os
import re
import shutil
import subprocess

import permissions

from shared import clamp_int, PDF_EXTRACT_MAX_OUTPUT, DOCX_MAX_DEFAULT, KV_VALUE_MAX_BYTES  # D4: 参数校验辅助
from toolkit import tool  # noqa: F401  # 装饰器 + 工具名 re-export
import deepseek_client as _dc  # 可变注入配置动态访问（dc.X 注入后立即生效）
from deepseek_client import (

    _TABLE_CELL_MAX,
    _atomic_write,
    _db_conn,
    _db_execute_mysql,
    _db_execute_postgres,
    _db_execute_sqlite,
    _load_secrets,
    _md_inline_html,
    _md_table_rows,
    _parse_page_range,
    _read_optional_text,
    _readonly_stmt,
    _register_cjk_font,
    _save_secrets,
    _strip_html_tags,
    _table_to_md,
)
from db_utils import force_limit  # L3: SQL 层强制 LIMIT（防无界查询）

# L3: SQLite 只读查询语句级超时（progress handler 中断慢查询，防占住共享工具线程池）
_SQLITE_QUERY_TIMEOUT_S = 15.0



@tool(
        {
            "type": "function",
            "function": {
                "name": "database_query_mysql",
                "description": "MySQL 只读查询（SELECT/SHOW/DESC）。连接在数据目录 db_config.json 的 mysql.<connection> 配置",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "connection": {"type": "string", "description": "可选：连接名（默认 default）"},
                        "sql": {"type": "string", "description": "只读 SQL 语句"},
                        "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 20）"},
                    },
                    "required": ["sql"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='MySQL 只读查询',
    preactivate=(('数据库', 'sql', 'mysql', 'postgres'),),
)
def database_query_mysql(connection="default", sql="", max_rows=20):
    """MySQL 只读查询（SELECT/SHOW/DESC；连接在 db_config.json 配置）。"""
    try:
        import pymysql
    except ImportError:
        return "错误：需要 pymysql（pip install pymysql）"
    if not str(sql or "").strip():
        return "错误：sql 必填"
    cfg, err = _db_conn("mysql", connection)
    if cfg is None:
        return err
    if not _readonly_stmt(sql):
        return "错误：仅允许只读查询（SELECT / SHOW / DESC）"
    limit = clamp_int(max_rows, 20, lo=1, hi=200)
    sql = force_limit(sql, limit)  # L3: SQL 层强制 LIMIT
    try:
        conn = pymysql.connect(
            host=str(cfg.get("host") or "127.0.0.1"),
            port=int(cfg.get("port") or 3306),
            user=str(cfg.get("user") or ""),
            password=str(cfg.get("password") or ""),
            database=str(cfg.get("database") or ""),
            charset="utf8mb4", connect_timeout=5, read_timeout=15,
        )
        try:
            cur = conn.cursor()
            try:
                cur.execute("SET SESSION max_execution_time=15000")
            except Exception:
                pass
            cur.execute(str(sql))
            cols = [d[0] for d in cur.description] if cur.description else []
            rows = cur.fetchmany(limit)
            lines = [" | ".join(cols)] if cols else []
            for r in rows:
                cells = [str(x) if x is not None else "" for x in r]
                cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
                lines.append(" | ".join(cells))
            extra = "" if len(rows) < limit else " [已截断]"
            return "\n".join(lines) + extra if lines else "执行成功（无结果集）"
        finally:
            try:
                conn.close()
            except Exception:
                pass
    except Exception as e:
        return f"错误：MySQL 查询失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "database_query_postgres",
                "description": "PostgreSQL 只读查询（SELECT/SHOW/DESC）。连接在数据目录 db_config.json 的 postgres.<connection> 配置",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "connection": {"type": "string", "description": "可选：连接名（默认 default）"},
                        "sql": {"type": "string", "description": "只读 SQL 语句"},
                        "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 20）"},
                    },
                    "required": ["sql"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='PostgreSQL 只读查询',
    preactivate=(('数据库', 'sql', 'mysql', 'postgres'),),
)
def database_query_postgres(connection="default", sql="", max_rows=20):
    """PostgreSQL 只读查询（SELECT/SHOW/DESC；连接在 db_config.json 配置）。"""
    try:
        import psycopg2
    except ImportError:
        return "错误：需要 psycopg2（pip install psycopg2）"
    if not str(sql or "").strip():
        return "错误：sql 必填"
    cfg, err = _db_conn("postgres", connection)
    if cfg is None:
        return err
    if not _readonly_stmt(sql):
        return "错误：仅允许只读查询（SELECT / SHOW / DESC）"
    limit = clamp_int(max_rows, 20, lo=1, hi=200)
    sql = force_limit(sql, limit)  # L3: SQL 层强制 LIMIT
    try:
        conn = psycopg2.connect(
            host=str(cfg.get("host") or "127.0.0.1"),
            port=int(cfg.get("port") or 5432),
            user=str(cfg.get("user") or ""),
            password=str(cfg.get("password") or ""),
            dbname=str(cfg.get("database") or ""),
            connect_timeout=5,
            # 语句超时：只读查询最长 15 秒，防慢查询占住共享工具线程池
            options="-c statement_timeout=15000",
        )
        try:
            cur = conn.cursor()
            cur.execute(str(sql))
            cols = [d[0] for d in cur.description] if cur.description else []
            rows = cur.fetchmany(limit)
            lines = [" | ".join(cols)] if cols else []
            for r in rows:
                cells = [str(x) if x is not None else "" for x in r]
                cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
                lines.append(" | ".join(cells))
            extra = "" if len(rows) < limit else " [已截断]"
            return "\n".join(lines) + extra if lines else "执行成功（无结果集）"
        finally:
            try:
                conn.close()
            except Exception:
                pass
    except Exception as e:
        return f"错误：PostgreSQL 查询失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "read_excel",
                "description": "读取 Excel 文件（.xlsx，openpyxl），返回表格文本",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Excel 文件绝对路径"},
                        "sheet": {"type": "string", "description": "可选：工作表名或序号（默认第一个）"},
                        "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 100）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='读 Excel',
    preactivate=(('表格', 'excel', 'csv', '报表'),),
)
def read_excel(path, sheet=0, max_rows=100):
    """读取 Excel 文件（openpyxl，.xlsx）。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "错误：需要 openpyxl（pip install openpyxl）"
    if not path or not str(path).strip():
        return "错误：path 必填"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{p}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    try:
        limit = clamp_int(max_rows, 100, lo=1, hi=500)
    except (TypeError, ValueError):
        limit = 100
    try:
        wb = load_workbook(p, read_only=True, data_only=True)
        try:
            if isinstance(sheet, int):
                ws = wb.worksheets[min(sheet, len(wb.worksheets) - 1)] if wb.worksheets else wb.active
            else:
                ws = wb[sheet]
        except KeyError:
            return f"错误：工作表不存在：{sheet}"
        lines = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= limit:
                break
            cells = ["" if v is None else str(v) for v in row]
            # 单元格截断：超长文本（minified JSON 等）撑爆上下文
            cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
            lines.append(" | ".join(cells))
        if not lines:
            return "（空工作表）"
        return "\n".join(lines) + (f"\n[前 {limit} 行…]" if len(lines) >= limit else "")
    except Exception as e:
        return f"错误：读取 Excel 失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "epub_read",
                "description": "读取 EPUB 电子书正文为纯文本（依赖 ebooklib，缺失时返回安装指引）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "EPUB 文件绝对路径"},
                        "max_chars": {"type": "integer", "description": "可选：最多返回字符数（默认 20000）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='读取 epub 电子书',
    preactivate=(('电子书', 'epub', 'mobi', 'kindle'),),
)
def epub_read(path, max_chars=20000):
    """读取 EPUB 电子书正文为纯文本（依赖 ebooklib，缺失时返回安装指引）。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        return "错误：需要 ebooklib（pip install ebooklib）"
    try:
        book = epub.read_epub(p)
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            body = item.get_body_content() or b""
            parts.append(_strip_html_tags(body.decode("utf-8", errors="replace")))
        text = "\n\n".join(x for x in parts if x)
        if not text:
            return "（EPUB 无正文内容）"
        if len(text) > limit:
            text = text[:limit] + f"\n[正文已截断前 {limit} 字符]"
        return text
    except Exception as e:
        return f"错误：读取 EPUB 失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "mobi_read",
                "description": "读取 MOBI 电子书正文为纯文本（依赖 mobi 库，缺失时返回安装指引）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "MOBI 文件绝对路径"},
                        "max_chars": {"type": "integer", "description": "可选：最多返回字符数（默认 20000）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='读取 mobi 电子书',
    preactivate=(('电子书', 'epub', 'mobi', 'kindle'),),
)
def mobi_read(path, max_chars=20000):
    """读取 MOBI 电子书正文为纯文本（依赖 mobi 库，缺失时返回安装指引）。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    try:
        from mobi import Mobi
    except ImportError:
        return "错误：需要 mobi（pip install mobi）"
    try:
        book = Mobi(p)
        book.parse()
        text = str(book) if hasattr(book, "__str__") else ""
        if not text:
            text = "\n\n".join(str(getattr(book, field, "")) for field in ("title", "author", "publisher", "description"))
        if len(text) > limit:
            text = text[:limit] + f"\n[正文已截断前 {limit} 字符]"
        return text or "（MOBI 解析无文本）"
    except Exception as e:
        return f"错误：读取 MOBI 失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "doc_read",
                "description": "读取旧版 .doc 二进制文档正文（依赖本机 antiword/catdoc 命令行工具）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".doc 文件绝对路径"},
                        "max_chars": {"type": "integer", "description": "可选：最多返回字符数（默认 20000）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='读取 doc/rtf 等旧格式',
    preactivate=(('outlook', 'msg邮件', 'msg文件', '旧版doc', 'rtf'),),
)
def doc_read(path, max_chars=20000):
    """读取旧版 .doc 二进制文档（优先 antiword，其次 catdoc）。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    import shutil as _shutil
    for exe in ("antiword", "catdoc"):
        if _shutil.which(exe):
            try:
                proc = subprocess.run(
                    [exe, p], capture_output=True, text=True, timeout=30,
                    errors="replace", creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
                )
                text = (proc.stdout or "").strip() or (proc.stderr or "").strip()
                if text:
                    if len(text) > limit:
                        text = text[:limit] + f"\n[正文已截断前 {limit} 字符]"
                    return text
            except Exception as e:
                return f"错误：读取 .doc 失败（{exe}）: {e}"
    return "错误：读取 .doc 需要 antiword 或 catdoc（Windows 可安装 antiword 或改用 .docx）"


@tool(
        {
            "type": "function",
            "function": {
                "name": "msg_read",
                "description": "读取 .msg Outlook 邮件（主题/发件人/正文/附件清单；依赖 extract_msg，缺失时返回安装指引）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".msg 文件绝对路径"},
                        "max_chars": {"type": "integer", "description": "可选：正文最多返回字符数（默认 20000）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📧 邮件与消息'],
    phrases='读取邮件消息',
    preactivate=(('outlook', 'msg邮件', 'msg文件', '旧版doc', 'rtf'),),
)
def msg_read(path, max_chars=20000):
    """读取 .msg Outlook 邮件（依赖 extract_msg，缺失时返回安装指引），返回主题/发件人/正文/附件清单。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    try:
        import extract_msg
    except ImportError:
        return "错误：需要 extract_msg（pip install extract-msg）"
    try:
        msg = extract_msg.Message(p)
        attachments = [a.longFilename or a.shortFilename or "?" for a in (msg.attachments or [])]
        lines = [
            f"主题：{msg.subject}",
            f"发件人：{msg.sender}",
            f"收件人：{msg.to}",
            f"日期：{msg.date}",
        ]
        if attachments:
            lines.append(f"附件（{len(attachments)}）：{', '.join(attachments[:20])}")
        body = str(msg.body or "").strip()
        if body:
            lines.append("正文：")
            if len(body) > limit:
                body = body[:limit] + f"\n[正文已截断前 {limit} 字符]"
            lines.append(body)
        return "\n".join(lines)
    except Exception as e:
        return f"错误：读取 .msg 失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "archive_list",
                "description": "列出压缩包内容：.zip / .tar / .gz / .7z / .rar（7z/rar 依赖 py7zr/rarfile 库）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "压缩包绝对路径"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='列出归档内容',
    preactivate=(('打包', '压缩成', '归档文件', '压缩包'), ('解压', '解包', '解压缩', '解压到')),
)
def archive_list(path):
    """列出压缩包内容：.zip / .tar / .gz / .7z / .rar（7z/rar 依赖 py7zr/rarfile 库）。"""
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：压缩包不存在：{path}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    ext = os.path.splitext(p)[1].lower()
    try:
        if ext in (".zip",):
            import zipfile
            with zipfile.ZipFile(p) as zf:
                infos = zf.infolist()
                lines = [f"压缩包 {path} 共 {len(infos)} 个条目："]
                for info in infos[:100]:
                    lines.append(f"· {info.filename}（{info.file_size} 字节）")
                if len(infos) > 100:
                    lines.append(f"… 其余 {len(infos) - 100} 个条目略")
                return "\n".join(lines)
        if ext in (".tar", ".gz", ".tgz"):
            import tarfile
            with tarfile.open(p) as tf:
                members = tf.getmembers()
                lines = [f"压缩包 {path} 共 {len(members)} 个条目："]
                for m in members[:100]:
                    lines.append(f"· {m.name}（{m.size} 字节）")
                if len(members) > 100:
                    lines.append(f"… 其余 {len(members) - 100} 个条目略")
                return "\n".join(lines)
        if ext == ".7z":
            import py7zr
            with py7zr.SevenZipFile(p, "r") as z:
                items = z.getnames()
                lines = [f"压缩包 {path} 共 {len(items)} 个条目："]
                for name in items[:100]:
                    lines.append(f"· {name}")
                if len(items) > 100:
                    lines.append(f"… 其余 {len(items) - 100} 个条目略")
                return "\n".join(lines)
        if ext == ".rar":
            import rarfile
            with rarfile.RarFile(p) as rf:
                infos = rf.infolist()
                lines = [f"压缩包 {path} 共 {len(infos)} 个条目："]
                for info in infos[:100]:
                    lines.append(f"· {info.filename}（{info.file_size} 字节）")
                if len(infos) > 100:
                    lines.append(f"… 其余 {len(infos) - 100} 个条目略")
                return "\n".join(lines)
        return f"错误：不支持的压缩格式 {ext or '（无扩展名）'}（支持 zip/tar/gz/7z/rar）"
    except ImportError as e:
        return f"错误：读取该格式需要额外依赖：{e}"
    except Exception as e:
        return f"错误：读取压缩包失败: {e}"


def _excel_rows_columns(data_rows):
    """C2: 从行数据推断列名：dict 行取首行 keys，否则返回 None（数组行无表头）。"""
    if data_rows and isinstance(data_rows[0], dict):
        return list(data_rows[0].keys())
    return None


def _excel_append_rows(ws, data_rows, cols):
    """C2: 把数据行写入 worksheet；dict 行按 cols 取字段（缺键补空），混合行健壮化。
    数组行保留原始类型（数字不转字符串），与 dict 行行为一致。"""
    if cols is not None:
        for row in data_rows:
            vals = [row.get(c, "") for c in cols] if isinstance(row, dict) else [""] * len(cols)
            ws.append(vals)
    else:
        for row in data_rows:
            ws.append(list(row) if isinstance(row, (list, tuple)) else [row])


@tool(
        {
            "type": "function",
            "function": {
                "name": "write_excel",
                "description": "写入/追加 Excel 文件（.xlsx）。data 传 JSON 数组（行数组或对象数组）；mode=append 追加到已有文件；sheets 传 {表名: 数据行} 一次写多表（此时 data 可传空数组）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "输出文件绝对路径"},
                        "data": {"type": "array", "items": {}, "description": "数据行（sheets 提供时可传 []）"},
                        "sheet": {"type": "string", "description": "可选：工作表名（默认 Sheet1）"},
                        "mode": {"type": "string", "description": "可选：overwrite 覆盖（默认）/ append 追加到已有文件"},
                        "sheets": {"type": "object", "description": "可选：{表名: 数据行} 多表一次写入，与 data 二选一"},
                    },
                    "required": ["path", "data"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='写 Excel',
    preactivate=(('表格', 'excel', 'csv', '报表'),),
)
def write_excel(path, data, sheet="Sheet1", mode="overwrite", sheets=None):
    """写入或追加 Excel 文件（.xlsx）。data 为 JSON 数组（行数组或对象数组）；
    mode=append 追加到已有文件；sheets 为 {表名: 数据行} 一次写多表。"""
    try:
        from openpyxl import Workbook, load_workbook
    except ImportError:
        return "错误：需要 openpyxl（pip install openpyxl）"
    if not path or not str(path).strip():
        return "错误：path 必填"
    if mode not in ("overwrite", "append"):
        return "错误：mode 仅支持 overwrite（覆盖）/ append（追加）"
    p = permissions.resolve(path)
    if not p:
        return "错误：路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    # 归一化写入计划：sheets（多表）优先；否则 data+sheet 单表
    if sheets is not None:
        if not isinstance(sheets, dict) or not sheets:
            return "错误：sheets 必须是 {表名: 数据行} 的非空对象"
        plan = {str(k)[:31]: v for k, v in sheets.items()}
        for sname, rows in plan.items():
            if not isinstance(rows, list):
                return f"错误：表 {sname} 的数据必须是数组"
    else:
        if not isinstance(data, list):
            return "错误：data 必须是非空数组"
        plan = {str(sheet or "Sheet1")[:31]: data}
    try:
        if mode == "append":
            if not os.path.isfile(p):
                return f"错误：追加模式要求文件已存在：{p}"
            wb = load_workbook(p)
            total = 0
            for sname, rows in plan.items():
                if sname in wb.sheetnames:
                    ws = wb[sname]
                    add_header = False  # 已有表直接追加数据，不重复写表头
                else:
                    ws = wb.create_sheet(sname)
                    add_header = True
                cols = _excel_rows_columns(rows)
                if add_header and cols is not None:
                    ws.append(cols)
                _excel_append_rows(ws, rows, cols)
                total += len(rows)
            wb.save(p)
            return f"已追加 Excel 至 {p}（{total} 行，{len(plan)} 个工作表）"
        # overwrite：多表/单表全量重写
        wb = Workbook()
        wb.remove(wb.active)
        total = 0
        for sname, rows in plan.items():
            ws = wb.create_sheet(sname)
            cols = _excel_rows_columns(rows)
            if cols is not None:
                ws.append(cols)
            _excel_append_rows(ws, rows, cols)
            total += len(rows)
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        wb.save(p)
        return f"已写入 Excel 至 {p}（{total} 行，{len(plan)} 个工作表）"
    except Exception as e:
        return f"错误：写入 Excel 失败: {e}"


def _chart_cjk_fonts():
    """C3: 探测系统可用的中文字体（Windows/macOS/Linux），返回 matplotlib sans-serif 候选。
    硬编码单字体在 Linux/无 YaHei 环境会中文变方块；这里按系统字体表实测筛选。"""
    try:
        from matplotlib import font_manager
        available = {f.name for f in font_manager.fontManager.ttflist}
        candidates = [
            "Microsoft YaHei", "SimHei",  # Windows
            "PingFang SC", "Hiragino Sans GB", "Heiti SC", "STHeiti",  # macOS
            "Noto Sans CJK SC", "Source Han Sans SC", "Source Han Sans CN",  # Linux
            "WenQuanYi Zen Hei", "WenQuanYi Micro Hei", "AR PL UMing CN",
            "DejaVu Sans",  # 兜底（无中文但保证不崩）
        ]
        hits = [c for c in candidates if c in available]
        return hits or ["DejaVu Sans"]
    except Exception:
        return ["DejaVu Sans"]


def _chart_parse_series(data):
    """C3: 把 chart_data 的 data 归一化为 [{"name", "xs", "ys"}] 系列列表（兼容旧单系列格式）。
    多系列：data=[{"name":"A","data":[1,2,3]}, {"name":"B","data":[[x,y],...]}]（data 支持数值数组 /
    [x,y] 对数组 / {"x":[], "y":[]} 对象）。"""
    if isinstance(data, dict) and isinstance(data.get("series"), list):
        data = data["series"]
    if not isinstance(data, list) or not data:
        raise ValueError("data 必须是非空数组")

    def to_float(v):
        try:
            return float(v or 0)
        except (TypeError, ValueError):
            raise ValueError(f"非数值数据：{v!r}（请只传数字）")

    first = data[0]
    if isinstance(first, dict) and "data" in first:
        # 多系列：每项 {"name","data":[...]} 或 {"name","data":{"x":[],"y":[]}}
        series = []
        for item in data:
            if not isinstance(item, dict) or "data" not in item:
                raise ValueError(f"多系列项必须是 {{name, data}} 对象：{item!r}")
            name = str(item.get("name", f"系列{len(series) + 1}"))
            rows = item["data"]
            if isinstance(rows, dict) and "y" in rows:
                xs = [to_float(v) for v in rows.get("x", [])] or list(range(len(rows["y"])))
                ys = [to_float(v) for v in rows["y"]]
            elif rows and isinstance(rows[0], (list, tuple)) and len(rows[0]) >= 2:
                xs = [str(r[0]) for r in rows]
                ys = [to_float(r[1]) for r in rows]
            elif rows and isinstance(rows[0], dict) and "y" in rows[0]:
                xs = [str(d.get("x", "")) for d in rows]
                ys = [to_float(d.get("y", 0)) for d in rows]
            else:
                xs = list(range(len(rows)))
                ys = [to_float(v) for v in rows]
            series.append({"name": name, "xs": xs, "ys": ys})
        return series
    # 旧单系列格式
    if isinstance(first, dict) and "x" in first:
        xs = [str(d.get("x", "")) for d in data]
        ys = [to_float(d.get("y", 0)) for d in data]
    elif isinstance(first, (list, tuple)) and len(first) >= 2:
        xs = [str(r[0]) for r in data]
        ys = [to_float(r[1]) for r in data]
    else:
        xs = list(range(len(data)))
        ys = [to_float(x) for x in data]
    return [{"name": "数据", "xs": xs, "ys": ys}]


@tool(
        {
            "type": "function",
            "function": {
                "name": "chart_data",
                "description": "数据可视化：生成图表 PNG（matplotlib）。单系列 data 传 [x,y] 数组/对象数组/数值数组；多系列传 [{\"name\":\"A\",\"data\":[...]},...]；kind: line/bar/pie/scatter（pie 仅单系列）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "data": {"type": "array", "items": {}, "description": "数据：单系列 [[x,y],...] 或 [{\"x\":..,\"y\":..},...] 或 [数值,...]；多系列 [{\"name\":\"A\",\"data\":[数值 或 [x,y] 或 {\"x\":[],\"y\":[]}]},...]"},
                        "path": {"type": "string", "description": "输出 PNG 绝对路径"},
                        "kind": {"type": "string", "description": "可选：line/bar/pie/scatter（默认 line）"},
                        "title": {"type": "string", "description": "可选：图表标题"},
                        "x_label": {"type": "string", "description": "可选：X 轴标签"},
                        "y_label": {"type": "string", "description": "可选：Y 轴标签"},
                    },
                    "required": ["data", "path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='数据可视化图表（线/柱/饼/散点）',
    preactivate=(('图片', '图像', '截图', '看图', '图表', '视觉执行', '视觉闭环', '屏幕操作'), ('表格', 'excel', 'csv', '报表')),
)
def chart_data(data, path, kind="line", title="", x_label="", y_label=""):
    """数据可视化：生成图表 PNG（matplotlib）。单系列 data 为 [x1,x2,...] / [[x,y],...] /
    [{"x":..,"y":..}]；多系列为 [{"name":"A","data":[...]},...]。kind: line/bar/pie/scatter。"""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        # C3: 系统字体探测（Windows/macOS/Linux 通用），避免中文变方块
        plt.rcParams["font.sans-serif"] = _chart_cjk_fonts()
        plt.rcParams["axes.unicode_minus"] = False
    except ImportError:
        return "错误：需要 matplotlib（pip install matplotlib）"
    if not path or not str(path).strip():
        return "错误：path 必填"
    if not isinstance(data, list) or not data:
        return "错误：data 必须是非空数组"
    p = permissions.resolve(path)
    if not p:
        return "错误：路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    try:
        series_list = _chart_parse_series(data)
        # NaN/inf 无法绘制：matplotlib 静默出空图，先挡掉
        for s in series_list:
            if any(v != v or v in (float("inf"), float("-inf")) for v in s["ys"]):
                return "错误：数据包含 NaN 或无穷值，请清洗后重试"
        k = str(kind or "line").lower()
        if k not in ("line", "bar", "pie", "scatter"):
            return f"错误：kind 非法：{kind}（支持 line/bar/pie/scatter）"
        multi = len(series_list) > 1
        if k == "pie":
            if multi:
                return "错误：饼图仅支持单系列数据，请合并为一份后重试"
            if len(series_list[0]["ys"]) > 20:
                return "错误：饼图最多支持 20 个数据点，请聚合后重试"
            if not any(v > 0 for v in series_list[0]["ys"]):
                return "错误：饼图需要至少一个正值数据"
        if _dc.CHART_THEME == "light":
            face = "#ffffff"
            grid = "#d5e4ec"
            tick = "#5c7a96"
            title_c = "#14283f"
            chart_bg = "#f5f9fc"
            palette = ["#00a3c8", "#ff7f50", "#9acd32", "#9370db", "#ffb800",
                       "#ff6b81", "#5cb85c", "#6c8ebf", "#e07020", "#b554c8"]
        else:
            face = "#0a101f"
            grid = "#14203a"
            tick = "#9db0d1"
            title_c = "#e9f1ff"
            chart_bg = "#0a101f"
            palette = ["#00d4ff", "#ff9e6d", "#b8e986", "#c9a8ff", "#ffd700",
                       "#ff8fa3", "#7ee07e", "#8fb8e8", "#f09a55", "#dd88f0"]
        fig, ax = plt.subplots(figsize=(8, 5), dpi=110, facecolor=face)
        ax.set_facecolor(chart_bg)
        for spine in ax.spines.values():
            spine.set_color(grid)
        ax.tick_params(colors=tick)
        ax.xaxis.label.set_color(tick)
        ax.yaxis.label.set_color(tick)
        ax.title.set_color(title_c)
        ax.grid(color=grid)
        total_pts = 0
        for idx, s in enumerate(series_list):
            color = palette[idx % len(palette)]
            xs, ys = s["xs"], s["ys"]
            total_pts += len(ys)
            if k == "bar":
                ax.bar(xs, ys, color=color, label=s["name"] if multi else None)
            elif k == "pie":
                ax.pie(ys, labels=xs, autopct="%1.1f%%", textprops={"color": title_c})
            elif k == "scatter":
                # 保留用户传入的 x 坐标（此前 range(len(ys)) 会把 x 丢弃成序号）
                ax.scatter(xs, ys, color=color, label=s["name"] if multi else None)
            else:
                ax.plot(xs, ys, color=color, marker="o", markersize=4, label=s["name"] if multi else None)
        if multi:
            ax.legend()
        if title:
            ax.set_title(str(title))
        if x_label:
            ax.set_xlabel(str(x_label))
        if y_label:
            ax.set_ylabel(str(y_label))
        fig.tight_layout()
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        fig.savefig(p)
        plt.close(fig)
        size = os.path.getsize(p) if os.path.exists(p) else 0
        return f"已生成图表至 {p}（{size} 字节，{k} 图，{len(series_list)} 系列，{total_pts} 个数据点）"
    except Exception as e:
        return f"错误：生成图表失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "database_query",
                "description": "对本地 SQLite 数据库执行只读查询（仅 SELECT/PRAGMA，最多 200 行），数据库文件须在允许目录内",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "db_path": {"type": "string", "description": "SQLite 数据库文件绝对路径"},
                        "sql": {"type": "string", "description": "只读 SQL 语句（SELECT / PRAGMA）"},
                        "max_rows": {"type": "integer", "description": "可选：返回行数上限（默认 20）"},
                    },
                    "required": ["db_path", "sql"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='SQLite 只读查询',
    preactivate=(('数据库', 'sql', 'mysql', 'postgres'),),
)
def database_query(db_path, sql, max_rows=20):
    """对本地 SQLite 数据库执行只读查询（SELECT/PRAGMA），路径需在允许目录内。"""
    if not db_path or not str(db_path).strip():
        return "错误：请提供数据库文件路径"
    if not sql or not str(sql).strip():
        return "错误：请提供 SQL 查询语句"
    ok, reason = permissions.check_filesystem(str(db_path), write=False)
    if not ok:
        return reason
    stmt = str(sql).strip()
    if not _readonly_stmt(stmt):
        return "错误：仅允许只读查询（SELECT / PRAGMA）"
    p = permissions.resolve(str(db_path))
    if not p or not os.path.exists(p):
        return f"错误：数据库文件不存在：{p}"
    try:
        import sqlite3
        import time as _t
        import urllib.parse

        # 路径含 ?/# 时按 URI 查询参数解析，需先 percent-encode
        conn = sqlite3.connect(
            f"file:{urllib.parse.quote(p)}?mode=ro", uri=True, timeout=5
        )
        try:
            # L3: SQL 层强制 LIMIT——仅 fetchmany 截断不够，无 LIMIT 的查询仍会全量执行
            limit = 20
            try:
                limit = clamp_int(max_rows, 20, lo=1, hi=200)
            except (TypeError, ValueError):
                pass
            stmt = force_limit(stmt, limit)
            # L3: 语句级超时（progress handler 每 500 条虚拟机指令检查一次，超时中断）
            _q_start = _t.monotonic()
            def _q_timeout():
                return 1 if _t.monotonic() - _q_start > _SQLITE_QUERY_TIMEOUT_S else 0
            conn.set_progress_handler(_q_timeout, 500)
            cur = conn.cursor()
            cur.execute(stmt)
            if cur.description is None:
                return "执行成功（无结果集）"
            cols = [d[0] for d in cur.description]
            rows = cur.fetchmany(limit)
            lines = [f"查询结果（{len(rows)} 行）:", " | ".join(str(c) for c in cols)]
            for r in rows:
                cells = ["" if v is None else str(v) for v in r]
                cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
                lines.append(" | ".join(cells))
            if len(rows) >= limit:
                lines.append("⚠ 已达行数上限，如需更多请缩小范围后分页查询")
            return "\n".join(lines)
        finally:
            try:
                conn.set_progress_handler(None, 0)
            except Exception:
                pass
            conn.close()
    except Exception as e:
        return f"错误：查询失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "database_execute",
                "description": "数据库写操作（UPDATE/INSERT/DELETE/DDL）。高危：变更前自动备份 + 审计；SQLite 的 connection 为数据库文件路径，mysql/postgres 用 db_config.json 的连接名",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "db_type": {"type": "string", "description": "sqlite / mysql / postgres"},
                        "connection": {"type": "string", "description": "sqlite=数据库文件绝对路径；mysql/postgres=连接名（默认 default）"},
                        "sql": {"type": "string", "description": "写操作 SQL 语句"},
                        "backup": {"type": "boolean", "description": "可选：变更前备份（默认 true）"},
                    },
                    "required": ["db_type", "sql", "connection"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='数据库写操作（SQLite/MySQL/PG，带审批）',
    preactivate=(('数据库写', '插入数据', '改数据库', '删除记录', 'update语句'),),
)
def database_execute(db_type="sqlite", connection="default", sql="", backup=True):
    """数据库写操作（UPDATE/INSERT/DELETE/DDL）。高危工具，走审批流 + 审计。
    db_type: sqlite / mysql / postgres；sqlite 的 connection 为数据库文件绝对路径。"""
    stmt = str(sql or "").strip()
    if not stmt:
        return "错误：sql 必填"
    if not stmt.lstrip().upper().startswith(("UPDATE", "INSERT", "DELETE", "CREATE", "DROP", "ALTER", "REPLACE")):
        return "错误：database_execute 仅用于写操作；只读查询请用 database_query*"
    # L4: 无 WHERE 的 UPDATE/DELETE 直接拒绝（防全表误操作；全表清空应分批带条件）
    if stmt.lstrip().upper().startswith(("UPDATE", "DELETE")) and not re.search(r"\bWHERE\b", stmt, re.I):
        return "错误：UPDATE/DELETE 必须带 WHERE 条件（防止全表误操作）；如需清空整表请分步删除并确认"
    dbtype = str(db_type or "sqlite").lower()
    try:
        if dbtype == "sqlite":
            return _db_execute_sqlite(str(connection or "").strip(), stmt, bool(backup))
        if dbtype == "mysql":
            return _db_execute_mysql(str(connection or "default"), stmt, bool(backup))
        if dbtype == "postgres":
            return _db_execute_postgres(str(connection or "default"), stmt, bool(backup))
        return f"错误：不支持的数据库类型：{db_type}（支持 sqlite / mysql / postgres）"
    except Exception as e:
        return f"错误：数据库执行失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "pdf_extract",
                "description": "从 PDF 提取文本（按页）、表格（Markdown 格式）或元数据；支持页码范围（如 1-5）。扫描件会提示改用 ocr_image",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "PDF 文件绝对路径（须在允许目录内）"},
                        "pages": {"type": "string", "description": "可选：页码范围，如 '1-5' / '3' / 'all'（默认 all）"},
                        "mode": {"type": "string", "description": "可选：text（文本，默认）/ table（表格）/ meta（元数据）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='提取 PDF 文本',
    preactivate=(('pdf', '转pdf', 'pdf提取', 'pdf生成', '读pdf'),),
)
def pdf_extract(path, pages="all", mode="text"):
    """从 PDF 提取文本（按页）/ 表格（Markdown）/ 元数据；支持页码范围与扫描件提示。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{path}"
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return "未安装 PyMuPDF，请先执行 pip_install PyMuPDF 后重试"
    m = str(mode or "text").lower()
    if m not in ("text", "table", "meta"):
        return f"错误：mode 仅支持 text / table / meta，收到：{mode}"
    try:
        doc = fitz.open(p)
        try:
            if doc.is_encrypted:
                return "PDF 已加密，暂不支持密码解密，请先移除密码后重试"
            total = doc.page_count
            if total == 0:
                return "PDF 为空文档"
            page_list = _parse_page_range(pages, total)
            if page_list is None:
                return f"错误：页码范围非法（应为 1-{total} 内逗号/连字符组合或 all）：{pages}"
            if m == "meta":
                md = doc.metadata or {}
                return "\n".join([
                    f"文件名: {os.path.basename(p)}",
                    f"页数: {total}",
                    f"标题: {md.get('title') or '（无）'}",
                    f"作者: {md.get('author') or '（无）'}",
                    f"主题: {md.get('subject') or '（无）'}",
                    f"创建时间: {md.get('creationDate') or '（无）'}",
                    f"修改时间: {md.get('modDate') or '（无）'}",
                    f"格式: {md.get('format') or '? '}",
                ])
            out = [f"文件名: {os.path.basename(p)}", f"页数: {total}"]
            out_len = sum(len(x) for x in out)
            truncated_hint = "\n[输出较长已截断，可用 pages= 指定页码范围分段提取]"
            if m == "text":
                for i in page_list:
                    page = doc.load_page(i - 1)
                    seg = f"\n--- 第{i}页 ---\n"
                    text = page.get_text("text").strip()
                    if not text:
                        seg += "（本页无文本层，疑似扫描件；可先用 web_screenshot/pdf 导出页面图片再用 ocr_image 识别）"
                    else:
                        seg += text
                    if out_len + len(seg) > PDF_EXTRACT_MAX_OUTPUT:
                        out.append(truncated_hint)
                        break
                    out.append(seg)
                    out_len += len(seg)
            else:  # table
                if not hasattr(doc.load_page(0), "find_tables"):
                    return "错误：当前 PyMuPDF 版本过低，表格提取需要 PyMuPDF 1.23+（可升级：pip_install PyMuPDF --upgrade，或改用 mode=text）"
                for i in page_list:
                    page = doc.load_page(i - 1)
                    seg_parts = [f"\n--- 第{i}页 表格 ---"]
                    try:
                        tables = page.find_tables()
                        found = False
                        for ti, tb in enumerate(tables.tables, 1):
                            data = tb.extract()
                            if not data:
                                continue
                            found = True
                            seg_parts.append(f"表格 {ti}:")
                            seg_parts.append(_table_to_md(data))
                        if not found:
                            seg_parts.append("（本页未检测到表格）")
                    except Exception:
                        seg_parts.append("（表格提取失败，可改用 mode=text 提取文本）")
                    seg = "\n".join(seg_parts)
                    if out_len + len(seg) > PDF_EXTRACT_MAX_OUTPUT:
                        out.append(truncated_hint)
                        break
                    out.append(seg)
                    out_len += len(seg)
            result = "\n".join(out)
            if len(result) > PDF_EXTRACT_MAX_OUTPUT:
                result = result[:PDF_EXTRACT_MAX_OUTPUT] + truncated_hint
            return result
        finally:
            doc.close()
    except Exception as e:
        return f"错误：PDF 读取失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "pdf_create",
                "description": "把文本或 Markdown 内容生成 PDF 文件（自动嵌入中文字体；支持标题/列表/代码块/表格排版）。长文档请分段生成",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "content": {"type": "string", "description": "PDF 正文（文本或 Markdown），与 source_path 二选一"},
                        "source_path": {"type": "string", "description": "可选：从本地 md/txt 文件读取内容（与 content 二选一）"},
                        "output": {"type": "string", "description": "输出 PDF 绝对路径（须在允许目录内）"},
                        "title": {"type": "string", "description": "可选：文档标题（默认取内容首行）"},
                    },
                    "required": ["output"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='生成 PDF',
    preactivate=(('pdf', '转pdf', 'pdf提取', 'pdf生成', '读pdf'),),
)
def pdf_create(content="", source_path="", output="", title=""):
    """把文本/Markdown 内容生成 PDF（中文字体嵌入；支持标题/列表/代码块/表格）。"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Preformatted, Table, TableStyle,
        )
    except ImportError:
        return "未安装 reportlab，请先执行 pip_install reportlab 后重试"
    if not str(output or "").strip():
        return "错误：output 必填"
    if str(content or "").strip() and str(source_path or "").strip():
        return "错误：content 与 source_path 只能二选一"
    if not str(content or "").strip():
        if not str(source_path or "").strip():
            return "错误：content 或 source_path 必填"
        src = permissions.resolve(source_path)
        if not src or not os.path.isfile(src):
            return f"错误：源文件不存在：{source_path}"
        ok, reason = permissions.check_filesystem(src, write=False)
        if not ok:
            return reason
        try:
            with open(src, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(2_000_000)
        except Exception as e:
            return f"错误：读取源文件失败: {e}"
    if not str(content or "").strip():
        return "错误：内容为空"
    out = permissions.resolve(output)
    if not out:
        return "错误：输出路径无效"
    if not out.lower().endswith(".pdf"):
        out += ".pdf"
    ok, reason = permissions.check_filesystem(out, write=True)
    if not ok:
        return reason
    try:
        import mdparse  # 复用项目自有 Markdown 块解析（无第三方依赖）
    except Exception:
        return "错误：mdparse 不可用"
    try:
        os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
        font = _register_cjk_font()
        mono = "Courier" if font == "Helvetica" else font
        styles = {
            "title": ParagraphStyle("title", fontName=font, fontSize=18, leading=24, spaceAfter=14),
            "h1": ParagraphStyle("h1", fontName=font, fontSize=16, leading=20, spaceBefore=10, spaceAfter=6),
            "h2": ParagraphStyle("h2", fontName=font, fontSize=14, leading=18, spaceBefore=8, spaceAfter=5),
            "h3": ParagraphStyle("h3", fontName=font, fontSize=13, leading=17, spaceBefore=6, spaceAfter=4),
            "h4": ParagraphStyle("h4", fontName=font, fontSize=12, leading=16, spaceBefore=5, spaceAfter=3),
            "h5": ParagraphStyle("h5", fontName=font, fontSize=11, leading=15, spaceBefore=4, spaceAfter=2),
            "h6": ParagraphStyle("h6", fontName=font, fontSize=10, leading=14, spaceBefore=4, spaceAfter=2),
            "body": ParagraphStyle("body", fontName=font, fontSize=10.5, leading=16, spaceAfter=6),
            "code": ParagraphStyle("code", fontName=mono, fontSize=8.5, leading=11.5,
                                   backColor=colors.Color(0.95, 0.95, 0.95), borderPadding=6,
                                   spaceBefore=4, spaceAfter=8),
        }
        flow = []
        # 标题：显式 title 优先；否则取内容首行作为文档元数据标题（PRD 规范）
        doc_title = str(title or "").strip()
        if not doc_title:
            first_line = next((ln.strip() for ln in str(content).splitlines() if ln.strip()), "")
            doc_title = first_line[:200]
        if str(title or "").strip():
            flow.append(Paragraph(_md_inline_html(title), styles["title"]))
        blocks = mdparse.parse_blocks(str(content))
        for blk in blocks:
            kind = blk[0]
            body = blk[1]
            if kind == "code":
                flow.append(Preformatted(body, styles["code"]))
            elif kind == "table":
                rows = _md_table_rows(body)
                if len(rows) >= 2:
                    t = Table(rows, repeatRows=1)
                    t.setStyle(TableStyle([
                        ("FONTNAME", (0, 0), (-1, 0), font),
                        ("FONTNAME", (0, 1), (-1, -1), font),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.9, 0.93, 0.97)),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ]))
                    flow.append(t)
                    flow.append(Spacer(1, 8))
            elif kind == "list":
                for ln in str(body).split("\n"):
                    if ln.strip():
                        flow.append(Paragraph("• " + _md_inline_html(ln), styles["body"]))
            elif kind.startswith("h"):
                level = int(kind[1])
                flow.append(Paragraph(_md_inline_html(body), styles[f"h{min(level, 6)}"]))
            else:
                flow.append(Paragraph(_md_inline_html(body), styles["body"]))
        if not flow:
            return "错误：内容未能解析为可排版元素"
        doc = SimpleDocTemplate(out, pagesize=A4,
                                leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54,
                                title=doc_title)
        doc.build(flow)
        size = os.path.getsize(out) if os.path.exists(out) else 0
        permissions.audit("pdf_create", out, f"{size} 字节")
        return f"已生成 PDF: {out}（{size / 1024:.1f} KB，中文字体 {'已嵌入' if font != 'Helvetica' else '未找到（可能乱码，请安装中文字体）'}）"
    except Exception as e:
        return f"错误：PDF 生成失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "docx_read",
                "description": "读取 Word .docx 文档为 Markdown 结构（标题层级/段落/列表/表格），旧版 .doc 会提示转换",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".docx 文件绝对路径"},
                        "max_chars": {"type": "integer", "description": "可选：输出字符上限（默认 50000，防超长文档撑爆上下文）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='读取 Word 文档',
    preactivate=(('word', 'docx', '读word', '读取文档'),),
)
def docx_read(path, max_chars=50000):
    """读取 Word .docx 为 Markdown 结构（标题/段落/列表/表格，保持文档顺序）。"""
    try:
        from docx import Document
    except ImportError:
        return "未安装 python-docx，请先执行 pip_install python-docx 后重试"
    if not str(path or "").strip():
        return "错误：path 必填"
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{path}"
    if p.lower().endswith(".doc"):
        return "错误：暂不支持旧版 .doc 格式，请先用 Word 另存为 .docx 后重试"
    try:
        limit = max(200, min(500000, int(max_chars or DOCX_MAX_DEFAULT)))
    except (TypeError, ValueError):
        limit = DOCX_MAX_DEFAULT
    try:
        from docx.table import Table as _Table
        from docx.text.paragraph import Paragraph as _Para

        doc = Document(p)
        parts = []
        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1]
            if tag == "p":
                para = _Para(child, doc)
                text = para.text.strip()
                if not text:
                    continue
                style = str(para.style.name or "")
                sl = style.lower()
                if sl.startswith("title"):
                    parts.append("# " + text)
                elif sl.startswith("heading"):
                    try:
                        level = int(style[-1])
                    except (TypeError, ValueError):
                        level = 1
                    parts.append("#" * min(level, 6) + " " + text)
                elif "list bullet" in sl:
                    parts.append("- " + text)
                elif "list number" in sl:
                    parts.append("1. " + text)
                else:
                    parts.append(text)
            elif tag == "tbl":
                table = _Table(child, doc)
                rows = [[c.text.strip() for c in r.cells] for r in table.rows]
                if rows and any(any(c for c in r) for r in rows):
                    parts.append(_table_to_md(rows))
        if not parts:
            return "（文档无可见内容）"
        result = "\n\n".join(parts)
        if len(result) > limit:
            result = result[:limit] + f"\n[内容较长已截断前 {limit} 字符]"
        return result
    except Exception as e:
        return f"错误：Word 读取失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "pptx_read",
                "description": "提取 PowerPoint .pptx 每页幻灯片的标题、正文要点与演讲者备注；图片以占位标注",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": ".pptx 文件绝对路径"},
                        "include_notes": {"type": "boolean", "description": "可选：是否包含演讲者备注（默认 true）"},
                    },
                    "required": ["path"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='读取 PPT',
    preactivate=(('ppt', 'pptx', '演示文稿', '读ppt'),),
)
def pptx_read(path, include_notes=True):
    """提取 PowerPoint 每页幻灯片的标题、正文要点与备注；图片占位标注。"""
    try:
        from pptx import Presentation
    except ImportError:
        return "未安装 python-pptx，请先执行 pip_install python-pptx 后重试"
    if not str(path or "").strip():
        return "错误：path 必填"
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{path}"
    if p.lower().endswith(".ppt"):
        return "错误：暂不支持旧版 .ppt 格式，请先用 PowerPoint 另存为 .pptx 后重试"
    try:
        prs = Presentation(p)
        out = [f"幻灯片数: {len(prs.slides)}"]

        def _walk_shapes(shapes, title_holder, body_lines, img_count, table_count, depth=0):
            """递归遍历形状（组合形状 GROUP 内文本/图片不丢），深度上限防畸形文件死循环。"""
            if depth > 10:
                return
            for shape in shapes:
                is_title_ph = (
                    getattr(shape, "is_placeholder", False)
                    and shape.placeholder_format.idx == 0
                    and shape.has_text_frame
                )
                if is_title_ph:
                    t = shape.text_frame.text.strip()
                    if t and not title_holder[0]:
                        title_holder[0] = t
                    continue
                if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    if hasattr(shape, "shapes"):
                        _walk_shapes(shape.shapes, title_holder, body_lines, img_count, table_count, depth + 1)
                elif shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        body_lines.extend(ln for ln in t.splitlines() if ln.strip())
                elif getattr(shape, "has_table", False):
                    table_count[0] += 1
                    rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                    if rows:
                        body_lines.append(f"[表格 {table_count[0]}]")
                        body_lines.append(_table_to_md(rows))
                elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_count[0] += 1

        for idx, slide in enumerate(prs.slides, 1):
            out.append(f"\n--- 第{idx}页 ---")
            title_holder = [None]
            body_lines = []
            img_count = [0]
            table_count = [0]
            _walk_shapes(slide.shapes, title_holder, body_lines, img_count, table_count)
            if title_holder[0]:
                out.append(f"标题: {title_holder[0]}")
            if body_lines:
                out.append("\n".join("- " + ln for ln in body_lines[:40]))
            elif not title_holder[0]:
                out.append("（本页无文本）")
            if img_count[0]:
                out.append(f"[图片占位: {img_count[0]} 张]")
            if include_notes and slide.has_notes_slide:
                try:
                    nt = slide.notes_slide.notes_text_frame.text.strip()
                    if nt:
                        out.append(f"备注: {nt[:500]}")
                except Exception:
                    pass
        return "\n".join(out)
    except Exception as e:
        return f"错误：PPT 读取失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "secret_store",
                "description": "密钥保险箱：DPAPI 加密托管 API key/令牌等敏感值。action=set 保存（value 只写不显示）/ get 按名取用 / delete 删除 / list 仅列出名称",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "action": {"type": "string", "description": "set / get / delete / list"},
                        "name": {"type": "string", "description": "密钥名称（如 openai_key）"},
                        "value": {"type": "string", "description": "set 必填：要托管的敏感值"},
                    },
                    "required": ["action"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='加密密钥存储',
    preactivate=(('密钥', 'api key', '令牌', '保险箱', '托管密码'),),
)
def secret_store(action="get", name="", value=""):
    """密钥保险箱：set / get / delete / list。value 只写不显示；list 只返回名称。"""
    act = str(action or "get").strip().lower()
    name = str(name or "").strip()
    data = _load_secrets()
    if act == "set":
        if not name:
            return "错误：name 必填"
        if not str(value or ""):
            return "错误：value 必填"
        data[name] = str(value)
        if not _save_secrets(data):
            return "错误：保存密钥失败"
        permissions.audit("secret_set", name, "值已加密保存", result="ok")
        return f"密钥「{name}」已加密保存（调用 secret_get(name) 取用）"
    if act == "get":
        if not name:
            return "错误：name 必填"
        if name not in data:
            return f"未找到密钥「{name}」"
        return str(data[name])
    if act == "delete":
        if not name:
            return "错误：name 必填"
        if name not in data:
            return f"未找到密钥「{name}」"
        data.pop(name, None)
        _save_secrets(data)
        permissions.audit("secret_delete", name, "删除")
        return f"已删除密钥「{name}」"
    if act == "list":
        return "已保存密钥：" + ("、".join(sorted(data.keys())) if data else "（空）")
    return "错误：action 仅支持 set / get / delete / list"


@tool(
        {
            "type": "function",
            "function": {
                "name": "kv_store",
                "description": "嵌入式键值存储：set 写入（可选 TTL 过期）/ get 读取 / delete 删除 / keys 列出全部 / search 按键或值模糊检索。适合缓存、配置、轻量状态（Redis 的零部署替代）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "action": {"type": "string", "description": "set / get / delete / keys / search"},
                        "key": {"type": "string", "description": "set/get/delete 必填：键"},
                        "value": {"type": "string", "description": "set 必填：值（上限 1MB）"},
                        "ttl_seconds": {"type": "integer", "description": "可选：set 时有效秒数（0=长期）"},
                        "pattern": {"type": "string", "description": "search 必填：键或值的模糊检索关键词"},
                    },
                    "required": ["action"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='轻量键值存储（缓存/状态）',
    preactivate=(('键值', 'kv存储', '缓存读写', '轻量状态'),),
)
def kv_store(action="get", key="", value="", pattern="", ttl_seconds=0):
    """嵌入式键值存储：set（可选 TTL）/ get / delete / keys / search。"""
    act = str(action or "get").strip().lower()
    if act not in ("set", "get", "delete", "keys", "search"):
        return "错误：action 仅支持 set / get / delete / keys / search"
    try:
        import diskcache
    except ImportError:
        return "未安装 diskcache，请先执行 pip_install diskcache 后重试"
    if not _dc.KV_CACHE_DIR:
        return "错误：KV 存储未初始化"
    try:
        os.makedirs(_dc.KV_CACHE_DIR, exist_ok=True)
    except Exception:
        return "错误：KV 目录创建失败"
    try:
        with diskcache.Cache(_dc.KV_CACHE_DIR) as cache:
            if act == "set":
                k = str(key or "").strip()
                if not k:
                    return "错误：set 需要 key"
                if len(k) > 256:
                    return "错误：key 过长（上限 256 字符）"
                # 区分 None 与 0/False：None 存空串，0/False 保留字面值
                v = str(value if value is not None else "")
                if len(v.encode("utf-8", "replace")) > KV_VALUE_MAX_BYTES:
                    return "错误：value 超过 1MB 上限"
                try:
                    ttl = max(0, min(365 * 24 * 3600, int(ttl_seconds or 0)))
                except (TypeError, ValueError):
                    ttl = 0
                cache.set(k, v, expire=ttl or None)
                n = len(cache)
                ttl_txt = f"TTL {ttl}s" if ttl else "长期有效"
                return f"已写入 key={k}（{ttl_txt}，当前共 {n} 个键）"
            if act == "get":
                k = str(key or "").strip()
                if not k:
                    return "错误：get 需要 key"
                v = cache.get(k)
                if v is None:
                    return f"key={k} 不存在或已过期"
                return f"key={k}: {v}"
            if act == "delete":
                k = str(key or "").strip()
                if not k:
                    return "错误：delete 需要 key"
                if k in cache:
                    del cache[k]
                    return f"已删除 key={k}"
                return f"key={k} 不存在"
            if act == "keys":
                keys = [k for k in cache.iterkeys() if not k.startswith("_")]
                if not keys:
                    return "KV 存储为空"
                lines = [f"共 {len(keys)} 个键："]
                for i, k in enumerate(sorted(keys), 1):
                    try:
                        v = cache.get(k) or ""
                        lines.append(f"{i}. {k} = {str(v)[:80]}")
                    except Exception:
                        lines.append(f"{i}. {k} = ?")
                return "\n".join(lines)
            # search
            pat = str(pattern or "").strip().lower()
            if not pat:
                return "错误：search 需要 pattern"
            hits = []
            for k in cache.iterkeys():
                if k.startswith("_"):
                    continue
                try:
                    v = str(cache.get(k) or "")
                except Exception:
                    continue
                if pat in k.lower() or pat in v.lower():
                    hits.append(f"{k}: {v[:120]}")
                    if len(hits) >= 50:
                        break
            if not hits:
                return f"未找到包含「{pattern}」的键或值"
            return f"命中 {len(hits)} 项：\n" + "\n".join(hits)
    except Exception as e:
        return f"错误：KV 操作失败: {e}"


@tool(
        {
            "type": "function",
            "function": {
                "name": "create_doc",
                "description": "创建文档（.md/.html 原生支持；.docx 依赖 python-docx），需 write 权限",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "文档绝对路径"},
                        "content": {"type": "string", "description": "文档内容（Markdown 或 HTML 文本）"},
                        "doc_type": {"type": "string", "description": "md / html / docx，默认按扩展名推断"},
                    },
                    "required": ["path", "content"],
                },
            },
        },
    groups=['📊 数据与文档'],
    phrases='创建 Office 文档（docx/pptx/pdf）',
    preactivate=(('写', '保存', '创建', '生成'),),
)
def create_doc(path, content, doc_type=""):
    """创建文档：.md/.html 原生；.docx 依赖 python-docx（缺失时返回安装指引）。"""
    ok, reason = permissions.check_filesystem(path, write=True)
    if not ok:
        return reason
    p = permissions.resolve(path)
    ext = (doc_type or "").lower() or os.path.splitext(p)[1].lstrip(".").lower()
    try:
        if ext == "docx":
            try:
                from docx import Document
            except ImportError:
                return "错误：生成 .docx 需要 python-docx：pip install python-docx"
            doc = Document()
            for para in (content or "").splitlines():
                if para.strip():
                    doc.add_paragraph(para)
            created = not os.path.exists(p)
            if not created:
                try:
                    shutil.copy2(p, p + ".bak")
                except Exception:
                    pass
            doc.save(p)
        else:
            if ext != "html":
                ext = "md"
            body = content or ""
            if ext == "html" and not body.lstrip().startswith("<"):
                body = (
                    "<!DOCTYPE html><html lang='zh'><head><meta charset='utf-8'>"
                    "<title>鲸语文档</title></head><body><pre>"
                    + body.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    + "</pre></body></html>"
                )
            created, _ = _atomic_write(p, body)
        permissions.audit("create_doc", p, f"{ext} {len(content or '')} 字符")
        return f"已创建文档 {p}（{'新建' if created else '覆盖并备份 .bak'}，{ext} 格式）"
    except Exception as e:
        return f"错误：文档创建失败: {e}"


__all__ = ['database_query_mysql', 'database_query_postgres', 'read_excel', 'epub_read', 'mobi_read', 'doc_read', 'msg_read', 'archive_list', 'write_excel', 'chart_data', 'database_query', 'database_execute', 'pdf_extract', 'pdf_create', 'docx_read', 'pptx_read', 'secret_store', 'kv_store', 'create_doc']
