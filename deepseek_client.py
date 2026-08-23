import contextlib
import copy
import itertools
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import threading
import time
import weakref
import ast
from collections import deque
from concurrent.futures import FIRST_COMPLETED, ThreadPoolExecutor, wait
from concurrent.futures.thread import _threads_queues, _worker  # 内部接口：daemon 化所需
from datetime import datetime
from urllib.parse import quote

import httpx
from openai import OpenAI
from openai import (
    APIError,
    APIConnectionError,
    APITimeoutError,
    RateLimitError,
)

import permissions
import crypto
from shared import cron_field_ok, OCR_IMAGE_PS, is_peak_hour  # noqa: F401  # is_peak_hour 由 main.py 经本模块 re-export
from security import (  # noqa: F401  # 供 main/tests 继续经 deepseek_client 访问
    SSRF_TRUSTED,
    set_ssrf_trusted,
    _trusted_host,
    _is_private_host,
    _safe_url,
    _url_host,
)
from db_utils import (  # noqa: F401  # 兼容旧访问名
    TABLE_CELL_MAX as _TABLE_CELL_MAX,
    DB_FORBIDDEN_KEYWORDS as _DB_FORBIDDEN_KEYWORDS,
    readonly_stmt as _readonly_stmt,
    db_preview_sql as _db_preview_sql,
    table_to_md as _table_to_md,
)
from pdf_utils import (  # noqa: F401  # 兼容旧访问名
    parse_page_range as _parse_page_range,
    find_cjk_font as _find_cjk_font,
    register_cjk_font as _register_cjk_font,
    md_inline_html as _md_inline_html,
    md_table_rows as _md_table_rows,
)
from proc_utils import kill_tree as _kill_tree  # noqa: F401  # 兼容旧访问名
from net_utils import (  # noqa: F401  # 兼容旧访问名
    DEFAULT_UA as _SEARCH_UA,
    FETCH_URL_TIMEOUT,
    _http_client,
    _shutdown_http_client,
    _safe_redirect_url,
)
from search_utils import (  # noqa: F401  # 兼容旧访问名
    BING_RESULT_RE as _BING_RESULT_RE,
    DDG_RESULT_RE as _DDG_RESULT_RE,
    SO360_RESULT_RE as _SO360_RESULT_RE,
    strip_tags as _strip_tags,
    decode_ddg_url as _decode_ddg_url,
    search_dedup as _search_dedup,
    search_safe as _search_safe,
)
import plugins as plugins_mod

# 按需加载能力：fetch_blocked（机场代理访问被墙站点）。独立模块按用户需要放
# 入项目目录并启用后才生效；文件缺失/被剔除时功能静默降级（不阻塞主程序）。
try:
    from fetch_blocked import fetch_blocked as _fetch_blocked_impl
except ImportError:
    _fetch_blocked_impl = None

logger = logging.getLogger("whaletalk")


def _decrypt_secret(v):
    """解密外部配置中的 dpapi: 密文字段；明文或空串原样返回。"""
    if isinstance(v, str) and v.startswith("dpapi:"):
        try:
            return crypto.decrypt(v)
        except Exception:
            logger.exception("外部配置敏感字段解密失败")
            return ""
    return v


class _DaemonThreadPool(ThreadPoolExecutor):
    """工具并行执行器：daemon 线程 + 模块级复用。

    CPython 3.9+ 的 ThreadPoolExecutor 在创建线程后立即 start，无法事后设置
    daemon（会抛 "cannot set daemon status of active thread"——用户实测崩溃）。
    这里复刻 _adjust_thread_count 并在 start 前设置 daemon，同时避免两个问题：
    - 非 daemon worker 在解释器退出时阻塞（进程挂起）；
    - 每次工具轮新建线程池的创建开销。
    """

    def _adjust_thread_count(self):
        if self._idle_semaphore.acquire(timeout=0):
            return

        def weakref_cb(_, q=self._work_queue):
            q.put(None)

        num_threads = len(self._threads)
        if num_threads < self._max_workers:
            t = threading.Thread(
                name="%s_%d" % (self._thread_name_prefix or self, num_threads),
                target=_worker,
                args=(weakref.ref(self, weakref_cb), self._work_queue,
                      self._initializer, self._initargs),
            )
            t.daemon = True  # 必须在 start() 之前
            t.start()
            self._threads.add(t)
            _threads_queues[t] = self._work_queue


_TOOL_EXECUTOR = _DaemonThreadPool(max_workers=4, thread_name_prefix="tool")
# 长耗时工具独立线程池：避免 run_wechat_writer / pip_install / 数据库大查询等
# 长时间占用普通工具池的全部 4 个 worker，导致同轮快速工具被慢任务排队拖死。
_LONG_TOOL_EXECUTOR = _DaemonThreadPool(max_workers=2, thread_name_prefix="longtool")
_LONG_TOOL_NAMES = frozenset({
    "run_wechat_writer",
    "pip_install",
    "database_execute",
    "database_query_mysql",
    "database_query_postgres",
    "webdav",
    "fetch_blocked",
    "download_file",
    "media_ffmpeg",
    "run_python",
    "run_command",
    "browser_navigate",
    "web_screenshot",
    "image_generate",
    "knowledge_index",
    "run_workflow",
    "schedule_task",
    "subagent_run",
    "run_tests",
    "verify_output",
    "read_email",
    "rss_fetch",
    "daily_brief",
    "search_web",
    "search_github",
    "search_realtime",
    "screen_see",
    "chart_read",
    "screenshot_to_html",
    "debug_screenshot",
    "scan_read",
    "image_batch",
})


def _tool_executor_for(name):
    """按工具名选择执行池：长任务走独立池，普通任务走快速池。"""
    return _LONG_TOOL_EXECUTOR if name in _LONG_TOOL_NAMES else _TOOL_EXECUTOR

DEFAULT_BASE_URL = "https://api.deepseek.com"
DEFAULT_MODEL = "deepseek-v4-flash"

SCENARIOS = {
    "通用": {"temperature": 1.0, "top_p": 1.0, "reasoning_effort": "high"},
    "编程": {"temperature": 0.15, "top_p": 0.95, "reasoning_effort": "max"},
    "Agent": {"temperature": 1.0, "top_p": 0.95, "reasoning_effort": "max"},
    "自定义": {"temperature": 1.0, "top_p": 1.0, "reasoning_effort": "high"},
}

MODELS = {
    "deepseek-v4-flash": {
        "label": "DeepSeek V4 Flash",
        "version": "DeepSeek-V4-Flash-0731",
        "max_context_tokens": 1_000_000,
        "max_output_tokens": 384 * 1024,
    },
    "deepseek-v4-pro": {
        "label": "DeepSeek V4 Pro",
        "version": "DeepSeek-V4-Pro-0813",
        "max_context_tokens": 1_000_000,
        "max_output_tokens": 384 * 1024,
    },
    "deepseek-v4-flash-vision-exp": {
        "label": "DeepSeek V4 Flash Vision (实验)",
        "version": "DeepSeek-V4-Flash-Vision-Exp",
        "max_context_tokens": 1_000_000,
        "max_output_tokens": 384 * 1024,
        "vision": True,
    },
}

VISION_MODEL = "deepseek-v4-flash-vision-exp"

# 图片内联限制（官方图像理解文档）：单张 ≤ 32 MiB，请求体 ≤ 48 MiB。
# base64 开销 4/3，留出文本/工具 schema 余量后按 40 MiB 计入。
IMAGE_MAX_BYTES = 32 * 1024 * 1024
IMAGE_INLINE_TOTAL_BASE64 = 40 * 1024 * 1024
IMAGE_EXTENSIONS = (".jpg", ".jpeg", ".png", ".gif", ".webp")


def is_vision_model(model):
    """判断模型是否支持图片输入（视觉模型）。"""
    if not model:
        return False
    if "vision" in str(model).lower():
        return True
    return bool((MODELS.get(model) or {}).get("vision"))


def _detect_image_mime(buf):
    """按文件实际内容（魔数）识别图片 MIME，不依赖文件名/声明类型。"""
    if buf[:3] == b"\xff\xd8\xff":
        return "image/jpeg"
    if buf[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if buf[:6] in (b"GIF87a", b"GIF89a"):
        return "image/gif"
    if buf[:4] == b"RIFF" and buf[8:12] == b"WEBP":
        return "image/webp"
    return "image/png"


def _is_image_path(p):
    low = str(p or "").lower()
    return low.endswith(IMAGE_EXTENSIONS) or low.startswith(("http://", "https://"))


def embed_message_images(messages, model, _log=None, detail="auto"):
        """把 user 消息中的 images（本地路径或 http(s) 图片 URL）内联为 image_url 内容块。

        - 仅对消息的浅拷贝生效，不修改调用方内存中的消息对象（UI/存档保持文本 content）；
        - 模型不支持视觉时抛 ValueError（附中文提示，UI 层会切换视觉模型）；
        - 单张 ≤32 MiB、base64 总量 ≤40 MiB，超限抛 ValueError 提示压缩；
        - 图片仅出现在 user 消息中（官方限制：system/assistant 携带图片返回 400）；
        - detail 控制官方图片处理级别：low（缩放 512 省 token）/ high / original / auto。
        """
        import base64

        if _log is None:
            _log = logger
        has_image = any(
            isinstance(m, dict) and m.get("images") and m.get("role") == "user"
            for m in messages
        )
        if not has_image:
            return messages
        if not is_vision_model(model):
            raise ValueError(
                f"当前模型 {model} 不支持图片输入，请切换到视觉模型 {VISION_MODEL}"
            )
        out = []
        total_b64 = 0
        for m in messages:
            imgs = m.get("images")
            if not (imgs and m.get("role") == "user"):
                out.append(m)
                continue
            m = dict(m)
            blocks = []
            text = m.get("content")
            if text:
                blocks.append({"type": "text", "text": str(text)})
            for p in imgs:
                try:
                    if str(p).strip().lower().startswith(("http://", "https://")):
                        url = p
                        err = _safe_url(url)
                        if err:
                            raise ValueError(f"图片 URL 不安全：{err}")
                        with _safe_stream("GET", url, timeout=20) as resp:
                            resp.raise_for_status()
                            raw = b""
                            for chunk in resp.iter_bytes(64 * 1024):
                                raw += chunk
                                if len(raw) > IMAGE_MAX_BYTES:
                                    raise ValueError(
                                        f"图片下载超过 {IMAGE_MAX_BYTES // (1024 * 1024)}MB，请先压缩"
                                    )
                            mime = _detect_image_mime(raw[:16])
                            b64 = base64.b64encode(raw).decode("ascii")
                        blocks.append(
                            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}", "detail": detail}}
                        )
                        total_b64 += len(b64)
                    else:
                        local = p
                        if os.path.isfile(local) and not os.path.isabs(local):
                            local = os.path.abspath(local)
                        if not os.path.isfile(local):
                            raise ValueError(f"图片文件不存在：{p}")
                        try:
                            size = os.path.getsize(local)
                        except OSError as e:
                            raise ValueError(f"图片文件读取失败：{p}: {e}")
                        if size > IMAGE_MAX_BYTES:
                            raise ValueError(
                                f"图片超过 {IMAGE_MAX_BYTES // (1024 * 1024)}MB，请先用 image_process 压缩：{os.path.basename(p)}"
                            )
                        with open(local, "rb") as f:
                            raw = f.read()
                        b64 = base64.b64encode(raw).decode("ascii")
                        blocks.append(
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{_detect_image_mime(raw[:16])};base64,{b64}",
                                    "detail": detail,
                                },
                            }
                        )
                        total_b64 += len(b64)
                except Exception as e:
                    raise ValueError(f"图片 {p} 处理失败: {e}")
            if total_b64 > IMAGE_INLINE_TOTAL_BASE64:
                raise ValueError("本轮图片总量过大（超过请求体限制），请减少图片或先压缩")
            if not blocks:
                blocks.append({"type": "text", "text": ""})
            m["content"] = blocks
            out.append(m)
        _log.info("已内联 %d 条消息的图片（base64 共 %.1f MB）", len(out), total_b64 / 1024 / 1024)
        return out

THINKING_MODES = {
    "none": "禁用思考 (none)",
    "low": "低思考 (low)",
    "medium": "中等思考 (medium)",
    "high": "高思考 (high)",
    "xhigh": "极高思考 (xhigh)",
    "max": "最大思考 (max)",
    "auto": "智能路由 (auto)",
}

EFFORT_BY_THINKING = {
    # 官方 effort 映射表（deepseek-v4-flash / pro 一致）：
    #   low→low · medium→high · high→high · xhigh→high · max→max
    # 全部档位 API 均接受；UI 如实展示档位与映射结果。
    "low": "low",
    "medium": "high",
    "high": "high",
    "xhigh": "high",
    "max": "max",
    "auto": None,  # 鲸语智能路由：启发式映射到 none/low/high/max
}

_AUTO_COMPLEX_WORDS = ("分析", "设计", "审查", "解释", "重构", "优化", "实现", "编写", "创建", "对比")
_AUTO_SIMPLE_WORDS = ("你好", "在吗", "谢谢", "再见", "ok", "yes", "no", "哈哈", "好的")


def _auto_effort(work):
    """auto 思考档：按任务复杂度启发式路由到 none/high/max（无额外 API 成本）。

    评分维度：内容长度 / 代码块 / 复杂词 / 多步骤结构 / 简单寒暄扣分。
    """
    text = ""
    for m in reversed(work):
        if m.get("role") == "user" and m.get("content"):
            c = m["content"]
            if isinstance(c, list):
                text = " ".join(
                    str(b.get("text", ""))
                    for b in c
                    if isinstance(b, dict) and b.get("type") == "text"
                )
            else:
                text = str(c)
            break
    score = 0
    if len(text) > 300:
        score += 1
    if "```" in text:
        score += 1
    if any(w in text for w in _AUTO_COMPLEX_WORDS):
        score += 1
    if any(w in text for w in _AUTO_SIMPLE_WORDS):
        score -= 1
    # 多步骤任务：编号/步骤式指令 ≥3 条，或段落数较多 → 升级思考深度
    step_lines = [
        ln for ln in text.splitlines()
        if ln.strip() and (
            ln.lstrip()[:2].rstrip(".").isdigit()
            or ln.lstrip().startswith(("步骤", "第一步", "然后", "接着", "- "))
        )
    ]
    if len(step_lines) >= 3:
        score += 1
    elif text.count("\n") >= 5:
        score += 1
    if score >= 2:
        return "max"
    if score == 1:
        return "high"
    return "none"

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "get_date",
            "description": "获取当前日期、具体时间与本地时区（如 2026-08-03 15:30:00 CST）",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "ask_user",
            "description": "向用户提问（遇到歧义、缺少关键信息、需确认高风险操作时使用）。阻塞等待用户回答后继续",
            "parameters": {
                "type": "object",
                "properties": {"prompt": {"type": "string", "description": "向用户提出的问题（简洁明确，可给出选项）"}},
                "required": ["prompt"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "request_permission",
            "description": "当工具执行被权限拒绝时，请求用户将操作加入白名单（允许目录 / 命令白名单 / 开启文件写权限）。用户同意后立即生效，可重试原操作",
            "parameters": {
                "type": "object",
                "properties": {
                    "action_type": {"type": "string", "description": "白名单类型：dir=加入允许目录 / command=加入命令白名单 / write=开启文件写权限"},
                    "value": {"type": "string", "description": "要加入白名单的值：目录绝对路径 或 命令名；write 类型可留空"},
                },
                "required": ["action_type"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_memory",
            "description": "写入一条长期记忆（用户偏好、关键结论、重要事实），自动去重，最多 2000 条；可附带类型、实体与关系三元组形成知识图谱",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "要记住的内容"},
                    "tags": {"type": "string", "description": "可选：逗号分隔的标签，便于检索"},
                    "type": {"type": "string", "description": "可选：记忆类型（偏好/事实/项目/联系/规则 等）"},
                    "entities": {"type": "string", "description": "可选：涉及的实体列表，逗号分隔，如 张三,项目A"},
                    "relations": {"type": "string", "description": "可选：关系三元组，分号分隔的 实体-关系-实体，如 张三-负责-项目A"},
                },
                "required": ["text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_memory",
            "description": "读取长期记忆：关键词支持语义相似度检索（不含关键词也能匹配相关记忆）；可按类型/实体过滤",
            "parameters": {
                "type": "object",
                "properties": {
                    "keyword": {"type": "string", "description": "可选：检索关键词（按语义相似度排序）"},
                    "max_items": {"type": "integer", "description": "可选：返回条数上限（默认 20）"},
                    "type": {"type": "string", "description": "可选：按记忆类型过滤（偏好/事实/项目/联系/规则 等）"},
                    "entity": {"type": "string", "description": "可选：按实体过滤（知识图谱节点）"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "query_memory_graph",
            "description": "知识图谱查询：按实体或关系检索关联记忆（返回结构化图谱片段），适合查找人与项目/任务间的关联",
            "parameters": {
                "type": "object",
                "properties": {
                    "entity": {"type": "string", "description": "可选：实体名（如 张三 / 项目A）"},
                    "relation": {"type": "string", "description": "可选：关系名（如 负责 / 参与）"},
                    "max_items": {"type": "integer", "description": "可选：返回条数上限（默认 20）"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_weather",
            "description": "获取指定城市某日期的天气（date 仅支持今天与近 3 天预报）",
            "parameters": {
                "type": "object",
                "properties": {
                    "location": {"type": "string", "description": "城市名称"},
                    "date": {"type": "string", "description": "日期，格式 YYYY-mm-dd（今天或未来 3 天内；留空为今天）"},
                },
                "required": ["location", "date"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_python",
            "description": "在隔离的 Python 子进程中执行代码（默认 -S 不加载第三方库、无网络库）；with_site=true 时加载已安装的第三方库并可访问外网（httpx/requests 等），需要新库时先调用 pip_install 安装",
            "parameters": {
                "type": "object",
                "properties": {
                    "code": {"type": "string", "description": "Python 代码"},
                    "with_site": {"type": "boolean", "description": "可选：true 时加载第三方库并允许网络请求（需已安装，如 pip_install 安装的库）"},
                },
                "required": ["code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "读取本地文本文件（UTF-8，默认前 100KB，须在允许目录内）；可指定 start_line/max_lines 按行读取超大文件",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "文件绝对路径（须在允许目录内）"},
                    "start_line": {"type": "integer", "description": "可选：起始行号（从 1 开始，与 max_lines 配合按行读取）"},
                    "max_lines": {"type": "integer", "description": "可选：读取行数上限（默认 200，最大 2000）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "fetch_url",
            "description": "抓取指定 URL 的文本/JSON 内容（超时 10 秒，最多 500KB）",
            "parameters": {
                "type": "object",
                "properties": {"url": {"type": "string", "description": "完整 URL，含 http(s)://"}},
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "fetch_blocked",
            "description": "抓取被墙/国际站点（linux.do、Google、archive.org 等）的文本/JSON 内容，自动使用本机机场节点代理 + 浏览器指纹绕过封锁。适用于 fetch_url 超时/失败或被墙的场景",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "目标网址，完整 URL 含 http(s)://"},
                    "proxy": {"type": "string", "description": "可选：代理节点，形如 https://user:pass@server:port，留空自动发现"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_web",
            "description": "联网搜索最新信息，返回相关网页的标题、链接与摘要（Bing + 360 + DuckDuckGo 并行聚合去重，默认最多 5 条，可指定 num 最多 20 条）。site 与 offset 参数保证生效；since/until 时间过滤依赖搜索引擎支持（可能不严格）。适合查询实时新闻、最新资讯、不熟悉的事实等；找到有用链接后可配合 fetch_url 抓取全文",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "搜索关键词（建议简洁明确）"},
                    "num": {"type": "integer", "description": "可选：返回条数（1-20，默认 5）"},
                    "offset": {"type": "integer", "description": "可选：翻页偏移（0 起，如 5 表示跳过前 5 条看第 6-10 条）"},
                    "since": {"type": "string", "description": "可选：起始日期过滤（YYYY-MM-DD，依赖引擎支持，可能不严格）"},
                    "until": {"type": "string", "description": "可选：截止日期过滤（YYYY-MM-DD，依赖引擎支持，可能不严格）"},
                    "site": {"type": "string", "description": "可选：限定站点域名（如 openai.com），只返回该站结果（保证生效）"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_github",
            "description": "搜索 GitHub 开源仓库（按 Star 排序）。支持 GitHub 原生搜索语法：org:（组织）、topic:、language:、stars:、in:readme 等，例如 org:deepseek-ai 精确查官方组织",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "搜索关键词（支持 org:/topic:/language:/stars: 等原生语法）"},
                    "num": {"type": "integer", "description": "可选：返回条数（1-20，默认 5）"},
                    "language": {"type": "string", "description": "可选：限定编程语言（如 python、javascript）"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_realtime",
            "description": "实时信息通道（Hacker News）：不传 query 返回当前热点榜（含点赞数），传 query 走全文搜索（含点赞/评论数）。适合查询正在发生的热点、技术社区讨论、实时新闻（弥补 search_web 实时性短板）",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "可选：搜索关键词；留空返回实时热点榜"},
                    "num": {"type": "integer", "description": "可选：返回条数（1-20，默认 5）"},
                    "source": {"type": "string", "description": "可选：数据源（当前仅支持 hn）"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "call_api",
            "description": "通用外部 API 调用（万能接口）：GET/POST/PUT/DELETE/PATCH，支持自定义查询参数/JSON 体/表单体/请求头。可对接任意开放 API（天气/翻译/大模型/企业服务等）。安全限制：仅公网 http(s) 地址（禁内网/回环，除非在配置 call_api_allowed_hosts 白名单中显式放行），响应 ≤500KB，超时 ≤180s",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "完整 API 地址（http/https）"},
                    "method": {"type": "string", "description": "请求方法：GET/POST/PUT/DELETE/PATCH/HEAD（默认 GET）"},
                    "params": {"type": "object", "description": "可选：查询参数对象（如 {\"limit\": 10}）"},
                    "json_body": {"type": "object", "description": "可选：JSON 请求体对象"},
                    "data": {"type": "string", "description": "可选：表单/原始请求体"},
                    "headers": {"type": "object", "description": "可选：自定义请求头（≤16 个，如 {\"Authorization\": \"Bearer xxx\"}）"},
                    "timeout": {"type": "integer", "description": "可选：超时秒数（1-180，默认 15）"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "system_status",
            "description": "系统资源自检：CPU 使用率、内存占用、工作区磁盘剩余、网络连通性（api.github.com/bing/api.deepseek.com）。适合任务前环境体检、排查网络/资源问题",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    # ===== 只读查询（需文件在允许目录内）=====
    {
        "type": "function",
        "function": {
            "name": "database_query",
            "description": "对本地 SQLite 数据库执行只读查询（仅 SELECT/PRAGMA，最多 200 行），数据库文件须在允许目录内",
            "parameters": {
                "type": "object",
                "properties": {
                    "db_path": {"type": "string", "description": "SQLite 数据库文件绝对路径"},
                    "sql": {"type": "string", "description": "只读 SQL 语句（SELECT / PRAGMA）"},
                    "max_rows": {"type": "integer", "description": "可选：返回行数上限（默认 20）"},
                },
                "required": ["db_path", "sql"],
            },
        },
    },
    # ===== 高风险工具（需权限设置开启 + 配置）=====
    {
        "type": "function",
        "function": {
            "name": "send_email",
            "description": "发送邮件（需要先在数据目录配置 email_config.json 的 SMTP 信息，未配置会提示）",
            "parameters": {
                "type": "object",
                "properties": {
                    "to": {"type": "string", "description": "收件人邮箱"},
                    "subject": {"type": "string", "description": "邮件主题"},
                    "body": {"type": "string", "description": "邮件正文"},
                },
                "required": ["to", "subject", "body"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "pip_install",
            "description": "安装 Python 库（完全体模式放行任意包），安装后配合 run_python(with_site=true) 使用；已装常用库：openpyxl/matplotlib/pymysql/psycopg2/Pillow 等",
            "parameters": {
                "type": "object",
                "properties": {"package": {"type": "string", "description": "要安装的包名（完全体模式不限白名单）"}},
                "required": ["package"],
            },
        },
    },
    # ===== L1 行动层（全部默认不启用，需权限设置开启）=====
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "写文件（自动建目录，目录白名单校验 + 大小限制 + 原子写 + 自动 .bak）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "目标文件绝对路径（须在允许目录内）"},
                    "content": {"type": "string", "description": "文件完整内容"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_file",
            "description": "编辑文件：按文本替换或正则替换（自动备份 .bak），需 write 权限",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "文件绝对路径"},
                    "old": {"type": "string", "description": "要替换的原文（与 regex 二选一）"},
                    "new": {"type": "string", "description": "替换后的新文本"},
                    "regex": {"type": "string", "description": "可选：正则表达式模式（Python re 语法）"},
                },
                "required": ["path", "new"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_dir",
            "description": "列出目录内容（只读，默认允许，目录须在允许目录内）",
            "parameters": {
                "type": "object",
                "properties": {"path": {"type": "string", "description": "目录绝对路径"}},
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": "执行白名单命令（python/pip/pytest/git 等，禁止 shell 拼接），需开启 shell 权限并可能需确认",
            "parameters": {
                "type": "object",
                "properties": {"command": {"type": "string", "description": "完整命令行，如 python hello.py"}},
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_local",
            "description": "在允许目录内检索文本文件内容（只读，支持常见文本格式）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "要检索的目录绝对路径"},
                    "query": {"type": "string", "description": "检索关键词"},
                    "max_results": {"type": "integer", "description": "最多返回条数，默认 20"},
                },
                "required": ["path", "query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_doc",
            "description": "创建文档（.md/.html 原生支持；.docx 需安装 python-docx），需 write 权限",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "文档绝对路径"},
                    "content": {"type": "string", "description": "文档内容（Markdown 或 HTML 文本）"},
                    "doc_type": {"type": "string", "description": "md / html / docx，默认按扩展名推断"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_code_project",
            "description": "创建多文件代码工程（批量写文件，自动建目录，需 write 权限，单次最多 50 个文件）",
            "parameters": {
                "type": "object",
                "properties": {
                    "project_dir": {"type": "string", "description": "工程根目录绝对路径（须在允许目录内）"},
                    "files": {
                        "type": "array",
                        "description": "文件列表",
                        "items": {
                            "type": "object",
                            "properties": {
                                "path": {"type": "string", "description": "相对工程根目录的路径，如 src/main.py"},
                                "content": {"type": "string", "description": "文件内容"},
                            },
                            "required": ["path", "content"],
                        },
                    },
                },
                "required": ["project_dir", "files"],
            },
        },
    },
    # ===== 桌面 RPA（P0）=====
    {
        "type": "function",
        "function": {
            "name": "rpa_screen_size",
            "description": "获取当前屏幕分辨率（桌面 RPA 坐标用，需 pyautogui）",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rpa_click",
            "description": "桌面 RPA：模拟鼠标点击屏幕坐标 (x,y)，button=left/right/middle",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer", "description": "屏幕 X 坐标"},
                    "y": {"type": "integer", "description": "屏幕 Y 坐标"},
                    "button": {"type": "string", "description": "可选：left/right/middle，默认 left"},
                    "clicks": {"type": "integer", "description": "可选：连击次数 1-5，默认 1"},
                },
                "required": ["x", "y"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rpa_type",
            "description": "桌面 RPA：模拟键盘输入文本（需先点击目标输入框）",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "要输入的文本"},
                    "interval": {"type": "number", "description": "可选：每个字符间隔秒数，默认 0.02"},
                },
                "required": ["text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rpa_hotkey",
            "description": "桌面 RPA：模拟组合键，如 ctrl+c / alt+tab / ctrl+shift+esc",
            "parameters": {
                "type": "object",
                "properties": {"keys": {"type": "string", "description": "组合键串，+ 分隔"}},
                "required": ["keys"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rpa_move",
            "description": "桌面 RPA：把鼠标移动到屏幕坐标 (x,y)",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer", "description": "X 坐标"},
                    "y": {"type": "integer", "description": "Y 坐标"},
                    "duration": {"type": "number", "description": "可选：移动耗时秒，默认 0.2"},
                },
                "required": ["x", "y"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rpa_scroll",
            "description": "桌面 RPA：滚动鼠标滚轮（正数向上，负数向下）",
            "parameters": {
                "type": "object",
                "properties": {
                    "clicks": {"type": "integer", "description": "滚动格数 -50~50"},
                    "x": {"type": "integer", "description": "可选：滚动位置 X"},
                    "y": {"type": "integer", "description": "可选：滚动位置 Y"},
                },
                "required": ["clicks"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rpa_screenshot",
            "description": "桌面 RPA：截取整个屏幕保存 PNG（默认工作区）",
            "parameters": {
                "type": "object",
                "properties": {"path": {"type": "string", "description": "可选：输出 PNG 绝对路径"}},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_navigate",
            "description": "浏览器可视操作（open/click/type/fill/submit/select/get_text），需安装 playwright。浏览器实例复用：连续调用共享同一页面，登录态保持，click/type/submit 不重新导航（多步操作有效）",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "目标网址（非 open 动作时若已在同页则不重复导航）"},
                    "action": {"type": "string", "description": "open / click / type / fill / submit / select / get_text"},
                    "selector": {"type": "string", "description": "CSS 选择器（click/type/fill/select/get_text 需要）"},
                    "text": {"type": "string", "description": "要输入的文本（type/fill）或要选择的选项（select）"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_screenshot",
            "description": "网页截图并保存到工作区，需安装 playwright（可选依赖）",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "目标网址"},
                    "width": {"type": "integer", "description": "视口宽度，默认 1280"},
                    "height": {"type": "integer", "description": "视口高度，默认 800"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "publish_draft",
            "description": "保存发布草稿到本地草稿箱（只建草稿不发布，发布权始终在用户手中）",
            "parameters": {
                "type": "object",
                "properties": {
                    "platform": {"type": "string", "description": "目标平台，如 公众号/博客/小红书"},
                    "title": {"type": "string", "description": "草稿标题"},
                    "content": {"type": "string", "description": "草稿正文"},
                },
                "required": ["platform", "title", "content"],
            },
        },
    },
    # ===== 推送通知（A6）：Webhook 一键推送 =====
    {
        "type": "function",
        "function": {
            "name": "send_webhook",
            "description": "推送通知到配置的 Webhook（钉钉/ServerChan/Slack/通用，webhooks.json 配置）。适合任务完成提醒、定时巡检结果推送",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "可选：通知标题（默认 鲸语提醒）"},
                    "text": {"type": "string", "description": "通知正文"},
                    "channel": {"type": "string", "description": "可选：指定通道（dingtalk/serverchan/slack/generic），留空推送全部已配置通道"},
                },
                "required": ["text"],
            },
        },
    },
    # ===== IM 主动触达（P1）=====
    {
        "type": "function",
        "function": {
            "name": "im_send",
            "description": "发送消息到 IM 通道（Telegram Bot / 企业微信群机器人，im_config.json 配置）。用于任务完成主动汇报、定时提醒。",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "消息正文"},
                    "title": {"type": "string", "description": "可选：消息标题（默认 鲸语提醒）"},
                    "channel": {"type": "string", "description": "可选：指定通道 telegram/wecom，留空推送全部已配置通道"},
                },
                "required": ["text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "telegram_poll_updates",
            "description": "接收 Telegram Bot 新消息（长轮询，游标自动去重）。AI 可定期调用检查用户是否通过 Telegram 召唤/下达新指令。",
            "parameters": {
                "type": "object",
                "properties": {
                    "timeout": {"type": "integer", "description": "可选：长轮询秒数 1-60，默认 15"},
                    "limit": {"type": "integer", "description": "可选：最多返回条数 1-20，默认 5"},
                },
                "required": [],
            },
        },
    },
    # ===== 二进制下载 / 文件格式扩展（P2）=====
    {
        "type": "function",
        "function": {
            "name": "download_file",
            "description": "下载二进制文件（图片/附件/文档/安装包等任意格式）到工作区或指定目录；流式写盘，单文件 200MB 上限。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "文件完整 URL（http/https）"},
                    "local_path": {"type": "string", "description": "可选：本地保存路径（留空保存到工作区 downloads/）"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "epub_read",
            "description": "读取 EPUB 电子书正文（ebooklib 可选依赖）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "EPUB 文件绝对路径"},
                    "max_chars": {"type": "integer", "description": "可选：最多返回字符数（默认 20000）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "mobi_read",
            "description": "读取 MOBI 电子书正文（mobi 可选依赖）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "MOBI 文件绝对路径"},
                    "max_chars": {"type": "integer", "description": "可选：最多返回字符数（默认 20000）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "doc_read",
            "description": "读取旧版 .doc 二进制文档正文（antiword/catdoc 可选依赖）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": ".doc 文件绝对路径"},
                    "max_chars": {"type": "integer", "description": "可选：最多返回字符数（默认 20000）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "msg_read",
            "description": "读取 .msg Outlook 邮件（主题/发件人/正文/附件清单，extract_msg 可选依赖）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": ".msg 文件绝对路径"},
                    "max_chars": {"type": "integer", "description": "可选：正文最多返回字符数（默认 20000）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "archive_list",
            "description": "列出压缩包内容：.zip / .tar / .gz / .7z / .rar（7z/rar 需可选依赖）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "压缩包绝对路径"},
                },
                "required": ["path"],
            },
        },
    },
    # ===== 多模态（A3）：语音合成 / 图像处理 / OCR =====
    {
        "type": "function",
        "function": {
            "name": "tts_save",
            "description": "把文本合成为语音 WAV 文件（Windows SAPI 中文语音），可调语速（rate -10~10）",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "要合成的文本"},
                    "path": {"type": "string", "description": "输出 WAV 文件绝对路径（须在允许目录内）"},
                    "rate": {"type": "integer", "description": "可选：语速 -10~10，默认 0"},
                },
                "required": ["text", "path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "image_process",
            "description": "图像处理：缩放/裁剪/旋转/格式转换/加水印（PIL）。ops 用分号分隔多个操作，如 resize=800x600; water=测试水印",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "源图片绝对路径"},
                    "output": {"type": "string", "description": "输出图片绝对路径"},
                    "ops": {"type": "string", "description": "可选：操作串。resize=宽x高; crop=x1,y1,x2,y2; rotate=度数; convert=PNG/JPEG; quality=1-100; water=水印文本"},
                },
                "required": ["path", "output"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "ocr_image",
            "description": "从图片文件提取文字（Windows OCR，适合截图/扫描件，需系统安装中文语言包）",
            "parameters": {
                "type": "object",
                "properties": {"path": {"type": "string", "description": "图片文件绝对路径"}},
                "required": ["path"],
            },
        },
    },
    # ===== 数据工具（A4）：CSV / Excel / 图表 / MySQL / PostgreSQL =====
    {
        "type": "function",
        "function": {
            "name": "read_csv",
            "description": "读取 CSV 文件（允许目录内），返回表格文本",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "CSV 文件绝对路径"},
                    "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 100）"},
                    "delimiter": {"type": "string", "description": "可选：分隔符（默认逗号）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_csv",
            "description": "写入 CSV 文件。rows 传 JSON 数组：[[v,v],...] 或 [{\"列\":值},...]",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "输出文件绝对路径"},
                    "rows": {"type": "array", "items": {}, "description": "数据行（数组的数组，或对象数组）"},
                    "headers": {"type": "string", "description": "可选：表头，逗号分隔"},
                },
                "required": ["path", "rows"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_excel",
            "description": "读取 Excel 文件（.xlsx，openpyxl），返回表格文本",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Excel 文件绝对路径"},
                    "sheet": {"type": "string", "description": "可选：工作表名或序号（默认第一个）"},
                    "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 100）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_excel",
            "description": "写入 Excel 文件（.xlsx）。data 传 JSON 数组（行数组或对象数组）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "输出文件绝对路径"},
                    "data": {"type": "array", "items": {}, "description": "数据行"},
                    "sheet": {"type": "string", "description": "可选：工作表名（默认 Sheet1）"},
                },
                "required": ["path", "data"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "chart_data",
            "description": "数据可视化：生成图表 PNG（matplotlib）。data 传 [x,y] 数组或对象数组或数值数组；kind: line/bar/pie/scatter",
            "parameters": {
                "type": "object",
                "properties": {
                    "data": {"type": "array", "items": {}, "description": "数据：[[x,y],...] 或 [{\"x\":..,\"y\":..},...] 或 [数值,...]"},
                    "path": {"type": "string", "description": "输出 PNG 绝对路径"},
                    "kind": {"type": "string", "description": "可选：line/bar/pie/scatter（默认 line）"},
                    "title": {"type": "string", "description": "可选：图表标题"},
                    "x_label": {"type": "string", "description": "可选：X 轴标签"},
                    "y_label": {"type": "string", "description": "可选：Y 轴标签"},
                },
                "required": ["data", "path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "database_query_mysql",
            "description": "MySQL 只读查询（SELECT/SHOW/DESC）。连接在数据目录 db_config.json 的 mysql.<connection> 配置",
            "parameters": {
                "type": "object",
                "properties": {
                    "connection": {"type": "string", "description": "可选：连接名（默认 default）"},
                    "sql": {"type": "string", "description": "只读 SQL 语句"},
                    "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 20）"},
                },
                "required": ["sql"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "database_query_postgres",
            "description": "PostgreSQL 只读查询（SELECT/SHOW/DESC）。连接在数据目录 db_config.json 的 postgres.<connection> 配置",
            "parameters": {
                "type": "object",
                "properties": {
                    "connection": {"type": "string", "description": "可选：连接名（默认 default）"},
                    "sql": {"type": "string", "description": "只读 SQL 语句"},
                    "max_rows": {"type": "integer", "description": "可选：最多返回行数（默认 20）"},
                },
                "required": ["sql"],
            },
        },
    },
    # ===== 子代理并行编排（A9）与自我验证闭环（A8）=====
    {
        "type": "function",
        "function": {
            "name": "subagent_run",
            "description": "并行子代理：把大任务拆成多个子任务，交给多个并发 AI 子代理分别完成并汇总结果（适合并行调研/多方案对比/多文件并行处理）",
            "parameters": {
                "type": "object",
                "properties": {
                    "tasks": {"type": "array", "items": {"type": "string"}, "description": "子任务列表（字符串数组，最多 8 个），如 [\"总结文件A\", \"总结文件B\"]"},
                    "parallel": {"type": "integer", "description": "可选：并行数 1-4（默认 2）"},
                    "context": {"type": "string", "description": "可选：共享背景上下文（注入每个子代理）"},
                },
                "required": ["tasks"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_tests",
            "description": "运行测试（pytest/unittest）并返回结果摘要：自我验证闭环第一步（写代码后必须自测）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "可选：测试文件或目录（留空自动扫描允许目录内的 test_*.py）"},
                    "framework": {"type": "string", "description": "可选：auto/pytest/unittest（默认 auto）"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "verify_output",
            "description": "对照标准答案自评：计算语义相似度（F1/覆盖率），指出缺失要点。自我验证闭环第二步：完成任务后与预期结果对照检查",
            "parameters": {
                "type": "object",
                "properties": {
                    "expected": {"type": "string", "description": "预期答案/标准要点"},
                    "actual": {"type": "string", "description": "实际输出/实现结果"},
                },
                "required": ["expected", "actual"],
            },
        },
    },
    # ===== 后台进程（服务器/长驻任务，实时输出见「进程终端」）=====
    {
        "type": "function",
        "function": {
            "name": "start_process",
            "description": "在后台启动长驻进程（如网站服务器），实时输出显示在进程终端，返回进程名与 pid",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "完整命令行，如 python -m http.server 8000 或 uvicorn app:app"},
                    "name": {"type": "string", "description": "可选：进程名（便于后续停止/查询）"},
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "stop_process",
            "description": "停止后台进程（按名称或 pid）",
            "parameters": {
                "type": "object",
                "properties": {
                    "target": {"type": "string", "description": "进程名或 pid，如 http.server 或 12345"},
                },
                "required": ["target"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_processes",
            "description": "列出所有后台进程的运行状态与最近输出（运行中/已退出）",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "environment_info",
            "description": "获取运行环境信息：Python 版本、已安装的常用包、工作区磁盘空间（避免重复安装已存在的东西）",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    # ===== 自我进化（感知自身代码 → 分支提案，不修改原文件）=====
    {
        "type": "function",
        "function": {
            "name": "project_info",
            "description": "感知鲸语自身代码库：版本、项目文件清单与规模（只读，自我进化分析用）",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_project_file",
            "description": "读取鲸语自身源码文件（仅限项目目录内 .py/.md/.json/.txt/.bat/.html，只读；支持 offset/limit 分页读取大文件）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "项目内文件路径，如 main.py 或 deepseek_client.py"},
                    "offset": {"type": "integer", "description": "可选：起始字符偏移（分页读取）"},
                    "limit": {"type": "integer", "description": "可选：本次读取字符数上限"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_evolution",
            "description": "自我进化：把改进后的代码作为提案写入 evolutions/<名称>_<时间戳>/ 分支（绝不修改原文件），用户可在「工具 → 自我进化」查看/采纳/忽略",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "提案名称，如 fix_typo / optimize_render"},
                    "files": {
                        "type": "array",
                        "description": "修改后的文件列表",
                        "items": {
                            "type": "object",
                            "properties": {
                                "path": {"type": "string", "description": "相对项目根目录的路径，如 main.py"},
                                "content": {"type": "string", "description": "修改后的完整文件内容"},
                            },
                            "required": ["path", "content"],
                        },
                    },
                },
                "required": ["name", "files"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "verify_files",
            "description": "批量核验文件是否存在及其大小（写文件/建工程后必须调用自检，防止幻觉；相对路径基于工作目录）",
            "parameters": {
                "type": "object",
                "properties": {
                    "paths": {
                        "type": "array",
                        "description": "要核验的文件路径列表（绝对路径或相对工作目录的路径）",
                        "items": {"type": "string"},
                    }
                },
                "required": ["paths"],
            },
        },
    },
    # ===== v2 能力层 · 任务调度 =====
    {
        "type": "function",
        "function": {
            "name": "schedule_task",
            "description": "创建定时任务（可主动安排：每周五周报、每小时巡检、生日提醒等）。expr_type=cron（5字段 分 时 日 月 周）/time（HH:MM 每日）/every（每 N 分钟）；action=message 到点自动发送指令执行 / notify 状态栏提醒 / backup 项目备份",
            "parameters": {
                "type": "object",
                "properties": {
                    "expr_type": {"type": "string", "description": "cron / time / every（默认 cron）"},
                    "expr": {"type": "string", "description": "表达式：cron 如 '30 9 * * 1'；time 如 '09:00'；every 如 '60'"},
                    "content": {"type": "string", "description": "到点要执行的内容（message 发指令 / notify 提醒文本）"},
                    "action": {"type": "string", "description": "message / notify / backup（默认 message）"},
                    "name": {"type": "string", "description": "可选：任务名称（便于后续取消/查看）"},
                    "enabled": {"type": "boolean", "description": "可选：是否启用（默认 true）"},
                },
                "required": ["expr", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_schedules",
            "description": "列出全部定时任务（id/时间/动作/内容/启用状态），配合 cancel_schedule 管理",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "cancel_schedule",
            "description": "取消定时任务（按 list_schedules 返回的 id 或名称）",
            "parameters": {
                "type": "object",
                "properties": {"target": {"type": "string", "description": "任务 id 或名称"}},
                "required": ["target"],
            },
        },
    },
    # ===== v2 能力层 · 桌面通知 =====
    {
        "type": "function",
        "function": {
            "name": "notify_desktop",
            "description": "发送 Windows 桌面 Toast 通知（离线可用）：任务完成、定时任务触发、长任务结束时提醒",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "可选：通知标题（默认 鲸语提醒）"},
                    "text": {"type": "string", "description": "通知正文"},
                },
                "required": ["text"],
            },
        },
    },
    # ===== v2 能力层 · 剪贴板 =====
    {
        "type": "function",
        "function": {
            "name": "clipboard_get",
            "description": "读取用户剪贴板文本（敏感：需用户确认授权），适合用户复制内容后直接处理",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "clipboard_set",
            "description": "把整理好的内容写入剪贴板，用户可直接粘贴使用",
            "parameters": {
                "type": "object",
                "properties": {"text": {"type": "string", "description": "要写入剪贴板的内容"}},
                "required": ["text"],
            },
        },
    },
    # ===== v2 能力层 · 文件闭环（删除/压缩/解压/批量重命名） =====
    {
        "type": "function",
        "function": {
            "name": "delete_file",
            "description": "删除文件或目录（默认移入回收站可恢复；permanent=true 才物理删除）。高危：需审批",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "要删除的文件/目录绝对路径（须在允许目录内）"},
                    "permanent": {"type": "boolean", "description": "可选：true 物理删除不可恢复（默认 false 移入回收站）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "archive_files",
            "description": "把多个文件/目录打包为 zip（工作区内，自动建目录）",
            "parameters": {
                "type": "object",
                "properties": {
                    "paths": {"type": "array", "description": "要打包的文件/目录绝对路径列表", "items": {"type": "string"}},
                    "output": {"type": "string", "description": "输出 zip 文件绝对路径"},
                },
                "required": ["paths", "output"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "extract_archive",
            "description": "解压 zip 压缩包到目标目录（自动越界防护）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "zip 文件绝对路径"},
                    "dest_dir": {"type": "string", "description": "解压目标目录绝对路径"},
                },
                "required": ["path", "dest_dir"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "batch_rename",
            "description": "批量重命名：把目录内所有文件名中的 pattern 替换为 replacement（dry_run=true 预览不实际改名）",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "目录绝对路径"},
                    "pattern": {"type": "string", "description": "要替换的字符串"},
                    "replacement": {"type": "string", "description": "替换后的字符串"},
                    "dry_run": {"type": "boolean", "description": "可选：true 仅预览"},
                },
                "required": ["directory", "pattern", "replacement"],
            },
        },
    },
    # ===== v2 能力层 · 媒体感知（图片理解/屏幕截图/语音识别） =====
    {
        "type": "function",
        "function": {
            "name": "image_understand",
            "description": "用多模态模型理解图片（本地文件路径或 http(s) 图片 URL）。当前模型不支持视觉时自动改用 deepseek-v4-flash-vision-exp，无需手动切换",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "图片文件绝对路径或 http(s) 图片 URL"},
                    "question": {"type": "string", "description": "可选：要问的问题（默认描述图片内容）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "screen_see",
            "description": "截图并让视觉模型解读当前屏幕（一步完成 截图+看图）。RPA/浏览器操作后自查首选：看清界面后决定下一步（点击/输入/验证）",
            "parameters": {
                "type": "object",
                "properties": {
                    "question": {"type": "string", "description": "可选：要看什么（默认描述屏幕内容）"},
                    "area": {"type": "string", "description": "可选：区域 left,top,right,bottom（默认全屏）"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "chart_read",
            "description": "图表截图 → 结构化数据 + 解读（折线/柱状/饼图/散点等，适合读报表/仪表盘截图）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "图表图片文件绝对路径"},
                    "question": {"type": "string", "description": "可选：针对图表的具体问题"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "screenshot_to_html",
            "description": "UI/网页截图 → 还原为 HTML+CSS 页面（前端还原），可保存到文件",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "截图文件绝对路径"},
                    "out_path": {"type": "string", "description": "可选：输出 HTML 绝对路径（默认仅返回代码）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "debug_screenshot",
            "description": "报错/异常截图 → 识别错误并给出诊断与修复建议（错误码/文案/行号/原因/修复）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "报错截图文件绝对路径"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "scan_read",
            "description": "扫描件/文档图片读取（图表、公式、手写、印刷混排），返回 Markdown 结构化内容",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "扫描件/文档图片绝对路径"},
                    "question": {"type": "string", "description": "可选：要提取/回答的内容"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "image_batch",
            "description": "批量视觉分析文件夹内图片：逐张理解后汇总报告（小并发，适合图库/截图/素材批量整理）",
            "parameters": {
                "type": "object",
                "properties": {
                    "folder": {"type": "string", "description": "图片所在目录绝对路径"},
                    "question": {"type": "string", "description": "可选：每张图要回答的问题（默认描述）"},
                    "pattern": {"type": "string", "description": "可选：文件通配符（默认 *.png，支持 *.jpg）"},
                    "max": {"type": "integer", "description": "可选：最多分析张数（1-200，默认 100）"},
                },
                "required": ["folder"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "screen_capture",
            "description": "截取当前屏幕保存到工作区（敏感：需审批）。配合 ocr_image/image_understand 描述屏幕内容",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "可选：输出 PNG 绝对路径（默认工作区 screenshots/）"},
                    "area": {"type": "string", "description": "可选：区域 left,top,right,bottom（默认全屏）"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "speech_to_text",
            "description": "本地语音转文字（faster-whisper 离线识别，未安装时提示安装；首次运行自动下载所选模型）",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "音频文件绝对路径（wav/mp3/m4a 等）"},
                    "model": {"type": "string", "description": "可选：tiny/base/small/medium/large-v3（默认 base，tiny 最快）"},
                },
                "required": ["path"],
            },
        },
    },
    # ===== v2 能力层 · 本地知识库 RAG =====
    {
        "type": "function",
        "function": {
            "name": "knowledge_index",
            "description": "对目录内文本文件建立语义检索索引（TF-IDF+bigram，零依赖）。建索引后可语义检索，措辞不同也能命中",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "可选：要索引的目录（默认工作区）"},
                    "force": {"type": "boolean", "description": "可选：强制重建"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "knowledge_search",
            "description": "语义检索知识库（先 knowledge_index 建索引）：找『之前写过的关于预算的文档』这类措辞模糊的问题",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "检索关键词/语义描述"},
                    "top_k": {"type": "integer", "description": "可选：返回条数（默认 5，最大 10）"},
                },
                "required": ["query"],
            },
        },
    },
    # ===== v2 能力层 · 数据库写操作（高危：审批 + 备份） =====
    {
        "type": "function",
        "function": {
            "name": "database_execute",
            "description": "数据库写操作（UPDATE/INSERT/DELETE/DDL）。高危：走审批流 + 变更前备份 + 审计；SQLite 的 connection 为数据库文件路径，mysql/postgres 用 db_config.json 的连接名",
            "parameters": {
                "type": "object",
                "properties": {
                    "db_type": {"type": "string", "description": "sqlite / mysql / postgres"},
                    "connection": {"type": "string", "description": "sqlite=数据库文件绝对路径；mysql/postgres=连接名（默认 default）"},
                    "sql": {"type": "string", "description": "写操作 SQL 语句"},
                    "backup": {"type": "boolean", "description": "可选：变更前备份（默认 true）"},
                },
                "required": ["db_type", "sql", "connection"],
            },
        },
    },
    # ===== v2 能力层 · 收邮件 =====
    {
        "type": "function",
        "function": {
            "name": "read_email",
            "description": "读取邮箱近期邮件（IMAP，email_config.json 配置 imap 段）。敏感：需审批",
            "parameters": {
                "type": "object",
                "properties": {
                    "limit": {"type": "integer", "description": "可选：最多返回封数（默认 10，最大 50）"},
                    "since_days": {"type": "integer", "description": "可选：最近 N 天（默认 3，0=全部）"},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "email_summary",
            "description": "读取近期邮件并整理为清单，供 AI 生成新邮件摘要（IMAP 配置同 read_email）",
            "parameters": {
                "type": "object",
                "properties": {
                    "limit": {"type": "integer", "description": "可选：最多返回封数（默认 10，最大 50）"},
                    "since_days": {"type": "integer", "description": "可选：最近 N 天（默认 1）"},
                },
            },
        },
    },
    # ===== Agent Mail（agently-cli，可选配置） =====
    {
        "type": "function",
        "function": {
            "name": "agent_mail",
            "description": "Agent 原生邮箱（agently-cli）：me 查看身份 / list 列邮件 / search 搜索 / read 读取 / send 发送 / reply 回复 / forward 转发 / trash 移回收站 / delete 永久删除 / download 下载附件。写操作需两阶段确认：首次调用返回 confirmation-token，向用户确认后再次调用",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "me / list / search / read / send / reply / forward / trash / delete / download"},
                    "q": {"type": "string", "description": "search：关键词"},
                    "id": {"type": "string", "description": "read/reply/forward/trash：msg_xxx"},
                    "to": {"type": "string", "description": "send/forward：收件人，多个用逗号分隔"},
                    "subject": {"type": "string", "description": "send：主题"},
                    "body": {"type": "string", "description": "send/reply/forward：正文"},
                    "dir": {"type": "string", "description": "list/search：inbox/sent/trash/spam"},
                    "limit": {"type": "integer", "description": "list/search：返回条数（默认 10）"},
                    "cursor": {"type": "string", "description": "list/search：翻页游标"},
                    "confirmation_token": {"type": "string", "description": "写操作二次确认：首次调用返回的 ctk_xxx"},
                    "attachment": {"type": "string", "description": "send/reply：附件路径，多个逗号分隔"},
                    "msg": {"type": "string", "description": "download：msg_xxx"},
                    "att": {"type": "string", "description": "download：att_xxx"},
                    "output": {"type": "string", "description": "download：保存目录"},
                },
                "required": ["action"],
            },
        },
    },
    # ===== v2 能力层 · 任务检查点（断点续跑） =====
    {
        "type": "function",
        "function": {
            "name": "task_checkpoint_save",
            "description": "保存任务进度检查点（长任务每完成一步就保存，崩溃/重启后可用 task_checkpoint_load 从断点继续）",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "任务名称"},
                    "status": {"type": "string", "description": "可选：状态（进行中/已完成/阻塞等）"},
                    "pending": {"type": "array", "description": "可选：剩余待办步骤列表", "items": {"type": "string"}},
                    "notes": {"type": "string", "description": "可选：进度备注"},
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "task_checkpoint_load",
            "description": "读取任务检查点（恢复未完成任务上下文）",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    # ===== v2 能力层 · 流程编排 =====
    {
        "type": "function",
        "function": {
            "name": "run_workflow",
            "description": "运行已保存的流程模板（workflows.json）：按顺序逐条发送指令，上一步完成后自动执行下一步",
            "parameters": {
                "type": "object",
                "properties": {"name": {"type": "string", "description": "流程名称"}},
                "required": ["name"],
            },
        },
    },
    # ===== v2 能力层 · 图片生成 =====
    {
        "type": "function",
        "function": {
            "name": "image_generate",
            "description": "生成图片（需在 config.json 配置 image_api_key/image_base_url/image_model，OpenAI 兼容 images API）",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {"type": "string", "description": "图片描述提示词"},
                    "path": {"type": "string", "description": "可选：输出路径（默认工作区 images/）"},
                    "size": {"type": "string", "description": "可选：尺寸如 1024x1024"},
                },
                "required": ["prompt"],
            },
        },
    },
    # ===== v2 能力层 · 用量洞察 =====
    {
        "type": "function",
        "function": {
            "name": "usage_report",
            "description": "生成用量洞察报告（近 N 天 token/费用/缓存命中/逐日明细），可配合定时任务每周自动生成",
            "parameters": {
                "type": "object",
                "properties": {"days": {"type": "integer", "description": "可选：统计最近 N 天（默认 7，最大 90）"}},
            },
        },
    },
    # ===== 文档处理：PDF 提取 / PDF 生成 / Word 读取 / PPT 读取 =====
    {
        "type": "function",
        "function": {
            "name": "pdf_extract",
            "description": "从 PDF 提取文本（按页）、表格（Markdown 格式）或元数据；支持页码范围（如 1-5）。扫描件会提示改用 ocr_image",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "PDF 文件绝对路径（须在允许目录内）"},
                    "pages": {"type": "string", "description": "可选：页码范围，如 '1-5' / '3' / 'all'（默认 all）"},
                    "mode": {"type": "string", "description": "可选：text（文本，默认）/ table（表格）/ meta（元数据）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "pdf_create",
            "description": "把文本或 Markdown 内容生成 PDF 文件（自动嵌入中文字体；支持标题/列表/代码块/表格排版）。长文档请分段生成",
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {"type": "string", "description": "PDF 正文（文本或 Markdown），与 source_path 二选一"},
                    "source_path": {"type": "string", "description": "可选：从本地 md/txt 文件读取内容（与 content 二选一）"},
                    "output": {"type": "string", "description": "输出 PDF 绝对路径（须在允许目录内）"},
                    "title": {"type": "string", "description": "可选：文档标题（默认取内容首行）"},
                },
                "required": ["output"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "docx_read",
            "description": "读取 Word .docx 文档为 Markdown 结构（标题层级/段落/列表/表格），旧版 .doc 会提示转换",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": ".docx 文件绝对路径"},
                    "max_chars": {"type": "integer", "description": "可选：输出字符上限（默认 50000，防超长文档撑爆上下文）"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "pptx_read",
            "description": "提取 PowerPoint .pptx 每页幻灯片的标题、正文要点与演讲者备注；图片以占位标注",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": ".pptx 文件绝对路径"},
                    "include_notes": {"type": "boolean", "description": "可选：是否包含演讲者备注（默认 true）"},
                },
                "required": ["path"],
            },
        },
    },
    # ===== 资讯聚合：RSS 订阅 / 抓取 / 简报 =====
    {
        "type": "function",
        "function": {
            "name": "rss_fetch",
            "description": "RSS 订阅管理：list 列出订阅 / preset 一键添加精选源（机器之心/量子位/少数派/IT之家/开源中国/HN）/ add 添加源 / remove 移除源 / fetch 抓取最新条目（含标题/链接/时间/摘要）。可配合 schedule_task 生成每日简报",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "list / preset / add / remove / fetch"},
                    "url": {"type": "string", "description": "add / fetch 时必填：RSS 源地址（http(s)）"},
                    "limit": {"type": "integer", "description": "可选：返回条数上限（默认 10，最大 20）"},
                    "since_hours": {"type": "integer", "description": "可选：只返回最近 N 小时的新条目（默认 24，0=全部）"},
                },
                "required": ["action"],
            },
        },
    },
    # ===== 二维码：生成 / 识别 =====
    {
        "type": "function",
        "function": {
            "name": "qrcode",
            "description": "二维码：generate 把文本/链接生成 PNG 二维码；read 识别本地图片中的二维码（可识别多个）。识别需 pyzbar，缺失时降级提示",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "generate / read"},
                    "text": {"type": "string", "description": "generate 必填：要编码的内容（链接/文本）"},
                    "output": {"type": "string", "description": "generate 必填：输出 PNG 路径"},
                    "image_path": {"type": "string", "description": "read 必填：待识别图片路径"},
                    "size": {"type": "integer", "description": "可选：生成边长像素（默认 300，64-1024）"},
                    "error_correction": {"type": "string", "description": "可选：纠错等级 L/M/Q/H（默认 M）"},
                },
                "required": ["action"],
            },
        },
    },
    # ===== 密钥保险箱（P2 信任基建） =====
    {
        "type": "function",
        "function": {
            "name": "secret_store",
            "description": "密钥保险箱：DPAPI 加密托管 API key/令牌等敏感值。action=set 保存（value 只写不显示）/ get 按名取用 / delete 删除 / list 仅列出名称",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "set / get / delete / list"},
                    "name": {"type": "string", "description": "密钥名称（如 openai_key）"},
                    "value": {"type": "string", "description": "set 必填：要托管的敏感值"},
                },
                "required": ["action"],
            },
        },
    },
    # ===== 嵌入式 KV 存储（轻量状态/缓存） =====
    {
        "type": "function",
        "function": {
            "name": "kv_store",
            "description": "嵌入式键值存储：set 写入（可选 TTL 过期）/ get 读取 / delete 删除 / keys 列出全部 / search 按键或值模糊检索。适合缓存、配置、轻量状态（Redis 的零部署替代）",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "set / get / delete / keys / search"},
                    "key": {"type": "string", "description": "set/get/delete 必填：键"},
                    "value": {"type": "string", "description": "set 必填：值（上限 1MB）"},
                    "ttl_seconds": {"type": "integer", "description": "可选：set 时有效秒数（0=长期）"},
                    "pattern": {"type": "string", "description": "search 必填：键或值的模糊检索关键词"},
                },
                "required": ["action"],
            },
        },
    },
    # ===== 音视频处理 =====
    {
        "type": "function",
        "function": {
            "name": "media_ffmpeg",
            "description": "音视频处理：info 读取时长/分辨率/码率/音频信息；thumbnail 指定时间点截图；transcode 转码（mp4/mp3 等）；extract_audio 提取音频。输入超 2GB 或耗时超 300 秒会拒绝",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "info / thumbnail / transcode / extract_audio"},
                    "input": {"type": "string", "description": "源文件绝对路径"},
                    "output": {"type": "string", "description": "thumbnail/transcode/extract_audio 必填：输出路径"},
                    "time": {"type": "string", "description": "可选：截图时间点，如 00:01:30（默认取开头 1 秒）"},
                    "width": {"type": "integer", "description": "可选：转码输出宽度（16-7680，保持宽高比）"},
                    "format": {"type": "string", "description": "可选：转码/提取输出格式：mp4/mp3/webm/mkv/avi/mov/ogg/flac/wav"},
                },
                "required": ["action", "input"],
            },
        },
    },
    # ===== WebDAV 云盘同步 =====
    {
        "type": "function",
        "function": {
            "name": "webdav",
            "description": "WebDAV 云盘同步（坚果云/Nextcloud/群晖等）：list 列目录 / upload 上传 / download 下载 / delete 删除。连接在 webdav_config.json 配置（密码可加密）",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "list / upload / download / delete"},
                    "remote_path": {"type": "string", "description": "远端路径，如 /Documents/report.pdf（默认 /）"},
                    "local_path": {"type": "string", "description": "upload/download 必填：本地文件绝对路径"},
                },
                "required": ["action"],
            },
        },
    },
    # ===== 公众号自动写作 =====
    {
        "type": "function",
        "function": {
            "name": "run_wechat_writer",
            "description": "运行公众号自动写作工具（WeChat Writer）：采集当日 AI 资讯（RSS+搜索+论坛）→ 选题（历史去重）→ LLM 写作（大纲/正文/润色）→ 质量门禁 → 存草稿箱（只产草稿不发布）。适合『写一篇今天的 AI 公众号文章』等请求；可 dry_run=true 只预览不落盘，topic= 指定主题，use_blocked=true 时被墙信源（linux.do 等）自动走代理通道",
            "parameters": {
                "type": "object",
                "properties": {
                    "dry_run": {"type": "boolean", "description": "可选：true 只预览不写草稿（默认 false）"},
                    "topic": {"type": "string", "description": "可选：指定主题，跳过自动选题"},
                    "use_blocked": {"type": "boolean", "description": "可选：true 时被墙信源自动经代理通道采集（需 fetch_blocked 能力就绪）"},
                },
                "required": [],
            },
        },
    },
    # ===== 每日简报（主动助手：采集当日资讯 → LLM 提炼 → 落盘） =====
    {
        "type": "function",
        "function": {
            "name": "daily_brief",
            "description": "生成每日简报：采集当日 AI/科技资讯（RSS+搜索）→ LLM 提炼要点与点评 → 保存到工作区 briefs/。适合『今天的资讯有什么』『生成今日简报』等请求；可配合 schedule_task 定时生成晨报",
            "parameters": {
                "type": "object",
                "properties": {
                    "topic": {"type": "string", "description": "可选：主题关键词（仅保留相关素材）"},
                    "max_items": {"type": "integer", "description": "可选：素材上限（默认 8，最大 15）"},
                },
                "required": [],
            },
        },
    },
    # ===== 插件工坊：AI 生成并安装插件（零代码能力扩展） =====
    {
        "type": "function",
        "function": {
            "name": "create_plugin",
            "description": "根据用户需求生成并安装鲸语插件：组合自定义工具/技能模板/自动化流程/场景配置，生成后立即生效。适合『添加一个XX工具』『创建一个XX流程』『帮我加个小红书文案技能』等需求",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "插件名称（简短，如 小红书文案助手）"},
                    "description": {"type": "string", "description": "可选：插件说明"},
                    "tools": {
                        "type": "array",
                        "description": "可选：自定义 HTTP 工具列表，每项 {name, endpoint, description, method, params}（params 为逗号分隔的参数名）",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string", "description": "工具名（英文标识）"},
                                "endpoint": {"type": "string", "description": "HTTP 地址（http/https）"},
                                "description": {"type": "string", "description": "工具描述（AI 何时调用）"},
                                "method": {"type": "string", "description": "可选：POST/GET（默认 POST）"},
                                "params": {"type": "string", "description": "可选：参数名，逗号分隔，如 topic, style"},
                            },
                        },
                    },
                    "skills": {
                        "type": "array",
                        "description": "可选：技能/提示词模板，每项 {name, text}（text 中 {{TEXT}} 会被输入框内容替换）",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string", "description": "技能名"},
                                "text": {"type": "string", "description": "提示词模板内容"},
                            },
                        },
                    },
                    "workflows": {
                        "type": "object",
                        "description": "可选：自动化流程 {流程名: {steps: [{text: 指令}]}}",
                        "additionalProperties": {"type": "object"},
                    },
                    "scenario": {
                        "type": "object",
                        "description": "可选：一键场景配置 {name, thinking, system_prompt, enabled_tools}",
                    },
                    "requires": {
                        "type": "array",
                        "description": "可选：依赖的 pip 包名列表",
                        "items": {"type": "string"},
                    },
                },
                "required": ["name"],
            },
        },
    },
]

RUN_PY_TIMEOUT = 10
RUN_PY_MAX_CHARS = 8000
RUN_PY_MAX_OUTPUT = 20000
READ_FILE_MAX_BYTES = 102400
_READ_LINE_MAX = 102400  # 按行读取的每行上限（防单行数百 MB 撑爆内存）
FETCH_URL_MAX_CHARS = 500000
WEATHER_TIMEOUT = 5
EDIT_FILE_MAX_SIZE = 20 * 1024 * 1024  # edit_file 全量读入上限（20MB）
EDIT_FILE_REGEX_MAX = 1000  # 正则长度上限（防灾难性回溯挂死工具线程的粗略防线）
EXTRACT_MAX_ENTRIES = 10000  # 解压条目数上限（防 zip 海量小文件 DoS）
EXTRACT_MAX_TOTAL_BYTES = 2 * 1024 * 1024 * 1024  # 解压总字节上限（防磁盘写满）
EXTRACT_MAX_SINGLE_BYTES = 2 * 1024 * 1024 * 1024  # 单文件解压大小上限

# ===== run_python 执行模式 =====
# 完全体模式：沙箱降为用户权限级（默认 -S 不加载第三方库），
# 由用户的授权决策与权限模型（阻止目录/审批）兜底。

# run_python 静态危险拦截（仅非完全智能模式生效；完全智能 = 用户显式授权任意代码）。
# 词边界正则避免 1.10.0 版"误拦合法代码"的教训（\|a\|x 分支误伤字母）。
_RUN_PY_FORBIDDEN = (
    (re.compile(r"\bos\.(?:system|popen|spawn\w*|kill|remove|unlink|rmdir|removedirs)\b"),
     "os 系统调用/文件删除"),
    (re.compile(r"\bsubprocess\b"), "subprocess"),
    (re.compile(r"\bshutil\.(?:rmtree|rmdir|move|copy\w*|make_archive)\b"), "shutil 文件操作"),
    (re.compile(r"\b(?:eval|exec)\s*\("), "eval/exec"),
    (re.compile(r"\b__import__\s*\("), "__import__"),
    (re.compile(r"\bctypes\b"), "ctypes"),
    (re.compile(r"\bsocket\b"), "socket 网络"),
)

# ast 深度检查（修复正则可绕过的攻击面）：
# - from os import system / importlib.import_module('subprocess') 等动态/别名导入
# - os["system"] / getattr(os, "system") 索引与反射式调用
# - open(..., "w"/"a"/"x"/"+") 写模式（-I -S 沙箱内写文件绕过权限模型）
_RUN_PY_FORBIDDEN_MODULES = {
    "subprocess", "ctypes", "socket", "pty", "winreg", "win32api",
    "win32process", "win32pipe", "msvcrt",
}
_RUN_PY_FROM_FORBIDDEN = {
    "system", "popen", "popen2", "spawn", "spawnl", "spawnv", "spawnle",
    "remove", "unlink", "rmdir", "removedirs", "kill", "killpg",
    "exec", "execv", "execl", "startfile", "pthread_kill",
}
_RUN_PY_DANGEROUS_ATTRS = {
    "system", "popen", "spawn", "spawnl", "spawnv", "spawnle",
    "kill", "killpg", "remove", "unlink", "rmdir", "removedirs",
    "execv", "execl", "execve", "startfile", "startfilepath",
}


def _call_open_mode(node):
    """提取 open() 调用中的 mode 参数（仅字符串字面量；变量无法静态判断则放行）。"""
    args = node.args
    kw = {}
    for k in node.keywords or ():
        if k.arg:
            kw[k.arg] = k.value
    mode_node = kw.get("mode")
    if mode_node is None and len(args) >= 2:
        mode_node = args[1]
    if isinstance(mode_node, ast.Constant) and isinstance(mode_node.value, str):
        return mode_node.value
    return None


def _run_python_ast_blocked(code):
    """ast 深度静态检查：拦截别名导入 / 反射式调用 / 写模式 open。

    返回拦截原因字符串，通过则返回 ""。语法错误不拦截（由子进程报告，
    避免把"用户代码有语法错误"误报为安全拦截）。
    """
    try:
        tree = ast.parse(code)
    except SyntaxError:
        return ""
    # 第一遍：收集 os / importlib 的别名，并拦截星号导入与危险模块
    os_aliases = set()
    importlib_aliases = set()
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                root = (alias.name or "").split(".")[0]
                if root in _RUN_PY_FORBIDDEN_MODULES or root == "builtins":
                    return f"静态拦截：禁止导入模块 {alias.name}（完全智能模式可放行）"
                if root == "os":
                    os_aliases.add(alias.asname or alias.name)
                if root == "importlib":
                    importlib_aliases.add(alias.asname or alias.name)
        elif isinstance(node, ast.ImportFrom):
            mod = (node.module or "").split(".")[0]
            if mod in _RUN_PY_FORBIDDEN_MODULES or mod == "builtins":
                return f"静态拦截：禁止导入模块 {node.module}（完全智能模式可放行）"
            if mod == "os":
                for a in node.names or ():
                    if a.name == "*":
                        return "静态拦截：禁止 from os import *（完全智能模式可放行）"
                    if a.name in _RUN_PY_FROM_FORBIDDEN:
                        return f"静态拦截：禁止 from os import {a.name}（完全智能模式可放行）"
            if mod == "importlib":
                for a in node.names or ():
                    if a.name == "import_module":
                        return "静态拦截：禁止 importlib 动态导入（完全智能模式可放行）"
    # 第二遍：检查调用/反射/索引
    for node in ast.walk(tree):
        if isinstance(node, ast.Call):
            f = node.func
            if isinstance(f, ast.Name):
                if f.id in ("eval", "exec", "__import__", "compile"):
                    return f"静态拦截：禁止调用 {f.id}()（完全智能模式可放行）"
                if f.id == "open":
                    mode = _call_open_mode(node)
                    if mode and any(ch in mode for ch in "wax+"):
                        return "静态拦截：禁止以写模式打开文件（写文件请用 write_file 工具）"
                if f.id in ("getattr", "setattr", "delattr", "vars", "globals", "locals"):
                    return f"静态拦截：禁止反射/内省调用 {f.id}()（完全智能模式可放行）"
            elif isinstance(f, ast.Attribute):
                attr = f.attr
                base = f.value
                if isinstance(base, ast.Name) and base.id in os_aliases and attr in _RUN_PY_DANGEROUS_ATTRS:
                    return f"静态拦截：禁止调用 {base.id}.{attr}()（完全智能模式可放行）"
                if attr == "import_module" and isinstance(base, ast.Name) and base.id in importlib_aliases:
                    return "静态拦截：禁止 importlib 动态导入（完全智能模式可放行）"
                if attr in ("write_text", "write_bytes"):
                    # Path('f').write_text(...) / pathlib.Path('f').write_text(...)
                    p_func = base.func if isinstance(base, ast.Call) else None
                    p_name = (
                        p_func.id
                        if isinstance(p_func, ast.Name)
                        else (p_func.attr if isinstance(p_func, ast.Attribute) else "")
                    )
                    if p_name in ("Path", "PurePath"):
                        return "静态拦截：禁止 pathlib 写文件（写文件请用 write_file 工具）"
        elif isinstance(node, ast.Subscript):
            # os["system"] / 别名索引式调用绕过
            if (
                isinstance(node.value, ast.Name)
                and node.value.id in os_aliases
                and isinstance(node.slice, ast.Constant)
                and isinstance(node.slice.value, str)
                and node.slice.value in _RUN_PY_FROM_FORBIDDEN
            ):
                return f"静态拦截：禁止索引调用 {node.value.id}[{node.slice.value!r}]（完全智能模式可放行）"
    return ""


def _run_python_blocked(code):
    """静态危险检查（非完全智能模式下拦截高危操作，完全智能模式仅校验非空）。

    双层防线：① 正则快速匹配（保留历史拦截面）；② ast 深度检查（拦截
    from-import 别名、importlib 动态导入、getattr/下标反射、写模式 open 等
    正则无法覆盖的绕过手法）。完全智能模式 = 用户显式授权，仅校验非空。
    """
    if not code:
        return "代码为空"
    if permissions.is_full_auto():
        return ""
    for pat, label in _RUN_PY_FORBIDDEN:
        if pat.search(code):
            return f"静态拦截：代码包含 {label} 操作（完全智能模式可放行，或删除该语句后重试）"
    ast_err = _run_python_ast_blocked(code)
    if ast_err:
        return ast_err
    return ""

# 工具结果"失败"前缀统一判定（main/taskpanel 共享，防散落魔法字符串漂移）
TOOL_RESULT_FAIL_PREFIXES = ("错误", "权限拒绝", "超时", "（用户停止")

JSON_HINT_MESSAGE = (
    "[JSON 输出模式] 请严格输出合法的 JSON 对象（已启用 response_format），"
    "不要输出任何 JSON 以外的内容。"
)

def get_date():
    """获取当前日期、具体时间与本地时区。"""
    now = datetime.now().astimezone()
    tz_name = now.tzinfo.tzname(now) if now.tzinfo else "?"
    return f"{now:%Y-%m-%d %H:%M:%S} {tz_name}"


def get_weather(location, date):
    """查询城市天气。date（YYYY-mm-dd）传给 wttr.in（仅支持今天与近 3 天预报，
    更早/更晚的日期返回错误，避免把过期预报当真）。"""
    loc = str(location or "").strip()
    if not loc:
        return "错误：location 必填（城市名称）"
    d = str(date or "").strip()
    base = f"{loc} {d}" if d else loc
    url = f"https://wttr.in/{quote(loc)}?format=j1&lang=zh"
    if d:
        url += f"&date={quote(d)}"
    try:
        resp = _http_client().get(url, timeout=WEATHER_TIMEOUT)
        resp.raise_for_status()
        data = resp.json()
        cur = (data.get("current_condition") or [{}])[0]
        temp = cur.get("temp_C") or "?"
        desc = cur.get("lang_zh") or []
        if desc and desc[0].get("value"):
            text = desc[0]["value"]
        else:
            text = ((cur.get("weatherDesc") or [{}])[0].get("value") or "未知")
        return f"{base} 天气：{text}，气温 {temp}°C"
    except Exception as e:
        # 不返回编造的"模拟数据"——AI 会把假天气当真用于决策
        logging.warning("天气查询失败: %s", e)
        return f"错误：天气查询失败（{e}），请稍后重试或改用其他方式获取天气"


def run_python(code, with_site=False):
    if not code or len(code) > RUN_PY_MAX_CHARS:
        return f"错误：代码为空或超过 {RUN_PY_MAX_CHARS} 字符"
    block = _run_python_blocked(code)
    if block:
        permissions.audit("run_python_blocked", "static_check", block[:200], result="denied")
        return f"权限拒绝：{block}"
    try:
        argv = [sys.executable, "-I"]
        if not with_site:
            argv.append("-S")
        argv += ["-c", code]
        # SpooledTemporaryFile 重定向输出：进程刷屏打印时内存峰值限 1MB，
        # 超时 kill 后读取截断，不再全量 buffered 进内存（GB 级打印防 OOM）
        import tempfile

        with tempfile.SpooledTemporaryFile(
            max_size=1 << 20, mode="w+t", encoding="utf-8", errors="replace"
        ) as out:
            proc = subprocess.Popen(
                argv,
                stdout=out,
                stderr=subprocess.STDOUT,
                text=True,
                encoding="utf-8",
                errors="replace",
                cwd=permissions.WORKSPACE_DIR or None,
            )
            try:
                proc.wait(timeout=RUN_PY_TIMEOUT)
            except subprocess.TimeoutExpired:
                _kill_tree(proc)
                try:
                    proc.wait(timeout=3)
                except Exception:
                    pass
                return f"错误：执行超时（>{RUN_PY_TIMEOUT}秒）"
            out.seek(0)
            out_data = out.read(RUN_PY_MAX_OUTPUT)
            out.seek(0, os.SEEK_END)
            if out.tell() > RUN_PY_MAX_OUTPUT:
                out_data += "\n[输出已截断]"
        if not out_data.strip():
            return f"执行成功（无输出），工作目录：{permissions.WORKSPACE_DIR or '（当前目录）'}"
        permissions.audit("run_python", "python -I -S -c <code>", f"{len(code)} 字符, rc={proc.returncode}")
        return (
            out_data
            + f"\n[工作目录：{permissions.WORKSPACE_DIR or '（当前目录）'}，"
            + ("加载第三方库" if with_site else "未加载第三方库（-S 隔离）")
            + "]"
        )
    except Exception as e:
        return f"错误：{e}"


def read_file(path, start_line=None, max_lines=None):
    if not path or len(str(path)) > 512:
        return "错误：路径为空或过长"
    # 与 list_dir / write_file 等一致：所有路径操作先经权限模型判定
    # （默认仅允许工作区；读取工作区外文件请在「权限设置 → allowed_dirs」加入目录）
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    try:
        if start_line is not None or max_lines is not None:
            # 按行读取（适合超大文件）：start_line 从 1 开始，max_lines 默认 200
            try:
                start = max(1, int(start_line or 1))
                count = max(1, min(2000, int(max_lines or 200)))
            except (TypeError, ValueError):
                return "错误：start_line / max_lines 必须是正整数"
            if start > 1_000_000:
                return "错误：start_line 过大（超过 100 万行，请缩小范围）"
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                for _ in range(start - 1):
                    f.readline()
                # readline(上限)：单行可达数百 MB（minified JSON/日志），
                # 不设上限会一次撑爆内存；超长行截断到 100KB 并标注
                lines = []
                for _ in range(count):
                    ln = f.readline(_READ_LINE_MAX)
                    if ln == "":
                        break
                    if len(ln) >= _READ_LINE_MAX and not ln.endswith("\n"):
                        ln = ln.rstrip("\n") + "…[超长行已截断]\n"
                    lines.append(ln)
            lines = [ln for ln in lines if ln != ""]
            if not lines:
                return f"（第 {start} 行起无内容）"
            body = "".join(lines)
            if not body.endswith("\n"):
                body += "\n"
            prefix = f"[按行读取 {path} 第 {start}-{start + len(lines) - 1} 行]\n"
            return prefix + body
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read(READ_FILE_MAX_BYTES)
        if len(content) >= READ_FILE_MAX_BYTES:
            content += "\n[文件较大，已截断前 100KB]"
        return content
    except Exception as e:
        return f"错误：无法读取文件 {path}: {e}"


def _patch_array_items(tools):
    """递归补齐 array 参数的 items（API 要求 type=array 必须带 items，缺则 400）。"""
    def fix(node):
        if not isinstance(node, dict):
            return
        if node.get("type") == "array" and "items" not in node:
            node["items"] = {}
        for v in node.values():
            if isinstance(v, dict):
                fix(v)
            elif isinstance(v, list):
                for item in v:
                    fix(item)

    for tool in tools or []:
        fix((tool.get("function") or {}).get("parameters"))


def _run_fetch_blocked(url, proxy=None, **kwargs):
    """工具分发：fetch_blocked（按需能力，模块缺失时明确提示）。

    与其他工具一致使用具名参数签名（分发器以 fn(**args) 调用）——
    此前误写成单个 dict 参数导致 unexpected keyword argument 'url'。
    """
    if _fetch_blocked_impl is None:
        return "错误: fetch_blocked 能力未安装（需要将 fetch_blocked.py 放入程序目录并启用后可用）"
    if not str(url or "").startswith(("http://", "https://")):
        return "错误: URL 必须以 http:// 或 https:// 开头"
    # SSRF 校验由 fetch_blocked.py 内部实现执行（含内网/元数据拦截）；
    # 包装层只做协议与参数分发，避免在 DNS 被测试/网络环境临时改写时误伤。
    return _fetch_blocked_impl(url, proxy)


def _safe_request(method, url, *, allow_loopback=True, max_redirects=5,
                  validate=None, **kwargs):
    """发起 HTTP 请求并逐跳校验重定向（防 SSRF 重定向绕过）。

    validate 可传入自定义校验函数：返回空串表示允许，返回字符串表示拒绝原因。
    """
    current = url
    for _ in range(max_redirects + 1):
        err = validate(current) if validate is not None else _safe_url(
            current, allow_loopback=allow_loopback
        )
        if err:
            raise ValueError(err)
        resp = _http_client().request(
            method, current, follow_redirects=False, **kwargs
        )
        resp_headers = getattr(resp, "headers", {}) or {}
        location = resp_headers.get("location") if hasattr(resp_headers, "get") else None
        if getattr(resp, "is_redirect", False) is True and location:
            loc = str(location)
            next_url = _safe_redirect_url(current, loc, allow_loopback=allow_loopback)
            if next_url is None:
                resp.close()
                raise ValueError("重定向目标被 SSRF 防护拦截")
            resp.close()
            current = next_url
            continue
        return resp
    raise ValueError("重定向次数过多")


@contextlib.contextmanager
def _safe_stream(method, url, *, allow_loopback=True, max_redirects=5,
                 validate=None, **kwargs):
    """流式请求并逐跳校验重定向（防 SSRF 重定向绕过）。"""
    current = url
    for _ in range(max_redirects + 1):
        err = validate(current) if validate is not None else _safe_url(
            current, allow_loopback=allow_loopback
        )
        if err:
            raise ValueError(err)
        try:
            _stream_cm = _http_client().stream(
                method, current, follow_redirects=False, **kwargs
            )
        except TypeError:
            # 兼容旧测试/自定义 mock 的 stream 签名不接受 follow_redirects
            _stream_cm = _http_client().stream(method, current, **kwargs)
        except AttributeError:
            # 兼容旧测试/自定义 mock 没有 stream 方法：退化为普通请求（非流式）
            resp = _safe_request(
                method, current,
                allow_loopback=allow_loopback,
                max_redirects=max_redirects,
                validate=validate,
                **kwargs,
            )
            yield resp
            return
        with _stream_cm as resp:
            resp_headers = getattr(resp, "headers", {}) or {}
            location = resp_headers.get("location") if hasattr(resp_headers, "get") else None
            if getattr(resp, "is_redirect", False) is True and location:
                loc = str(location)
                next_url = _safe_redirect_url(current, loc, allow_loopback=allow_loopback)
                if next_url is None:
                    raise ValueError("重定向目标被 SSRF 防护拦截")
                current = next_url
                continue
            yield resp
            return
    raise ValueError("重定向次数过多")


def fetch_url(url):
    err = _safe_url(url)
    if err:
        return f"错误：{err}"
    try:
        # stream 边读边断：大响应不再全量下载进内存后才截断（防 GB 级内存峰值）
        with _safe_stream("GET", url, timeout=FETCH_URL_TIMEOUT) as resp:
            resp.raise_for_status()
            raw = b""
            truncated = False
            for chunk in resp.iter_bytes(64 * 1024):
                raw += chunk
                if len(raw) >= FETCH_URL_MAX_CHARS * 3:  # UTF-8 中文 1 字 3 字节
                    truncated = True
                    break  # 提前断开，连接随 with 释放
        # 编码自适应：charset 声明优先；无声明时尝试 utf-8 → gb18030（中文网页常见）→ 兜底
        charset = (resp.headers.get("content-type") or "").split("charset=")[-1].strip() or ""
        text = None
        for enc in (charset or "utf-8", "utf-8", "gb18030", "latin-1"):
            if not enc:
                continue
            try:
                text = raw.decode(enc, errors="strict")
                break
            except (LookupError, UnicodeDecodeError):
                continue
        if text is None:
            text = raw.decode("utf-8", errors="replace")
        if truncated:
            text = text[:FETCH_URL_MAX_CHARS] + "\n[内容较大，已截断前 500KB]"
        return text
    except Exception as e:
        return f"错误：{e}"


# ===== 二进制下载（P2）：图片/附件/安装包等任意文件 =====
DOWNLOAD_MAX_BYTES = 200 * 1024 * 1024  # 单文件 200MB 上限（与 WebDAV 对齐）


def download_file(url, local_path=""):
    """下载二进制文件（图片/附件/文档/安装包等）到工作区或指定目录。

    黑名单模式默认放行任意 URL（network.blocklist 除外）；流式写盘，超过
    200MB 自动中止并删除半成品。
    """
    if not str(url or "").startswith(("http://", "https://")):
        return "错误：url 必须以 http:// 或 https:// 开头"
    err = _safe_url(url)
    if err:
        return f"错误：{err}"
    try:
        from urllib.parse import urlparse, unquote
        fn = os.path.basename(unquote(urlparse(str(url)).path)) or f"download_{datetime.now():%Y%m%d_%H%M%S}.bin"
        fn = re.sub(r"[\\/:*?\"<>|]", "_", fn)[:120]
    except Exception:
        fn = f"download_{datetime.now():%Y%m%d_%H%M%S}.bin"
    if str(local_path or "").strip():
        p = permissions.resolve(local_path)
    else:
        base = os.path.join(permissions.WORKSPACE_DIR or "", "downloads")
        p = permissions.resolve(os.path.join(base, fn))
    if not p:
        return "错误：本地路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    try:
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        total = 0
        too_large = False
        with _safe_stream("GET", url, timeout=60) as resp:
            resp.raise_for_status()
            with open(p, "wb") as f:
                for chunk in resp.iter_bytes(64 * 1024):
                    total += len(chunk)
                    if total > DOWNLOAD_MAX_BYTES:
                        too_large = True
                        break
                    f.write(chunk)
        if too_large:
            try:
                os.remove(p)
            except OSError:
                pass
            return f"错误：文件超过 {DOWNLOAD_MAX_BYTES // 1024 // 1024}MB 上限，已中止"
        permissions.audit("download_file", url, f"{p} {total} 字节")
        return f"已下载 {url} → {p}（{total} 字节）"
    except Exception as e:
        try:
            if os.path.exists(p):
                os.remove(p)
        except OSError:
            pass
        return f"错误：下载失败: {e}"


# ===== 推送通知（A6）：钉钉 / ServerChan / Slack / 通用 Webhook =====
WEBHOOK_CONFIG_FILE = None  # 由 main 注入（DATA_DIR/webhooks.json）

CHART_THEME = "dark"  # 图表配色跟随主题，由 main 注入（"dark"/"light"）


def _load_webhooks():
    if not WEBHOOK_CONFIG_FILE or not os.path.exists(WEBHOOK_CONFIG_FILE):
        return {}
    try:
        with open(WEBHOOK_CONFIG_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict):
            return {k: _decrypt_secret(v) for k, v in data.items()}
        return {}
    except Exception:
        logging.exception("读取 webhook 配置失败")
        return {}


def send_webhook_notify(text, title="鲸语提醒", channel=""):
    """发送 Webhook 推送（钉钉/ServerChan/Slack/通用）。

    webhooks.json 格式：{"dingtalk": "https://oapi.dingtalk.com/robot/send?access_token=xxx",
    "serverchan": "https://sctapi.ftqq.com/KEY.send", "slack": "https://hooks.slack.com/services/...",
    "generic": "https://example.com/hook"}
    """
    if not str(text or "").strip():
        return "错误：text 必填"
    cfgs = _load_webhooks()
    if not cfgs:
        return "错误：未配置 Webhook（数据目录 webhooks.json 为空）"
    if channel:
        candidates = {str(channel).strip().lower(): cfgs.get(channel)} if cfgs.get(channel) else {}
    else:
        candidates = cfgs
    sent = []
    for name, url in candidates.items():
        if not url:
            continue
        try:
            ch = str(name).lower()
            if ch == "serverchan":
                # ServerChan 期望 application/x-www-form-urlencoded（title/desp）
                resp = _http_client().post(
                    str(url),
                    data=_webhook_payload(ch, str(title), str(text)),
                    timeout=10,
                )
            else:
                resp = _http_client().post(
                    str(url),
                    json=_webhook_payload(ch, str(title), str(text)),
                    timeout=10,
                )
            ok = resp.status_code < 400
            sent.append(f"{name}:{'✅' if ok else '❌' + str(resp.status_code)}")
        except Exception as e:
            sent.append(f"{name}:❌ {e}")
    return "；".join(sent) if sent else "错误：没有可用的 Webhook 通道"


def _webhook_payload(channel, title, text):
    if channel == "dingtalk":
        return {"msgtype": "text", "text": {"content": f"{title}\n{text}"}}
    if channel == "slack":
        return {"text": f"*{title}*\n{text}"}
    if channel == "generic":
        return {"title": title, "text": text}
    return {"title": title, "desp": text}  # serverchan 及其他


def send_webhook(title="", text="", channel=""):
    """主动推送通知到配置的 Webhook（钉钉/ServerChan/Slack/通用）。"""
    return send_webhook_notify(str(text or "").strip() or "（无内容）", str(title or "鲸语提醒"), channel)


# ===== IM 主动触达（P1）：Telegram Bot / 企业微信群机器人 =====
IM_CONFIG_FILE = None  # 由 main 注入（DATA_DIR/im_config.json）
AGENT_MAIL_ENABLED = False  # 由 main 按 config.agent_mail_enabled 注入（默认关闭，不配置不启用）
AGENT_MAIL_CLI = "agently-cli"
_TELEGRAM_OFFSET = 0   # Telegram getUpdates 游标（进程内去重）


def _load_im_config():
    """读取 im_config.json（敏感字段支持 dpapi: 密文）。返回 (dict, error)。"""
    if not IM_CONFIG_FILE or not os.path.exists(IM_CONFIG_FILE):
        return {}, (
            "IM 通道未配置（如不需要推送可忽略；如需开启：系统菜单 → IM 通道配置）。"
            "支持 wecom_webhook / wecom_aibot_bot_id+secret / telegram_bot_token+chat_id"
        )
    try:
        with open(IM_CONFIG_FILE, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        if not isinstance(cfg, dict):
            return {}, "错误：im_config.json 格式不是对象"
        return {k: _decrypt_secret(v) for k, v in cfg.items()}, ""
    except Exception:
        logging.exception("读取 IM 配置失败")
        return {}, "错误：读取 IM 配置失败"


def im_send(text, title="", channel=""):
    """主动触达：发送消息到 Telegram / 企业微信群机器人（可同时推送多通道）。"""
    if not str(text or "").strip():
        return "错误：text 必填"
    cfg, err = _load_im_config()
    if not cfg:
        return err
    title = str(title or "鲸语提醒").strip()
    body = f"{title}\n{text}" if title else str(text)
    ch = str(channel or "").strip().lower()
    targets = {}
    if not ch or ch == "telegram":
        if cfg.get("telegram_bot_token") and cfg.get("telegram_chat_id"):
            targets["telegram"] = (cfg["telegram_bot_token"], str(cfg["telegram_chat_id"]))
    if not ch or ch in ("wecom", "wechat", "weixin"):
        if cfg.get("wecom_webhook"):
            targets["wecom"] = (cfg["wecom_webhook"],)
    if not targets:
        return "错误：未配置可用的 IM 通道（telegram_bot_token/telegram_chat_id 或 wecom_webhook）"
    sent = []
    for name, val in targets.items():
        try:
            if name == "telegram":
                token, chat_id = val
                url = f"https://api.telegram.org/bot{token}/sendMessage"
                resp = _http_client().post(url, json={"chat_id": chat_id, "text": body[:4000]}, timeout=15)
            else:
                resp = _http_client().post(val[0], json={"msgtype": "text", "text": {"content": body[:4000]}}, timeout=15)
            ok = resp.status_code < 400
            sent.append(f"{name}:{'✅' if ok else '❌' + str(resp.status_code)}")
        except Exception as e:
            sent.append(f"{name}:❌ {e}")
    permissions.audit("im_send", ",".join(sent), body[:80], result="ok")
    return "；".join(sent)


def telegram_poll_updates(timeout=15, limit=5):
    """接收 Telegram 消息（供 AI 定期检查或用户召唤）。返回最近消息；游标自动前移去重。"""
    global _TELEGRAM_OFFSET
    cfg, err = _load_im_config()
    if not cfg:
        return err
    token = cfg.get("telegram_bot_token")
    chat_id = str(cfg.get("telegram_chat_id") or "").strip()
    if not token:
        return "未配置 telegram_bot_token（系统菜单 → IM 通道配置 可开启）"
    try:
        timeout = max(1, min(60, int(timeout or 15)))
        limit = max(1, min(20, int(limit or 5)))
    except (TypeError, ValueError):
        timeout, limit = 15, 5
    try:
        url = f"https://api.telegram.org/bot{token}/getUpdates"
        params = {"timeout": timeout, "limit": limit}
        if _TELEGRAM_OFFSET:
            params["offset"] = _TELEGRAM_OFFSET + 1
        resp = _http_client().post(url, json=params, timeout=timeout + 15)
        resp.raise_for_status()
        data = resp.json()
    except Exception as e:
        return f"错误：Telegram 接收失败: {e}"
    updates = data.get("result") or []
    if not updates:
        return "（暂无新消息）"
    lines = []
    for u in updates:
        msg = u.get("message") or {}
        if chat_id and str(msg.get("chat", {}).get("id")) != chat_id:
            continue
        sender = (msg.get("from") or {}).get("username") or (msg.get("from") or {}).get("first_name") or "?"
        text = str(msg.get("text") or "[非文本消息]")[:500]
        lines.append(f"@{sender}: {text}")
        if int(u.get("update_id") or 0) > _TELEGRAM_OFFSET:
            _TELEGRAM_OFFSET = int(u["update_id"])
    if not lines:
        return "（暂无来自配置 chat_id 的新消息）"
    return "\n".join(lines)


# ===== 多模态（A3）：语音合成 / 图像处理 / 文件 OCR =====
def tts_save(text, path, rate=0):
    """语音合成保存为 WAV 文件（Windows SAPI，可选 pywin32；无则用 PowerShell）。"""
    if not text or not str(text).strip():
        return "错误：text 必填"
    if not path or not str(path).strip():
        return "错误：path 必填"
    p = permissions.resolve(path)
    if not p:
        return "错误：路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    try:
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
    except Exception:
        pass
    try:
        import pythoncom
        import win32com.client

        synth = str(text)[:8000]  # 文本上限：超长合成会长时间占住共享工具线程池
        result = {"err": None}

        def _speak():
            pythoncom.CoInitialize()
            stream = None
            try:
                speaker = win32com.client.Dispatch("SAPI.SpVoice")
                stream = win32com.client.Dispatch("SAPI.SpFileStream")
                stream.Open(p, 3)  # SSFMCreateForWrite
                speaker.AudioOutputStream = stream
                try:
                    speaker.Rate = max(-10, min(10, int(rate or 0)))
                except (TypeError, ValueError):
                    pass
                # Speak 同步阻塞且无法安全强杀：放后台线程执行，主路径只等 60s
                speaker.Speak(synth)
            except Exception as e:
                result["err"] = e
            finally:
                # COM 资源成对释放：Speak 抛异常也要关流 + CoUninitialize（防单元泄漏）
                if stream is not None:
                    try:
                        stream.Close()
                    except Exception:
                        pass
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

        t = threading.Thread(target=_speak, daemon=True)
        t.start()
        t.join(timeout=60.0)
        if t.is_alive():
            return (
                f"语音合成进行中（文本较长，后台继续生成），稍后可在 {p} 查看。"
                "需要控制时长请缩短文本。"
            )
        if result["err"]:
            raise result["err"]
        try:
            size = os.path.getsize(p) if os.path.exists(p) else 0
        except OSError:
            size = 0
        if size < 100:
            return (
                f"已生成语音文件 {p}（{size} 字节）但内容可能为空："
                "系统未安装中文语音包（设置 → 时间和语言 → 语音）时 SAPI 无可用音色"
            )
        return f"已合成语音保存至 {p}"
    except ImportError:
        return "错误：需要 pywin32（pip install pywin32）"
    except Exception as e:
        return f"错误：语音合成失败: {e}"


_CLIENT_HOLDER = {"client": None}  # main 在 ensure_client 时注入


def set_active_client(client):
    _CLIENT_HOLDER["client"] = client


def subagent_run(tasks, parallel=2, context=""):
    """并行子代理：把大任务拆给多个并发 LLM 子代理，各自输出结论后汇总。

    tasks：任务数组（字符串列表，最多 8 个）；parallel：并行数 1-4。
    context：可选背景上下文（注入每个子代理）。
    """
    if not isinstance(tasks, list) or not tasks:
        return "错误：tasks 必须是非空数组（每个元素是一个子任务目标）"
    tasks = [str(t) for t in tasks][:8]
    try:
        parallel = max(1, min(4, int(parallel or 2)))
    except (TypeError, ValueError):
        parallel = 2
    client = _CLIENT_HOLDER.get("client")
    if client is None:
        return "错误：没有可用客户端（请先完成一次对话建立连接）"
    base = "你是并行子代理，专注完成分配的子任务，输出简洁、可执行的结论（不要提及子代理身份）。"
    if str(context or "").strip():
        base += f"\n\n【共享背景上下文】\n{context}"
    results = [None] * len(tasks)

    def run(i, task):
        last_err = None
        for attempt in range(2):
            try:
                resp = client.client.chat.completions.create(
                    model=client.model,
                    messages=[
                        {"role": "system", "content": base},
                        {"role": "user", "content": str(task)},
                    ],
                    max_tokens=2048,
                    stream=False,
                    timeout=120.0,
                    extra_body={"thinking": {"type": "disabled"}},
                )
                results[i] = (resp.choices[0].message.content or "").strip() or "（子代理无输出）"
                return
            except Exception as e:
                last_err = e
                if attempt == 0:
                    time.sleep(1)
        results[i] = f"（子任务失败：{last_err}）"

    import concurrent.futures as cf

    with cf.ThreadPoolExecutor(max_workers=parallel) as ex:
        futures = [ex.submit(run, i, t) for i, t in enumerate(tasks)]
        for _f in cf.as_completed(futures):
            pass
    lines = []
    for i, t in enumerate(tasks):
        lines.append(f"## 子任务 {i + 1}：{t[:80]}\n{results[i]}")
    return "\n\n".join(lines)


# ===== 自我验证闭环（A8）：跑测试 / 对照标准答案自评 =====
def run_tests(path=None, framework="auto"):
    """在允许目录内运行测试（pytest/unittest），返回结果摘要。

    path：测试文件或目录（留空则扫描允许目录内 *_test.py / test_*.py）。
    """
    target = None
    if str(path or "").strip():
        target = permissions.resolve(path)
        if not target or not os.path.exists(target):
            return f"错误：路径不存在：{path}"
        ok, reason = permissions.check_filesystem(target, write=False)
        if not ok:
            return reason
    import glob as _glob

    if target is None:
        base = permissions.WORKSPACE_DIR
        if not base:
            return "错误：未配置工作目录"
        found = _glob.glob(os.path.join(base, "**", "test_*.py"), recursive=True)[:20] + \
                _glob.glob(os.path.join(base, "**", "*_test.py"), recursive=True)[:20]
        if not found:
            return "错误：允许目录内未找到测试文件（test_*.py / *_test.py）"
        target = found[0]
    fw = str(framework or "auto").lower()
    if fw == "unittest":
        cmd = [sys.executable, "-m", "unittest", "discover", "-v"]
        if target and os.path.isfile(target):
            cmd = [sys.executable, "-m", "unittest", "-v", str(target)]
    elif fw == "pytest":
        cmd = [sys.executable, "-m", "pytest", "-q"]
        if target:
            cmd.append(str(target))
    else:
        cmd = [sys.executable, "-m", "pytest", "-q", str(target)]
        if not os.path.isfile(target):
            cmd = [sys.executable, "-m", "unittest", "discover", "-v", str(target)]
    try:
        # SpooledTemporaryFile 限流：pytest -v / unittest 输出可达 MB 级，
        # capture_output 全量进内存会 OOM；内存峰值限 1MB 后自动转磁盘
        import tempfile

        with tempfile.SpooledTemporaryFile(
            max_size=1 << 20, mode="w+t", encoding="utf-8", errors="replace"
        ) as out:
            proc = subprocess.Popen(
                cmd, stdout=out, stderr=subprocess.STDOUT, text=True,
                encoding="utf-8", errors="replace",
                cwd=os.path.dirname(target) if os.path.isfile(target) else target,
            )
            try:
                proc.wait(timeout=180)
            except subprocess.TimeoutExpired:
                _kill_tree(proc)
                try:
                    proc.wait(timeout=3)
                except Exception:
                    pass
                return "错误：测试执行超时（180 秒）"
            out.seek(0)
            out_data = out.read(12000)
            out.seek(0, os.SEEK_END)
            if out.tell() > 12000:
                out_data += "\n[输出已截断]"
        return f"退出码 {proc.returncode}\n{out_data}"
    except subprocess.TimeoutExpired:
        return "错误：测试超时（>180 秒）"
    except Exception as e:
        return f"错误：运行测试失败: {e}"


def verify_output(expected, actual):
    """对照标准答案自评：计算语义相似度并指出差异要点（自我验证闭环）。"""
    e = str(expected or "")
    a = str(actual or "")
    if not e.strip():
        return "错误：expected 必填"
    if not a.strip():
        return "评估：实际输出为空（0% 匹配）"
    et, at = set(_mem_tokens(e)), set(_mem_tokens(a))
    if not et or not at:
        return "评估：无法分词比较（内容过短）"
    inter = et & at
    recall = len(inter) / len(et)          # 预期要点覆盖率
    precision = len(inter) / len(at)       # 输出聚焦度
    f1 = 2 * recall * precision / (recall + precision) if (recall + precision) else 0.0
    missing = [t for t in et if t not in at][:10]
    verdict = "通过" if f1 >= 0.7 else ("基本通过" if f1 >= 0.5 else "未通过")
    lines = [
        f"评估：{verdict}（F1={f1:.2f}，覆盖率 {recall:.0%}，聚焦度 {precision:.0%}）",
    ]
    if missing:
        lines.append("缺失要点：" + "、".join(missing))
    if len(e) > 0 and len(a) > 0 and a.strip().startswith(TOOL_RESULT_FAIL_PREFIXES):
        lines.append("提示：实际输出以错误开头，请检查执行是否成功")
    return "\n".join(lines)



def image_process(path, output, ops=""):
    """PIL 图像处理：resize=宽x高; crop=x1,y1,x2,y2; rotate=度数;
    convert=PNG/JPEG; quality=1-100; water=水印文本（右下角）。"""
    if not path or not output:
        return "错误：path 与 output 必填"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：源图片不存在：{p}"
    out = permissions.resolve(output)
    if not out:
        return "错误：输出路径无效"
    ok, reason = permissions.check_filesystem(out, write=True)
    if not ok:
        return reason
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return "错误：需要 Pillow（pip install pillow）"
    try:
        img = Image.open(p)
        # 防 decompression bomb：先检查像素尺寸，拒绝超大图再进入解码/处理
        try:
            if img.width * img.height > 100_000_000:
                return "错误：图片像素过大（超过 1 亿像素），请先压缩后再处理"
        except Exception:
            pass
        applied = 0
        quality = None
        for op in str(ops or "").split(";"):
            op = op.strip()
            if not op:
                continue
            key, _, val = op.partition("=")
            key = key.strip().lower()
            val = val.strip()
            # 每个操作独立容错：单个操作参数错误明确报错，不吞掉整个处理
            try:
                if key == "resize" and val:
                    w, _, h = val.lower().partition("x")
                    if not (w.isdigit() and h.isdigit()):
                        return f"错误：resize 格式应为 宽x高（如 800x600），收到：{op}"
                    img = img.resize((max(1, int(w)), max(1, int(h))))
                elif key == "crop" and val:
                    parts = [v.strip() for v in val.split(",")]
                    if len(parts) != 4:
                        return f"错误：crop 需要 4 个坐标 x1,y1,x2,y2（如 0,0,100,100），收到：{op}"
                    x1, y1, x2, y2 = (int(v) for v in parts)
                    img = img.crop((x1, y1, x2, y2))
                elif key == "rotate" and val:
                    try:
                        deg = float(val)
                    except ValueError:
                        return f"错误：rotate 需要数字角度（如 90），收到：{op}"
                    img = img.rotate(deg, expand=True)
                elif key == "convert" and val:
                    # PIL convert 需要模式名（RGB/RGBA/L），"JPEG" 这类格式名需映射
                    mode_map = {
                        "JPEG": "RGB", "JPG": "RGB", "PNG": "RGBA",
                        "GRAY": "L", "GREY": "L", "BMP": "RGB", "WEBP": "RGB",
                    }
                    target = mode_map.get(val.upper(), val.upper())
                    img = img.convert(target)
                elif key == "water" and val:
                    draw = ImageDraw.Draw(img)
                    try:
                        font = ImageFont.truetype(
                            "C:/Windows/Fonts/msyh.ttc", max(12, img.width // 20)
                        )
                    except Exception:
                        font = ImageFont.load_default()
                    w, h = img.size
                    tw, th = draw.textbbox((0, 0), val, font=font)[2:4]
                    draw.text((w - tw - 10, h - th - 10), val, fill=(255, 255, 255, 200), font=font)
                elif key == "quality" and val:
                    try:
                        quality = max(1, min(100, int(val)))
                    except ValueError:
                        return f"错误：quality 应为 1-100 的数字，收到：{op}"
                    applied += 1  # quality 不算图像变换，但记录已生效
                    continue
                else:
                    continue  # 未知操作静默跳过（保持向后兼容）
                applied += 1
            except (ValueError, IndexError) as e:
                return f"错误：操作 {op} 参数非法：{e}"
        os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
        save_kw = {"quality": quality} if quality else {}
        img.save(out, **save_kw)
        size = os.path.getsize(out) if os.path.exists(out) else 0
        if applied == 0:
            return f"已复制图像至 {out}（{size} 字节）未做处理（ops 为空或未识别）。支持：resize/crop/rotate/convert/quality/water"
        return f"已处理图像并保存至 {out}（{size} 字节，{applied} 项操作生效）"
    except Exception as e:
        return f"错误：图像处理失败: {e}"


def ocr_image(path):
    """从图片文件提取文字（Windows OCR，需系统语言包支持）。"""
    if not path or not str(path).strip():
        return "错误：path 必填"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：图片不存在：{p}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    try:
        import tempfile

        fd, ps_path = tempfile.mkstemp(suffix=".ps1")
        os.close(fd)
        try:
            script = OCR_IMAGE_PS.replace("@PATH@", "'" + str(p).replace("'", "''") + "'")
            with open(ps_path, "w", encoding="utf-8-sig") as f:
                f.write(script)
            proc = subprocess.run(
                ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-File", ps_path],
                capture_output=True, text=True, timeout=60,
                encoding="utf-8", errors="replace",
            )
            out = (proc.stdout or "").strip()
            return out or "未能识别出文字"
        finally:
            try:
                os.remove(ps_path)
            except OSError:
                pass
    except Exception as e:
        return f"错误：OCR 失败: {e}"


# ===== 数据工具（A4）：CSV / Excel / 图表 / MySQL / PostgreSQL =====
DB_CONFIG_FILE = None  # 由 main 注入（DATA_DIR/db_config.json）


def _db_conn(kind, name):
    """从 db_config.json 读取连接配置。格式：{"mysql": {"default": {host,port,user,password,database}}}。"""
    if not DB_CONFIG_FILE:
        return None, "错误：数据库配置未初始化"
    try:
        if not os.path.exists(DB_CONFIG_FILE):
            return None, "错误：未找到数据库配置文件 db_config.json（需先在数据目录配置）"
        with open(DB_CONFIG_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        conns = data.get(kind) or {}
        cfg = conns.get(str(name or "default"))
        if not cfg or not isinstance(cfg, dict):
            return None, f"错误：未找到 {kind} 连接「{name}」（可用：{list(conns) or '无'}）"
        cfg = dict(cfg)
        if cfg.get("password"):
            cfg["password"] = _decrypt_secret(cfg["password"])
        return cfg, ""
    except Exception as e:
        return None, f"错误：读取数据库配置失败: {e}"


def database_query_mysql(connection="default", sql="", max_rows=20):
    """MySQL 只读查询（SELECT/SHOW/DESC；连接在 db_config.json 配置）。"""
    try:
        import pymysql
    except ImportError:
        return "错误：需要 pymysql（pip install pymysql）"
    if not str(sql or "").strip():
        return "错误：sql 必填"
    cfg, err = _db_conn("mysql", connection)
    if cfg is None:
        return err
    if not _readonly_stmt(sql):
        return "错误：仅允许只读查询（SELECT / SHOW / DESC）"
    try:
        conn = pymysql.connect(
            host=str(cfg.get("host") or "127.0.0.1"),
            port=int(cfg.get("port") or 3306),
            user=str(cfg.get("user") or ""),
            password=str(cfg.get("password") or ""),
            database=str(cfg.get("database") or ""),
            charset="utf8mb4", connect_timeout=5, read_timeout=15,
        )
        try:
            cur = conn.cursor()
            try:
                cur.execute("SET SESSION max_execution_time=15000")
            except Exception:
                pass
            cur.execute(str(sql))
            cols = [d[0] for d in cur.description] if cur.description else []
            rows = cur.fetchmany(max(1, min(200, int(max_rows or 20))))
            lines = [" | ".join(cols)] if cols else []
            for r in rows:
                cells = [str(x) if x is not None else "" for x in r]
                cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
                lines.append(" | ".join(cells))
            extra = "" if len(rows) < max(1, min(200, int(max_rows or 20))) else " [已截断]"
            return "\n".join(lines) + extra if lines else "执行成功（无结果集）"
        finally:
            try:
                conn.close()
            except Exception:
                pass
    except Exception as e:
        return f"错误：MySQL 查询失败: {e}"


def database_query_postgres(connection="default", sql="", max_rows=20):
    """PostgreSQL 只读查询（SELECT/SHOW/DESC；连接在 db_config.json 配置）。"""
    try:
        import psycopg2
    except ImportError:
        return "错误：需要 psycopg2（pip install psycopg2）"
    if not str(sql or "").strip():
        return "错误：sql 必填"
    cfg, err = _db_conn("postgres", connection)
    if cfg is None:
        return err
    if not _readonly_stmt(sql):
        return "错误：仅允许只读查询（SELECT / SHOW / DESC）"
    try:
        conn = psycopg2.connect(
            host=str(cfg.get("host") or "127.0.0.1"),
            port=int(cfg.get("port") or 5432),
            user=str(cfg.get("user") or ""),
            password=str(cfg.get("password") or ""),
            dbname=str(cfg.get("database") or ""),
            connect_timeout=5,
            # 语句超时：只读查询最长 15 秒，防慢查询占住共享工具线程池
            options="-c statement_timeout=15000",
        )
        try:
            cur = conn.cursor()
            cur.execute(str(sql))
            cols = [d[0] for d in cur.description] if cur.description else []
            rows = cur.fetchmany(max(1, min(200, int(max_rows or 20))))
            lines = [" | ".join(cols)] if cols else []
            for r in rows:
                cells = [str(x) if x is not None else "" for x in r]
                cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
                lines.append(" | ".join(cells))
            extra = "" if len(rows) < max(1, min(200, int(max_rows or 20))) else " [已截断]"
            return "\n".join(lines) + extra if lines else "执行成功（无结果集）"
        finally:
            try:
                conn.close()
            except Exception:
                pass
    except Exception as e:
        return f"错误：PostgreSQL 查询失败: {e}"


def read_csv(path, max_rows=100, delimiter=","):
    """读取 CSV 文件（允许目录内），返回表格文本。"""
    if not path or not str(path).strip():
        return "错误：path 必填"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{p}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    import csv as _csv

    try:
        limit = max(1, min(500, int(max_rows or 100)))
    except (TypeError, ValueError):
        limit = 100
    try:
        delim = str(delimiter or ",")
        if len(delim) != 1:
            return "错误：delimiter 必须是单个字符（如 , ; | \\t）"
        if delim == "\\t":
            delim = "\t"
        with open(p, "r", encoding="utf-8-sig", errors="replace", newline="") as f:
            rows = list(itertools.islice(_csv.reader(f, delimiter=delim), limit))
        if not rows:
            return "（空文件）"
        # 列宽截断：超宽单元格（minified JSON/长 URL）会撑爆上下文，单格限 100 字符
        rows = [
            [str(c)[:_TABLE_CELL_MAX] + ("…" if len(str(c)) > _TABLE_CELL_MAX else "") for c in r]
            for r in rows
        ]
        widths = []
        for c in range(max(len(r) for r in rows)):
            widths.append(max(len(r[c]) if c < len(r) else 0 for r in rows))
        lines = []
        for r in rows:
            cells = [r[c].ljust(widths[c]) if c < len(r) else "" for c in range(len(widths))]
            lines.append(" | ".join(cells))
        total = "…" if len(rows) >= limit else ""
        return "\n".join(lines) + (f"\n[前 {limit} 行{total}]" if len(rows) >= limit else "")
    except Exception as e:
        return f"错误：读取 CSV 失败: {e}"


def write_csv(path, rows, headers=""):
    """写入 CSV 文件。rows 为 JSON 数组：[[v,v],...] 或 [{"col":v},...]；headers 逗号分隔。"""
    if not path or not str(path).strip():
        return "错误：path 必填"
    if not isinstance(rows, list):
        return "错误：rows 必须是非空数组"
    p = permissions.resolve(path)
    if not p:
        return "错误：路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    import csv as _csv

    try:
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        if rows and isinstance(rows[0], dict):
            cols = [h.strip() for h in str(headers or "").split(",") if h.strip()] or list(rows[0].keys())
            # 混合 dict/非 dict 行健壮化：非 dict 行按空字典处理（防中间夹杂标量崩溃）
            data = [
                [row.get(c, "") for c in cols] if isinstance(row, dict) else [""] * len(cols)
                for row in rows
            ]
        else:
            cols = [h.strip() for h in str(headers or "").split(",") if h.strip()]
            data = [
                [str(x) for x in row] if isinstance(row, (list, tuple)) else [str(row)]
                for row in rows
            ]
        with open(p, "w", encoding="utf-8-sig", newline="") as f:
            w = _csv.writer(f)
            if cols:
                w.writerow(cols)
            w.writerows(data)
        return f"已写入 CSV 至 {p}（{len(data)} 行）"
    except Exception as e:
        return f"错误：写入 CSV 失败: {e}"


def read_excel(path, sheet=0, max_rows=100):
    """读取 Excel 文件（openpyxl，.xlsx）。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "错误：需要 openpyxl（pip install openpyxl）"
    if not path or not str(path).strip():
        return "错误：path 必填"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{p}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    try:
        limit = max(1, min(500, int(max_rows or 100)))
    except (TypeError, ValueError):
        limit = 100
    try:
        wb = load_workbook(p, read_only=True, data_only=True)
        try:
            if isinstance(sheet, int):
                ws = wb.worksheets[min(sheet, len(wb.worksheets) - 1)] if wb.worksheets else wb.active
            else:
                ws = wb[sheet]
        except KeyError:
            return f"错误：工作表不存在：{sheet}"
        lines = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= limit:
                break
            cells = ["" if v is None else str(v) for v in row]
            # 单元格截断：超长文本（minified JSON 等）撑爆上下文
            cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
            lines.append(" | ".join(cells))
        if not lines:
            return "（空工作表）"
        return "\n".join(lines) + (f"\n[前 {limit} 行…]" if len(lines) >= limit else "")
    except Exception as e:
        return f"错误：读取 Excel 失败: {e}"


def _read_optional_text(path, max_chars):
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return None, f"错误：文件不存在：{path}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return None, reason
    try:
        max_chars = max(1000, min(100000, int(max_chars or 20000)))
    except (TypeError, ValueError):
        max_chars = 20000
    return p, max_chars


def _strip_html_tags(html_text):
    import html as _html
    txt = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", html_text or "")
    txt = re.sub(r"(?s)<[^>]+>", "\n", txt)
    txt = _html.unescape(txt)
    txt = re.sub(r"\n{3,}", "\n\n", txt)
    return txt.strip()


def epub_read(path, max_chars=20000):
    """读取 EPUB 电子书正文（ebooklib 可选依赖）。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        return "错误：需要 ebooklib（pip install ebooklib）"
    try:
        book = epub.read_epub(p)
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            body = item.get_body_content() or b""
            parts.append(_strip_html_tags(body.decode("utf-8", errors="replace")))
        text = "\n\n".join(x for x in parts if x)
        if not text:
            return "（EPUB 无正文内容）"
        if len(text) > limit:
            text = text[:limit] + f"\n[正文已截断前 {limit} 字符]"
        return text
    except Exception as e:
        return f"错误：读取 EPUB 失败: {e}"


def mobi_read(path, max_chars=20000):
    """读取 MOBI 电子书正文（mobi 可选依赖）。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    try:
        from mobi import Mobi
    except ImportError:
        return "错误：需要 mobi（pip install mobi）"
    try:
        book = Mobi(p)
        book.parse()
        text = str(book) if hasattr(book, "__str__") else ""
        if not text:
            text = "\n\n".join(str(getattr(book, field, "")) for field in ("title", "author", "publisher", "description"))
        if len(text) > limit:
            text = text[:limit] + f"\n[正文已截断前 {limit} 字符]"
        return text or "（MOBI 解析无文本）"
    except Exception as e:
        return f"错误：读取 MOBI 失败: {e}"


def doc_read(path, max_chars=20000):
    """读取旧版 .doc 二进制文档（优先 antiword，其次 catdoc）。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    import shutil as _shutil
    for exe in ("antiword", "catdoc"):
        if _shutil.which(exe):
            try:
                proc = subprocess.run(
                    [exe, p], capture_output=True, text=True, timeout=30,
                    errors="replace", creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
                )
                text = (proc.stdout or "").strip() or (proc.stderr or "").strip()
                if text:
                    if len(text) > limit:
                        text = text[:limit] + f"\n[正文已截断前 {limit} 字符]"
                    return text
            except Exception as e:
                return f"错误：读取 .doc 失败（{exe}）: {e}"
    return "错误：读取 .doc 需要 antiword 或 catdoc（Windows 可安装 antiword 或改用 .docx）"


def msg_read(path, max_chars=20000):
    """读取 .msg Outlook 邮件（extract_msg 可选依赖），返回主题/发件人/正文/附件清单。"""
    p_or_err = _read_optional_text(path, max_chars)
    if p_or_err[0] is None:
        return p_or_err[1]
    p, limit = p_or_err
    try:
        import extract_msg
    except ImportError:
        return "错误：需要 extract_msg（pip install extract-msg）"
    try:
        msg = extract_msg.Message(p)
        attachments = [a.longFilename or a.shortFilename or "?" for a in (msg.attachments or [])]
        lines = [
            f"主题：{msg.subject}",
            f"发件人：{msg.sender}",
            f"收件人：{msg.to}",
            f"日期：{msg.date}",
        ]
        if attachments:
            lines.append(f"附件（{len(attachments)}）：{', '.join(attachments[:20])}")
        body = str(msg.body or "").strip()
        if body:
            lines.append("正文：")
            if len(body) > limit:
                body = body[:limit] + f"\n[正文已截断前 {limit} 字符]"
            lines.append(body)
        return "\n".join(lines)
    except Exception as e:
        return f"错误：读取 .msg 失败: {e}"


def archive_list(path):
    """列出压缩包内容：.zip / .tar / .gz / .7z / .rar（可选依赖）。"""
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：压缩包不存在：{path}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    ext = os.path.splitext(p)[1].lower()
    try:
        if ext in (".zip",):
            import zipfile
            with zipfile.ZipFile(p) as zf:
                infos = zf.infolist()
                lines = [f"压缩包 {path} 共 {len(infos)} 个条目："]
                for info in infos[:100]:
                    lines.append(f"· {info.filename}（{info.file_size} 字节）")
                if len(infos) > 100:
                    lines.append(f"… 其余 {len(infos) - 100} 个条目略")
                return "\n".join(lines)
        if ext in (".tar", ".gz", ".tgz"):
            import tarfile
            with tarfile.open(p) as tf:
                members = tf.getmembers()
                lines = [f"压缩包 {path} 共 {len(members)} 个条目："]
                for m in members[:100]:
                    lines.append(f"· {m.name}（{m.size} 字节）")
                if len(members) > 100:
                    lines.append(f"… 其余 {len(members) - 100} 个条目略")
                return "\n".join(lines)
        if ext == ".7z":
            import py7zr
            with py7zr.SevenZipFile(p, "r") as z:
                items = z.getnames()
                lines = [f"压缩包 {path} 共 {len(items)} 个条目："]
                for name in items[:100]:
                    lines.append(f"· {name}")
                if len(items) > 100:
                    lines.append(f"… 其余 {len(items) - 100} 个条目略")
                return "\n".join(lines)
        if ext == ".rar":
            import rarfile
            with rarfile.RarFile(p) as rf:
                infos = rf.infolist()
                lines = [f"压缩包 {path} 共 {len(infos)} 个条目："]
                for info in infos[:100]:
                    lines.append(f"· {info.filename}（{info.file_size} 字节）")
                if len(infos) > 100:
                    lines.append(f"… 其余 {len(infos) - 100} 个条目略")
                return "\n".join(lines)
        return f"错误：不支持的压缩格式 {ext or '（无扩展名）'}（支持 zip/tar/gz/7z/rar）"
    except ImportError as e:
        return f"错误：读取该格式需要额外依赖：{e}"
    except Exception as e:
        return f"错误：读取压缩包失败: {e}"


def write_excel(path, data, sheet="Sheet1"):
    """写入 Excel 文件（.xlsx）。data 为 JSON 数组（行数组或对象数组）。"""
    try:
        from openpyxl import Workbook
    except ImportError:
        return "错误：需要 openpyxl（pip install openpyxl）"
    if not path or not str(path).strip():
        return "错误：path 必填"
    if not isinstance(data, list):
        return "错误：data 必须是非空数组"
    p = permissions.resolve(path)
    if not p:
        return "错误：路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    try:
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        wb = Workbook()
        ws = wb.active
        ws.title = str(sheet or "Sheet1")[:31]
        if data and isinstance(data[0], dict):
            cols = list(data[0].keys())
            ws.append(cols)
            for row in data:
                # 混合 dict/非 dict 行健壮化：非 dict 行按空字典处理
                vals = [row.get(c, "") for c in cols] if isinstance(row, dict) else [""] * len(cols)
                ws.append(vals)
        else:
            for row in data:
                ws.append([str(x) for x in row] if isinstance(row, (list, tuple)) else [str(row)])
        wb.save(p)
        return f"已写入 Excel 至 {p}（{len(data)} 行）"
    except Exception as e:
        return f"错误：写入 Excel 失败: {e}"


def chart_data(data, path, kind="line", title="", x_label="", y_label=""):
    """数据可视化：生成图表 PNG（matplotlib）。data 为 JSON 数组：
    [x1, x2, ...]（单系列）或 [[x,y],...] 或 [{"x":..,"y":..}]。kind: line/bar/pie/scatter。"""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        # 中文字体（Windows 环境）：避免标题/标签中文变方块
        plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
        plt.rcParams["axes.unicode_minus"] = False
    except ImportError:
        return "错误：需要 matplotlib（pip install matplotlib）"
    if not path or not str(path).strip():
        return "错误：path 必填"
    if not isinstance(data, list) or not data:
        return "错误：data 必须是非空数组"
    p = permissions.resolve(path)
    if not p:
        return "错误：路径无效"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    try:
        def _to_float(v):
            try:
                return float(v or 0)
            except (TypeError, ValueError):
                raise ValueError(f"非数值数据：{v!r}（请只传数字）")

        if isinstance(data[0], dict) and "x" in data[0]:
            xs = [str(d.get("x", "")) for d in data]
            ys = [_to_float(d.get("y", 0)) for d in data]
        elif isinstance(data[0], (list, tuple)) and len(data[0]) >= 2:
            xs = [str(r[0]) for r in data]
            ys = [_to_float(r[1]) for r in data]
        else:
            xs = list(range(len(data)))
            ys = [_to_float(x) for x in data]
        # 非数值字符串在 float() 处会抛 ValueError：给出明确报错而不是笼统的失败
        # （NaN/inf 也无法绘制：matplotlib 静默出空图，先挡掉）
        if any(v != v or v in (float("inf"), float("-inf")) for v in ys):
            return "错误：数据包含 NaN 或无穷值，请清洗后重试"
        k = str(kind or "line").lower()
        if k not in ("line", "bar", "pie", "scatter"):
            return f"错误：kind 非法：{kind}（支持 line/bar/pie/scatter）"
        if k == "pie":
            if len(ys) > 20:
                return "错误：饼图最多支持 20 个数据点，请聚合后重试"
            if not any(v > 0 for v in ys):
                return "错误：饼图需要至少一个正值数据"
        if CHART_THEME == "light":
            face = "#ffffff"
            grid = "#d5e4ec"
            tick = "#5c7a96"
            title_c = "#14283f"
            series = "#00a3c8"
            chart_bg = "#f5f9fc"
        else:
            face = "#0a101f"
            grid = "#14203a"
            tick = "#9db0d1"
            title_c = "#e9f1ff"
            series = "#00d4ff"
            chart_bg = "#0a101f"
        fig, ax = plt.subplots(figsize=(8, 5), dpi=110, facecolor=face)
        ax.set_facecolor(chart_bg)
        for spine in ax.spines.values():
            spine.set_color(grid)
        ax.tick_params(colors=tick)
        ax.xaxis.label.set_color(tick)
        ax.yaxis.label.set_color(tick)
        ax.title.set_color(title_c)
        ax.grid(color=grid)
        if k == "bar":
            ax.bar(xs, ys, color=series)
        elif k == "pie":
            ax.pie(ys, labels=xs, autopct="%1.1f%%", textprops={"color": title_c})
        elif k == "scatter":
            ax.scatter(list(range(len(ys))), ys, color=series)
        else:
            ax.plot(xs, ys, color=series, marker="o", markersize=4)
        if title:
            ax.set_title(str(title))
        if x_label:
            ax.set_xlabel(str(x_label))
        if y_label:
            ax.set_ylabel(str(y_label))
        fig.tight_layout()
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        fig.savefig(p)
        plt.close(fig)
        size = os.path.getsize(p) if os.path.exists(p) else 0
        return f"已生成图表至 {p}（{size} 字节，{k} 图，{len(ys)} 个数据点）"
    except Exception as e:
        return f"错误：生成图表失败: {e}"


# ===== 长期记忆（与 main 的 memory.json 兼容：{"enabled", "facts":[{key,value}], "notes":[{text,tags,ts}]}）=====
MEMORY_FILE = None  # 由 main 初始化时注入（DATA_DIR/memory.json）
MEMORY_MAX_ITEMS = 2000  # v2.16.2 起扩容：伙伴需要记住的更多
MEMORY_MAX_TEXT = 2000
_MEMORY_LOCK = threading.Lock()  # 并行 write_memory 读-改-写串行化，防丢失更新


def _load_memory():
    if not MEMORY_FILE or not os.path.exists(MEMORY_FILE):
        return {"enabled": False, "facts": [], "notes": []}
    try:
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return {"enabled": False, "facts": [], "notes": []}
        data.setdefault("enabled", False)
        data.setdefault("facts", [])
        data.setdefault("notes", [])
        return data
    except Exception:
        logging.exception("读取记忆失败")
        return {"enabled": False, "facts": [], "notes": []}


def _save_memory(data):
    if not MEMORY_FILE:
        return False
    try:
        os.makedirs(os.path.dirname(MEMORY_FILE) or ".", exist_ok=True)
        tmp = MEMORY_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=1)
        os.replace(tmp, MEMORY_FILE)  # 原子写：崩溃不损坏 memory.json
        return True
    except Exception:
        logging.exception("保存记忆失败")
        return False


# ---------- 语义检索（纯标准库：字符 bigram + 词重叠 + IDF 权重） ----------
_BIGRAM_RE = re.compile(r"[\w\u4e00-\u9fff]+")


def _mem_tokens(text):
    """中文按双字符 bigram（含数字混排段也切分）、英文按词切分，用于相似度计算。"""
    text = str(text or "").lower()
    tokens = []
    for seg in _BIGRAM_RE.findall(text):
        if not seg:
            continue
        if re.search(r"[\u4e00-\u9fff]", seg):
            # 含中文的段（可能混入数字/字母，如 凌晨3点、备份v2）：
            # 按字符 bigram 切分，兼容数字/中文字形差异
            tokens.extend(seg[i : i + 2] for i in range(max(0, len(seg) - 1)))
            if len(seg) <= 4:
                tokens.append(seg)
        else:
            tokens.append(seg)
    return tokens


def _mem_idf(facts):
    """IDF 权重：低频特征更有区分度。"""
    import math

    n = max(1, len(facts))
    doc_count = {}
    for f in facts:
        seen = set(_mem_tokens((f.get("value") or "") + " " + (f.get("key") or "")))
        for tok in seen:
            doc_count[tok] = doc_count.get(tok, 0) + 1
    return {tok: math.log((n + 1) / (cnt + 1)) + 1 for tok, cnt in doc_count.items()}


def _mem_score(query_tokens, idf, text):
    """查询与记忆文本的加权相似度（0~1）。

    余弦（IDF 加权）+ 查询要点覆盖率混合：短查询（如"谁在搞备份"）只有
    一两个共同 bigram 时，覆盖率维度能正确反映"核心词命中"。"""
    if not query_tokens or not text:
        return 0.0
    toks = _mem_tokens(text)
    if not toks:
        return 0.0
    q_set, t_set = set(query_tokens), set(toks)
    common = q_set & t_set
    if not common:
        return 0.0
    w = sum(idf.get(t, 1.0) for t in common)
    cosine = w / (len(q_set) ** 0.5 * len(t_set) ** 0.5)
    recall = len(common) / len(q_set)
    return 0.55 * cosine + 0.45 * recall


def write_memory(text, tags="", type="", entities="", relations=""):
    """写入一条长期记忆（Agent 自动写入，与手动维护的 facts 同文件）。

    type：记忆类型（偏好/事实/项目/联系/规则 等，便于分类检索）；
    entities：涉及的实体列表（逗号分隔，如 张三,项目A），构成知识图谱节点；
    relations：关系三元组（分号分隔的 "实体-关系-实体"，如 张三-负责-项目A）。
    """
    text = str(text or "").strip()
    if not text:
        return "错误：记忆内容为空"
    if len(text) > MEMORY_MAX_TEXT:
        text = text[:MEMORY_MAX_TEXT] + "…"
    with _MEMORY_LOCK:
        data = _load_memory()
        key = str(tags or "").strip().split(",")[0].strip() or "自动记忆"
        facts = data.get("facts") or []
        for f in facts:
            if f.get("value") == text:
                return "该内容已存在，未重复写入"
        entry = {"key": key[:40], "value": text}
        if str(type or "").strip():
            entry["type"] = str(type).strip()[:20]
        ent = [e.strip()[:30] for e in str(entities or "").split(",") if e.strip()]
        if ent:
            entry["entities"] = ent
        rels = []
        for r in str(relations or "").split(";"):
            parts = [p.strip() for p in str(r).split("-") if p.strip()]
            if len(parts) == 3:
                rels.append({"rel": parts[1][:20], "to": parts[2][:30]})
        if rels:
            entry["relations"] = rels
        entry["ts"] = datetime.now().isoformat(timespec="seconds")
        facts.append(entry)
        if len(facts) > MEMORY_MAX_ITEMS:
            del facts[: len(facts) - MEMORY_MAX_ITEMS]
        data["facts"] = facts
        if _save_memory(data):
            return f"已写入记忆（当前共 {len(facts)} 条）"
        return "错误：记忆写入失败"


def read_memory(keyword="", max_items=20, type="", entity=""):
    """读取长期记忆（facts + notes）。

    keyword 为空时按时间倒序返回；非空时按语义相似度排序（TF-IDF + bigram，
    即使不含关键词也能检索到相关记忆）。
    type：按记忆类型过滤；entity：按实体过滤（知识图谱节点检索）。
    """
    data = _load_memory()
    facts = data.get("facts") or []
    if str(type or "").strip():
        t = str(type).strip()
        facts = [f for f in facts if str(f.get("type") or "") == t]
    if str(entity or "").strip():
        e = str(entity).strip().lower()
        facts = [f for f in facts if e in [x.lower() for x in (f.get("entities") or [])]]
    entries = []
    for f in facts:
        k = str(f.get("key") or "").strip()
        v = str(f.get("value") or "").strip()
        if k or v:
            entries.append((f"{k}: {v}".strip(": "), v, f))
    for n in data.get("notes") or []:
        t = str(n.get("text") or "").strip()
        if t:
            entries.append((t, t, {}))
    try:
        limit = max(1, min(100, int(max_items or 20)))
    except (TypeError, ValueError):
        limit = 20
    kw = str(keyword or "").strip()
    if kw:
        kwl = kw.lower()
        q_tokens = _mem_tokens(kw)
        idf = _mem_idf(facts)
        scored = []
        for label, v, f in entries:
            text = label + " " + v
            # 强匹配（包含关键词）优先；语义相似度作为补充（低相似度不返回，
            # 避免短关键词把无关条目带出来）
            exact = kwl in text.lower()
            sim = _mem_score(q_tokens, idf, text)
            if exact or sim >= 0.15:
                scored.append((1.0 if exact else sim, label, v, f))
        scored.sort(key=lambda x: -x[0])
        entries = [(label, v, f) for _, label, v, f in scored[:limit]]
        if not entries:
            return "（无匹配记忆）"
    else:
        entries = entries[-limit:][::-1]  # 最新在前
    if not entries:
        return "（暂无记忆）" if not kw else "（无匹配记忆）"
    lines = []
    for label, v, f in entries:
        meta = []
        if f.get("type"):
            meta.append(f"类型:{f['type']}")
        if f.get("entities"):
            meta.append(f"实体:{','.join(f['entities'])}")
        if f.get("relations"):
            meta.append(f"关系:{';'.join(r['rel'] + '→' + r['to'] for r in f['relations'])}")
        suffix = f" [{', '.join(meta)}]" if meta else ""
        lines.append(f"- {label}{suffix}")
    return "\n".join(lines)


def query_memory_graph(entity=None, relation="", max_items=20):
    """知识图谱查询：按实体/关系检索关联记忆（返回结构化的图谱片段）。

    示例：query_memory_graph(entity='张三') 返回所有涉及张三的记忆；
    query_memory_graph(relation='负责') 返回所有"负责"关系。
    """
    data = _load_memory()
    facts = data.get("facts") or []
    e = str(entity or "").strip().lower() if entity else ""
    r = str(relation or "").strip() if relation else ""
    hits = []
    for f in facts:
        match_e = not e or e in [x.lower() for x in (f.get("entities") or [])]
        rels = f.get("relations") or []
        match_r = not r or any(rel.get("rel") == r for rel in rels)
        if match_e and match_r:
            hits.append(f)
    try:
        limit = max(1, min(100, int(max_items or 20)))
    except (TypeError, ValueError):
        limit = 20
    hits = hits[-limit:]
    if not hits:
        return "（图谱中无匹配记忆）"
    lines = []
    for f in hits:
        v = str(f.get("value") or "").strip()
        lines.append(f"- {v}")
        for rel in f.get("relations") or []:
            lines.append(f"    {rel.get('rel')} → {rel.get('to')}")
    return "\n".join(lines)


# ===== SQLite 只读查询 =====
def database_query(db_path, sql, max_rows=20):
    """对本地 SQLite 数据库执行只读查询（SELECT/PRAGMA），路径需在允许目录内。"""
    if not db_path or not str(db_path).strip():
        return "错误：请提供数据库文件路径"
    if not sql or not str(sql).strip():
        return "错误：请提供 SQL 查询语句"
    ok, reason = permissions.check_filesystem(str(db_path), write=False)
    if not ok:
        return reason
    stmt = str(sql).strip()
    if not _readonly_stmt(stmt):
        return "错误：仅允许只读查询（SELECT / PRAGMA）"
    p = permissions.resolve(str(db_path))
    if not p or not os.path.exists(p):
        return f"错误：数据库文件不存在：{p}"
    try:
        import sqlite3
        import urllib.parse

        # 路径含 ?/# 时按 URI 查询参数解析，需先 percent-encode
        conn = sqlite3.connect(
            f"file:{urllib.parse.quote(p)}?mode=ro", uri=True, timeout=5
        )
        try:
            cur = conn.cursor()
            cur.execute(stmt)
            if cur.description is None:
                return "执行成功（无结果集）"
            cols = [d[0] for d in cur.description]
            try:
                limit = max(1, min(200, int(max_rows or 20)))
            except (TypeError, ValueError):
                limit = 20
            rows = cur.fetchmany(limit)
            lines = [f"查询结果（{len(rows)} 行）:", " | ".join(str(c) for c in cols)]
            for r in rows:
                cells = ["" if v is None else str(v) for v in r]
                cells = [c[:_TABLE_CELL_MAX] + ("…" if len(c) > _TABLE_CELL_MAX else "") for c in cells]
                lines.append(" | ".join(cells))
            return "\n".join(lines)
        finally:
            conn.close()
    except Exception as e:
        return f"错误：查询失败: {e}"


# ===== 邮件（需配置 SMTP）=====
EMAIL_CONFIG_FILE = None  # 由 main 注入（DATA_DIR/email_config.json）


def send_email(to, subject, body):
    """发送邮件：需要先配置 SMTP（email_config.json：smtp_host/smtp_port/user/password/from）。"""
    if not EMAIL_CONFIG_FILE or not os.path.exists(EMAIL_CONFIG_FILE):
        return (
            "错误：未配置邮件。请在数据目录创建 email_config.json，格式：\n"
            '{"smtp_host": "smtp.example.com", "smtp_port": 465, '
            '"user": "you@example.com", "password": "***", "from": "you@example.com"}'
        )
    try:
        with open(EMAIL_CONFIG_FILE, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        if isinstance(cfg, dict) and cfg.get("password"):
            cfg["password"] = _decrypt_secret(cfg["password"])
        smtp_host = str(cfg.get("smtp_host", "")).strip()
        smtp_port = int(cfg.get("smtp_port", 465))
        user = str(cfg.get("user", "")).strip()
        password = str(cfg.get("password", ""))
        from_addr = str(cfg.get("from") or user).strip()
        if not (smtp_host and user and password):
            return "错误：email_config.json 缺少 smtp_host / user / password"
        to = str(to or "").strip()
        # 严格校验：parseaddr + 无换行（CRLF 注入面），多收件人逗号分隔逐个校验
        import email.utils

        recipients = [r.strip() for r in to.split(",") if r.strip()]
        if not recipients or any(
            "@" not in email.utils.parseaddr(r)[1] or
            "\n" in r or "\r" in r
            for r in recipients
        ):
            return "错误：收件人邮箱格式不正确（支持逗号分隔多个地址）"
        to = ", ".join(recipients)
        import smtplib
        from email.header import Header
        from email.mime.text import MIMEText

        msg = MIMEText(str(body or ""), "plain", "utf-8")
        msg["Subject"] = Header(str(subject or ""), "utf-8")
        msg["From"] = from_addr
        msg["To"] = to
        try:
            server = smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=10)
        except Exception:
            # SSL 失败回退普通 SMTP+STARTTLS：465 端口是 SSL 专用，回退必须用 587
            fallback_port = 587 if smtp_port in (465, 0) else smtp_port
            server = smtplib.SMTP(smtp_host, fallback_port, timeout=10)
            try:
                server.starttls()
            except Exception:
                pass
        try:
            server.login(user, password)
            # sendmail 第二参必须是收件人列表：此前把逗号拼接串当单个收件人，
            # 多收件人时对方收不到信（rcpt 被当成一个非法地址）
            server.sendmail(from_addr, recipients, msg.as_string())
        finally:
            try:
                server.quit()
            except Exception:
                pass
        return f"邮件已发送至 {', '.join(recipients)}"
    except Exception as e:
        return f"错误：邮件发送失败: {e}"


# ===== 受限 pip 安装 =====
# 完全体模式：None = 全部放行（由用户授权决定）；如需恢复白名单，改为列表即可
PIP_ALLOWLIST = None
PIP_ALLOWLIST_NOTICE = (
    "注意：run_python 沙箱默认隔离（不加载第三方库），"
    "如需使用请调用 run_python 时设置 with_site=true。"
)


def pip_install(package):
    """安装 Python 库到当前环境（配合 run_python(with_site=true) 使用）。

    完全体模式下不限制包名；若 PIP_ALLOWLIST 为列表则仅允许白名单内安装。
    """
    pkg = str(package or "").strip()
    if not pkg:
        return "错误：请提供要安装的包名"
    base = re.split(r"[<>=!~]", pkg)[0].strip().lower()
    if PIP_ALLOWLIST is not None and base not in PIP_ALLOWLIST:
        return f"权限拒绝：仅允许安装白名单库：{PIP_ALLOWLIST}"
    try:
        # SpooledTemporaryFile 限流：--quiet 下 pip 错误输出仍可能 MB 级，防 OOM
        import tempfile

        with tempfile.SpooledTemporaryFile(
            max_size=1 << 20, mode="w+t", encoding="utf-8", errors="replace"
        ) as out:
            proc = subprocess.Popen(
                [sys.executable, "-m", "pip", "install", "--quiet", "--disable-pip-version-check", pkg],
                stdout=out, stderr=subprocess.STDOUT, text=True,
                encoding="utf-8", errors="replace",
            )
            try:
                proc.wait(timeout=300)
            except subprocess.TimeoutExpired:
                _kill_tree(proc)
                try:
                    proc.wait(timeout=3)
                except Exception:
                    pass
                return "错误：安装超时（300 秒）"
            out.seek(0)
            out_data = out.read()
            if len(out_data) > 20000:
                out_data = out_data[-20000:] + "\n[较早输出已省略]"
        if proc.returncode == 0:
            return f"已安装 {pkg}。\n{PIP_ALLOWLIST_NOTICE}"
        return f"安装失败：{out_data[-800:]}"
    except subprocess.TimeoutExpired:
        return "错误：安装超时（300 秒）"
    except Exception as e:
        return f"错误：{e}"


SEARCH_TIMEOUT = 8
SEARCH_MAX_RESULTS = 5
def _search_bing(query, num=SEARCH_MAX_RESULTS, offset=0, since="", until=""):
    url = (
        f"https://www.bing.com/search?q={quote(query)}"
        f"&count={num}&first={offset + 1}&setlang=zh-CN"
    )
    # 时间范围过滤（Bing filters 语法：ex1:"起..止"，可只给一端）
    if since or until:
        daterange = f"{since}..{until}" if since and until else f"{since or ''}..{until or ''}"
        url += f"&filters=ex1:%22{daterange}%22"
    resp = _http_client().get(
        url,
        headers={"User-Agent": _SEARCH_UA, "Accept-Language": "zh-CN,zh;q=0.9"},
        timeout=SEARCH_TIMEOUT,
    )
    resp.raise_for_status()
    results = []
    for m in _BING_RESULT_RE.finditer(resp.text):
        link, title, rest = m.groups()
        title = _strip_tags(title)
        if not title or link.startswith(("javascript:", "//")):
            continue
        snip_m = re.search(r"<p[^>]*>(.*?)</p>", rest, re.S)
        snippet = _strip_tags(snip_m.group(1)) if snip_m else ""
        results.append({"title": title, "url": link, "snippet": snippet})
        if len(results) >= num:
            break
    return results


def _search_duckduckgo(query, num=SEARCH_MAX_RESULTS, since=""):
    url = f"https://html.duckduckgo.com/html/?q={quote(query)}"
    if since:
        url += f"&df={since}"
    resp = _http_client().get(
        url,
        headers={"User-Agent": _SEARCH_UA},
        timeout=SEARCH_TIMEOUT,
    )
    resp.raise_for_status()
    results = []
    for m in _DDG_RESULT_RE.finditer(resp.text):
        link, title, snippet = m.groups()
        title = _strip_tags(title)
        if not title:
            continue
        results.append(
            {
                "title": title,
                "url": _decode_ddg_url(link),
                "snippet": _strip_tags(snippet or ""),
            }
        )
        if len(results) >= num:
            break
    return results


def _search_so360(query, num=SEARCH_MAX_RESULTS):
    """360 搜索：国内可达的稳定源（结果多、反爬弱），链接可能为 /link 加密跳转。"""
    url = f"https://www.so.com/s?q={quote(query)}"
    resp = _http_client().get(
        url,
        headers={"User-Agent": _SEARCH_UA, "Accept-Language": "zh-CN,zh;q=0.9"},
        timeout=SEARCH_TIMEOUT,
    )
    resp.raise_for_status()
    results = []
    for m in _SO360_RESULT_RE.finditer(resp.text):
        link, title_html = m.groups()
        if link.startswith("javascript:"):
            continue  # 置顶功能入口/推广位，非真实结果
        title = _strip_tags(title_html).strip()
        if not title:
            continue
        if link.startswith("/"):
            link = "https://www.so.com" + link  # /link?m= 加密跳转补全
        results.append({"title": title, "url": link, "snippet": ""})
        if len(results) >= num:
            break
    return results


# 搜索引擎注册表：(名称, 质量权重)。函数名为 "_search_<名称>"，调用时经
# globals() 动态查找——测试可 mock.patch 模块属性替换实现。权重决定聚合
# 输出顺序（数值大的优先展示）。实测结论：bing/so360 国内稳定；duckduckgo
# 时好时坏（健康度机制自动跳过）；baidu/sogou/yandex 反爬；google 等不可达。
_SEARCH_ENGINES = (
    ("bing", 3),
    ("so360", 2),
    ("duckduckgo", 1),
)

# 引擎健康度：连续失败 3 次暂停 10 分钟，成功一次即恢复
_SEARCH_HEALTH = {}  # name -> {"fails": int, "skip_until": float}
_SEARCH_HEALTH_FAIL_LIMIT = 3
_SEARCH_HEALTH_COOLDOWN = 600.0
_SEARCH_HEALTH_LOCK = threading.Lock()


def _search_healthy(name):
    with _SEARCH_HEALTH_LOCK:
        h = _SEARCH_HEALTH.get(name)
        if not h:
            return True
        if h["fails"] >= _SEARCH_HEALTH_FAIL_LIMIT:
            if time.time() >= h.get("skip_until", 0):
                h["fails"] = 0  # 冷却结束，重新尝试
                return True
            return False
        return True


def _search_report(name, ok):
    with _SEARCH_HEALTH_LOCK:
        h = _SEARCH_HEALTH.setdefault(name, {"fails": 0, "skip_until": 0.0})
        if ok:
            h["fails"] = 0
        else:
            h["fails"] += 1
            if h["fails"] >= _SEARCH_HEALTH_FAIL_LIMIT:
                h["skip_until"] = time.time() + _SEARCH_HEALTH_COOLDOWN
                logger.warning("搜索源 %s 连续 %d 次失败，暂停 %d 分钟", name,
                               _SEARCH_HEALTH_FAIL_LIMIT, _SEARCH_HEALTH_COOLDOWN // 60)


def search_web(query, num=SEARCH_MAX_RESULTS, offset=0, since="", until="", site=""):
    """联网搜索：多引擎并行聚合（bing/360/duckduckgo），支持条数/翻页/时间/站点过滤。

    引擎健康度：连续失败 3 次的引擎自动暂停 10 分钟（进程内），可用引擎互补；
    不可用/质量差的源（baidu/sogou/yandex 反爬、google 等不可达）不在注册表中。

    Args:
        query: 搜索关键词
        num: 返回条数（1-20，默认 5）
        offset: 翻页偏移（0 起，如 5 表示第 6-10 条；仅 Bing 支持）
        since/until: 时间范围过滤（YYYY-MM-DD，可只给一端；DDG 仅支持 since）
        site: 限定站点域名（如 "openai.com"，自动追加 site:）
    """
    if not query or not str(query).strip():
        return "错误：搜索词为空"
    try:
        num = max(1, min(20, int(num)))
        offset = max(0, min(200, int(offset)))
    except (TypeError, ValueError):
        num, offset = SEARCH_MAX_RESULTS, 0
    for tag, val in (("since", since), ("until", until)):
        if val and not re.match(r"^\d{4}-\d{2}-\d{2}$", str(val)):
            return f"错误：{tag} 日期格式应为 YYYY-MM-DD"
    q = str(query).strip()
    if len(q) > 200:
        return "错误：搜索词过长（上限 200 字符）"
    site = str(site or "").strip()
    if site:
        if len(site) > 100 or not re.match(r"^[A-Za-z0-9.\-]+$", site):
            return "错误：site 参数应为域名（如 openai.com）"
        q = f"{q} site:{site}"
    since = str(since or "").strip()
    until = str(until or "").strip()

    # 并行调用所有健康引擎；每引擎请求 num+offset 条供聚合后手动翻页
    import concurrent.futures as _cf

    def _run(entry):
        name, _weight = entry
        fn = globals().get("_search_" + name)  # 动态查找：支持测试 mock 替换
        if fn is None:
            return name, [], None
        try:
            kw = {"num": max(num + offset, 10)}
            if name == "bing":
                kw.update(offset=offset, since=since, until=until)
            elif name == "duckduckgo":
                kw["since"] = since
            results = fn(q, **kw)
            return name, results, None
        except Exception as e:
            return name, [], e

    engines = [e for e in _SEARCH_ENGINES if _search_healthy(e[0])]
    with _cf.ThreadPoolExecutor(max_workers=len(engines) or 1) as ex:
        outcomes = list(ex.map(_run, engines))

    merged, last_err = [], None
    for name, results, err in outcomes:
        if err is not None or not results:
            _search_report(name, False)
            if err is not None:
                last_err = err
            continue
        _search_report(name, True)
        merged.extend(_search_safe(results))
    merged = _search_dedup(merged)
    # site 硬过滤：搜索引擎可能忽略 site: 语法，聚合后按域名兜底保证生效
    pre_site = merged
    if site:
        merged = [
            r for r in merged
            if str(r.get("url") or "").split("/")[2].lower() in (site, "www." + site)
            or str(r.get("url") or "").split("/")[2].lower().endswith("." + site)
        ]
    # offset 手动翻页：请求时已多取，这里直接切片（引擎不支持 first= 也生效）
    merged = merged[offset:offset + num]
    if merged:
        lines = [f"搜索结果（{len(merged)} 条）:"]
        for i, r in enumerate(merged, 1):
            lines.append(f"{i}. {r['title']}\n   {r['url']}\n   {r['snippet']}".rstrip())
        return "\n\n".join(lines)
    if site and pre_site:
        return f"未找到限定站点 {site} 的结果（搜索引擎未返回该站点内容，可尝试去掉 site 参数）"
    detail = f": {last_err}" if last_err is not None else ""
    return f"错误：搜索失败（可用搜索源均不可用{detail}）"


def search_github(query, num=5, language=""):
    """GitHub 仓库搜索（代码/开源项目垂直源，实测国内可达）。

    GitHub API 未认证限流 60 次/小时，适合低频垂直检索。
    """
    if not query or not str(query).strip():
        return "错误：搜索词为空"
    try:
        num = max(1, min(20, int(num)))
    except (TypeError, ValueError):
        num = 5
    q = str(query).strip()
    if len(q) > 200:
        return "错误：搜索词过长（上限 200 字符）"
    language = str(language or "").strip()
    if language:
        if len(language) > 40 or not re.match(r"^[A-Za-z0-9+#.\-]+$", language):
            return "错误：language 参数不合法"
        q = f"{q} language:{language}"
    try:
        resp = _http_client().get(
            "https://api.github.com/search/repositories",
            params={"q": q, "per_page": num, "sort": "stars"},
            headers={"Accept": "application/vnd.github+json", "User-Agent": _SEARCH_UA},
            timeout=10,
        )
        if resp.status_code == 403:
            return "错误：GitHub API 限流（每小时 60 次），请稍后再试"
        resp.raise_for_status()
        items = (resp.json() or {}).get("items") or []
    except Exception as e:
        return f"错误：GitHub 搜索失败: {e}"
    if not items:
        return "未找到相关仓库"
    lines = [f"GitHub 仓库（{len(items)} 个，按 Star 排序）:"]
    for i, it in enumerate(items, 1):
        desc = (it.get("description") or "").strip()[:120]
        stars = it.get("stargazers_count", 0)
        lines.append(f"{i}. {it.get('full_name', '?')} ⭐{stars}\n   {it.get('html_url', '')}\n   {desc}".rstrip())
    return "\n\n".join(lines)


def search_realtime(query="", num=5, source="hn"):
    """实时信息通道：Hacker News（热点/搜索），绕开通用搜索引擎的实时性短板。

    HN 无 query 时返回 Top Stories 热点；有 query 时走 Algolia 全文搜索
    （含 points/评论数/时间）。Reddit 等源在当前网络不可达，未接入。
    """
    try:
        num = max(1, min(20, int(num)))
    except (TypeError, ValueError):
        num = 5
    src = str(source or "hn").strip().lower()
    try:
        if src != "hn":
            return f"错误：暂不支持数据源 {src}（当前仅 hn）"
        q = str(query or "").strip()
        if q:
            if len(q) > 120:
                return "错误：搜索词过长（上限 120 字符）"
            resp = _http_client().get(
                "https://hn.algolia.com/api/v1/search",
                params={"query": q, "hitsPerPage": num},
                headers={"User-Agent": _SEARCH_UA},
                timeout=10,
            )
            resp.raise_for_status()
            hits = (resp.json() or {}).get("hits") or []
            if not hits:
                return f"Hacker News 未找到与「{q}」相关的结果"
            lines = [f"Hacker News 搜索结果（{len(hits)} 条）:"]
            for i, h in enumerate(hits, 1):
                title = str(h.get("title") or "").strip()[:120]
                url = str(h.get("url") or "").strip() or (
                    f"https://news.ycombinator.com/item?id={h.get('objectID')}"
                )
                pts = h.get("points") or 0
                cmts = h.get("num_comments") or 0
                lines.append(f"{i}. {title}（👍{pts} 💬{cmts}）\n   {url}")
            return "\n\n".join(lines)
        # 无 query：实时热点榜
        resp = _http_client().get(
            "https://hacker-news.firebaseio.com/v0/topstories.json",
            headers={"User-Agent": _SEARCH_UA},
            timeout=10,
        )
        resp.raise_for_status()
        ids = (resp.json() or [])[:num]
        if not ids:
            return "Hacker News 热点暂时为空"
        items = []
        for sid in ids:
            try:
                r = _http_client().get(
                    f"https://hacker-news.firebaseio.com/v0/item/{sid}.json",
                    headers={"User-Agent": _SEARCH_UA},
                    timeout=8,
                )
                it = r.json()
                if it and it.get("title"):
                    items.append(it)
            except Exception:
                continue
            if len(items) >= num:
                break
        if not items:
            return "Hacker News 热点获取失败"
        lines = [f"Hacker News 实时热点（{len(items)} 条）:"]
        for i, it in enumerate(items, 1):
            title = str(it.get("title") or "").strip()[:120]
            url = str(it.get("url") or "").strip() or (
                f"https://news.ycombinator.com/item?id={it.get('id')}"
            )
            pts = it.get("score") or 0
            lines.append(f"{i}. {title}（👍{pts}）\n   {url}")
        return "\n\n".join(lines)
    except Exception as e:
        return f"错误：实时信息获取失败: {e}"


CALL_API_MAX_BYTES = 500 * 1024  # 响应体上限 500KB（与 fetch_url 输出对齐）
CALL_API_METHODS = ("GET", "POST", "PUT", "DELETE", "PATCH", "HEAD")
CALL_API_MAX_HEADERS = 16
# 内网/回环白名单（main 从 config 的 call_api_allowed_hosts 注入）：
# 命中精确主机名的请求跳过 SSRF 拦截——用于用户显式放行的本地服务
# （如 127.0.0.1:8000 的本地 API），其余内网地址照常拦截。仅精确匹配，
# 建议填 IP 而非域名（避免 DNS 重绑定绕过）。
CALL_API_ALLOWED_HOSTS = []


def _call_api_host_allowed(url):
    try:
        host = _url_host(url)
        if not host:
            return False
        allow = {str(h).strip().lower() for h in CALL_API_ALLOWED_HOSTS if str(h).strip()}
        return host.lower() in allow
    except Exception:
        return False


def call_api(url, method="GET", params=None, json_body=None, data=None,
             headers=None, timeout=15):
    """通用外部 API 调用（自主 AI 的"万能接口"）。

    支持 GET/POST/PUT/DELETE/PATCH/HEAD，自定义查询参数/JSON 体/表单体/请求头。
    安全约束：仅 http(s) 公网地址（SSRF 防护，与 fetch_url 同规则）；请求头
    禁止 CRLF 注入；响应体 ≤200KB 截断；超时上限 60s。

    Args:
        url: 完整 API 地址（http/https）
        method: 请求方法（默认 GET）
        params: 可选，查询参数 dict（如 {"limit": 10}）
        json_body: 可选，JSON 请求体 dict
        data: 可选，表单/原始请求体
        headers: 可选，自定义请求头 dict（≤8 个）
        timeout: 可选，超时秒数（1-60，默认 15）
    """
    if not url or not str(url).startswith(("http://", "https://")):
        return "错误：url 必须以 http:// 或 https:// 开头"
    if not _call_api_host_allowed(url):
        err = _safe_url(url, allow_loopback=False)
        if err:
            return f"错误：{err}（如需访问本地/内网服务，可在配置 call_api_allowed_hosts 中加入该主机白名单）"
    method = str(method or "GET").strip().upper()
    if method not in CALL_API_METHODS:
        return f"错误：method 仅支持 {'/'.join(CALL_API_METHODS)}"
    try:
        timeout = max(1, min(180, int(timeout or 15)))
    except (TypeError, ValueError):
        timeout = 15
    hdrs = {}
    if headers:
        if not isinstance(headers, dict):
            return "错误：headers 必须是键值对象"
        if len(headers) > CALL_API_MAX_HEADERS:
            return f"错误：headers 最多 {CALL_API_MAX_HEADERS} 个"
        for k, v in headers.items():
            k, v = str(k).strip(), str(v or "").strip()
            if not k or not re.match(r"^[A-Za-z0-9\-]+$", k):
                return f"错误：请求头名称不合法：{k}"
            if "\r" in v or "\n" in v:
                return "错误：请求头值禁止包含换行（防 CRLF 注入）"
            hdrs[k] = v
    try:
        kw = {"params": params} if params else {}
        if json_body is not None:
            kw["json"] = json_body
        if data is not None:
            kw["data"] = data
        def _validate(u):
            if _call_api_host_allowed(u):
                return ""
            return _safe_url(u, allow_loopback=False)

        raw = b""
        truncated = False
        status_code = 0
        content_type = ""
        # 流式读取：大响应不再全量进内存，超过上限立即断开连接
        with _safe_stream(
            method, url, validate=_validate,
            headers=hdrs or None, timeout=timeout, **kw
        ) as resp:
            resp.raise_for_status()
            status_code = resp.status_code
            content_type = (resp.headers or {}).get("content-type", "") if hasattr(resp.headers, "get") else ""
            if hasattr(resp, "iter_bytes"):
                for chunk in resp.iter_bytes(64 * 1024):
                    raw += chunk
                    if len(raw) >= CALL_API_MAX_BYTES:
                        truncated = True
                        break
            else:
                # 兼容旧测试/自定义 mock 的普通响应对象（无流式接口）
                raw = getattr(resp, "content", b"") or b""
                truncated = len(raw) > CALL_API_MAX_BYTES
                raw = raw[:CALL_API_MAX_BYTES]
        body = raw
        text = body[:CALL_API_MAX_BYTES].decode("utf-8", errors="replace")
        # JSON 美化输出（若可解析），便于阅读
        try:
            if content_type.startswith("application/json"):
                import json as _json
                text = _json.dumps(_json.loads(text), ensure_ascii=False, indent=2)
        except Exception:
            pass
        head = f"HTTP {status_code} · {method} {url.split('?')[0][:80]}"
        if truncated:
            head += f" · 响应已截断（>{CALL_API_MAX_BYTES // 1024}KB，显示前 {CALL_API_MAX_BYTES // 1024}KB）"
        return f"{head}\n\n{text}" if text.strip() else head
    except Exception as e:
        return f"错误：API 调用失败（{type(e).__name__}: {str(e)[:120]}）"


def system_status():
    """系统资源自检：CPU / 内存 / 磁盘 / 网络连通性（自主 AI 的"体检"能力）。

    psutil 可选（缺失时 CPU/内存降级提示）；网络连通性探测
    api.github.com / bing / api.deepseek.com 三个关键端点。
    """
    lines = ["系统状态："]
    try:
        import psutil

        cpu = psutil.cpu_percent(interval=0.5)
        mem = psutil.virtual_memory()
        lines.append(f"- CPU：{cpu:.0f}% 使用（{psutil.cpu_count()} 核）")
        lines.append(
            f"- 内存：{mem.used / 1024 ** 3:.1f}GB / {mem.total / 1024 ** 3:.1f}GB"
            f"（{mem.percent:.0f}%）"
        )
    except ImportError:
        lines.append("- CPU/内存：需安装 psutil（pip install psutil）获得详情")
    except Exception as e:
        lines.append(f"- CPU/内存：读取失败（{e}）")
    try:
        base = (
            permissions.WORKSPACE_DIR
            if permissions.WORKSPACE_DIR and os.path.isdir(permissions.WORKSPACE_DIR)
            else os.getcwd()
        )
        du = shutil.disk_usage(base)
        lines.append(
            f"- 磁盘（工作区）：剩余 {du.free / 1024 ** 3:.1f}GB / 总 {du.total / 1024 ** 3:.1f}GB"
        )
    except Exception:
        pass
    reach = []
    for host, port in (
        ("api.github.com", 443),
        ("www.bing.com", 443),
        ("api.deepseek.com", 443),
    ):
        try:
            import socket

            s = socket.create_connection((host, port), timeout=1.5)
            s.close()
            reach.append(f"{host} ✓")
        except Exception:
            reach.append(f"{host} ✗")
    lines.append("- 网络：" + " ".join(reach))
    return "\n".join(lines)


def _atomic_write(path, content):
    """原子写：唯一临时文件（防并行写同一路径互相截断）+ os.replace，
    覆盖前自动备份 .bak。返回 (created, real_size)。"""
    created = not os.path.exists(path)
    if not created:
        try:
            shutil.copy2(path, path + ".bak")
        except Exception:
            pass
    import tempfile

    dirname = os.path.dirname(path) or "."
    os.makedirs(dirname, exist_ok=True)
    fd, tmp = tempfile.mkstemp(
        dir=dirname, prefix=os.path.basename(path) + ".", suffix=".tmp"
    )
    try:
        # newline="" 关闭文本翻译：Windows 默认把 \n 写成 \r\n，既改变换行风格，
        # 也使"按 UTF-8 字节校验"的大小与实际落盘字节不一致
        with os.fdopen(fd, "w", encoding="utf-8", newline="") as f:
            f.write(content)
        os.replace(tmp, path)
        tmp = None
    finally:
        if tmp is not None:
            try:
                os.unlink(tmp)
            except OSError:
                pass
    try:
        real_size = os.path.getsize(path)
    except OSError:
        real_size = len(content.encode("utf-8", "replace"))
    return created, real_size


def write_file(path, content):
    """写文件：目录白名单 + 大小限制（按 UTF-8 字节计）+ 原子写 + 自动 .bak + 写入后真实核验。"""
    ok, reason = permissions.check_filesystem(path, write=True)
    if not ok:
        return reason
    if content is None:
        return "错误：内容为空"
    # 中文按字符数会超限 3 倍：统一按 UTF-8 字节数校验
    content_bytes = len(str(content).encode("utf-8", "ignore"))
    if content_bytes > permissions.max_write_size():
        return (
            f"错误：内容 {content_bytes} 字节超过大小限制 "
            f"{permissions.max_write_size()}"
        )
    p = permissions.resolve(path)
    try:
        created, real_size = _atomic_write(p, str(content))
        if not os.path.exists(p):
            return f"错误：写入后核验失败，文件不存在：{p}"
        permissions.audit("write_file", p, f"{real_size} 字节")
        return (
            f"已写入 {p}（{'新建' if created else '覆盖并备份 .bak'}，"
            f"实际 {real_size} 字节，已核验存在）"
        )
    except Exception as e:
        return f"错误：写入失败: {e}"


def edit_file(path, old="", new="", regex=None):
    """编辑文件：按文本或正则替换（自动备份 .bak）。"""
    ok, reason = permissions.check_filesystem(path, write=True)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not os.path.isfile(p):
        return f"错误：文件不存在：{p}"
    try:
        # 读入上限：允许目录内也可能有 GB 级文件，全量读入内存会 OOM
        if os.path.getsize(p) > EDIT_FILE_MAX_SIZE:
            return f"错误：文件超过 {EDIT_FILE_MAX_SIZE // 1024 // 1024}MB 上限，请改用其他方式处理"
        with open(p, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
    except Exception as e:
        return f"错误：读取失败: {e}"
    if regex:
        if len(str(regex)) > EDIT_FILE_REGEX_MAX:
            return f"错误：正则过长（>{EDIT_FILE_REGEX_MAX} 字符）"
        try:
            # lambda 返回 new 原样：re.sub 的字符串替换会把 new 中的 \1 / \g<1> /
            # \\ 解释为分组引用与转义，模型生成的替换文本含反斜杠时会被静默改写
            new_content, n = re.subn(regex, lambda m: new or "", content)
        except re.error as e:
            return f"错误：正则无效: {e}"
    else:
        if not old:
            return "错误：需要提供 old（原文）或 regex（正则）"
        if old not in content:
            return "错误：目标文本未找到"
        new_content = content.replace(old, new or "")
        n = content.count(old)
    if n == 0:
        return "错误：无匹配内容，未做修改"
    try:
        _atomic_write(p, new_content)
        if not os.path.exists(p):
            return f"错误：写入后核验失败，文件不存在：{p}"
        permissions.audit("edit_file", p, f"替换 {n} 处")
        return f"已替换 {n} 处，写入 {p}（已备份 .bak，已核验存在）"
    except Exception as e:
        return f"错误：写入失败: {e}"


def list_dir(path):
    """列目录（只读）。"""
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not os.path.isdir(p):
        return f"错误：目录不存在：{p}"
    # scandir 惰性迭代 + islice：百万条目目录不再全量 listdir（1-2s + 数百 MB 内存）
    import itertools

    try:
        entries = list(itertools.islice(os.scandir(p), 200))
        entries.sort(key=lambda e: e.name)  # 保持确定性排序（只排前 200 条）
    except Exception as e:
        return f"错误：{e}"
    lines = []
    for e in entries:
        try:
            is_dir = e.is_dir()
            kind = "DIR " if is_dir else "FILE"
            size = ""
            if not is_dir:
                try:
                    size = f" {e.stat().st_size}B"
                except OSError:
                    pass
            lines.append(f"{kind} {e.name}{size}")
        except OSError:
            continue
    # 统计总数需要完整遍历，代价大：仅当目录小才统计（大目录直接显示前 200）
    total = "200+" if len(entries) >= 200 else str(len(entries))
    tail = f"\n[共 {total} 项，仅显示前 200]" if len(entries) >= 200 else f"\n共 {total} 项"
    return "\n".join(lines) + tail


def run_command(command):
    """执行白名单命令（argv 直传，禁止 shell 拼接）。"""
    ok, reason, argv = permissions.check_shell(command)
    if not ok:
        return reason
    timeout = permissions.shell_timeout()
    try:
        # 输出 spool 到临时文件：命令刷屏（type 大日志）时内存峰值限 1MB
        import tempfile

        with tempfile.SpooledTemporaryFile(
            max_size=1 << 20, mode="w+t", encoding="utf-8", errors="replace"
        ) as out:
            proc = subprocess.Popen(
                argv,
                stdout=out,
                stderr=subprocess.STDOUT,
                text=True,
                encoding="utf-8",
                errors="replace",
                # 命令在工作目录执行（跟随 📁 目录设置），相对路径引用不再漂移
                cwd=WORKING_DIR or permissions.WORKSPACE_DIR or None,
            )
            try:
                proc.wait(timeout=timeout)
            except subprocess.TimeoutExpired:
                _kill_tree(proc)
                try:
                    proc.wait(timeout=3)
                except Exception:
                    pass
                return f"错误：命令超时（>{timeout} 秒）"
            out.seek(0)
            out_data = out.read(20000)
            out.seek(0, os.SEEK_END)
            if out.tell() > 20000:
                out_data += "\n[输出已截断]"
        permissions.audit("run_command", " ".join(argv), f"rc={proc.returncode}")
        if not out_data.strip():
            return f"执行成功（无输出），退出码 {proc.returncode}"
        return f"退出码 {proc.returncode}\n{out_data}"
    except Exception as e:
        return f"错误：{e}"


# ===== 后台进程管理（服务器/长驻任务）=====
MAX_PROCESSES = 8
WORKING_DIR = None  # 由 main 注入（工作目录：run_command/start_process 的 cwd）
PROCESSES = {}  # name -> {"proc", "pid", "name", "started", "exited", "code", "lines": deque}
_PROCESSES_LOCK = threading.Lock()
_process_cb = None  # (name, line) -> None（主线程经队列消费）


def snapshot_processes():
    """返回进程表快照（线程安全，供主线程/面板读取）。"""
    with _PROCESSES_LOCK:
        return list(PROCESSES.items())


def get_process(name):
    """读取单个进程条目（线程安全）。"""
    with _PROCESSES_LOCK:
        return PROCESSES.get(name)


def set_process_output_callback(cb):
    global _process_cb
    _process_cb = cb


def _emit_process(name, line):
    if _process_cb:
        try:
            _process_cb(name, line)
        except Exception:
            logger.debug("进程输出回调异常", exc_info=True)


def _process_reader(name):
    entry = get_process(name)
    if not entry:
        return
    proc = entry["proc"]
    batch = []
    last = time.monotonic()

    def flush():
        if not batch:
            return
        joined = "\n".join(batch)
        try:
            # 持锁 extend：与 list_processes 的拷贝迭代互斥，防 deque 迭代竞态
            with _PROCESSES_LOCK:
                entry["lines"].extend(batch)  # deque(maxlen=2000) 自动裁剪
        except Exception:
            pass
        _emit_process(name, joined)  # 整批回调：队列条目降为 1，UI 批量渲染
        batch.clear()

    try:
        for raw in proc.stdout:
            line = raw.rstrip("\n")
            if len(line) > 2000:
                line = line[:2000] + "…[截断]"  # 单行限长：防单行 100KB 撑爆 deque 与 UI
            batch.append(line)
            now = time.monotonic()
            if len(batch) >= 50 or now - last >= 0.1:
                flush()
                last = now
    except Exception:
        pass
    flush()
    entry["exited"] = True
    entry["code"] = proc.poll()
    _emit_process(name, f"── 进程已退出（code={entry['code']}）──")


def start_process(command, name=""):
    """后台启动长驻进程（服务器等），输出实时推送终端面板。"""
    ok, reason, argv = permissions.check_shell(command)
    if not ok:
        return reason
    with _PROCESSES_LOCK:
        for k in [k for k, v in PROCESSES.items() if v.get("exited")]:
            PROCESSES.pop(k, None)
        if len(PROCESSES) >= MAX_PROCESSES:
            return f"错误：后台进程数已达上限（{MAX_PROCESSES} 个），请先 stop_process 停止部分进程"
    try:
        # Python 子进程输出到管道为全缓冲，实时日志需 -u 无缓冲
        if os.path.basename(argv[0]).lower().startswith("python") and "-u" not in argv:
            argv = [argv[0], "-u"] + argv[1:]
        proc = subprocess.Popen(
            argv,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            encoding="utf-8",
            errors="replace",
            creationflags=(
                subprocess.CREATE_NO_WINDOW
                if hasattr(subprocess, "CREATE_NO_WINDOW")
                else 0
            ),
            # 与 run_command 一致：在工作目录启动（服务器/脚本的相对路径引用落地）
            cwd=WORKING_DIR or permissions.WORKSPACE_DIR or None,
        )
    except Exception as e:
        return f"错误：进程启动失败: {e}"
    # 命名与插入在启动后同一临界区完成：杜绝并发 start 同 base 名互相覆盖
    # （此前查重与插入分开持锁，两并行调用可在间隙抢到同名 → 进程失管无法 stop）
    with _PROCESSES_LOCK:
        for k in [k for k, v in PROCESSES.items() if v.get("exited")]:
            PROCESSES.pop(k, None)
        if len(PROCESSES) >= MAX_PROCESSES:
            _kill_tree(proc)  # 启动后才发现达上限：进程必须回收，不能留孤儿
            return f"错误：后台进程数已达上限（{MAX_PROCESSES} 个），请先 stop_process 停止部分进程"
        base = str(name or "").strip()[:40] or os.path.basename(argv[0]) or f"proc{len(PROCESSES) + 1}"
        proc_name = base
        i = 2
        while proc_name in PROCESSES:
            proc_name = f"{base}_{i}"
            i += 1
        PROCESSES[proc_name] = {
            "proc": proc,
            "pid": proc.pid,
            "name": proc_name,
            "started": datetime.now().strftime("%H:%M:%S"),
            "exited": False,
            "code": None,
            "lines": deque(maxlen=2000),
        }
    _emit_process(proc_name, f"── 进程启动 pid={proc.pid}（{' '.join(argv)}）──")
    threading.Thread(target=_process_reader, args=(proc_name,), daemon=True).start()
    permissions.audit("start_process", proc_name, " ".join(argv))
    return (
        f"已启动后台进程「{proc_name}」（pid={proc.pid}）\n"
        "实时输出见「工具 → 进程终端」。可用 stop_process 停止，list_processes 查询状态。"
    )


def stop_process(target):
    """停止后台进程（按名称或 pid）。"""
    target = str(target or "").strip()
    if not target:
        return "错误：需要进程名或 pid"
    for name, entry in snapshot_processes():
        if name == target or str(entry["pid"]) == target:
            if not entry["exited"]:
                try:
                    _kill_tree(entry["proc"])
                    try:
                        entry["proc"].wait(timeout=3)
                    except subprocess.TimeoutExpired:
                        pass
                except Exception:
                    pass
                with _PROCESSES_LOCK:
                    entry["exited"] = True
                    PROCESSES.pop(name, None)  # 立即回收条目，防失管/重复 stop
                _emit_process(name, "── 已停止 ──")
                return f"已停止进程「{name}」（pid={entry['pid']}）"
            with _PROCESSES_LOCK:
                PROCESSES.pop(name, None)
            return f"进程「{name}」已退出（code={entry['code']}）"
    running = [f"{n}({e['pid']})" for n, e in snapshot_processes()]
    return f"未找到进程：{target}（运行中：{running or '无'}）"


def list_processes():
    """列出后台进程状态与最近输出。"""
    entries = snapshot_processes()
    if not entries:
        return "当前没有后台进程"
    lines = []
    for name, e in entries:
        status = "运行中" if not e["exited"] else f"已退出(code={e['code']})"
        lines.append(f"· {name}（pid={e['pid']}，{e['started']}，{status}）")
        # 持锁拷贝：reader 线程 extend 与这里迭代并发可能抛 "deque mutated during iteration"
        with _PROCESSES_LOCK:
            recent = list(e["lines"])[-3:]
        for ln in recent:
            lines.append(f"    {ln[:120]}")
    return "\n".join(lines)


def stop_all_processes():
    """退出程序时终止所有后台进程（防孤儿进程，含进程树）。"""
    for entry in snapshot_processes():
        try:
            _kill_tree(entry["proc"])
            try:
                entry["proc"].wait(timeout=2)
            except subprocess.TimeoutExpired:
                pass
        except Exception:
            pass


_COMMON_PACKAGES = (
    "flask", "django", "fastapi", "uvicorn", "requests", "bs4", "pandas",
    "numpy", "matplotlib", "playwright", "docx", "pytest", "httpx",
    "openai", "tiktoken", "pillow", "tqdm", "yaml", "jinja2",
)


def environment_info():
    """运行环境信息：避免 AI 重复安装/做无用假设。"""
    import importlib.util

    lines = [f"平台: {sys.platform}", f"Python: {sys.version.split()[0]}"]
    # find_spec 按导入名检测：pillow 的导入名是 PIL（pip 包名 ≠ 导入名）
    _spec_names = {
        "pillow": "PIL", "bs4": "bs4", "yaml": "yaml", "docx": "docx",
        "opencv-python": "cv2", "xlrd": "xlrd", "openpyxl": "openpyxl",
    }
    installed = [
        p if importlib.util.find_spec(_spec_names.get(p, p)) else None
        for p in _COMMON_PACKAGES
    ]
    installed = [p for p in installed if p]
    lines.append(f"已安装常用包: {', '.join(installed) if installed else '（无）'}")
    if permissions.WORKSPACE_DIR:
        try:
            usage = shutil.disk_usage(permissions.WORKSPACE_DIR)
            lines.append(
                f"工作区磁盘: 剩余 {usage.free / 1024 ** 3:.1f}GB / 总 {usage.total / 1024 ** 3:.1f}GB"
            )
        except Exception:
            pass
    return "\n".join(lines)


# ===== 自我进化（感知自身代码 → 分支提案）=====
PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
EVOLUTIONS_DIR = os.path.join(PROJECT_DIR, "evolutions")
PROJECT_READ_EXTS = (".py", ".md", ".json", ".txt", ".bat", ".html")
EVO_WRITE_EXTS = (".py", ".md", ".json", ".txt", ".html")


def _current_version():
    try:
        with open(os.path.join(PROJECT_DIR, "main.py"), encoding="utf-8") as f:
            src = f.read()
        m = re.search(r'^VERSION\s*=\s*"([\d.]+)"', src, re.MULTILINE)
        if m:
            return m.group(1)
    except Exception:
        pass
    return "unknown"


def _py_stats(path):
    """统计 py 文件函数/类数量（粗略）。"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            src = f.read()
        fn = len(re.findall(r"^    def |^def ", src, re.MULTILINE))
        cls = len(re.findall(r"^class ", src, re.MULTILINE))
        return f"{len(src.splitlines())} 行 · {fn} 函数 · {cls} 类"
    except Exception:
        return "?"


def project_info():
    """感知鲸语自身代码库（只读）。"""
    lines = [f"鲸语版本: {_current_version()}", "项目文件："]
    for fn in sorted(os.listdir(PROJECT_DIR)):
        full = os.path.join(PROJECT_DIR, fn)
        if os.path.isfile(full) and fn.endswith(PROJECT_READ_EXTS):
            try:
                size = os.path.getsize(full)
                size_txt = f"{size / 1024:.1f}KB" if size >= 1024 else f"{size}B"
                extra = f" · {_py_stats(full)}" if fn.endswith(".py") else ""
                lines.append(f"- {fn}（{size_txt}{extra}）")
            except Exception:
                lines.append(f"- {fn}（读取失败）")
    lines.append("说明：自我改进请用 create_evolution 写入 evolutions/ 分支，勿修改原文件。")
    return "\n".join(lines)


def read_project_file(path, offset=0, limit=0):
    """读取鲸语自身源码（仅项目目录内白名单扩展名，只读）。

    offset/limit：按字符分页（大型文件如 main.py 320KB 需分页读取），
    limit=0 表示读取到 offset+80000 或文件尾。
    """
    p = os.path.abspath(os.path.expanduser(str(path or "")))
    base = os.path.abspath(PROJECT_DIR)
    # Windows 路径大小写不敏感：normcase 后比较，防合法路径被误拒
    if os.path.normcase(p) != os.path.normcase(base) and not os.path.normcase(p).startswith(
        os.path.normcase(base).rstrip("\\/") + os.sep
    ):
        return "权限拒绝：只能读取项目目录内的文件"
    if not os.path.isfile(p):
        return f"错误：文件不存在：{p}"
    p_lower = p.lower()
    if not any(p_lower.endswith(ext) for ext in PROJECT_READ_EXTS):
        return f"错误：不支持的文件类型（仅 {'/'.join(PROJECT_READ_EXTS)}）"
    try:
        with open(p, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        try:
            off = max(0, int(offset or 0))
            lim = max(0, int(limit or 0))
        except (TypeError, ValueError):
            off, lim = 0, 0
        total = len(content)
        if off >= total:
            return f"[已到达文件末尾] {p}（共 {total} 字符）"
        if lim > 0:
            chunk = content[off : off + lim]
        else:
            chunk = content[off : off + 80000]
        head = f"[{p} 第 {off}-{off + len(chunk)} 字符 / 共 {total} 字符]\n"
        return head + chunk + ("\n[已截断，可继续用 offset 读取后续]" if off + len(chunk) < total else "")
    except Exception as e:
        return f"错误：读取失败: {e}"


def create_evolution(name, files):
    """自我进化提案：写入 evolutions/<name>_<ts>/ 分支，绝不修改原文件。

    返回分支路径；EVOLUTION.md 缺失时自动生成基础说明（AI 应尽量在
    files 中包含完整的 EVOLUTION.md：改动内容/原因/风险/验证方式）。
    """
    name = re.sub(r'[\\/:*?"<>|]', "_", str(name or "evolution").strip())[:40] or "evolution"
    if not isinstance(files, list) or not files:
        return '错误：files 必须是非空数组 [{"path": "main.py", "content": "..."}]'
    if len(files) > 20:
        return "错误：文件数超过 20 上限"
    total_bytes = sum(len(str(f.get("content") or "")) for f in files if isinstance(f, dict))
    if total_bytes > 50 * 1024 * 1024:
        return "错误：提案内容超过 50MB 总上限"
    # 校验前置：非法路径/类型在创建任何目录前拒绝（避免空分支残留）
    for f in files[:20]:
        if not isinstance(f, dict):
            return "错误：files 元素必须是对象"
        rel = str(f.get("path") or "").strip().replace("\\", "/")
        if not rel or rel in (".", "..") or ".." in rel.split("/"):
            return f"错误：非法相对路径：{rel}"
        if not rel.endswith(EVO_WRITE_EXTS):
            return f"错误：不支持的文件类型：{rel}"
        branch_preview = os.path.join(EVOLUTIONS_DIR, "_preview")
        full = os.path.normpath(os.path.join(branch_preview, rel))
        if full != branch_preview and not full.startswith(branch_preview.rstrip("\\/") + os.sep):
            return f"错误：路径越界：{rel}"
    ts = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
    branch = os.path.join(EVOLUTIONS_DIR, f"{name}_{ts}")
    try:
        os.makedirs(branch, exist_ok=True)
    except Exception as e:
        return f"错误：创建分支失败: {e}"
    written = []
    has_md = False
    for f in files[:20]:
        rel = str(f.get("path") or "").strip().replace("\\", "/")
        content = f.get("content") or ""
        full = os.path.normpath(os.path.join(branch, rel))
        try:
            os.makedirs(os.path.dirname(full) or branch, exist_ok=True)
            with open(full, "w", encoding="utf-8") as fh:
                fh.write(content)
            written.append(rel)
            if rel == "EVOLUTION.md":
                has_md = True
        except Exception as e:
            return f"错误：写入 {rel} 失败: {e}"
    if not has_md:
        try:
            with open(os.path.join(branch, "EVOLUTION.md"), "w", encoding="utf-8") as fh:
                fh.write(
                    f"# 进化提案：{name}\n\n"
                    f"- 时间：{datetime.now():%Y-%m-%d %H:%M:%S}\n"
                    f"- 修改文件：{', '.join(written)}\n\n"
                    "## 说明\n（鲸语补充：改动内容、原因、风险与验证方式）\n"
                )
        except Exception:
            pass
    permissions.audit("create_evolution", name, ", ".join(written))
    return (
        f"自我进化提案已创建：{branch}\n"
        f"文件：{', '.join(written)}\n"
        "请在「工具 → 自我进化」中查看差异、采纳或忽略。"
    )


_SEARCH_EXTS = (
    ".py", ".md", ".txt", ".json", ".html", ".css", ".js", ".ts",
    ".yaml", ".yml", ".csv", ".log", ".ini", ".cfg", ".toml",
)
_SEARCH_SKIP_DIRS = {".git", "__pycache__", ".venv", "node_modules", "dist", "build"}


def search_local(path, query, max_results=20):
    """在允许目录内检索文本文件内容（只读）。"""
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not os.path.isdir(p):
        return f"错误：目录不存在：{p}"
    try:
        limit = max(1, min(200, int(max_results or 20)))
    except (TypeError, ValueError):
        limit = 20
    q = str(query or "").lower()
    if not q:
        return "错误：查询关键词为空"
    hits = []
    scanned = 0
    _MAX_SCAN_FILES = 2000  # 扫描预算：大工作区防单次工具调用卡分钟级
    try:
        for root, dirs, files in os.walk(p):
            dirs[:] = [d for d in dirs if d not in _SEARCH_SKIP_DIRS]
            for fn in files:
                if len(hits) >= limit or scanned >= _MAX_SCAN_FILES:
                    return _search_local_result(hits, scanned, limit, q)
                scanned += 1
                if not fn.lower().endswith(_SEARCH_EXTS):
                    continue
                full = os.path.join(root, fn)
                try:
                    if os.path.getsize(full) > 512 * 1024:
                        continue
                    with open(full, "r", encoding="utf-8", errors="replace") as f:
                        for ln in f:
                            if q in ln.lower():
                                rel = os.path.relpath(full, p)
                                hits.append(f"{rel}: {ln.strip()[:150]}")
                                if len(hits) >= limit:
                                    return _search_local_result(hits, scanned, limit, q)
                except Exception:
                    continue
    except Exception as e:
        return f"错误：检索失败: {e}"
    return _search_local_result(hits, scanned, limit, q)


def _search_local_result(hits, scanned, limit, query):
    """search_local 结果格式化（命中已满/扫描预算耗尽/正常结束共用出口）。"""
    if not hits:
        return f"未找到包含「{query}」的文件（已扫描 {scanned} 个文件）"
    note = ""
    if len(hits) >= limit or scanned >= 2000:
        note = f"\n[已限制显示前 {limit} 条]"
    return f"找到 {len(hits)} 个匹配文件：\n" + "\n".join(hits) + note


def verify_files(paths):
    """批量核验文件存在性与大小（写文件后自检，防幻觉）。

    相对路径基于工作目录解析；只读操作。
    """
    if not isinstance(paths, list) or not paths:
        return '错误：paths 必须是非空数组，如 ["src/main.py", "C:/x/app.py"]'
    if len(paths) > 30:
        return "错误：文件数超过 30 上限"
    lines = []
    exist = 0
    missing = 0
    for raw in paths[:30]:
        p = str(raw or "").strip()
        if not p:
            continue
        if not os.path.isabs(p) and permissions.WORKSPACE_DIR:
            # 相对路径限定在工作区内，防 ../ 越界探测工作区外文件
            ws = permissions.WORKSPACE_DIR.rstrip("\\/")
            full = os.path.normpath(os.path.join(permissions.WORKSPACE_DIR, p))
            if full != ws and not full.startswith(ws + os.sep):
                lines.append(f"❌ 越界路径被拒绝 {p}")
                missing += 1
                continue
            if os.path.exists(full):
                p = full
        elif os.path.isabs(p):
            # 绝对路径同样走权限判定：防探测磁盘任意文件的存在性与大小
            ok_abs, _ = permissions.check_filesystem(p, write=False)
            if not ok_abs:
                lines.append(f"❌ 越界路径被拒绝 {p}")
                missing += 1
                continue
        if os.path.isfile(p):
            try:
                size = os.path.getsize(p)
                lines.append(f"✅ 存在 {p}（{size} 字节）")
            except OSError:
                lines.append(f"✅ 存在 {p}（大小未知）")
            exist += 1
        else:
            lines.append(f"❌ 缺失 {p}")
            missing += 1
    lines.append(f"核验结果：{exist} 个存在 / {missing} 个缺失")
    return "\n".join(lines)


def create_doc(path, content, doc_type=""):
    """创建文档：.md/.html 原生；.docx 需 python-docx（可选）。"""
    ok, reason = permissions.check_filesystem(path, write=True)
    if not ok:
        return reason
    p = permissions.resolve(path)
    ext = (doc_type or "").lower() or os.path.splitext(p)[1].lstrip(".").lower()
    try:
        if ext == "docx":
            try:
                from docx import Document
            except ImportError:
                return "错误：生成 .docx 需要 python-docx：pip install python-docx"
            doc = Document()
            for para in (content or "").splitlines():
                if para.strip():
                    doc.add_paragraph(para)
            created = not os.path.exists(p)
            if not created:
                try:
                    shutil.copy2(p, p + ".bak")
                except Exception:
                    pass
            doc.save(p)
        else:
            if ext != "html":
                ext = "md"
            body = content or ""
            if ext == "html" and not body.lstrip().startswith("<"):
                body = (
                    "<!DOCTYPE html><html lang='zh'><head><meta charset='utf-8'>"
                    "<title>鲸语文档</title></head><body><pre>"
                    + body.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    + "</pre></body></html>"
                )
            created, _ = _atomic_write(p, body)
        permissions.audit("create_doc", p, f"{ext} {len(content or '')} 字符")
        return f"已创建文档 {p}（{'新建' if created else '覆盖并备份 .bak'}，{ext} 格式）"
    except Exception as e:
        return f"错误：文档创建失败: {e}"


def write_code_project(project_dir, files):
    """创建多文件代码工程：批量写文件（逐文件原子写 + 越界防护）。"""
    ok, reason = permissions.check_filesystem(project_dir, write=True)
    if not ok:
        return reason
    if not isinstance(files, list) or not files:
        return '错误：files 必须是非空数组 [{"path": "...", "content": "..."}]'
    if len(files) > 50:
        return "错误：文件数超过 50 上限"
    base = permissions.resolve(project_dir)
    created = []
    failed = []
    total = 0
    for f in files[:50]:
        if not isinstance(f, dict):
            failed.append(("?", "元素必须是对象"))
            continue
        rel = str(f.get("path") or "").strip().replace("\\", "/")
        content = f.get("content") or ""
        if not rel or rel in (".", "..") or ".." in rel.split("/"):
            failed.append((rel or "?", "非法相对路径"))
            continue
        # 与 write_file 同规则：按 UTF-8 字节校验（中文 3 字节/字）
        if len(str(content).encode("utf-8", "ignore")) > permissions.max_write_size():
            failed.append((rel, "内容超过大小限制"))
            continue
        full = os.path.normpath(os.path.join(base, rel))
        if full != base and not full.startswith(base.rstrip("\\/") + os.sep):
            failed.append((rel, "路径越界"))
            continue
        try:
            os.makedirs(os.path.dirname(full) or base, exist_ok=True)
            _atomic_write(full, content)
            if not os.path.exists(full):
                failed.append((rel, "写入后核验失败"))
                continue
            created.append(rel)
            total += len(content)
        except Exception as e:
            failed.append((rel, str(e)))
    if not created:
        return "错误：全部文件写入失败：" + "；".join(f"{r}({why})" for r, why in failed)
    permissions.audit("write_code_project", base, f"{len(created)} 个文件")
    lines = [f"已创建代码工程 {base}（{len(created)} 个文件，均核验存在）", "文件清单："]
    lines += ["· " + c for c in created]
    if failed:
        lines.append(f"⚠ 失败 {len(failed)} 个：")
        lines += ["· " + r + "：" + why for r, why in failed]
    lines.append(f"共 {len(created)} 个文件，{total} 字符")
    return "\n".join(lines)


def _playwright_ready():
    try:
        import playwright  # noqa: F401

        return True, ""
    except ImportError:
        return False, "浏览器操作需要安装 playwright：pip install playwright && playwright install chromium"


# ===== 桌面 RPA（P0）：pyautogui 鼠标键盘，操作任意桌面软件 =====
RPA_FAILSAFE = True  # 鼠标移到屏幕左上角时立即中断 RPA（pyautogui failsafe）


def _rpa_ready():
    try:
        import pyautogui  # noqa: F401
        return True, ""
    except ImportError:
        return False, "桌面 RPA 需要安装 pyautogui：pip install pyautogui"


def rpa_screen_size():
    """当前屏幕分辨率（RPA 坐标用）。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    try:
        import pyautogui
        w, h = pyautogui.size()
        return f"屏幕分辨率：{w} x {h}"
    except Exception as e:
        return f"错误：获取屏幕尺寸失败: {e}"


def rpa_click(x, y, button="left", clicks=1):
    """模拟鼠标点击。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    try:
        import pyautogui
        x = int(x)
        y = int(y)
        button = str(button or "left").strip().lower()
        clicks = max(1, min(5, int(clicks or 1)))
        if button not in ("left", "right", "middle"):
            return "错误：button 仅支持 left/right/middle"
        pyautogui.FAILSAFE = RPA_FAILSAFE
        pyautogui.click(x, y, button=button, clicks=clicks)
        permissions.audit("rpa_click", f"{x},{y}", f"{button} x{clicks}")
        return f"已点击 ({x}, {y})，{button} 键 x{clicks}"
    except Exception as e:
        return f"错误：RPA 点击失败: {e}"


def rpa_type(text, interval=0.02):
    """模拟键盘输入文本。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    if not str(text or ""):
        return "错误：text 必填"
    try:
        import pyautogui
        interval = max(0.0, min(0.2, float(interval or 0.02)))
        pyautogui.FAILSAFE = RPA_FAILSAFE
        pyautogui.typewrite(str(text), interval=interval)
        permissions.audit("rpa_type", "键盘输入", str(text)[:60])
        return f"已输入 {len(str(text))} 个字符"
    except Exception as e:
        return f"错误：RPA 输入失败: {e}"


def rpa_hotkey(keys):
    """模拟组合键，如 ctrl+c / alt+tab / ctrl+shift+esc。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    if not str(keys or "").strip():
        return "错误：keys 必填"
    try:
        import pyautogui
        seq = [str(k).strip().lower() for k in str(keys).replace(" ", "").split("+") if str(k).strip()]
        if not seq:
            return "错误：keys 格式应为 ctrl+c 或 alt+tab"
        pyautogui.FAILSAFE = RPA_FAILSAFE
        pyautogui.hotkey(*seq)
        permissions.audit("rpa_hotkey", "+".join(seq), "组合键")
        return f"已按下组合键 {'+'.join(seq)}"
    except Exception as e:
        return f"错误：RPA 组合键失败: {e}"


def rpa_move(x, y, duration=0.2):
    """移动鼠标到坐标。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    try:
        import pyautogui
        x, y = int(x), int(y)
        duration = max(0.0, min(2.0, float(duration or 0.2)))
        pyautogui.FAILSAFE = RPA_FAILSAFE
        pyautogui.moveTo(x, y, duration=duration)
        return f"鼠标已移动到 ({x}, {y})"
    except Exception as e:
        return f"错误：RPA 移动失败: {e}"


def rpa_scroll(clicks, x=None, y=None):
    """滚动鼠标滚轮（正数向上，负数向下）。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    try:
        import pyautogui
        n = max(-50, min(50, int(clicks or 0)))
        if x is not None and y is not None:
            pyautogui.scroll(n, x=int(x), y=int(y))
        else:
            pyautogui.scroll(n)
        permissions.audit("rpa_scroll", str(n), "滚轮")
        return f"已滚动 {n} 格"
    except Exception as e:
        return f"错误：RPA 滚动失败: {e}"


def rpa_screenshot(path=""):
    """截取当前屏幕保存为 PNG（不指定路径保存到工作区）。"""
    ok, hint = _rpa_ready()
    if not ok:
        return hint
    if str(path or "").strip():
        p = permissions.resolve(path)
    else:
        p = permissions.resolve(os.path.join(
            permissions.WORKSPACE_DIR or "", f"rpa_screen_{datetime.now():%Y%m%d_%H%M%S}.png"
        ))
    if not p:
        return "错误：截图路径无效"
    try:
        import pyautogui
        os.makedirs(os.path.dirname(p) or ".", exist_ok=True)
        img = pyautogui.screenshot()
        img.save(p)
        permissions.audit("rpa_screenshot", p, "屏幕截图")
        return f"已截屏保存至 {p}"
    except Exception as e:
        return f"错误：RPA 截屏失败: {e}"


# ===== 浏览器模式（由 main 注入，True=无头静默，False=有头弹出窗口可实时预览）=====
BROWSER_HEADLESS = True
BROWSER_PROFILE_DIR = None  # 由 main 注入（DATA_DIR/browser_profile，登录态持久化）

_BROWSER_LOCK = threading.RLock()  # 可重入：browser_navigate 持锁内调用 _get_browser_page
_BROWSER_PW = None
_BROWSER = None          # Browser 或 persistent BrowserContext
_BROWSER_PAGE = None     # 共享页面（多步操作保持状态/登录态）


def _browser_headless():
    return BROWSER_HEADLESS


def _get_browser_page():
    """获取共享浏览器页面（实例复用 + 登录态持久 + 页面状态保持）。"""
    global _BROWSER_PW, _BROWSER, _BROWSER_PAGE
    with _BROWSER_LOCK:
        if _BROWSER is None:
            from playwright.sync_api import sync_playwright

            _BROWSER_PW = sync_playwright().start()
            kwargs = {"headless": _browser_headless()}
            if BROWSER_PROFILE_DIR:
                os.makedirs(BROWSER_PROFILE_DIR, exist_ok=True)
                ctx = _BROWSER_PW.chromium.launch_persistent_context(BROWSER_PROFILE_DIR, **kwargs)
                _BROWSER = ctx
                _BROWSER_PAGE = ctx.pages[0] if ctx.pages else ctx.new_page()
            else:
                _BROWSER = _BROWSER_PW.chromium.launch(**kwargs)
                _BROWSER_PAGE = _BROWSER.new_page()
        return _BROWSER_PAGE


def close_browser():
    """关闭共享浏览器（main 退出时调用，登录态已持久化不受影响）。"""
    global _BROWSER_PW, _BROWSER, _BROWSER_PAGE
    with _BROWSER_LOCK:
        try:
            if _BROWSER is not None:
                _BROWSER.close()
        except Exception:
            pass
        try:
            if _BROWSER_PW is not None:
                _BROWSER_PW.stop()
        except Exception:
            pass
        _BROWSER_PW = _BROWSER = _BROWSER_PAGE = None


def _browser_goto(page, url):
    if str(page.url or "").strip("/") != str(url).strip("/"):
        page.goto(url, timeout=15000, wait_until="domcontentloaded")
        return True
    return False


def browser_navigate(url, action="open", selector="", text=""):
    """浏览器可视操作（Playwright 可选依赖，未安装时返回安装提示）。

    浏览器实例复用（连续操作共享同一页面）：click/type/fill/submit 等动作
    不重新导航，保留当前页面状态与登录态；open 才会跳转新页面。
    有头/无头跟随全局开关 BROWSER_HEADLESS。
    """
    ok, hint = _playwright_ready()
    if not ok:
        return hint
    action = (action or "open").lower()
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except Exception as e:
        return f"错误：playwright 初始化失败: {e}"
    # SSRF 防护：与 fetch_url 同规则（模型可控 URL 禁止 file:// 与内网/回环）
    err = _safe_url(url)
    if err:
        return f"错误：{err}"
    try:
        with _BROWSER_LOCK:  # playwright 非线程安全：浏览器操作串行化
            page = _get_browser_page()
            if action == "open":
                _browser_goto(page, url)
                return f"已打开：{page.title() or url}\n当前 URL: {page.url}"
            # 非 open 动作：确保在目标页（已在此页则保持状态，不重复导航）
            _browser_goto(page, url)
            if action == "get_text":
                if not selector:
                    return "错误：get_text 需要 selector"
                els = page.query_selector_all(selector)
                if not els:
                    return f"未找到匹配 {selector} 的元素"
                texts = [e.inner_text()[:500] for e in els[:10]]
                return "\n".join(f"· {t}" for t in texts)
            if action == "click":
                if not selector:
                    return "错误：click 需要 selector"
                page.click(selector, timeout=5000)
                page.wait_for_timeout(1000)
                return f"已点击 {selector}，当前 URL: {page.url}"
            if action in ("type", "fill"):
                if not selector:
                    return "错误：type/fill 需要 selector"
                page.fill(selector, text or "", timeout=5000)
                return f"已在 {selector} 输入文本"
            if action == "submit":
                if selector:
                    page.click(selector, timeout=5000)
                else:
                    page.keyboard.press("Enter")
                page.wait_for_timeout(1500)
                return f"已提交表单，当前 URL: {page.url}"
            if action == "select":
                if not selector:
                    return "错误：select 需要 selector"
                page.select_option(selector, text or "")
                return f"已在 {selector} 选择 {text}"
            return f"错误：未知动作 {action}（open/click/type/fill/submit/select/get_text）"
    except Exception as e:
        return f"错误：浏览器操作失败: {e}"


def web_screenshot(url, width=1280, height=800):
    """网页截图并保存到工作区（Playwright 可选依赖，复用共享浏览器）。"""
    ok, hint = _playwright_ready()
    if not ok:
        return hint
    if not permissions.WORKSPACE_DIR:
        return "错误：工作区未初始化"
    try:
        w = max(320, min(2560, int(width or 1280)))
        h = max(240, min(1920, int(height or 800)))
    except (TypeError, ValueError):
        w, h = 1280, 800
    err = _safe_url(url)
    if err:
        return f"错误：{err}"
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe = re.sub(r"[\\/:*?\"<>|]", "_", str(url))[:30]
    path = os.path.join(permissions.WORKSPACE_DIR, f"screenshot_{safe}_{ts}.png")
    try:
        with _BROWSER_LOCK:
            page = _get_browser_page()
            if str(page.url or "").strip("/") != str(url).strip("/"):
                page.goto(url, timeout=15000, wait_until="domcontentloaded")
            page.set_viewport_size({"width": w, "height": h})
            page.wait_for_timeout(1200)
            page.screenshot(path=path, full_page=False)
        permissions.audit("web_screenshot", url, path)
        return f"截图已保存：{path}"
    except Exception as e:
        return f"错误：截图失败: {e}"


def publish_draft(platform, title, content):
    """保存发布草稿到本地草稿箱（只建草稿不发布，双确认由审批流保证）。"""
    if not permissions.WORKSPACE_DIR:
        return "错误：工作区未初始化"
    drafts = os.path.join(permissions.WORKSPACE_DIR, "drafts")
    # 路径穿越防护：写前必须经权限模型判定（草稿箱必须在工作区内）
    ok, reason = permissions.check_filesystem(drafts, write=True)
    if not ok:
        return reason
    try:
        os.makedirs(drafts, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        # platform 与 title 一律清洗，杜绝 '..' / 分隔符穿越
        safe_platform = re.sub(r'[\\/:*?"<>|]', "_", str(platform or "draft")).strip(" .")[:40] or "draft"
        safe = re.sub(r'[\\/:*?"<>|]', "_", str(title or "草稿"))[:40] or "草稿"
        path = os.path.join(drafts, f"{safe_platform}_{safe}_{ts}.md")
        # 二次兜底：规范化后必须仍位于草稿箱内
        if os.path.normpath(path) != path or not path.startswith(os.path.normpath(drafts) + os.sep):
            return f"错误：非法路径被拦截：{path}"
        with open(path, "w", encoding="utf-8") as f:
            f.write(f"# {title}\n\n{content}")
        permissions.audit("publish_draft", path, safe_platform)
        return f"草稿已保存到本地草稿箱（未发布）：{path}\n正式发布请在平台后台操作。"
    except Exception as e:
        return f"错误：草稿保存失败: {e}"


# ============================================================================
# 新工具分类（v2 能力层）：任务调度 / 桌面通知 / 剪贴板 / 文件闭环 / 媒体感知
#   / 数据写入 / 知识库 RAG / 流程编排 / 洞察报告 / 任务检查点
# ============================================================================

# 注入（main 启动时设置）：
SCHEDULES_FILE = None        # DATA_DIR/schedules.json（与定时任务面板同文件）
KNOWLEDGE_INDEX_FILE = None  # DATA_DIR/knowledge_index.json
WORKFLOWS_FILE = None        # DATA_DIR/workflows.json
CHECKPOINT_FILE = None       # DATA_DIR/task_checkpoint.json
STATS_FILE = None            # DATA_DIR/stats.json
PATTERNS_FILE = None         # DATA_DIR/patterns.json（成功模式配方，run_workflow 的 recipe 步骤用）
IMAGE_GEN_BASE = None        # 图片生成端点（默认 = base_url）
IMAGE_GEN_KEY = None         # 图片生成 API Key（默认 = api_key）
IMAGE_GEN_MODEL = "gpt-image-1"
VISION_SELF_REVIEW = False   # 视觉自审（由 main 注入）：工具产出图片时自动审图
RSS_SOURCES_FILE = None      # DATA_DIR/rss_sources.json（RSS 订阅列表）
KV_CACHE_DIR = None          # DATA_DIR/kv_cache（diskcache 存储目录）
WEBDAV_CONFIG_FILE = None    # DATA_DIR/webdav_config.json（WebDAV 连接）
PLUGIN_PATHS = None          # 插件体系路径（plugins_dir/user_tools/prompts/workflows，main 注入）

SCHEDULES_LOCK = threading.Lock()  # 与 main 的定时任务面板共享（防并发覆盖）
_SEND_CALLBACK = None              # run_workflow：向主线程投递要发送的消息
_BUSY_PROVIDER = None              # run_workflow：查询是否正在生成


def set_send_callback(cb):
    global _SEND_CALLBACK
    _SEND_CALLBACK = cb


def set_busy_provider(cb):
    global _BUSY_PROVIDER
    _BUSY_PROVIDER = cb


# ---------- 任务调度工具（复用 main 的 cron 引擎与 schedules.json） ----------
# cron 字段校验/匹配由 shared.py 统一提供（CRON_RANGES / cron_field_ok）


def _load_schedules_plain():
    if not SCHEDULES_FILE or not os.path.exists(SCHEDULES_FILE):
        return []
    try:
        with open(SCHEDULES_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, list) else []
    except Exception:
        logging.exception("读取定时任务失败")
        return []


def _save_schedules_plain(schedules):
    if not SCHEDULES_FILE:
        return False
    try:
        os.makedirs(os.path.dirname(SCHEDULES_FILE) or ".", exist_ok=True)
        tmp = SCHEDULES_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(schedules, f, ensure_ascii=False, indent=2)
        os.replace(tmp, SCHEDULES_FILE)
        return True
    except Exception:
        logging.exception("保存定时任务失败")
        return False


def schedule_task(expr_type="cron", expr="", content="", action="message", name="", enabled=True, off_peak=False):
    """创建定时任务（与手动「定时任务」面板同文件同引擎，AI 可主动安排）。

    expr_type: cron（5 字段：分 时 日 月 周）/ time（HH:MM 每日一次）/ every（每 N 分钟）
    action: message（到点自动发送指令执行任务）/ notify（状态栏提醒）/ backup（项目备份）
            / workflow（到点自动运行 workflows.json 中的流程，content 为流程名）
    off_peak: 高峰错峰——触发时刻处于高峰时段（9-12 / 14-18）时自动顺延到
            最近空闲时段开始执行（官方峰谷定价：空闲价格仅为高峰一半）
    """
    expr = str(expr or "").strip()
    if not expr:
        return "错误：expr 必填（cron 表达式 / HH:MM / 分钟数）"
    if not SCHEDULES_FILE:
        return "错误：定时任务模块未初始化"
    s = {"enabled": bool(enabled), "action": str(action or "message"), "last": "", "last_run": 0}
    if expr_type == "time":
        if not re.match(r"^\d{1,2}:\d{2}$", expr):
            return "错误：time 格式应为 HH:MM（如 09:00）"
        hh, _, mm = expr.partition(":")
        if not 0 <= int(hh) <= 23 or not 0 <= int(mm) <= 59:
            return "错误：time 时间非法（小时 0-23，分钟 0-59）"
        s["time"] = expr
    elif expr_type == "every":
        try:
            n = int(expr)
        except (TypeError, ValueError):
            return "错误：every 需要整数分钟数"
        if not 1 <= n <= 1440:
            return "错误：every 应在 1-1440 分钟之间"
        s["every"] = n
    else:
        fields = expr.split()
        if len(fields) != 5:
            return "错误：cron 需 5 个字段：分 时 日 月 周（如 30 9 * * 1）"
        if not all(cron_field_ok(f, i) for i, f in enumerate(fields)):
            return (
                "错误：cron 字段非法（值域：分 0-59，时 0-23，日 1-31，月 1-12，周 1-7；"
                "仅支持数字、*、,、-、/）"
            )
        s["cron"] = expr
    if str(action or "") not in ("message", "notify", "backup", "workflow"):
        return "错误：action 仅支持 message / notify / backup / workflow"
    if str(action) in ("message", "notify") and not str(content or "").strip():
        return "错误：message / notify 动作需要 content 内容"
    if str(action) == "workflow" and not str(content or "").strip():
        return "错误：workflow 动作需要流程名称（workflows.json 中的流程名）"
    if off_peak:
        s["off_peak"] = True
    if str(name or "").strip():
        s["name"] = str(name).strip()[:40]
    if str(content or "").strip():
        s["text"] = str(content).strip()[:2000]
    s["id"] = f"s{int(time.time() * 1000)}"
    with SCHEDULES_LOCK:
        schedules = _load_schedules_plain()
        schedules.append(s)
        _save_schedules_plain(schedules)
    when = expr
    permissions.audit("schedule_task", s["id"], f"{expr_type}:{expr} -> {action}")
    return f"已创建定时任务（id={s['id']}）：{when} 执行「{s.get('text', '')[:60]}」"


def list_schedules():
    """列出全部定时任务（含 id/时间/动作/内容/状态）。"""
    with SCHEDULES_LOCK:
        schedules = _load_schedules_plain()
    if not schedules:
        return "当前没有定时任务"
    act_map = {"message": "发指令", "notify": "提醒", "backup": "备份", "workflow": "流程"}
    lines = [f"共 {len(schedules)} 个定时任务："]
    for i, s in enumerate(schedules, 1):
        act = act_map.get(str(s.get("action") or "message"), str(s.get("action")))
        if s.get("cron"):
            when = f"cron:{s['cron']}"
        elif s.get("every"):
            when = f"每{s['every']}分钟"
        else:
            when = f"每日 {s.get('time', '')}"
        status = "启用" if s.get("enabled") else "停用"
        name = str(s.get("name") or "").strip()
        content = str(s.get("text") or "").strip()
        lines.append(
            f"{i}. [{status}] id={s.get('id', '-')} {name} | {when} | {act} | {content[:60]}"
        )
    return "\n".join(lines)


def cancel_schedule(target=""):
    """取消定时任务（按 id 或名称）。"""
    t = str(target or "").strip()
    if not t:
        return "错误：target 必填（任务 id 或名称，可用 list_schedules 查看）"
    with SCHEDULES_LOCK:
        schedules = _load_schedules_plain()
        if not schedules:
            return "当前没有定时任务"
        kept = []
        removed = None
        for s in schedules:
            sid = str(s.get("id") or "")
            sname = str(s.get("name") or "")
            if (sid and sid == t) or (sname and sname == t):
                removed = s
            else:
                kept.append(s)
        if removed is None:
            return f"错误：未找到定时任务：{t}（可用 list_schedules 查看）"
        _save_schedules_plain(kept)
    permissions.audit("cancel_schedule", t, "removed")
    return f"已取消定时任务：{removed.get('name') or t}（{removed.get('cron') or removed.get('every') or removed.get('time', '')}）"


# ---------- 桌面通知（Windows Toast，零依赖） ----------
# 占位符用 @TITLE@/@BODY@ 而非 $title/$body：用户内容若含字面 "$body" 会被
# 顺序 replace 二次替换污染脚本（$title 先替换成含 "$body" 的内容时同样被污染）。
_NOTIFY_PS = r"""
$ErrorActionPreference='Stop'
Add-Type -AssemblyName System.Runtime.WindowsRuntime
$asTaskGeneric = ([System.WindowsRuntimeSystemExtensions].GetMethods() | Where-Object { $_.Name -eq 'AsTask' -and $_.GetParameters().Count -eq 1 -and $_.GetParameters()[0].ParameterType.Name -eq 'IAsyncOperation`1' })[0]
function Await($WinRtTask, $ResultType) {
    $asTask = $asTaskGeneric.MakeGenericMethod($ResultType)
    $netTask = $asTask.Invoke($null, @($WinRtTask))
    $netTask.Wait(-1) | Out-Null
    $netTask.Result
}
[Windows.UI.Notifications.ToastNotificationManager, Windows.UI.Notifications, ContentType=WindowsRuntime] | Out-Null
[Windows.UI.Notifications.ToastNotification, Windows.UI.Notifications, ContentType=WindowsRuntime] | Out-Null
[Windows.Data.Xml.Dom.XmlDocument, Windows.Data.Xml.Dom.XmlDocument, ContentType=WindowsRuntime] | Out-Null
$template = [Windows.UI.Notifications.ToastNotificationManager]::GetTemplateContent([Windows.UI.Notifications.ToastTemplateType]::ToastText02)
$textNodes = $template.GetElementsByTagName("text")
$textNodes.Item(0).AppendChild($template.CreateTextNode('@TITLE@')) | Out-Null
$textNodes.Item(1).AppendChild($template.CreateTextNode('@BODY@')) | Out-Null
$toast = New-Object Windows.UI.Notifications.ToastNotification $template
[Windows.UI.Notifications.ToastNotificationManager]::CreateToastNotifier("鲸语 WhaleTalk").Show($toast)
"""


def notify_desktop(title="鲸语提醒", text=""):
    """Windows 桌面 Toast 通知（离线可用，任务完成/定时任务触发时使用）。"""
    if not str(text or "").strip():
        return "错误：text 必填"
    title = str(title or "鲸语提醒")[:60]
    body = str(text).strip()[:300]
    try:
        import tempfile

        fd, ps_path = tempfile.mkstemp(suffix=".ps1")
        os.close(fd)
        try:
            title_quoted = "'" + str(title).replace("'", "''") + "'"
            body_quoted = "'" + body.replace("'", "''") + "'"
            # 先替换标题为哨兵，再替换正文，最后回填标题：防止标题/正文互相包含对方占位符
            title_sentinel = "__WHALETALK_TITLE__"
            script = _NOTIFY_PS.replace("@TITLE@", title_sentinel)
            script = script.replace("@BODY@", body_quoted)
            script = script.replace(title_sentinel, title_quoted)
            with open(ps_path, "w", encoding="utf-8-sig") as f:
                f.write(script)
            proc = subprocess.run(
                ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-File", ps_path],
                capture_output=True, text=True, timeout=15,
                encoding="utf-8", errors="replace",
            )
            if proc.returncode != 0:
                # Toast 不可用（老系统/受限环境）时兜底为提示音
                try:
                    import winsound

                    winsound.Beep(880, 250)
                    winsound.Beep(660, 250)
                except Exception:
                    pass
                return f"通知显示失败（已播放提示音）：{(proc.stderr or '')[:150]}"
            return f"已发送桌面通知：{title}"
        finally:
            try:
                os.remove(ps_path)
            except OSError:
                pass
    except Exception as e:
        return f"错误：通知失败: {e}"


# ---------- 剪贴板读写（Win32，可在任意线程调用） ----------
def _win_clipboard_get():
    import ctypes

    try:
        user32 = ctypes.windll.user32
        kernel32 = ctypes.windll.kernel32
        CF_UNICODETEXT = 13
        if not user32.OpenClipboard(None):
            return None
        try:
            h = user32.GetClipboardData(CF_UNICODETEXT)
            if not h:
                return None
            size = kernel32.GlobalSize(h)
            lock = kernel32.GlobalLock(h)
            if not lock:
                return None
            try:
                buf = ctypes.create_string_buffer(size)
                ctypes.memmove(buf, lock, size)
            finally:
                kernel32.GlobalUnlock(h)
            raw = buf.raw
            return raw.decode("utf-16-le", errors="replace").lstrip("\ufeff").rstrip("\x00")
        finally:
            user32.CloseClipboard()
    except Exception:
        logging.exception("剪贴板读取失败")
        return None


def _win_clipboard_set(text):
    import ctypes

    try:
        user32 = ctypes.windll.user32
        kernel32 = ctypes.windll.kernel32
        CF_UNICODETEXT = 13
        data = (str(text) + "\x00").encode("utf-16-le")
        if not user32.OpenClipboard(None):
            return False
        try:
            user32.EmptyClipboard()
            h = kernel32.GlobalAlloc(0x0002, len(data))  # GMEM_MOVEABLE
            if not h:
                return False
            try:
                lock = kernel32.GlobalLock(h)
                if lock:
                    ctypes.memmove(lock, data, len(data))
                    kernel32.GlobalUnlock(h)
                ok = bool(user32.SetClipboardData(CF_UNICODETEXT, h))
                if not ok:
                    kernel32.GlobalFree(h)  # 失败时系统未接管，释放句柄防泄漏
                return ok
            except Exception:
                kernel32.GlobalFree(h)
                return False
        finally:
            user32.CloseClipboard()
    except Exception:
        logging.exception("剪贴板写入失败")
        return False


def clipboard_get():
    """读取用户剪贴板文本（敏感操作，走审批闸门默认需确认）。"""
    text = _win_clipboard_get()
    if text is None:
        return "错误：无法访问剪贴板"
    if not text.strip():
        return "剪贴板为空"
    return f"[剪贴板内容（{len(text)} 字符）]\n{text}"


def clipboard_set(text):
    """把内容写入剪贴板（用户可直接粘贴使用）。"""
    if not str(text or "").strip():
        return "错误：text 必填"
    t = str(text)
    if len(t) > 500000:
        t = t[:500000]
    if _win_clipboard_set(t):
        return f"已写入剪贴板（{len(t)} 字符）"
    return "错误：剪贴板写入失败"


# ---------- 文件闭环：删除（回收站优先）/ 压缩 / 解压 / 批量重命名 ----------
def _recycle_path(p):
    """把文件/目录移入 Windows 回收站（SHFileOperationW，可恢复）。"""
    import ctypes
    from ctypes import wintypes

    FO_DELETE = 3
    FOF_ALLOWUNDO = 0x40
    FOF_NOCONFIRMATION = 0x10
    FOF_SILENT = 0x4

    class SHFILEOPSTRUCT(ctypes.Structure):
        _fields_ = [
            ("hwnd", wintypes.HWND),
            ("wFunc", ctypes.c_uint),
            ("pFrom", wintypes.LPCWSTR),
            ("pTo", wintypes.LPCWSTR),
            ("fFlags", ctypes.c_ushort),
            ("fAnyOperationsAborted", wintypes.BOOL),
            ("hNameMappings", ctypes.c_void_p),
            ("lpszProgressTitle", wintypes.LPCWSTR),
        ]

    op = SHFILEOPSTRUCT(
        0, FO_DELETE, p + "\x00\x00", None,
        FOF_ALLOWUNDO | FOF_NOCONFIRMATION | FOF_SILENT, 0, None, None,
    )
    return ctypes.windll.shell32.SHFileOperationW(ctypes.byref(op)) == 0


def delete_file(path, permanent=False):
    """删除文件/目录。默认移入回收站（可恢复）；permanent=True 才物理删除。"""
    p = permissions.resolve(path)
    if not p or not os.path.exists(p):
        return f"错误：路径不存在：{path}"
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    try:
        if permanent:
            if os.path.isdir(p):
                shutil.rmtree(p)
            else:
                os.remove(p)
            permissions.audit("delete_file", p, "permanent")
            return f"已物理删除：{p}"
        if os.name == "nt" and _recycle_path(p):
            permissions.audit("delete_file", p, "recycle")
            return f"已移入回收站（可恢复）：{p}"
        if os.path.isdir(p):
            shutil.rmtree(p)
        else:
            os.remove(p)
        permissions.audit("delete_file", p, "permanent(fallback)")
        return f"已删除：{p}"
    except Exception as e:
        return f"错误：删除失败: {e}"


def archive_files(paths, output):
    """把多个文件/目录打包为 zip（工作区内，自动建目录，跳过 .git/__pycache__ 等）。"""
    if not isinstance(paths, list) or not paths:
        return "错误：paths 必须是非空数组（文件/目录路径列表）"
    if not str(output or "").strip():
        return "错误：output 必填"
    import zipfile

    out = permissions.resolve(output)
    if not out:
        return "错误：输出路径无效"
    if not out.lower().endswith(".zip"):
        out += ".zip"
    ok, reason = permissions.check_filesystem(out, write=True)
    if not ok:
        return reason
    resolved = []
    for raw in paths[:50]:
        p = permissions.resolve(raw)
        if not p or not os.path.exists(p):
            return f"错误：路径不存在：{raw}"
        ok, reason = permissions.check_filesystem(p, write=False)
        if not ok:
            return reason
        resolved.append(p)
    try:
        os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
        count = 0
        total_bytes = 0
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
            for p in resolved:
                if os.path.isdir(p):
                    for root, dirs, files in os.walk(p):
                        dirs[:] = [d for d in dirs if d not in _ARCHIVE_SKIP_DIRS]
                        for fn in files:
                            full = os.path.join(root, fn)
                            if count >= EXTRACT_MAX_ENTRIES:
                                raise ValueError(f"文件数超过上限（{EXTRACT_MAX_ENTRIES}）")
                            fsize = os.path.getsize(full)
                            if total_bytes + fsize > EXTRACT_MAX_TOTAL_BYTES:
                                raise ValueError("总大小超过打包上限")
                            zf.write(full, os.path.relpath(full, os.path.dirname(p)))
                            count += 1
                            total_bytes += fsize
                else:
                    if count >= EXTRACT_MAX_ENTRIES:
                        raise ValueError(f"文件数超过上限（{EXTRACT_MAX_ENTRIES}）")
                    fsize = os.path.getsize(p)
                    if total_bytes + fsize > EXTRACT_MAX_TOTAL_BYTES:
                        raise ValueError("总大小超过打包上限")
                    zf.write(p, os.path.basename(p))
                    count += 1
                    total_bytes += fsize
        size = os.path.getsize(out)
        permissions.audit("archive_files", out, f"{count} 个文件")
        return f"已打包 {count} 个文件到 {out}（{size / 1024:.1f} KB）"
    except Exception as e:
        try:
            if os.path.exists(out):
                os.remove(out)
        except OSError:
            pass
        return f"错误：打包失败: {e}"


def extract_archive(path, dest_dir):
    """解压 zip 到目标目录（zip-slip 越界防护）。"""
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：压缩包不存在：{path}"
    ok, reason = permissions.check_filesystem(p, write=False)
    if not ok:
        return reason
    dest = permissions.resolve(dest_dir) if str(dest_dir or "").strip() else None
    if not dest:
        return "错误：dest_dir 必填"
    ok, reason = permissions.check_filesystem(dest, write=True)
    if not ok:
        return reason
    try:
        os.makedirs(dest, exist_ok=True)
        base = os.path.normpath(dest)
        ext = os.path.splitext(p)[1].lower()
        if ext == ".zip":
            import zipfile
            count = 0
            with zipfile.ZipFile(p) as zf:
                infos = zf.infolist()
                if len(infos) > EXTRACT_MAX_ENTRIES:
                    return f"错误：压缩包条目数超过上限（{EXTRACT_MAX_ENTRIES}），已中止"
                total_size = 0
                for info in infos:
                    target = os.path.normpath(os.path.join(base, info.filename))
                    if not (target == base or target.startswith(base + os.sep)):
                        return f"错误：压缩包含越界条目，已中止：{info.filename}"
                    if info.file_size > EXTRACT_MAX_SINGLE_BYTES:
                        return f"错误：压缩包单文件超过大小上限：{info.filename}"
                    total_size += info.file_size
                    if total_size > EXTRACT_MAX_TOTAL_BYTES:
                        return f"错误：压缩包总解压大小超过上限，已中止"
                for info in infos:
                    zf.extract(info, dest)
                    count += 1
            permissions.audit("extract_archive", dest, f"{count} 个条目")
            return f"已解压 {count} 个条目到 {dest}"
        if ext in (".tar", ".gz", ".tgz"):
            import tarfile
            with tarfile.open(p) as tf:
                members = tf.getmembers()
                if len(members) > EXTRACT_MAX_ENTRIES:
                    return f"错误：压缩包条目数超过上限（{EXTRACT_MAX_ENTRIES}），已中止"
                total_size = sum(m.size for m in members)
                if total_size > EXTRACT_MAX_TOTAL_BYTES:
                    return f"错误：压缩包总解压大小超过上限，已中止"
                for m in members:
                    target = os.path.normpath(os.path.join(base, m.name))
                    if not (target == base or target.startswith(base + os.sep)):
                        return f"错误：压缩包含越界条目，已中止：{m.name}"
                tf.extractall(dest)
            permissions.audit("extract_archive", dest, f"{len(members)} 个条目")
            return f"已解压 {len(members)} 个条目到 {dest}"
        if ext == ".7z":
            import py7zr
            with py7zr.SevenZipFile(p, "r") as z:
                names = z.getnames()
                if len(names) > EXTRACT_MAX_ENTRIES:
                    return f"错误：压缩包条目数超过上限（{EXTRACT_MAX_ENTRIES}），已中止"
                for name in names:
                    target = os.path.normpath(os.path.join(base, name))
                    if not (target == base or target.startswith(base + os.sep)):
                        return f"错误：压缩包含越界条目，已中止：{name}"
                z.extractall(dest)
            permissions.audit("extract_archive", dest, f"{len(names)} 个条目")
            return f"已解压 {len(names)} 个条目到 {dest}"
        if ext == ".rar":
            import rarfile
            with rarfile.RarFile(p) as rf:
                infos = rf.infolist()
                if len(infos) > EXTRACT_MAX_ENTRIES:
                    return f"错误：压缩包条目数超过上限（{EXTRACT_MAX_ENTRIES}），已中止"
                total_size = 0
                for info in infos:
                    target = os.path.normpath(os.path.join(base, info.filename))
                    if not (target == base or target.startswith(base + os.sep)):
                        return f"错误：压缩包含越界条目，已中止：{info.filename}"
                    if info.file_size > EXTRACT_MAX_SINGLE_BYTES:
                        return f"错误：压缩包单文件超过大小上限：{info.filename}"
                    total_size += info.file_size
                    if total_size > EXTRACT_MAX_TOTAL_BYTES:
                        return f"错误：压缩包总解压大小超过上限，已中止"
                rf.extractall(dest)
            permissions.audit("extract_archive", dest, f"{len(infos)} 个条目")
            return f"已解压 {len(infos)} 个条目到 {dest}"
        return f"错误：不支持的压缩格式 {ext or '（无扩展名）'}（支持 zip/tar/gz/7z/rar）"
    except ImportError as e:
        return f"错误：解压该格式需要额外依赖：{e}"
    except Exception as e:
        return f"错误：解压失败: {e}"


def batch_rename(directory, pattern, replacement, dry_run=False):
    """批量重命名：把文件名中的 pattern 全部替换为 replacement（含扩展名）。"""
    d = permissions.resolve(directory)
    if not d or not os.path.isdir(d):
        return f"错误：目录不存在：{directory}"
    ok, reason = permissions.check_filesystem(d, write=True)
    if not ok:
        return reason
    pattern = str(pattern or "")
    if not pattern:
        return "错误：pattern 必填"
    replacement = str(replacement or "")
    try:
        renamed = []
        for fn in sorted(os.listdir(d)):
            if pattern in fn:
                new = fn.replace(pattern, replacement)
                if new == fn:
                    continue
                src = os.path.join(d, fn)
                dst = os.path.join(d, new)
                if os.path.exists(dst):
                    continue
                if not dry_run:
                    os.rename(src, dst)
                renamed.append(f"{fn} → {new}")
        if not renamed:
            return f"目录内没有包含「{pattern}」的文件名"
        note = "（预览，未实际重命名）" if dry_run else ""
        if not dry_run:
            permissions.audit("batch_rename", d, f"{len(renamed)} 个文件")
        return f"共 {len(renamed)} 个文件{note}：\n" + "\n".join(renamed[:20])
    except Exception as e:
        return f"错误：批量重命名失败: {e}"


# ---------- 媒体感知：图片理解 / 屏幕截图 / 语音识别 ----------
def image_understand(path, question=""):
    """用多模态模型理解图片（本地文件或 http(s) 图片 URL）。

    自动适配视觉模型：当前客户端模型不支持图片时，自动改用
    deepseek-v4-flash-vision-exp（同一 API Key / 端点），无需手动切换。
    """
    if not str(path or "").strip():
        return "错误：path 必填"
    import base64

    is_url = str(path).strip().lower().startswith(("http://", "https://"))
    if is_url:
        err = _safe_url(path)
        if err:
            return f"错误：{err}"
    else:
        p = permissions.resolve(path)
        if not p or not os.path.isfile(p):
            return f"错误：图片不存在：{path}"
        ok, reason = permissions.check_filesystem(p, write=False)
        if not ok:
            return reason
        try:
            if os.path.getsize(p) > 32 * 1024 * 1024:
                return "错误：图片超过 32MB，请先用 image_process 压缩"
        except OSError:
            pass
    try:
        if is_url:
            try:
                # stream 边读边断：URL 图片大小不可信，防恶意/超大图全量进内存
                with _safe_stream("GET", path, timeout=20) as resp:
                    resp.raise_for_status()
                    img_buf = b""
                    truncated = False
                    for chunk in resp.iter_bytes(64 * 1024):
                        img_buf += chunk
                        if len(img_buf) > 32 * 1024 * 1024:
                            truncated = True
                            break
                if truncated:
                    return "错误：图片下载超过 32MB，请先用 image_process 压缩"
                b64 = base64.b64encode(img_buf).decode("ascii")
                # 格式按文件实际内容（魔数）识别，而非声明的 MIME
                mime = _detect_image_mime(img_buf[:16])
            except Exception as e:
                return f"错误：图片下载失败: {e}"
        else:
            with open(p, "rb") as f:
                raw = f.read()
            b64 = base64.b64encode(raw).decode("ascii")
            # 格式按文件实际内容（魔数）识别，而非文件名/扩展名
            mime = _detect_image_mime(raw[:16])
        client = _CLIENT_HOLDER.get("client")
        if client is None:
            return "错误：没有可用客户端（请先完成一次对话建立连接）"
        model, switched = client.model, False
        if not is_vision_model(model):
            model = VISION_MODEL
            switched = True
        resp = client.client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": str(question or "请描述这张图片的内容")},
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                    ],
                }
            ],
            max_tokens=2048,
            stream=False,
            timeout=120.0,
        )
        out = (resp.choices[0].message.content or "").strip()
        if not out:
            return "模型未返回内容（当前模型可能不支持图片输入，可配置支持视觉的模型端点）"
        if switched:
            out += f"\n\n（注：当前模型不支持图片，已自动改用视觉模型 {VISION_MODEL}）"
        return out
    except Exception as e:
        return f"错误：图片理解失败: {e}（当前模型可能不支持视觉输入，可切换支持视觉的模型端点）"


def screen_capture(path="", area=""):
    """截取当前屏幕保存到工作区（默认全屏；area 形如 left,top,right,bottom）。"""
    if str(path or "").strip():
        out = permissions.resolve(path)
        if not out:
            return "错误：输出路径无效"
        if not out.lower().endswith(".png"):
            out += ".png"
        ok, reason = permissions.check_filesystem(out, write=True)
        if not ok:
            return reason
    else:
        base = os.path.join(permissions.WORKSPACE_DIR or ".", "screenshots")
        ok, reason = permissions.check_filesystem(base, write=True)
        if not ok:
            return reason
        try:
            os.makedirs(base, exist_ok=True)
        except Exception:
            pass
        out = os.path.join(base, f"screen_{datetime.now():%Y%m%d_%H%M%S}.png")
    try:
        from PIL import ImageGrab

        bbox = None
        if str(area or "").strip():
            try:
                parts = [int(x.strip()) for x in str(area).split(",")]
                if len(parts) == 4:
                    bbox = tuple(parts)
            except (TypeError, ValueError):
                bbox = None
        img = ImageGrab.grab(bbox=bbox)
        img.save(out, "PNG")
        size = os.path.getsize(out)
        permissions.audit("screen_capture", out, f"{size} 字节")
        return f"已截屏保存至 {out}（{size / 1024:.0f} KB），可用 ocr_image / image_understand 分析"
    except Exception as e:
        return f"错误：截屏失败: {e}"


# ---------- 视觉 Agent 能力（多模态闭环）：看图 / 自审 / 批量 ----------
# 图片路径匹配：优先绝对路径（允许路径含空格，遇括号/引号/换行即止），
# 其次相对路径（单 token，不含空格）。
_IMAGE_FILE_PATH_RE = re.compile(
    r"(?:[A-Za-z]:[\\/][^（）()\"'\r\n]*?|"
    r"[\\/][^（）()\"'\r\n]*?|"
    r"[^\s（）()\"'\r\n]+)\.(?:png|jpe?g|gif|webp)",
    re.I,
)
_IMAGE_PATH_TRAIL = "，。、,;；：:()（）\"'` \t\r\n"


def _extract_image_path(text):
    """从工具结果文本中提取已存在的图片文件路径（自审用，找不到返回 None）。

    兼容路径含空格与中文路径；候选串尾的标点会被剥离后再校验存在性。
    """
    text = str(text or "")
    ws = permissions.WORKSPACE_DIR
    for raw in re.findall(_IMAGE_FILE_PATH_RE, text):
        for cand in (raw, raw.rstrip(_IMAGE_PATH_TRAIL)):
            if not cand:
                continue
            if os.path.isfile(cand):
                return cand
            # 相对候选 → 相对工作区再试（工具一般返回绝对路径，这是兜底）
            if ws and not os.path.isabs(cand):
                joined = os.path.join(ws, cand)
                if os.path.isfile(joined):
                    return joined
    # 兜底：按空白拆分逐 token 校验（防正则漏网，如路径前带空格/标点）
    for tok in re.split(r"\s+", text):
        tok = tok.rstrip(_IMAGE_PATH_TRAIL)
        if not tok:
            continue
        if re.search(r"\.(?:png|jpe?g|gif|webp)$", tok, re.I):
            if os.path.isfile(tok):
                return tok
            if ws and not os.path.isabs(tok) and os.path.isfile(os.path.join(ws, tok)):
                return os.path.join(ws, tok)
    return None


# 产出图片文件、可触发视觉自审的工具（VISION_SELF_REVIEW 开启时自动审图）
_IMAGE_PRODUCING_TOOLS = frozenset({
    "image_generate",
    "chart_data",
    "screen_capture",
    "rpa_screenshot",
    "web_screenshot",
    "image_process",
    "qrcode",
})


def _capture_screen_png(area=""):
    """截取当前屏幕到工作区截图目录，返回 PNG 路径；失败返回 None。"""
    try:
        from PIL import ImageGrab

        bbox = None
        if str(area or "").strip():
            try:
                parts = [int(x.strip()) for x in str(area).split(",")]
                if len(parts) == 4:
                    bbox = tuple(parts)
            except (TypeError, ValueError):
                bbox = None
        img = ImageGrab.grab(bbox=bbox)
        base = os.path.join(permissions.WORKSPACE_DIR or ".", "screenshots")
        os.makedirs(base, exist_ok=True)
        out = os.path.join(base, f"screen_see_{datetime.now():%Y%m%d_%H%M%S}.png")
        img.save(out, "PNG")
        permissions.audit("screen_see", out, "屏幕截图（视觉解读）")
        return out
    except Exception:
        return None


def screen_see(question="", area=""):
    """截图并让视觉模型解读当前屏幕（一步完成 截图+看图）。

    RPA/浏览器操作后自查首选：看清界面后决定下一步操作（点击/输入/验证）。
    """
    q = str(question or "请描述当前屏幕内容，重点关注界面元素、按钮、文字与状态。").strip()
    try:
        from PIL import ImageGrab  # noqa: F401  # 提前校验依赖，给出明确安装提示
    except ImportError:
        return "错误：屏幕截图需要 Pillow，请先安装：pip install Pillow"
    path = _capture_screen_png(area)
    if not path:
        return "错误：屏幕截图失败"
    return image_understand(path, question=q)


def chart_read(path, question=""):
    """图表截图 → 结构化数据 + 解读（折线/柱状/饼图/散点等）。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    q = str(question or "").strip()
    if q:
        q += "；"
    q += (
        "请解读这张图表：1) 标题与图表类型；2) 坐标轴/图例/数据点，尽量精确给出数值；"
        "3) 关键趋势与结论。数据可用 Markdown 表格输出时请表格化。"
    )
    return image_understand(path, question=q)


def screenshot_to_html(path, out_path=""):
    """UI/网页截图 → 还原为 HTML+CSS 页面（可保存到文件）。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    q = (
        "请把这张截图还原成等价的 HTML+CSS 页面：像素级还原布局、配色、文字、间距与元素位置，"
        "输出完整可用的 HTML（CSS 内联，<html> 到 </html> 全量代码），只输出代码，不要解释。"
    )
    result = image_understand(path, question=q)
    if str(out_path or "").strip():
        # 视觉调用失败（错误文案）时绝不写入目标文件，原样返回错误
        if not result or str(result).startswith("错误"):
            return result or "错误：图片理解未返回内容"
        out = permissions.resolve(out_path)
        if not out:
            return "错误：输出路径无效"
        if not out.lower().endswith(".html"):
            out += ".html"
        ok, reason = permissions.check_filesystem(out, write=True)
        if not ok:
            return reason
        code = re.sub(r"^```(?:html|htm)?\s*|\s*```$", "", result.strip(), flags=re.I)
        with open(out, "w", encoding="utf-8") as f:
            f.write(code)
        permissions.audit("screenshot_to_html", out, str(path)[:80])
        return f"已根据截图生成 HTML 保存至 {out}\n\n{result}"
    return result


def debug_screenshot(path):
    """报错/异常截图 → 识别错误并给出诊断与修复建议。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    q = (
        "这是报错/异常截图。请：1) 识别错误类型与关键信息（报错文案、错误码、行号、堆栈线索）；"
        "2) 分析可能原因；3) 给出具体的修复建议（需要时可提到相关文件/函数）。"
    )
    return image_understand(path, question=q)


def scan_read(path, question=""):
    """扫描件/文档图片读取（图表、公式、手写、印刷体混排）。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    q = str(question or "").strip()
    if q:
        q += "；"
    q += (
        "这是扫描件/文档图片。请完整提取其中的文字、图表、公式与数据，保持原有结构，"
        "用 Markdown 呈现；手写内容按可辨识程度尽量转写，不确定处标注。"
    )
    return image_understand(path, question=q)


def image_batch(folder, question="", pattern="*.png", max=100):
    """批量视觉分析文件夹内图片：逐张理解后汇总报告（小并发）。"""
    if not str(folder or "").strip():
        return "错误：folder 必填"
    base = permissions.resolve(folder)
    if not base or not os.path.isdir(base):
        return f"错误：目录不存在：{folder}"
    ok, reason = permissions.check_filesystem(base, write=False)
    if not ok:
        return reason
    try:
        m = int(max or 100)
    except (TypeError, ValueError):
        m = 100
    limit = 1 if m < 1 else 200 if m > 200 else m
    # 防路径穿越：pattern 含 .. 等分隔符时，glob 可能越过允许目录返回外部文件。
    # 收集后用规范化路径强校验「必须位于 base 之内」，越界文件一律丢弃。
    import glob

    base_norm = os.path.normpath(base)
    files = []
    for f in sorted(glob.glob(os.path.join(base, str(pattern or "*.png")))):
        try:
            inside = os.path.commonpath([base_norm, os.path.normpath(f)]) == base_norm
        except ValueError:
            inside = False
        if inside and os.path.isfile(f) and f.lower().endswith(IMAGE_EXTENSIONS):
            files.append(f)
    if not files:
        return f"错误：目录 {base} 内没有匹配「{pattern}」的图片"
    files = files[:limit]
    q = str(question or "请描述这张图片的主要内容，并用一句话概括。").strip()
    results = [None] * len(files)

    def _one(i, p):
        results[i] = image_understand(p, question=q)

    import concurrent.futures as _cf

    with _cf.ThreadPoolExecutor(max_workers=min(4, len(files))) as ex:
        futures = [ex.submit(_one, i, p) for i, p in enumerate(files)]
        for _f in _cf.as_completed(futures):
            pass
    lines = []
    for i, p in enumerate(files):
        lines.append(f"### {os.path.basename(p)}\n{results[i] or '（分析失败）'}")
    lines.append(f"\n—— 共分析 {len(files)} 张图片 ——")
    return "\n\n".join(lines)


# Whisper 模型实例缓存：按模型名复用，避免每次调用重新加载（large-v3 可耗时数十秒）
_WHISPER_CACHE = {}
_WHISPER_CACHE_LOCK = threading.Lock()


def speech_to_text(path, model="base"):
    """本地语音转文字（faster-whisper，未安装时提示先安装）。
    model: tiny/base/small/medium/large-v3（首次运行需下载对应模型，tiny/base 较小）。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：音频文件不存在：{path}"
    model_name = str(model or "base").strip().lower()
    if model_name not in ("tiny", "base", "small", "medium", "large-v3"):
        model_name = "base"
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        return "错误：需要 faster-whisper（pip install faster-whisper），安装后重试"
    try:
        inst = _WHISPER_CACHE.get(model_name)
        if inst is None:
            with _WHISPER_CACHE_LOCK:
                inst = _WHISPER_CACHE.get(model_name)
                if inst is None:
                    inst = WhisperModel(model_name, device="cpu", compute_type="int8")
                    _WHISPER_CACHE[model_name] = inst
        segments, _info = inst.transcribe(p)
        text = "".join(seg.text for seg in segments).strip()
        return text or "（未识别出语音内容）"
    except Exception as e:
        return f"错误：语音识别失败: {e}"


# ---------- 本地知识库 RAG（TF-IDF + bigram 语义检索，零依赖） ----------
_KNOWLEDGE_EXTS = (
    ".md", ".txt", ".py", ".json", ".html", ".css", ".js", ".ts", ".csv",
    ".yml", ".yaml", ".ini", ".toml", ".xml", ".log",
)
_KNOWLEDGE_SKIP_DIRS = {".git", "__pycache__", ".venv", "node_modules", "dist", "build", "backups"}
_ARCHIVE_SKIP_DIRS = {".git", "__pycache__", ".venv", "node_modules", "dist", "build"}


def _knowledge_walk(root):
    docs = []
    for base, dirs, files in os.walk(root):
        dirs[:] = [d for d in dirs if d not in _KNOWLEDGE_SKIP_DIRS]
        for fn in files:
            if not fn.lower().endswith(_KNOWLEDGE_EXTS):
                continue
            full = os.path.join(base, fn)
            try:
                if os.path.getsize(full) > 1024 * 1024:
                    continue
            except OSError:
                continue
            docs.append(full)
        if len(docs) >= 200:
            break
    return docs


def _knowledge_snippet(text, query, width=800):
    idx = -1
    for tok in _mem_tokens(query):
        i = str(text).lower().find(tok)
        if i >= 0 and (idx < 0 or i < idx):
            idx = i
    if idx < 0:
        return str(text)[:width]
    start = max(0, idx - width // 4)
    end = min(len(text), idx + width * 3 // 4)
    seg = str(text)[start:end]
    if start > 0:
        seg = "…" + seg
    if end < len(text):
        seg += "…"
    return seg


def knowledge_index(directory="", force=False):
    """对目录内文本文件建立语义检索索引（增量：mtime/size 未变的文档直接复用）。"""
    root = permissions.resolve(str(directory or "").strip() or permissions.WORKSPACE_DIR or ".")
    if not root or not os.path.isdir(root):
        return f"错误：目录不存在：{directory}"
    ok, reason = permissions.check_filesystem(root, write=False)
    if not ok:
        return reason
    try:
        docs = _knowledge_walk(root)
        if not docs:
            return f"错误：目录内没有可索引的文本文件：{root}"
        # 增量复用：上次索引的文档 mtime/size 未变则沿用旧文本（省读取+解析）
        old_docs = {}
        if not force and KNOWLEDGE_INDEX_FILE and os.path.exists(KNOWLEDGE_INDEX_FILE):
            try:
                with open(KNOWLEDGE_INDEX_FILE, "r", encoding="utf-8") as f:
                    old = json.load(f)
                for d in old.get("docs") or []:
                    if isinstance(d, dict) and d.get("path"):
                        old_docs[d["path"]] = d
            except Exception:
                old_docs = {}
        entries = []
        reused = 0
        for full in sorted(docs):
            try:
                st = os.stat(full)
            except OSError:
                continue
            old = old_docs.get(full)
            # 复用判据用纳秒时间戳：秒级 st_mtime 在同一秒内快速改写
            # （size 恰好相同）时误判"未变化"，导致索引漏更新（真实 bug）
            if old and old.get("mtime_ns") == st.st_mtime_ns and old.get("size") == st.st_size:
                if old.get("text"):
                    entries.append({
                        "path": full,
                        "text": old["text"][:100000],
                        "mtime_ns": st.st_mtime_ns,
                        "mtime": st.st_mtime,
                        "size": st.st_size,
                    })
                    reused += 1
                continue
            try:
                with open(full, "r", encoding="utf-8", errors="replace") as f:
                    text = f.read(200000)
            except Exception:
                continue
            if text.strip():
                entries.append({
                    "path": full,
                    "text": text[:100000],
                    "mtime_ns": st.st_mtime_ns,
                    "mtime": st.st_mtime,
                    "size": st.st_size,
                })
        if not entries:
            return "错误：没有可索引的内容"
        idf = _mem_idf([{"value": e["text"][:50000], "key": e["path"]} for e in entries])
        index = {
            "root": root,
            "count": len(entries),
            "idf": idf,
            "docs": entries,
        }
        if KNOWLEDGE_INDEX_FILE:
            os.makedirs(os.path.dirname(KNOWLEDGE_INDEX_FILE) or ".", exist_ok=True)
            tmp = KNOWLEDGE_INDEX_FILE + ".tmp"
            with open(tmp, "w", encoding="utf-8") as f:
                json.dump(index, f, ensure_ascii=False)
            os.replace(tmp, KNOWLEDGE_INDEX_FILE)
        extra = f"（新增 {len(entries) - reused}，复用 {reused}）"
        return f"已索引 {len(entries)} 个文档（{root}）{extra}，可用 knowledge_search 检索"
    except Exception as e:
        return f"错误：建索引失败: {e}"


def knowledge_search(query, top_k=5):
    """语义检索知识库（TF-IDF + bigram，措辞不同也能命中）。"""
    q = str(query or "").strip()
    if not q:
        return "错误：query 必填"
    if not KNOWLEDGE_INDEX_FILE or not os.path.exists(KNOWLEDGE_INDEX_FILE):
        return "错误：知识库尚未建立索引（先用 knowledge_index 对目录建索引）"
    try:
        with open(KNOWLEDGE_INDEX_FILE, "r", encoding="utf-8") as f:
            index = json.load(f)
    except Exception:
        logging.exception("读取知识库索引失败")
        return "错误：索引读取失败，请重新 knowledge_index"
    idf = index.get("idf") or {}
    docs = index.get("docs") or []
    if not docs:
        return "知识库为空（请先用 knowledge_index 建索引）"
    try:
        k = max(1, min(10, int(top_k or 5)))
    except (TypeError, ValueError):
        k = 5
    qt = _mem_tokens(q)
    scored = []
    for d in docs:
        s = _mem_score(qt, idf, d.get("text") or "")
        if s > 0:
            scored.append((s, d))
    scored.sort(key=lambda x: x[0], reverse=True)
    top = scored[:k]
    if not top:
        return f"未找到相关内容（知识库共 {len(docs)} 个文档，查询：{q}）"
    lines = [f"知识库命中 {len(top)}/{len(docs)} 个文档（查询：{q}）："]
    for s, d in top:
        lines.append(f"\n【{s:.2f}】{d['path']}\n{_knowledge_snippet(d.get('text') or '', q)}")
    return "\n".join(lines)


# ---------- 数据库写操作（高危：审批闸门 + 变更前备份） ----------
def database_execute(db_type="sqlite", connection="default", sql="", backup=True):
    """数据库写操作（UPDATE/INSERT/DELETE/DDL）。高危工具，走审批流 + 审计。
    db_type: sqlite / mysql / postgres；sqlite 的 connection 为数据库文件绝对路径。"""
    stmt = str(sql or "").strip()
    if not stmt:
        return "错误：sql 必填"
    if not stmt.lstrip()[:6].upper().startswith(("UPDATE", "INSERT", "DELETE", "CREATE", "DROP", "ALTER", "REPLACE")):
        return "错误：database_execute 仅用于写操作；只读查询请用 database_query*"
    dbtype = str(db_type or "sqlite").lower()
    try:
        if dbtype == "sqlite":
            return _db_execute_sqlite(str(connection or "").strip(), stmt, bool(backup))
        if dbtype == "mysql":
            return _db_execute_mysql(str(connection or "default"), stmt, bool(backup))
        if dbtype == "postgres":
            return _db_execute_postgres(str(connection or "default"), stmt, bool(backup))
        return f"错误：不支持的数据库类型：{db_type}（支持 sqlite / mysql / postgres）"
    except Exception as e:
        return f"错误：数据库执行失败: {e}"


def _db_execute_sqlite(path, stmt, backup):
    import sqlite3

    if not path:
        return "错误：sqlite 的 connection 需为数据库文件绝对路径"
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：数据库文件不存在：{path}"
    # 高危写操作：与 write_file 同级，要求显式写权限 + 审批闸门
    ok, reason = permissions.check_filesystem(p, write=True)
    if not ok:
        return reason
    bak = ""
    if backup:
        try:
            bak_dir = os.path.join(permissions.WORKSPACE_DIR or ".", "db_backups")
            os.makedirs(bak_dir, exist_ok=True)
            bak = os.path.join(bak_dir, f"{os.path.basename(p)}_{datetime.now():%Y%m%d_%H%M%S}.bak")
            shutil.copy2(p, bak)
        except Exception:
            bak = ""
    conn = sqlite3.connect(p, timeout=10)
    try:
        cur = conn.cursor()
        preview = ""
        head = stmt.lstrip().upper()
        if head.startswith(("UPDATE", "DELETE")):
            sel = _db_preview_sql(stmt)
            if sel:
                try:
                    preview = f"\n变更预览：命中 {len(cur.execute(sel).fetchall())} 行"
                except Exception:
                    preview = ""
        cur.execute(stmt)
        affected = max(0, cur.rowcount)
        conn.commit()
        permissions.audit("database_execute", p, f"{affected} 行变更")
        return f"已执行（{affected} 行受影响）{preview}\n备份：{bak or '无'}\nSQL：{stmt[:200]}"
    finally:
        conn.close()


def _db_execute_mysql(connection, stmt, backup):
    try:
        import pymysql
    except ImportError:
        return "错误：需要 pymysql（pip install pymysql）"
    cfg, err = _db_conn("mysql", connection)
    if cfg is None:
        return err
    conn = pymysql.connect(
        host=str(cfg.get("host") or "127.0.0.1"),
        port=int(cfg.get("port") or 3306),
        user=str(cfg.get("user") or ""),
        password=str(cfg.get("password") or ""),
        database=str(cfg.get("database") or ""),
        charset="utf8mb4", connect_timeout=5, read_timeout=30,
    )
    try:
        cur = conn.cursor()
        try:
            cur.execute("SET SESSION max_execution_time=15000")
        except Exception:
            pass
        preview = ""
        head = stmt.lstrip().upper()
        if head.startswith(("UPDATE", "DELETE")):
            sel = _db_preview_sql(stmt)
            if sel:
                try:
                    cur.execute(sel)
                    preview = f"\n变更预览：命中 {len(cur.fetchall())} 行"
                except Exception:
                    preview = ""
        cur.execute(stmt)
        affected = max(0, cur.rowcount)
        conn.commit()
        permissions.audit("database_execute", f"mysql:{connection}", f"{affected} 行变更")
        backup_note = "\n⚠ 当前 MySQL 暂不支持自动备份，请自行确保数据安全" if backup else ""
        return f"已执行（{affected} 行受影响）{preview}{backup_note}\nSQL：{stmt[:200]}"
    finally:
        conn.close()


def _db_execute_postgres(connection, stmt, backup):
    try:
        import psycopg2
    except ImportError:
        return "错误：需要 psycopg2（pip install psycopg2）"
    cfg, err = _db_conn("postgres", connection)
    if cfg is None:
        return err
    conn = psycopg2.connect(
        host=str(cfg.get("host") or "127.0.0.1"),
        port=int(cfg.get("port") or 5432),
        user=str(cfg.get("user") or ""),
        password=str(cfg.get("password") or ""),
        dbname=str(cfg.get("database") or ""),
        connect_timeout=5,
    )
    try:
        cur = conn.cursor()
        preview = ""
        head = stmt.lstrip().upper()
        if head.startswith(("UPDATE", "DELETE")):
            sel = _db_preview_sql(stmt)
            if sel:
                try:
                    cur.execute(sel)
                    preview = f"\n变更预览：命中 {len(cur.fetchall())} 行"
                except Exception:
                    preview = ""
        cur.execute(stmt)
        affected = max(0, cur.rowcount)
        conn.commit()
        permissions.audit("database_execute", f"postgres:{connection}", f"{affected} 行变更")
        backup_note = "\n⚠ 当前 PostgreSQL 暂不支持自动备份，请自行确保数据安全" if backup else ""
        return f"已执行（{affected} 行受影响）{preview}{backup_note}\nSQL：{stmt[:200]}"
    finally:
        conn.close()


# ---------- 收邮件（IMAP，email_config.json 的 imap 段配置） ----------
def read_email(limit=10, since_days=3):
    """读取邮箱近期邮件（IMAP，email_config.json 配置 imap 段：host/port/user/password/ssl）。"""
    if not EMAIL_CONFIG_FILE or not os.path.exists(EMAIL_CONFIG_FILE):
        return "错误：未找到 email_config.json（需配置 imap 段）"
    try:
        with open(EMAIL_CONFIG_FILE, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        imap = cfg.get("imap") if isinstance(cfg, dict) else None
        if not isinstance(imap, dict):
            imap = {}  # 兼容扁平键格式：imap_host / imap_port / imap_user / imap_password / imap_ssl
        host = str(imap.get("host") or cfg.get("imap_host") or "")
        user = str(imap.get("user") or cfg.get("imap_user") or "")
        pwd = _decrypt_secret(str(imap.get("password") or cfg.get("imap_password") or ""))
        if not (host and user and pwd):
            return "错误：imap 配置不完整（host/user/password 必填）"
        try:
            lim = max(1, min(50, int(limit or 10)))
        except (TypeError, ValueError):
            lim = 10
        try:
            days = max(0, min(30, int(since_days or 3)))
        except (TypeError, ValueError):
            days = 3
        import imaplib
        from email.header import decode_header
        from email import message_from_bytes

        ssl_flag = imap.get("ssl", cfg.get("imap_ssl", "true"))
        if str(ssl_flag).lower() in ("true", "1", "yes"):
            conn = imaplib.IMAP4_SSL(host, int(imap.get("port") or cfg.get("imap_port") or 993), timeout=15)
        else:
            conn = imaplib.IMAP4(host, int(imap.get("port") or cfg.get("imap_port") or 143), timeout=15)
        try:
            conn.login(user, pwd)
            conn.select("INBOX")
            if days > 0:
                # IMAP SINCE 需要 d-MMM-yyyy（英文月份）。不能用 strftime("%b")：
                # 中文系统下输出"8月"导致服务器返回 BAD（真实 bug）
                import datetime as _dt

                _d = _dt.date.today() - _dt.timedelta(days=days)
                _months = ("Jan", "Feb", "Mar", "Apr", "May", "Jun",
                           "Jul", "Aug", "Sep", "Oct", "Nov", "Dec")
                since = f"{_d.day:02d}-{_months[_d.month - 1]}-{_d.year}"
                status, data = conn.search(None, "SINCE", since)
            else:
                status, data = conn.search(None, "ALL")
            ids = (data[0] or b"").split() if status == "OK" and data and data[0] else []
            ids = ids[-lim:]
            out = []

            def dec(v):
                parts = decode_header(str(v or ""))
                return "".join(
                    p.decode(ch or "utf-8", errors="replace") if isinstance(p, bytes) else str(p)
                    for p, ch in parts
                )

            for mid in reversed(ids):
                st, msg_data = conn.fetch(mid, "(BODY.PEEK[HEADER.FIELDS (SUBJECT FROM DATE)])")
                if st != "OK" or not msg_data or not msg_data[0]:
                    continue
                raw = msg_data[0][1] if isinstance(msg_data[0], tuple) else None
                if not raw:
                    continue
                try:
                    m = message_from_bytes(raw)
                    out.append(f"发件人: {dec(m.get('From'))}\n主题: {dec(m.get('Subject'))}\n日期: {m.get('Date')}")
                except Exception:
                    continue
            if not out:
                return f"邮箱（{user}）近 {days} 天没有邮件"
            return f"邮箱（{user}）最近 {len(out)} 封邮件：\n\n" + "\n\n".join(out)
        finally:
            try:
                conn.logout()
            except Exception:
                pass
    except Exception as e:
        return f"错误：读取邮件失败: {e}"


# ---------- 新邮件汇总（P1 收件箱模式） ----------
def email_summary(limit=10, since_days=1):
    """读取近期邮件并整理为可汇总的清单（供 AI 生成摘要）。"""
    raw = read_email(limit=limit, since_days=since_days)
    if str(raw).startswith("错误"):
        return raw
    return "新邮件汇总任务：请根据以下邮件清单生成要点摘要（发件人/主题/日期）：\n\n" + raw


# ---------- Agent Mail（agently-cli，可选集成） ----------
def _resolve_agent_mail_cli(cli):
    """把 cli 名解析为 subprocess 可直接执行的绝对路径（兼容 Windows .CMD/.BAT shim）。

    Windows 下 shell=False 不会解析 npm 全局安装的 .CMD shim，直接传命令名会
    FileNotFoundError；这里显式解析出完整路径后仍保持 shell=False 执行。
    """
    cli = str(cli or "agently-cli").strip() or "agently-cli"
    if os.path.isabs(cli):
        return cli
    found = shutil.which(cli)
    if not found:
        return cli  # 兜底：让 FileNotFoundError 走 127 提示
    if os.name == "nt" and found.lower().endswith((".cmd", ".bat")):
        return found
    return found


def _agent_mail_run(args, timeout=60):
    """调用 agently-cli 并返回 (exit_code, stdout_text)。"""
    cli = _resolve_agent_mail_cli(str(AGENT_MAIL_CLI or "agently-cli").strip() or "agently-cli")
    cmd = [cli] + args
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=timeout,
            errors="replace",
            creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
        )
        out = (proc.stdout or "").strip()
        if proc.returncode != 0 and (proc.stderr or "").strip():
            out += (("\n" + proc.stderr.strip()) if out else proc.stderr.strip())
        return proc.returncode, out
    except FileNotFoundError:
        return 127, "未找到 agently-cli（请先 npm install -g @tencent-qqmail/agently-cli，并在 系统 → 外部服务配置 → Agent Mail 开启）"
    except Exception as e:
        return 1, f"调用 agently-cli 失败: {e}"


def _agent_mail_tip():
    return (
        "Agent Mail 未启用（不需要可忽略）。启用方法：系统菜单 → 外部服务配置 → Agent Mail 页签，"
        "勾选启用并确认 CLI 已安装（npm install -g @tencent-qqmail/agently-cli）。"
    )


def agent_mail(action="list", q="", id="", to="", subject="", body="", dir="",
               limit=10, cursor="", confirmation_token="", attachment="", msg="", att="", output=""):
    """Agent 原生邮箱（通过 agently-cli）：me/list/search/read/send/reply/forward/trash/delete/download。

    写操作遵循 CLI 两阶段确认：首次调用不带 confirmation_token 会返回 ctk 与 summary，
    请向用户展示并等待明确许可后，再用相同参数 + confirmation_token 完成。
    """
    if not AGENT_MAIL_ENABLED:
        return _agent_mail_tip()
    act = str(action or "list").strip().lower()
    # 解析收件人与附件（逗号分隔，CLI 要求可重复参数）
    def _multi(value):
        return [x.strip() for x in str(value or "").split(",") if x.strip()]

    cmd = []
    try:
        limit = max(1, min(50, int(limit or 10)))
    except (TypeError, ValueError):
        limit = 10

    try:
        if act == "me":
            cmd += ["+me"]
        elif act == "list":
            cmd += ["message", "+list", "--limit", str(limit)]
            if str(dir or "").strip():
                cmd += ["--dir", str(dir).strip()]
            if str(cursor or "").strip():
                cmd += ["--cursor", str(cursor).strip()]
        elif act == "search":
            cmd += ["message", "+search", "--q", str(q or ""), "--limit", str(limit)]
            if str(dir or "").strip():
                cmd += ["--dir", str(dir).strip()]
            if str(cursor or "").strip():
                cmd += ["--cursor", str(cursor).strip()]
        elif act == "read":
            if not str(id or "").strip():
                return "错误：read 需要 id（msg_xxx）"
            cmd += ["message", "+read", "--id", str(id).strip()]
        elif act == "send":
            if not _multi(to) or not str(subject or "").strip():
                return "错误：send 需要 to（可逗号分隔多个）与 subject"
            cmd += ["message", "+send", "--subject", str(subject).strip(), "--body", str(body or "")]
            for x in _multi(to):
                cmd += ["--to", x]
            for x in _multi(attachment):
                cmd += ["--attachment", x]
            if str(confirmation_token or "").strip():
                cmd += ["--confirmation-token", str(confirmation_token).strip()]
        elif act == "reply":
            if not str(id or "").strip():
                return "错误：reply 需要 id（msg_xxx）"
            cmd += ["message", "+reply", "--id", str(id).strip(), "--body", str(body or "")]
            for x in _multi(attachment):
                cmd += ["--attachment", x]
            if str(confirmation_token or "").strip():
                cmd += ["--confirmation-token", str(confirmation_token).strip()]
        elif act == "forward":
            if not str(id or "").strip() or not _multi(to):
                return "错误：forward 需要 id 与 to"
            cmd += ["message", "+forward", "--id", str(id).strip(), "--body", str(body or "")]
            for x in _multi(to):
                cmd += ["--to", x]
            if str(confirmation_token or "").strip():
                cmd += ["--confirmation-token", str(confirmation_token).strip()]
        elif act == "trash":
            if not str(id or "").strip():
                return "错误：trash 需要 id"
            cmd += ["message", "+trash", "--id", str(id).strip()]
            if str(confirmation_token or "").strip():
                cmd += ["--confirmation-token", str(confirmation_token).strip()]
        elif act == "delete":
            if str(id or "").strip():
                cmd += ["message", "+delete", "--id", str(id).strip()]
            else:
                cmd += ["message", "+delete", "--all"]
        elif act == "download":
            if not str(msg or "").strip() or not str(att or "").strip():
                return "错误：download 需要 msg（msg_xxx）与 att（att_xxx）"
            cmd += ["attachment", "+download", "--msg", str(msg).strip(), "--att", str(att).strip()]
            if str(output or "").strip():
                cmd += ["--output", str(output).strip()]
        else:
            return "错误：action 仅支持 me/list/search/read/send/reply/forward/trash/delete/download"
    except Exception as e:
        return f"错误：参数构造失败: {e}"

    timeout = 120 if act in ("send", "reply", "forward") else 60
    code, out = _agent_mail_run(cmd, timeout=timeout)
    if code == 127:
        return out
    # 授权失效：提示用户重新 OAuth，不自动重试
    if code == 3 or "invalid_grant" in str(out).lower() or "unauthorized" in str(out).lower():
        return (
            "[授权失效] Agent Mail 登录状态已过期。请在系统终端运行 "
            "`agently-cli auth login` 重新授权后再试。\n\n" + out
        )
    # exit 8 = 缺少 confirmation-token：把 ctk/summary 原样交回，AI 必须停下等用户确认
    if code == 8 and not str(confirmation_token or "").strip():
        return (
            "[需要用户确认] 请把下面的 summary 展示给用户并等待明确许可；"
            "用户许可后，用相同参数加 confirmation_token 再次调用。\n\n" + out
        )
    if code == 0:
        return out or "（命令成功，无输出）"
    return f"[agently-cli exit {code}] " + out


# ---------- 任务检查点（断点续跑） ----------
def task_checkpoint_save(name="", status="进行中", pending=None, notes="", auto=False):
    """保存任务进度检查点（崩溃/重启后可从此继续）。

    auto=True：鲸语工具链执行中的自动断点（main 每步工具后写入），
    任务正常完成时由 task_checkpoint_clear 自动清除；手动断点不受影响。
    """
    if not str(name or "").strip() and not str(notes or "").strip():
        return "错误：name 或 notes 必填"
    data = {
        "name": str(name or "未命名任务")[:60],
        "status": str(status or "进行中")[:20],
        "pending": [str(p) for p in (pending or [])][:20],
        "notes": str(notes or "")[:2000],
        "saved_at": datetime.now().isoformat(timespec="seconds"),
    }
    if auto:
        data["auto"] = True
    if not CHECKPOINT_FILE:
        return "错误：检查点模块未初始化"
    try:
        os.makedirs(os.path.dirname(CHECKPOINT_FILE) or ".", exist_ok=True)
        tmp = CHECKPOINT_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        os.replace(tmp, CHECKPOINT_FILE)
        return f"已保存任务检查点：{data['name']}（{data['status']}）"
    except Exception as e:
        return f"错误：保存检查点失败: {e}"


def task_checkpoint_clear():
    """清除自动检查点（任务正常完成时调用）；手动断点保留。"""
    if not CHECKPOINT_FILE or not os.path.exists(CHECKPOINT_FILE):
        return "当前没有检查点"
    try:
        with open(CHECKPOINT_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict) and data.get("auto"):
            os.remove(CHECKPOINT_FILE)
            return "已清除自动断点"
        return "当前检查点为手动保存，未清除"
    except Exception:
        return "错误：清除检查点失败"


def task_checkpoint_load():
    """读取任务检查点（断点续跑时恢复任务上下文）。"""
    if not CHECKPOINT_FILE or not os.path.exists(CHECKPOINT_FILE):
        return "当前没有任务检查点"
    try:
        with open(CHECKPOINT_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return "检查点文件损坏"
        lines = [
            f"任务：{data.get('name', '')}（状态：{data.get('status', '')}）",
            f"保存时间：{data.get('saved_at', '')}",
        ]
        if data.get("auto"):
            lines.append("（自动断点：由鲸语在任务执行中自动保存）")
        if data.get("pending"):
            lines.append("待办步骤：\n" + "\n".join(f"- {p}" for p in data["pending"]))
        if data.get("notes"):
            lines.append(f"备注：{data['notes']}")
        return "\n".join(lines)
    except Exception:
        logging.exception("读取检查点失败")
        return "错误：读取检查点失败"


# ---------- 流程编排（workflows.json：步骤 = 依次发送的指令） ----------
_WORKFLOW_RUNNING = False  # 流程防重：同一时刻只允许一个流程运行
_WORKFLOW_LOCK = threading.Lock()  # 检查-置位原子化：并行工具调用下防双流程同时启动


def _recipe_chain(name):
    """按配方名读取成功模式工具链（patterns.json 的 name 字段）。"""
    if not PATTERNS_FILE or not os.path.exists(PATTERNS_FILE):
        return []
    try:
        with open(PATTERNS_FILE, "r", encoding="utf-8") as f:
            pats = json.load(f)
        for p in pats if isinstance(pats, list) else []:
            if isinstance(p, dict) and str(p.get("name") or "") == str(name or "").strip():
                return [str(c) for c in (p.get("chain") or []) if str(c).strip()]
    except Exception:
        logging.exception("读取配方失败")
    return []


def _workflow_step_text(st, name=""):
    """把单个流程步骤解析为待发送指令。

    dict 步骤支持：
    - text：指令原文（必填或与 recipe 二选一）
    - recipe：配方名（patterns.json）——发送时注入已验证的工具链，text 作为任务目标
    """
    if isinstance(st, dict):
        text = str(st.get("text") or "").strip()
        recipe = str(st.get("recipe") or "").strip()
        if recipe:
            chain = _recipe_chain(recipe)
            if chain:
                prefix = "请按以下已验证成功的工具链顺序执行任务：\n" + " → ".join(chain)
                return f"{prefix}\n\n任务目标：{text}" if text else prefix
            return f"[配方「{recipe}」不存在或为空，请直接完成以下任务]\n{text}" if text else ""
        return text
    return str(st or "").strip()


def run_workflow(name):
    """运行已保存的流程模板：按顺序逐条发送指令，上一步完成后自动执行下一步。

    步骤支持 {"text": "任务目标", "recipe": "配方名"}：自动注入配方工具链。
    """
    global _WORKFLOW_RUNNING
    if not WORKFLOWS_FILE:
        return "错误：流程模块未初始化"
    if not _SEND_CALLBACK:
        return "错误：发送通道不可用"
    try:
        if not os.path.exists(WORKFLOWS_FILE):
            return "错误：没有已保存的流程（workflows.json 为空）"
        with open(WORKFLOWS_FILE, "r", encoding="utf-8") as f:
            wf = json.load(f)
        steps = wf.get(str(name)) if isinstance(wf, dict) else None
        if not steps or not isinstance(steps, dict):
            avail = list(wf) if isinstance(wf, dict) else []
            return f"错误：未找到流程「{name}」（可用：{avail}）"
        step_list = steps.get("steps")
        if not isinstance(step_list, list) or not step_list:
            return f"错误：流程「{name}」没有步骤"
        texts = []
        for st in step_list:
            t = _workflow_step_text(st, name)
            if t:
                texts.append(t)
        if not texts:
            return f"错误：流程「{name}」的步骤均为空"
        desc = f"启动流程「{name}」（{len(texts)} 步）\n" + "\n".join(f"{i}. {t[:100]}" for i, t in enumerate(texts, 1))
        # 校验全部通过后才检查-置位（同一临界区：并行工具调用下防双流程同时启动；
        # 校验失败绝不占位，避免一次失败流程标记永久占用）
        with _WORKFLOW_LOCK:
            if _WORKFLOW_RUNNING:
                return "错误：已有流程正在运行，请等待完成后再启动新流程"
            _WORKFLOW_RUNNING = True
        # 异步执行：在后台线程逐条下发，等待上一步生成结束
        def _run():
            global _WORKFLOW_RUNNING
            try:
                for i, t in enumerate(texts, 1):
                    deadline = time.time() + 600
                    while _BUSY_PROVIDER and _BUSY_PROVIDER():
                        if time.time() > deadline:
                            return
                        time.sleep(1)
                    _SEND_CALLBACK(t)
                    time.sleep(2)  # 让主线程进入生成状态
            except Exception:
                logging.exception("流程执行异常")
            finally:
                with _WORKFLOW_LOCK:
                    _WORKFLOW_RUNNING = False
        threading.Thread(target=_run, daemon=True).start()
        return desc
    except Exception as e:
        with _WORKFLOW_LOCK:
            _WORKFLOW_RUNNING = False
        return f"错误：读取流程失败: {e}"


# ---------- 图片生成（OpenAI 兼容 images API） ----------
def image_generate(prompt, path="", size="1024x1024"):
    """生成图片（需配置 image_api_key / image_base_url / image_model，OpenAI 兼容接口）。"""
    p = str(prompt or "").strip()
    if not p:
        return "错误：prompt 必填"
    # size 白名单校验：非法尺寸让模型自纠（接口对任意字符串返回 400，报错不友好）
    sz = str(size or "1024x1024").strip().lower()
    if not re.match(r"^(256|512|768|1024|1536|2048)x(256|512|768|1024|1536|2048)$", sz):
        return (
            f"错误：size 非法：{size}（支持 256/512/768/1024/1536/2048 的正方形或 "
            "两者组合，如 1024x1024 / 1536x1024）"
        )
    key = str(IMAGE_GEN_KEY or "").strip()
    if not key:
        return "错误：未配置图片生成（config.json 的 image_api_key / image_base_url / image_model）"
    base = str(IMAGE_GEN_BASE or "").strip().rstrip("/") or DEFAULT_BASE_URL
    if str(path or "").strip():
        out = permissions.resolve(path)
        if not out:
            return "错误：输出路径无效"
        if not out.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
            out += ".png"
        ok, reason = permissions.check_filesystem(out, write=True)
        if not ok:
            return reason
    else:
        base_dir = os.path.join(permissions.WORKSPACE_DIR or ".", "images")
        ok, reason = permissions.check_filesystem(base_dir, write=True)
        if not ok:
            return reason
        try:
            os.makedirs(base_dir, exist_ok=True)
        except Exception:
            pass
        out = os.path.join(base_dir, f"gen_{datetime.now():%Y%m%d_%H%M%S}.png")
    try:
        resp = _http_client().post(
            f"{base}/images/generations",
            json={
                "model": IMAGE_GEN_MODEL,
                "prompt": p,
                "n": 1,
                "size": sz,
                "response_format": "b64_json",
            },
            headers={"Authorization": f"Bearer {key}"},
            timeout=120.0,
        )
        resp.raise_for_status()
        data = resp.json()
        items = data.get("data") or []
        if not items:
            return "错误：接口未返回图片"
        import base64

        if items[0].get("b64_json"):
            with open(out, "wb") as f:
                f.write(base64.b64decode(items[0]["b64_json"]))
        elif items[0].get("url"):
            # URL 图片大小不可信：20MB 上限，防写满磁盘；返回地址也做 SSRF 校验
            dl_url = str(items[0]["url"])
            err = _safe_url(dl_url, allow_loopback=False)
            if err:
                return f"错误：图片接口返回了不安全的下载地址（{err}）"
            try:
                with _safe_stream("GET", dl_url, allow_loopback=False, timeout=60) as r:
                    total = 0
                    truncated = False
                    with open(out, "wb") as f:
                        for chunk in r.iter_bytes(64 * 1024):
                            total += len(chunk)
                            if total > 20 * 1024 * 1024:
                                truncated = True
                                break
                            f.write(chunk)
            except Exception:
                try:
                    os.remove(out)
                except OSError:
                    pass
                raise
            if truncated:
                try:
                    os.remove(out)
                except OSError:
                    pass
                return "错误：图片下载超过 20MB 上限，已放弃保存"
        else:
            return "错误：接口返回格式无法解析"
        size_b = os.path.getsize(out)
        permissions.audit("image_generate", out, p[:80])
        return f"已生成图片保存至 {out}（{size_b / 1024:.0f} KB）"
    except Exception as e:
        return f"错误：图片生成失败: {e}"


# ---------- 用量洞察报告 ----------
def usage_report(days=7):
    """生成用量洞察报告（按天/模型汇总 token 与估算费用）。"""
    if not STATS_FILE or not os.path.exists(STATS_FILE):
        return "暂无用量统计数据"
    try:
        days = max(1, min(90, int(days or 7)))
    except (TypeError, ValueError):
        days = 7
    try:
        import stats as stats_mod
        from datetime import date as _date, timedelta

        data = stats_mod.load_stats(STATS_FILE)
        totals = stats_mod.empty_day()
        model_usage = {}
        per_day = []
        for i in range(days - 1, -1, -1):
            d = (_date.today() - timedelta(days=i)).isoformat()
            day_data = data.get(d)
            if not day_data:
                continue
            day_total = stats_mod.day_total(data, d)
            if not any(day_total.values()):
                continue
            per_day.append(
                f"{d}: 输入 {day_total['prompt']:,} / 输出 {day_total['completion']:,}"
                f" / 缓存命中 {day_total['cache_hit']:,}"
            )
            for k in totals:
                totals[k] += day_total[k]
            for model, usage in day_data.items():
                acc = model_usage.setdefault(model, stats_mod.empty_day())
                for k in acc:
                    acc[k] += usage.get(k, 0)
        if not any(totals.values()):
            return f"近 {days} 天没有使用记录"
        hit_ratio = totals["cache_hit"] / max(1, totals["prompt"])
        lines = [
            f"近 {days} 天用量报告：",
            f"输入 {totals['prompt']:,} / 输出 {totals['completion']:,} token，"
            f"缓存命中 {totals['cache_hit']:,}（{hit_ratio:.0%}）",
        ]
        for model, u in model_usage.items():
            lines.append(f"模型 {model}: 输入 {u['prompt']:,} / 输出 {u['completion']:,} / 费用约 ¥{stats_mod.estimate_cost(u, model):.2f}")
        if per_day:
            lines.append("逐日明细：\n" + "\n".join(per_day))
        return "\n".join(lines)
    except Exception as e:
        return f"错误：生成报告失败: {e}"


def _persist_long_result(name, text):
    """工具结果过长时自动落盘到工作区，上下文只保留路径 + 首尾摘要（省 token 不丢信息）。"""
    if len(str(text)) <= _RESULT_INTO_CONTEXT_MAX:
        return text
    if not permissions.WORKSPACE_DIR:
        return text
    try:
        d = os.path.join(permissions.WORKSPACE_DIR, "long_results")
        os.makedirs(d, exist_ok=True)
        fn = f"{name}_{datetime.now():%Y%m%d_%H%M%S%f}.txt"
        path = os.path.join(d, fn)
        with open(path, "w", encoding="utf-8", errors="replace") as f:
            f.write(str(text))
        keep = _RESULT_INTO_CONTEXT_MAX // 8
        body = str(text)
        head = body[:keep]
        tail = body[-keep:]
        return (
            f"[结果过长已落盘] 完整结果（{len(body)} 字符）已保存至：{path}\n"
            f"——开头摘要——\n{head}\n——结尾摘要——\n{tail}\n"
            f"如需查看全部内容，用 read_file 按行读取该文件。"
        )
    except Exception:
        logging.exception("长结果落盘失败，按截断处理")
        return str(text)[:_RESULT_INTO_CONTEXT_MAX] + "\n[结果过大，已截断，请缩小范围后重试]"


# ============================================================================
# 文档处理：PDF 提取 / PDF 生成 / Word 读取 / PPT 读取（可选依赖模式）
# ============================================================================
PDF_EXTRACT_MAX_OUTPUT = 60000   # pdf_extract 单次输出上限（防撑爆上下文）
DOCX_MAX_DEFAULT = 50000         # docx_read 默认输出上限


def pdf_extract(path, pages="all", mode="text"):
    """从 PDF 提取文本（按页）/ 表格（Markdown）/ 元数据；支持页码范围与扫描件提示。"""
    if not str(path or "").strip():
        return "错误：path 必填"
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{path}"
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return "未安装 PyMuPDF，请先执行 pip_install PyMuPDF 后重试"
    m = str(mode or "text").lower()
    if m not in ("text", "table", "meta"):
        return f"错误：mode 仅支持 text / table / meta，收到：{mode}"
    try:
        doc = fitz.open(p)
        try:
            if doc.is_encrypted:
                return "PDF 已加密，暂不支持密码解密，请先移除密码后重试"
            total = doc.page_count
            if total == 0:
                return "PDF 为空文档"
            page_list = _parse_page_range(pages, total)
            if page_list is None:
                return f"错误：页码范围非法（应为 1-{total} 内逗号/连字符组合或 all）：{pages}"
            if m == "meta":
                md = doc.metadata or {}
                return "\n".join([
                    f"文件名: {os.path.basename(p)}",
                    f"页数: {total}",
                    f"标题: {md.get('title') or '（无）'}",
                    f"作者: {md.get('author') or '（无）'}",
                    f"主题: {md.get('subject') or '（无）'}",
                    f"创建时间: {md.get('creationDate') or '（无）'}",
                    f"修改时间: {md.get('modDate') or '（无）'}",
                    f"格式: {md.get('format') or '? '}",
                ])
            out = [f"文件名: {os.path.basename(p)}", f"页数: {total}"]
            out_len = sum(len(x) for x in out)
            truncated_hint = "\n[输出较长已截断，可用 pages= 指定页码范围分段提取]"
            if m == "text":
                for i in page_list:
                    page = doc.load_page(i - 1)
                    seg = f"\n--- 第{i}页 ---\n"
                    text = page.get_text("text").strip()
                    if not text:
                        seg += "（本页无文本层，疑似扫描件；可先用 web_screenshot/pdf 导出页面图片再用 ocr_image 识别）"
                    else:
                        seg += text
                    if out_len + len(seg) > PDF_EXTRACT_MAX_OUTPUT:
                        out.append(truncated_hint)
                        break
                    out.append(seg)
                    out_len += len(seg)
            else:  # table
                if not hasattr(doc.load_page(0), "find_tables"):
                    return "错误：当前 PyMuPDF 版本过低，表格提取需要 PyMuPDF 1.23+（可升级：pip_install PyMuPDF --upgrade，或改用 mode=text）"
                for i in page_list:
                    page = doc.load_page(i - 1)
                    seg_parts = [f"\n--- 第{i}页 表格 ---"]
                    try:
                        tables = page.find_tables()
                        found = False
                        for ti, tb in enumerate(tables.tables, 1):
                            data = tb.extract()
                            if not data:
                                continue
                            found = True
                            seg_parts.append(f"表格 {ti}:")
                            seg_parts.append(_table_to_md(data))
                        if not found:
                            seg_parts.append("（本页未检测到表格）")
                    except Exception:
                        seg_parts.append("（表格提取失败，可改用 mode=text 提取文本）")
                    seg = "\n".join(seg_parts)
                    if out_len + len(seg) > PDF_EXTRACT_MAX_OUTPUT:
                        out.append(truncated_hint)
                        break
                    out.append(seg)
                    out_len += len(seg)
            result = "\n".join(out)
            if len(result) > PDF_EXTRACT_MAX_OUTPUT:
                result = result[:PDF_EXTRACT_MAX_OUTPUT] + truncated_hint
            return result
        finally:
            doc.close()
    except Exception as e:
        return f"错误：PDF 读取失败: {e}"


# ---------- PDF 生成（reportlab，中文字体自动嵌入） ----------
# _find_cjk_font / _register_cjk_font 已移至 pdf_utils.py
def pdf_create(content="", source_path="", output="", title=""):
    """把文本/Markdown 内容生成 PDF（中文字体嵌入；支持标题/列表/代码块/表格）。"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Preformatted, Table, TableStyle,
        )
    except ImportError:
        return "未安装 reportlab，请先执行 pip_install reportlab 后重试"
    if not str(output or "").strip():
        return "错误：output 必填"
    if str(content or "").strip() and str(source_path or "").strip():
        return "错误：content 与 source_path 只能二选一"
    if not str(content or "").strip():
        if not str(source_path or "").strip():
            return "错误：content 或 source_path 必填"
        src = permissions.resolve(source_path)
        if not src or not os.path.isfile(src):
            return f"错误：源文件不存在：{source_path}"
        ok, reason = permissions.check_filesystem(src, write=False)
        if not ok:
            return reason
        try:
            with open(src, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(2_000_000)
        except Exception as e:
            return f"错误：读取源文件失败: {e}"
    if not str(content or "").strip():
        return "错误：内容为空"
    out = permissions.resolve(output)
    if not out:
        return "错误：输出路径无效"
    if not out.lower().endswith(".pdf"):
        out += ".pdf"
    ok, reason = permissions.check_filesystem(out, write=True)
    if not ok:
        return reason
    try:
        import mdparse  # 复用项目自有 Markdown 块解析（无第三方依赖）
    except Exception:
        return "错误：mdparse 不可用"
    try:
        os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
        font = _register_cjk_font()
        mono = "Courier" if font == "Helvetica" else font
        styles = {
            "title": ParagraphStyle("title", fontName=font, fontSize=18, leading=24, spaceAfter=14),
            "h1": ParagraphStyle("h1", fontName=font, fontSize=16, leading=20, spaceBefore=10, spaceAfter=6),
            "h2": ParagraphStyle("h2", fontName=font, fontSize=14, leading=18, spaceBefore=8, spaceAfter=5),
            "h3": ParagraphStyle("h3", fontName=font, fontSize=13, leading=17, spaceBefore=6, spaceAfter=4),
            "h4": ParagraphStyle("h4", fontName=font, fontSize=12, leading=16, spaceBefore=5, spaceAfter=3),
            "h5": ParagraphStyle("h5", fontName=font, fontSize=11, leading=15, spaceBefore=4, spaceAfter=2),
            "h6": ParagraphStyle("h6", fontName=font, fontSize=10, leading=14, spaceBefore=4, spaceAfter=2),
            "body": ParagraphStyle("body", fontName=font, fontSize=10.5, leading=16, spaceAfter=6),
            "code": ParagraphStyle("code", fontName=mono, fontSize=8.5, leading=11.5,
                                   backColor=colors.Color(0.95, 0.95, 0.95), borderPadding=6,
                                   spaceBefore=4, spaceAfter=8),
        }
        flow = []
        # 标题：显式 title 优先；否则取内容首行作为文档元数据标题（PRD 规范）
        doc_title = str(title or "").strip()
        if not doc_title:
            first_line = next((ln.strip() for ln in str(content).splitlines() if ln.strip()), "")
            doc_title = first_line[:200]
        if str(title or "").strip():
            flow.append(Paragraph(_md_inline_html(title), styles["title"]))
        blocks = mdparse.parse_blocks(str(content))
        for blk in blocks:
            kind = blk[0]
            body = blk[1]
            if kind == "code":
                flow.append(Preformatted(body, styles["code"]))
            elif kind == "table":
                rows = _md_table_rows(body)
                if len(rows) >= 2:
                    t = Table(rows, repeatRows=1)
                    t.setStyle(TableStyle([
                        ("FONTNAME", (0, 0), (-1, 0), font),
                        ("FONTNAME", (0, 1), (-1, -1), font),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.9, 0.93, 0.97)),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ]))
                    flow.append(t)
                    flow.append(Spacer(1, 8))
            elif kind == "list":
                for ln in str(body).split("\n"):
                    if ln.strip():
                        flow.append(Paragraph("• " + _md_inline_html(ln), styles["body"]))
            elif kind.startswith("h"):
                level = int(kind[1])
                flow.append(Paragraph(_md_inline_html(body), styles[f"h{min(level, 6)}"]))
            else:
                flow.append(Paragraph(_md_inline_html(body), styles["body"]))
        if not flow:
            return "错误：内容未能解析为可排版元素"
        doc = SimpleDocTemplate(out, pagesize=A4,
                                leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54,
                                title=doc_title)
        doc.build(flow)
        size = os.path.getsize(out) if os.path.exists(out) else 0
        permissions.audit("pdf_create", out, f"{size} 字节")
        return f"已生成 PDF: {out}（{size / 1024:.1f} KB，中文字体 {'已嵌入' if font != 'Helvetica' else '未找到（可能乱码，请安装中文字体）'}）"
    except Exception as e:
        return f"错误：PDF 生成失败: {e}"


def docx_read(path, max_chars=50000):
    """读取 Word .docx 为 Markdown 结构（标题/段落/列表/表格，保持文档顺序）。"""
    try:
        from docx import Document
    except ImportError:
        return "未安装 python-docx，请先执行 pip_install python-docx 后重试"
    if not str(path or "").strip():
        return "错误：path 必填"
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{path}"
    if p.lower().endswith(".doc"):
        return "错误：暂不支持旧版 .doc 格式，请先用 Word 另存为 .docx 后重试"
    try:
        limit = max(200, min(500000, int(max_chars or DOCX_MAX_DEFAULT)))
    except (TypeError, ValueError):
        limit = DOCX_MAX_DEFAULT
    try:
        from docx.table import Table as _Table
        from docx.text.paragraph import Paragraph as _Para

        doc = Document(p)
        parts = []
        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1]
            if tag == "p":
                para = _Para(child, doc)
                text = para.text.strip()
                if not text:
                    continue
                style = str(para.style.name or "")
                sl = style.lower()
                if sl.startswith("title"):
                    parts.append("# " + text)
                elif sl.startswith("heading"):
                    try:
                        level = int(style[-1])
                    except (TypeError, ValueError):
                        level = 1
                    parts.append("#" * min(level, 6) + " " + text)
                elif "list bullet" in sl:
                    parts.append("- " + text)
                elif "list number" in sl:
                    parts.append("1. " + text)
                else:
                    parts.append(text)
            elif tag == "tbl":
                table = _Table(child, doc)
                rows = [[c.text.strip() for c in r.cells] for r in table.rows]
                if rows and any(any(c for c in r) for r in rows):
                    parts.append(_table_to_md(rows))
        if not parts:
            return "（文档无可见内容）"
        result = "\n\n".join(parts)
        if len(result) > limit:
            result = result[:limit] + f"\n[内容较长已截断前 {limit} 字符]"
        return result
    except Exception as e:
        return f"错误：Word 读取失败: {e}"


def pptx_read(path, include_notes=True):
    """提取 PowerPoint 每页幻灯片的标题、正文要点与备注；图片占位标注。"""
    try:
        from pptx import Presentation
    except ImportError:
        return "未安装 python-pptx，请先执行 pip_install python-pptx 后重试"
    if not str(path or "").strip():
        return "错误：path 必填"
    ok, reason = permissions.check_filesystem(path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(path)
    if not p or not os.path.isfile(p):
        return f"错误：文件不存在：{path}"
    if p.lower().endswith(".ppt"):
        return "错误：暂不支持旧版 .ppt 格式，请先用 PowerPoint 另存为 .pptx 后重试"
    try:
        prs = Presentation(p)
        out = [f"幻灯片数: {len(prs.slides)}"]

        def _walk_shapes(shapes, title_holder, body_lines, img_count, table_count, depth=0):
            """递归遍历形状（组合形状 GROUP 内文本/图片不丢），深度上限防畸形文件死循环。"""
            if depth > 10:
                return
            for shape in shapes:
                is_title_ph = (
                    getattr(shape, "is_placeholder", False)
                    and shape.placeholder_format.idx == 0
                    and shape.has_text_frame
                )
                if is_title_ph:
                    t = shape.text_frame.text.strip()
                    if t and not title_holder[0]:
                        title_holder[0] = t
                    continue
                if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    if hasattr(shape, "shapes"):
                        _walk_shapes(shape.shapes, title_holder, body_lines, img_count, table_count, depth + 1)
                elif shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        body_lines.extend(ln for ln in t.splitlines() if ln.strip())
                elif getattr(shape, "has_table", False):
                    table_count[0] += 1
                    rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                    if rows:
                        body_lines.append(f"[表格 {table_count[0]}]")
                        body_lines.append(_table_to_md(rows))
                elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_count[0] += 1

        for idx, slide in enumerate(prs.slides, 1):
            out.append(f"\n--- 第{idx}页 ---")
            title_holder = [None]
            body_lines = []
            img_count = [0]
            table_count = [0]
            _walk_shapes(slide.shapes, title_holder, body_lines, img_count, table_count)
            if title_holder[0]:
                out.append(f"标题: {title_holder[0]}")
            if body_lines:
                out.append("\n".join("- " + ln for ln in body_lines[:40]))
            elif not title_holder[0]:
                out.append("（本页无文本）")
            if img_count[0]:
                out.append(f"[图片占位: {img_count[0]} 张]")
            if include_notes and slide.has_notes_slide:
                try:
                    nt = slide.notes_slide.notes_text_frame.text.strip()
                    if nt:
                        out.append(f"备注: {nt[:500]}")
                except Exception:
                    pass
        return "\n".join(out)
    except Exception as e:
        return f"错误：PPT 读取失败: {e}"


# ============================================================================
# 资讯聚合：RSS 订阅管理 + 抓取（feedparser 可选依赖）
# ============================================================================
RSS_FETCH_TIMEOUT = 10
RSS_MAX_ITEMS = 20
RSS_SUMMARY_MAX = 300


def _load_rss_sources():
    if not RSS_SOURCES_FILE or not os.path.exists(RSS_SOURCES_FILE):
        return []
    try:
        with open(RSS_SOURCES_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, list) else []
    except Exception:
        logging.exception("读取 RSS 订阅失败")
        return []


def _save_rss_sources(sources):
    if not RSS_SOURCES_FILE:
        return False
    try:
        os.makedirs(os.path.dirname(RSS_SOURCES_FILE) or ".", exist_ok=True)
        tmp = RSS_SOURCES_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(sources, f, ensure_ascii=False, indent=2)
        os.replace(tmp, RSS_SOURCES_FILE)
        return True
    except Exception:
        logging.exception("保存 RSS 订阅失败")
        return False


# 精选 RSS 预置源（action=preset 一键添加）：中文 AI/科技/开发者为主
RSS_PRESET_SOURCES = [
    {"name": "机器之心", "url": "https://www.jiqizhixin.com/rss"},
    {"name": "量子位", "url": "https://www.qbitai.com/feed"},
    {"name": "少数派", "url": "https://sspai.com/feed"},
    {"name": "IT之家", "url": "https://www.ithome.com/rss/"},
    {"name": "开源中国", "url": "https://www.oschina.net/news/rss"},
    {"name": "Hacker News", "url": "https://news.ycombinator.com/rss"},
]


def rss_fetch(action="list", url="", limit=10, since_hours=24):
    """RSS 订阅管理（list/add/remove/preset）+ 抓取最新条目（标题/链接/时间/摘要）。

    action=preset 一键添加精选科技/AI 源（机器之心/量子位/少数派/IT之家/开源中国/HN）。
    """
    act = str(action or "list").strip().lower()
    if act not in ("list", "add", "remove", "fetch", "preset"):
        return "错误：action 仅支持 list / add / remove / fetch / preset"
    if act in ("add", "fetch") and not str(url or "").strip():
        return f"错误：{act} 需要 url（RSS 源地址）"
    try:
        lim = max(1, min(RSS_MAX_ITEMS, int(limit or 10)))
    except (TypeError, ValueError):
        lim = 10
    try:
        hours = max(0, min(24 * 30, int(since_hours) if since_hours not in (None, "") else 24))
    except (TypeError, ValueError):
        hours = 24
    if act == "list":
        sources = _load_rss_sources()
        if not sources:
            return "当前没有 RSS 订阅（用 action=preset 一键添加精选源，或 action=add url=... 手动添加）"
        lines = [f"共 {len(sources)} 个订阅源："]
        for i, s in enumerate(sources, 1):
            lines.append(f"{i}. {s.get('name') or s.get('url')} | {s.get('url')}")
        return "\n".join(lines)
    if act == "preset":
        sources = _load_rss_sources()
        existing = {s.get("url") for s in sources}
        added = [s for s in RSS_PRESET_SOURCES if s["url"] not in existing]
        if not added:
            return "精选源均已订阅"
        _save_rss_sources(sources + added)
        return f"已添加 {len(added)} 个精选源：" + "、".join(s["name"] for s in added)
    if act == "add":
        u = str(url).strip()
        if len(u) > 2048 or not u.startswith(("http://", "https://")):
            return "错误：url 必须是 http(s) 开头的 RSS 源地址"
        err = _safe_url(u, allow_loopback=False)
        if err:
            return f"错误：{err}"
        sources = _load_rss_sources()
        if any(s.get("url") == u for s in sources):
            return "该源已订阅"
        sources.append({"url": u, "name": "", "added": datetime.now().isoformat(timespec="seconds")})
        if _save_rss_sources(sources):
            return f"已添加订阅源（当前共 {len(sources)} 个）：{u}"
        return "错误：订阅保存失败"
    if act == "remove":
        u = str(url).strip()
        sources = _load_rss_sources()
        kept = [s for s in sources if s.get("url") != u]
        if len(kept) == len(sources):
            return f"未找到订阅源：{u}"
        if _save_rss_sources(kept):
            return f"已移除订阅源（剩余 {len(kept)} 个）：{u}"
        return "错误：订阅保存失败"
    # fetch
    try:
        import feedparser
    except ImportError:
        return "未安装 feedparser，请先执行 pip_install feedparser 后重试"
    u = str(url).strip()
    if not u.startswith(("http://", "https://")):
        return "错误：url 必须是 http(s) 开头的 RSS 源地址"
    err = _safe_url(u, allow_loopback=False)
    if err:
        return f"错误：{err}"
    # feedparser 6.x 的 parse() 不再支持 timeout 关键字（旧版支持）：
    # 统一用内部线程 + join 超时实现可靠超时，兼容所有版本
    box = {}

    def _parse():
        try:
            box["parsed"] = feedparser.parse(u, request_headers={"User-Agent": _SEARCH_UA})
        except Exception as e:
            box["err"] = e

    t = threading.Thread(target=_parse, daemon=True)
    t.start()
    t.join(RSS_FETCH_TIMEOUT)
    if t.is_alive():
        return "错误：RSS 抓取超时（>10 秒），请稍后重试或检查源地址"
    if "err" in box:
        return f"错误：RSS 抓取失败: {box['err']}"
    parsed = box["parsed"]
    # getattr 兼容 FeedParserDict 与普通对象
    if getattr(parsed, "bozo", 0) and not getattr(parsed, "entries", []):
        return f"错误：无效的 RSS 源（{getattr(parsed, 'bozo_exception', None) or '解析失败'}）"
    import calendar
    from datetime import datetime as _dt

    cutoff_ts = None
    if hours > 0:
        cutoff_ts = time.time() - hours * 3600
    picked = []
    seen = set()
    for e in parsed.entries:
        title = str(getattr(e, "title", "") or "").strip()
        link = str(getattr(e, "link", "") or "").strip()
        if not title and not link:
            continue
        fp = (link or title)
        if fp in seen:
            continue
        seen.add(fp)
        published = getattr(e, "published_parsed", None) or getattr(e, "updated_parsed", None)
        if published:
            try:
                ts = calendar.timegm(published)
            except Exception:
                ts = None
            if cutoff_ts is not None and ts is not None and ts < cutoff_ts:
                continue  # feed 不一定按时间排序：用 continue 而非 break，防乱序源丢条目
        picked.append(e)
        if len(picked) >= lim:
            break
    if not picked:
        return f"来源 {u} 最近 {hours} 小时没有新条目" if hours else f"来源 {u} 没有可显示的条目"
    feed_title = str(getattr(getattr(parsed, "feed", None), "title", "") or "") or u
    lines = [f"来源: {feed_title}（{len(picked)} 条，最近 {hours}h）"]
    for i, e in enumerate(picked, 1):
        title = str(getattr(e, "title", "") or "（无标题）")[:100]
        link = str(getattr(e, "link", "") or "").strip()
        pub = ""
        published = getattr(e, "published_parsed", None) or getattr(e, "updated_parsed", None)
        if published:
            try:
                pub = _dt.fromtimestamp(calendar.timegm(published), _dt.timezone.utc).strftime("%Y-%m-%d %H:%M")
            except Exception:
                pub = ""
        summary = re.sub(r"<[^>]+>", " ", str(getattr(e, "summary", "") or getattr(e, "description", "")))
        summary = re.sub(r"\s+", " ", summary).strip()[:RSS_SUMMARY_MAX]
        lines.append(f"{i}. {title} | {pub} | {link}")
        if summary:
            lines.append(f"   {summary}")
    return "\n".join(lines)


# ============================================================================
# 二维码：生成（qrcode 可选依赖）/ 识别（pyzbar 可选依赖，缺失降级提示）
# ============================================================================
def qrcode(action="generate", text="", output="", image_path="", size=300, error_correction="M"):
    """二维码生成与识别。"""
    act = str(action or "generate").strip().lower()
    if act not in ("generate", "read"):
        return "错误：action 仅支持 generate / read"
    if act == "generate":
        if not str(text or "").strip():
            return "错误：generate 需要 text（要编码的内容）"
        if not str(output or "").strip():
            return "错误：generate 需要 output（PNG 路径）"
        try:
            import qrcode
        except ImportError:
            return "未安装 qrcode，请先执行 pip_install qrcode 后重试"
        try:
            from qrcode.constants import ERROR_CORRECT_L, ERROR_CORRECT_M, ERROR_CORRECT_Q, ERROR_CORRECT_H
        except Exception:
            ERROR_CORRECT_L, ERROR_CORRECT_M, ERROR_CORRECT_Q, ERROR_CORRECT_H = 1, 0, 3, 2
        out = permissions.resolve(output)
        if not out:
            return "错误：输出路径无效"
        if not out.lower().endswith(".png"):
            out += ".png"
        ok, reason = permissions.check_filesystem(out, write=True)
        if not ok:
            return reason
        try:
            s = max(64, min(1024, int(size or 300)))
        except (TypeError, ValueError):
            s = 300
        ec_map = {"L": ERROR_CORRECT_L, "M": ERROR_CORRECT_M,
                  "Q": ERROR_CORRECT_Q, "H": ERROR_CORRECT_H}
        ec = ec_map.get(str(error_correction or "M").upper(), ERROR_CORRECT_M)
        try:
            os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
            qr = qrcode.QRCode(version=None, error_correction=ec, box_size=10, border=2)
            qr.add_data(str(text))
            qr.make(fit=True)
            img = qr.make_image(fill_color="black", back_color="white")
            if img.size != (s, s):
                img = img.resize((s, s))
            img.save(out)
            permissions.audit("qrcode_generate", out, f"{len(str(text))} 字符")
            return f"已生成二维码: {out}（{s}x{s}px）"
        except Exception as e:
            return f"错误：二维码生成失败: {e}"
    # read
    if not str(image_path or "").strip():
        return "错误：read 需要 image_path"
    ok, reason = permissions.check_filesystem(image_path, write=False)
    if not ok:
        return reason
    p = permissions.resolve(image_path)
    if not p or not os.path.isfile(p):
        return f"错误：图片不存在：{image_path}"
    try:
        from PIL import Image
        import pyzbar.pyzbar as pyzbar
    except Exception:
        # pyzbar 在 Windows 依赖系统 zbar DLL，缺失时 import 即抛异常 → 统一降级提示
        return (
            "未安装 pyzbar（Windows 需系统 zbar 库，pip_install pyzbar 后还需安装 "
            "zbar DLL）。降级方案：可先用 ocr_image 对图片做粗识别"
        )
    try:
        img = Image.open(p).convert("RGB")
        results = pyzbar.decode(img)
        if not results:
            return "未识别到二维码（可尝试 ocr_image 粗识别）"
        lines = [f"识别到 {len(results)} 个二维码："]
        for r in results:
            lines.append("· " + r.data.decode("utf-8", errors="replace"))
        return "\n".join(lines)
    except Exception as e:
        return f"错误：二维码识别失败: {e}"


# ============================================================================
# 密钥保险箱（P2 信任基建）：DPAPI 加密托管，按名取用，不落明文日志
# ============================================================================
SECRETS_FILE = None  # 由 main 注入（DATA_DIR/secrets.json）


def _load_secrets():
    if not SECRETS_FILE:
        return {}
    if not os.path.exists(SECRETS_FILE):
        return {}
    try:
        with open(SECRETS_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return {}
        out = {}
        for k, v in data.items():
            try:
                out[str(k)] = crypto.decrypt(str(v))
            except Exception:
                out[str(k)] = str(v)
        return out
    except Exception:
        logging.exception("读取密钥保险箱失败")
        return {}


def _save_secrets(data):
    if not SECRETS_FILE:
        return False
    try:
        os.makedirs(os.path.dirname(SECRETS_FILE) or ".", exist_ok=True)
        out = {str(k): crypto.encrypt(str(v)) for k, v in (data or {}).items()}
        from persistence import atomic_json_write
        return atomic_json_write(SECRETS_FILE, out, indent=2)
    except Exception:
        logging.exception("保存密钥保险箱失败")
        return False


def secret_store(action="get", name="", value=""):
    """密钥保险箱：set / get / delete / list。value 只写不显示；list 只返回名称。"""
    act = str(action or "get").strip().lower()
    name = str(name or "").strip()
    data = _load_secrets()
    if act == "set":
        if not name:
            return "错误：name 必填"
        if not str(value or ""):
            return "错误：value 必填"
        data[name] = str(value)
        if not _save_secrets(data):
            return "错误：保存密钥失败"
        permissions.audit("secret_set", name, "值已加密保存", result="ok")
        return f"密钥「{name}」已加密保存（调用 secret_get(name) 取用）"
    if act == "get":
        if not name:
            return "错误：name 必填"
        if name not in data:
            return f"未找到密钥「{name}」"
        return str(data[name])
    if act == "delete":
        if not name:
            return "错误：name 必填"
        if name not in data:
            return f"未找到密钥「{name}」"
        data.pop(name, None)
        _save_secrets(data)
        permissions.audit("secret_delete", name, "删除")
        return f"已删除密钥「{name}」"
    if act == "list":
        return "已保存密钥：" + ("、".join(sorted(data.keys())) if data else "（空）")
    return "错误：action 仅支持 set / get / delete / list"


# ============================================================================
# 嵌入式 KV 存储（diskcache 可选依赖；支持 TTL 与模糊检索）
# ============================================================================
KV_VALUE_MAX_BYTES = 1024 * 1024  # value 上限 1MB


def kv_store(action="get", key="", value="", pattern="", ttl_seconds=0):
    """嵌入式键值存储：set（可选 TTL）/ get / delete / keys / search。"""
    act = str(action or "get").strip().lower()
    if act not in ("set", "get", "delete", "keys", "search"):
        return "错误：action 仅支持 set / get / delete / keys / search"
    try:
        import diskcache
    except ImportError:
        return "未安装 diskcache，请先执行 pip_install diskcache 后重试"
    if not KV_CACHE_DIR:
        return "错误：KV 存储未初始化"
    try:
        os.makedirs(KV_CACHE_DIR, exist_ok=True)
    except Exception:
        return "错误：KV 目录创建失败"
    try:
        with diskcache.Cache(KV_CACHE_DIR) as cache:
            if act == "set":
                k = str(key or "").strip()
                if not k:
                    return "错误：set 需要 key"
                if len(k) > 256:
                    return "错误：key 过长（上限 256 字符）"
                v = str(value or "")
                if len(v.encode("utf-8", "replace")) > KV_VALUE_MAX_BYTES:
                    return "错误：value 超过 1MB 上限"
                try:
                    ttl = max(0, min(365 * 24 * 3600, int(ttl_seconds or 0)))
                except (TypeError, ValueError):
                    ttl = 0
                cache.set(k, v, expire=ttl or None)
                n = len(cache)
                ttl_txt = f"TTL {ttl}s" if ttl else "长期有效"
                return f"已写入 key={k}（{ttl_txt}，当前共 {n} 个键）"
            if act == "get":
                k = str(key or "").strip()
                if not k:
                    return "错误：get 需要 key"
                v = cache.get(k)
                if v is None:
                    return f"key={k} 不存在或已过期"
                return f"key={k}: {v}"
            if act == "delete":
                k = str(key or "").strip()
                if not k:
                    return "错误：delete 需要 key"
                if k in cache:
                    del cache[k]
                    return f"已删除 key={k}"
                return f"key={k} 不存在"
            if act == "keys":
                keys = [k for k in cache.iterkeys() if not k.startswith("_")]
                if not keys:
                    return "KV 存储为空"
                lines = [f"共 {len(keys)} 个键："]
                for i, k in enumerate(sorted(keys), 1):
                    try:
                        v = cache.get(k) or ""
                        lines.append(f"{i}. {k} = {str(v)[:80]}")
                    except Exception:
                        lines.append(f"{i}. {k} = ?")
                return "\n".join(lines)
            # search
            pat = str(pattern or "").strip().lower()
            if not pat:
                return "错误：search 需要 pattern"
            hits = []
            for k in cache.iterkeys():
                if k.startswith("_"):
                    continue
                try:
                    v = str(cache.get(k) or "")
                except Exception:
                    continue
                if pat in k.lower() or pat in v.lower():
                    hits.append(f"{k}: {v[:120]}")
                    if len(hits) >= 50:
                        break
            if not hits:
                return f"未找到包含「{pattern}」的键或值"
            return f"命中 {len(hits)} 项：\n" + "\n".join(hits)
    except Exception as e:
        return f"错误：KV 操作失败: {e}"


# ============================================================================
# 音视频处理（imageio-ffmpeg 自带 ffmpeg 二进制；参数白名单 + 超时/大小限制）
# ============================================================================
MEDIA_MAX_INPUT = 2 * 1024 * 1024 * 1024   # 输入 2GB 上限
MEDIA_TIMEOUT = 300                        # 单次转码/提取最长 300s
MEDIA_FORMATS = {"mp4", "mp3", "webm", "mkv", "avi", "mov", "ogg", "flac", "wav", "gif", "png", "jpg"}
_FFMPEG_BIN = None


def _ffmpeg_path():
    global _FFMPEG_BIN
    if _FFMPEG_BIN is None:
        try:
            import imageio_ffmpeg
            _FFMPEG_BIN = imageio_ffmpeg.get_ffmpeg_exe()
        except Exception:
            _FFMPEG_BIN = "ffmpeg"
    return _FFMPEG_BIN


def _ffmpeg_run(args, timeout=MEDIA_TIMEOUT):
    """执行 ffmpeg（argv 直传，禁止 shell 拼接）；超时杀进程树。"""
    try:
        proc = subprocess.Popen(
            [_ffmpeg_path()] + args,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True, encoding="utf-8", errors="replace",
            creationflags=subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0,
        )
    except Exception as e:
        return None, f"ffmpeg 启动失败: {e}"
    try:
        out, err = proc.communicate(timeout=timeout)
    except subprocess.TimeoutExpired:
        _kill_tree(proc)
        return None, f"执行超时（>{timeout} 秒）"
    return proc.returncode, (out or "") + (err or "")


def media_ffmpeg(action="info", input="", output="", time="", width=0, format=""):
    """音视频：info / thumbnail / transcode / extract_audio（参数白名单化）。"""
    act = str(action or "info").strip().lower()
    if act not in ("info", "thumbnail", "transcode", "extract_audio"):
        return "错误：action 仅支持 info / thumbnail / transcode / extract_audio"
    if not str(input or "").strip():
        return "错误：input 必填"
    ok, reason = permissions.check_filesystem(input, write=False)
    if not ok:
        return reason
    src = permissions.resolve(input)
    if not src or not os.path.isfile(src):
        return f"错误：源文件不存在：{input}"
    try:
        if os.path.getsize(src) > MEDIA_MAX_INPUT:
            return "错误：输入文件超过 2GB 上限"
    except OSError:
        pass
    if act == "info":
        code, text = _ffmpeg_run(["-hide_banner", "-i", src], timeout=20)
        if code is None:
            return f"错误：{text}"
        lines = []
        for ln in (text or "").splitlines():
            s = ln.strip()
            if s.startswith("Duration:"):
                dur = s.split("Duration:", 1)[1].split(",", 1)[0].strip()
                lines.append("时长: " + dur)
            elif s.startswith("Stream #"):
                lines.append("流: " + s.split("Stream #", 1)[1].strip())
        if not lines:
            return f"错误：无法解析媒体信息（ffmpeg 输出：{(text or '')[:200]}）"
        return "\n".join([f"文件名: {os.path.basename(src)}"] + lines)
    # thumbnail / transcode / extract_audio 需要 output
    if not str(output or "").strip():
        return f"错误：{act} 需要 output（输出路径）"
    out = permissions.resolve(output)
    if not out:
        return "错误：输出路径无效"
    if act == "thumbnail":
        if not out.lower().endswith((".png", ".jpg", ".jpeg")):
            out += ".png"
        ok, reason = permissions.check_filesystem(out, write=True)
        if not ok:
            return reason
        ts = str(time or "").strip()
        if ts and not re.match(r"^\d{1,2}:\d{2}:\d{2}(\.\d+)?$|^\d+(\.\d+)?$", ts):
            return "错误：time 格式应为 HH:MM:SS 或秒数（如 00:01:30）"
        args = ["-hide_banner", "-y", "-ss", ts or "1", "-i", src, "-frames:v", "1", "-q:v", "2", out]
        code, text = _ffmpeg_run(args, timeout=60)
        if code is None:
            return f"错误：{text}"
        if code != 0:
            return f"错误：截图失败：{(text or '')[-300:]}"
        size = os.path.getsize(out) if os.path.exists(out) else 0
        permissions.audit("media_ffmpeg_thumbnail", out, f"{size} 字节")
        return f"已截图保存至 {out}（{size / 1024:.0f} KB，时间点 {ts or '1s'}）"
    fmt = str(format or "").strip().lower().lstrip(".")
    if fmt not in MEDIA_FORMATS:
        return (
            f"错误：format 非法：{format or '（空）'}（支持 {'/'.join(sorted(MEDIA_FORMATS))}；"
            "如未指定可按输出扩展名推断）"
        )
    if not out.lower().endswith(("." + fmt, ".jpg")):
        out += "." + fmt
    ok, reason = permissions.check_filesystem(out, write=True)
    if not ok:
        return reason
    if act == "transcode":
        w = 0
        if width:
            try:
                w = max(16, min(7680, int(width)))
            except (TypeError, ValueError):
                return "错误：width 应为 16-7680 的整数"
        args = ["-hide_banner", "-y", "-i", src]
        if w:
            args += ["-vf", f"scale={w}:-2", "-c:v", "libx264", "-preset", "veryfast"]
        if fmt in ("mp3", "ogg", "flac", "wav"):
            args += ["-vn"]
        args += [out]
    else:  # extract_audio
        # 强制转码（copy 与目标容器可能不兼容；testsrc 无音频流时 mp3 也能正常产出空流）
        acodec = "libmp3lame" if fmt == "mp3" else ("flac" if fmt == "flac" else "pcm_s16le" if fmt == "wav" else "libmp3lame")
        args = ["-hide_banner", "-y", "-i", src, "-vn", "-acodec", acodec, out]
    code, text = _ffmpeg_run(args)
    if code is None:
        return f"错误：{text}"
    if code != 0:
        return f"错误：处理失败：{(text or '')[-300:]}"
    size = os.path.getsize(out) if os.path.exists(out) else 0
    permissions.audit("media_ffmpeg", out, f"{act} {size} 字节")
    return f"已{'转码' if act == 'transcode' else '提取音频'}保存至 {out}（{size / 1024 / 1024:.1f} MB）"


# ============================================================================
# WebDAV 云盘同步（httpx 原生 PROPFIND/GET/PUT/DELETE；凭据可 DPAPI 加密）
# ============================================================================
WEBDAV_MAX_SIZE = 200 * 1024 * 1024  # 单文件 200MB 上限（防全量进内存）


def _load_webdav_config():
    if not WEBDAV_CONFIG_FILE or not os.path.exists(WEBDAV_CONFIG_FILE):
        return None, (
            "未配置 WebDAV。请在数据目录创建 webdav_config.json："
            '{"url": "https://dav.example.com", "username": "you@example.com", "password": "***"}'
            "（password 可先用 dpapi: 前缀的密文，或直接明文）"
        )
    try:
        with open(WEBDAV_CONFIG_FILE, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        url = str(cfg.get("url") or "").strip().rstrip("/")
        user = str(cfg.get("username") or "").strip()
        pwd = crypto.decrypt(str(cfg.get("password") or ""))
        if not (url and user and pwd):
            return None, "错误：webdav_config.json 缺少 url / username / password"
        if not url.startswith(("http://", "https://")):
            return None, "错误：webdav url 必须为 http(s) 开头"
        return {"url": url, "username": user, "password": pwd}, ""
    except Exception:
        logging.exception("读取 WebDAV 配置失败")
        return None, "错误：读取 WebDAV 配置失败"


def _webdav_request(cfg, method, path, **kw):
    url = cfg["url"] + str(path)
    return _http_client().request(
        method, url,
        auth=(cfg["username"], cfg["password"]),
        timeout=30,
        **kw,
    )


def webdav(action="list", remote_path="/", local_path=""):
    """WebDAV 云盘操作：list / upload / download / delete。"""
    act = str(action or "list").strip().lower()
    if act not in ("list", "upload", "download", "delete"):
        return "错误：action 仅支持 list / upload / download / delete"
    cfg, err = _load_webdav_config()
    if cfg is None:
        return err
    remote = str(remote_path or "/").strip()
    if not remote.startswith("/"):
        remote = "/" + remote
    if ".." in remote.split("/") or any(ord(ch) < 32 for ch in remote):
        return "错误：remote_path 非法（禁止 .. 与控制字符）"
    if act == "list":
        body = (
            '<?xml version="1.0"?><d:propfind xmlns:d="DAV:">'
            "<d:prop><d:displayname/><d:getcontentlength/><d:getlastmodified/>"
            "<d:resourcetype/></d:prop></d:propfind>"
        )
        try:
            resp = _webdav_request(
                cfg, "PROPFIND", remote,
                headers={"Depth": "1", "Content-Type": "application/xml"},
                content=body,
            )
        except Exception as e:
            return f"错误：WebDAV 请求失败: {e}"
        if resp.status_code not in (200, 207):
            return f"错误：WebDAV 列表失败（HTTP {resp.status_code}）"
        import xml.etree.ElementTree as ET

        try:
            root = ET.fromstring(resp.text)
        except Exception:
            return "错误：WebDAV 响应解析失败"
        ns = {"d": "DAV:"}
        items = []
        for resp_el in root.findall(".//d:response", ns):
            href_el = resp_el.find("d:href", ns)
            href = (href_el.text if href_el is not None else "") or ""
            props = resp_el.find("d:propstat/d:prop", ns)
            if props is None:
                continue
            dn = props.find("d:displayname", ns)
            sz = props.find("d:getcontentlength", ns)
            mod = props.find("d:getlastmodified", ns)
            rt = props.find("d:resourcetype/d:collection", ns)
            name = (dn.text if dn is not None and dn.text else None) or href.rstrip("/").split("/")[-1] or "/"
            items.append((
                rt is not None,
                name,
                (sz.text if sz is not None else "") or "",
                (mod.text if mod is not None else "") or "",
            ))
        if not items:
            return f"远端目录为空：{remote}"
        lines = [f"远端目录 {remote}："]
        for is_dir, name, size, mod in sorted(items, key=lambda x: (not x[0], x[1].lower())):
            lines.append(f"{'DIR ' if is_dir else 'FILE'} {name} | {size}B | {mod}")
        return "\n".join(lines)
    if act in ("upload", "download"):
        if not str(local_path or "").strip():
            return f"错误：{act} 需要 local_path"
        if act == "download":
            ok, reason = permissions.check_filesystem(local_path, write=True)
            if not ok:
                return reason
            out = permissions.resolve(local_path)
            if not out:
                return "错误：本地路径无效"
            total = 0
            try:
                os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
                client = _http_client()
                if hasattr(client, "stream"):
                    with _safe_stream("GET", cfg["url"] + remote, timeout=30) as resp:
                        if resp.status_code != 200:
                            return f"错误：下载失败（HTTP {resp.status_code}）"
                        with open(out, "wb") as f:
                            for chunk in resp.iter_bytes(64 * 1024):
                                total += len(chunk)
                                if total > WEBDAV_MAX_SIZE:
                                    try:
                                        os.remove(out)
                                    except OSError:
                                        pass
                                    return f"错误：远端文件超过 {WEBDAV_MAX_SIZE // 1024 // 1024}MB 上限，请分段下载"
                                f.write(chunk)
                else:
                    resp = _webdav_request(cfg, "GET", remote)
                    if resp.status_code != 200:
                        return f"错误：下载失败（HTTP {resp.status_code}）"
                    if len(resp.content) > WEBDAV_MAX_SIZE:
                        return f"错误：远端文件超过 {WEBDAV_MAX_SIZE // 1024 // 1024}MB 上限，请分段下载"
                    total = len(resp.content)
                    with open(out, "wb") as f:
                        f.write(resp.content)
            except Exception as e:
                try:
                    if os.path.exists(out):
                        os.remove(out)
                except OSError:
                    pass
                return f"错误：WebDAV 下载失败: {e}"
            return f"已下载 {remote} → {out}（{total} 字节）"
        # upload
        ok, reason = permissions.check_filesystem(local_path, write=False)
        if not ok:
            return reason
        src = permissions.resolve(local_path)
        if not src or not os.path.isfile(src):
            return f"错误：本地文件不存在：{local_path}"
        try:
            if os.path.getsize(src) > WEBDAV_MAX_SIZE:
                return f"错误：本地文件超过 {WEBDAV_MAX_SIZE // 1024 // 1024}MB 上限，请压缩后上传"
        except Exception as e:
            return f"错误：读取本地文件信息失败: {e}"
        total = 0
        try:
            client = _http_client()
            if hasattr(client, "stream"):
                def _chunks():
                    nonlocal total
                    with open(src, "rb") as f:
                        while True:
                            chunk = f.read(64 * 1024)
                            if not chunk:
                                break
                            total += len(chunk)
                            yield chunk
                # 流式上传：大文件不再一次性读入内存
                with _safe_stream(
                    "PUT", cfg["url"] + remote,
                    auth=(cfg["username"], cfg["password"]),
                    timeout=30,
                    content=_chunks(),
                ) as resp:
                    status_code = resp.status_code
            else:
                # 兼容旧测试/自定义 mock：无 stream 时退化为普通请求
                with open(src, "rb") as f:
                    data = f.read()
                total = len(data)
                resp = _webdav_request(cfg, "PUT", remote, content=data)
                status_code = resp.status_code
        except Exception as e:
            return f"错误：WebDAV 上传失败: {e}"
        if status_code not in (200, 201, 204):
            return f"错误：上传失败（HTTP {status_code}）"
        permissions.audit("webdav_upload", remote, f"{total} 字节")
        return f"已上传 {src} → {remote}（{total} 字节）"
    # delete
    try:
        resp = _webdav_request(cfg, "DELETE", remote)
    except Exception as e:
        return f"错误：WebDAV 删除失败: {e}"
    if resp.status_code not in (200, 204, 404):
        return f"错误：删除失败（HTTP {resp.status_code}）"
    permissions.audit("webdav_delete", remote, "ok")
    return f"已删除远端路径：{remote}"


# ============================================================================
# 公众号自动写作（wechat_writer 独立包，薄封装）
# ============================================================================
def run_wechat_writer(dry_run=False, topic="", use_blocked=False):
    """运行公众号自动写作工具：采集→选题→写作→质检→存草稿箱（只产草稿）。

    耗时可能 1-3 分钟（多次 LLM 调用）；返回结构化摘要文本。
    草稿统一写到工作区 drafts/（与 publish_draft 同目录，用户可从产物面板直达）。
    use_blocked=True 时被墙信源（linux.do/hostloc 等）自动经代理通道采集。
    """
    try:
        from wechat_writer import run_once
    except ImportError:
        return "错误：wechat_writer 模块不可用（请确认项目目录完整）"
    try:
        drafts_dir = None
        archive_dir = None
        if permissions.WORKSPACE_DIR:
            drafts_dir = os.path.join(permissions.WORKSPACE_DIR, "drafts")
            archive_dir = os.path.join(permissions.WORKSPACE_DIR, "wechat_articles")
        result = run_once(
            dry_run=bool(dry_run),
            topic_override=str(topic or ""),
            drafts_dir=drafts_dir,
            archive_dir=archive_dir,
            use_blocked=bool(use_blocked),
        )
    except Exception as e:
        return f"错误：公众号写作工具运行失败: {e}"
    if not result.get("ok"):
        reasons = result.get("quality", {}).get("reasons") or result.get("errors") or ["未知原因"]
        return f"公众号写作未完成：{'；'.join(str(r) for r in reasons[:3])}"
    q = result.get("quality") or {}
    paths = result.get("paths") or {}
    lines = [
        f"✅ 公众号文章已完成（{result.get('chars', 0)} 字，质检分 {q.get('score', '?')}）：",
        f"主题：{result.get('topic', '')}",
        f"标题：{result.get('title', '')}",
    ]
    if paths.get("draft_path"):
        lines.append(f"草稿：{paths['draft_path']}")
    if paths.get("html_path"):
        lines.append(f"HTML：{paths['html_path']}")
    if paths.get("archive_path"):
        lines.append(f"存档：{paths['archive_path']}")
    if dry_run:
        lines.append("（dry-run 预览：未写入草稿箱，正式运行后草稿存草稿箱，请在公众号后台审阅后手动发布）")
    else:
        lines.append("草稿已存入草稿箱，请在公众号后台审阅后手动发布（工具不自动发布）。")
    return "\n".join(lines)


# ============================================================================
# 每日简报（主动助手：采集当日资讯 → LLM 提炼 → 落盘工作区 briefs/）
# ============================================================================
def daily_brief(topic="", max_items=8):
    """生成每日简报：采集当日 AI/科技资讯（复用 WeChat Writer 采集引擎）
    → LLM 提炼要点与点评 → 保存到工作区 briefs/brief_YYYYMMDD.md。

    topic：可选主题关键词（仅保留标题/摘要命中的素材）。
    返回简报正文 + 落盘路径。
    """
    try:
        from wechat_writer import config as _ww_config
        from wechat_writer import sources as _ww_sources
    except ImportError:
        return "错误：wechat_writer 模块不可用（请确认项目目录完整）"
    try:
        cfg = _ww_config.load_config()
        items = _ww_sources.collect_all(cfg)
    except Exception as e:
        return f"错误：资讯采集失败: {e}"
    kw = str(topic or "").strip()
    if kw:
        items = [it for it in items if kw in (f"{it.title} {it.summary}")]
    if not items:
        return "今日暂无资讯素材（RSS 与搜索均无结果），可稍后再试"
    try:
        limit = max(3, min(15, int(max_items or 8)))
    except (TypeError, ValueError):
        limit = 8
    items = items[:limit]
    client = _CLIENT_HOLDER.get("client")
    if client is None:
        return "错误：没有可用客户端（请先完成一次对话建立连接）"
    material = "\n\n".join(
        f"{i + 1}. {it.title}（{it.source}）\n   {it.url}\n   {it.summary[:200]}"
        for i, it in enumerate(items)
    )
    prompt = (
        "你是每日资讯主编。基于以下今日采集的资讯，生成一份精炼简报：\n"
        "1. 简报标题：一句话概括今日主题（## 开头）\n"
        "2. 3-6 条要点，每条用 - 前缀：主题 + 一句话点评\n"
        "3. 结尾「今日趋势」：2-3 句话总结值得关注的动向\n"
        "只输出简报正文（Markdown），不要任何说明文字。\n\n"
        f"素材（{len(items)} 条）：\n{material}"
    )
    try:
        resp = client.client.chat.completions.create(
            model=client.model,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2048,
            stream=False,
            timeout=120.0,
            extra_body={"thinking": {"type": "disabled"}},
        )
        brief = (resp.choices[0].message.content or "").strip()
    except Exception as e:
        return f"错误：简报生成失败: {e}"
    if not brief:
        return "简报生成失败：模型返回空内容，请重试"
    out = ""
    if permissions.WORKSPACE_DIR:
        d = os.path.join(permissions.WORKSPACE_DIR, "briefs")
        ok, reason = permissions.check_filesystem(d, write=True)
        if not ok:
            out = f"\n（简报落盘被权限拒绝：{reason}）"
        else:
            try:
                os.makedirs(d, exist_ok=True)
                path = os.path.join(d, f"brief_{datetime.now():%Y%m%d}.md")
                with open(path, "w", encoding="utf-8") as f:
                    f.write(f"# 每日简报 {datetime.now():%Y-%m-%d}\n\n{brief}\n")
                out = f"\n已保存：{path}"
            except Exception as e:
                out = f"\n（简报落盘失败：{e}）"
    return f"📰 今日简报（{len(items)} 条素材）：\n\n{brief}{out}"


# ============================================================================
# 插件工坊：AI 生成并安装插件（零代码能力扩展）
# ============================================================================
def _to_tool_schema(t):
    """把简化工具描述转成 user_tools.json 完整 schema（兼容已完整 schema 的输入）。"""
    if isinstance(t, dict) and t.get("function"):
        return t
    if not isinstance(t, dict):
        return None
    name = str(t.get("name") or "").strip()
    endpoint = str(t.get("endpoint") or "").strip()
    if not (name and endpoint):
        return None
    params = [p.strip() for p in str(t.get("params") or "").split(",") if p.strip()]
    return {
        "type": "function",
        "function": {
            "name": name,
            "description": str(t.get("description") or "")[:200],
            "parameters": {
                "type": "object",
                "properties": {p: {"type": "string", "description": p} for p in params},
                "required": params,
            },
            "endpoint": endpoint,
            "method": str(t.get("method") or "POST").upper(),
        },
    }


def create_plugin(name, description="", tools=None, skills=None, workflows=None,
                  scenario=None, requires=None):
    """AI 生成并安装插件：根据需求组合工具/技能/流程/场景，生成后立即生效。

    内部走审批闸门（ACTION_TOOLS 注册）；安装后可在「工具 → 插件管理」停用/卸载。
    """
    name = str(name or "").strip()
    if not name:
        return "错误：name 必填（插件名称）"
    contents = {}
    if tools:
        converted = [_to_tool_schema(t) for t in tools] if isinstance(tools, list) else [_to_tool_schema(tools)]
        converted = [t for t in converted if t]
        if converted:
            contents["tools"] = converted
    if skills:
        contents["skills"] = skills if isinstance(skills, list) else [skills]
    if workflows:
        contents["workflows"] = workflows
    if scenario:
        contents["scenario"] = scenario
    if not contents:
        return "错误：插件需要至少一项能力（tools / skills / workflows / scenario）"
    plugin = {
        "format": plugins_mod.PLUGIN_FORMAT,
        "version": 1,
        "meta": {
            "name": name[:40],
            "description": str(description or "")[:200],
            "author": "鲸语 AI",
            "version": "1.0.0",
        },
        "requires": [str(r).strip() for r in (requires or []) if str(r).strip()],
        "contents": contents,
    }
    ok, err = plugins_mod.validate_plugin(plugin)
    if not ok:
        return f"错误：插件校验失败：{err}"
    if not PLUGIN_PATHS:
        return "错误：插件模块未初始化"
    res = plugins_mod.apply_plugin(plugin, PLUGIN_PATHS)
    if not res.get("ok"):
        return f"错误：插件安装失败：{res.get('error')}"
    added = res.get("added") or {}
    parts = []
    if added.get("tools"):
        parts.append(f"工具 {'、'.join(added['tools'])}")
    if added.get("skills"):
        parts.append(f"技能 {'、'.join(added['skills'])}")
    if added.get("workflows"):
        parts.append(f"流程 {'、'.join(added['workflows'])}")
    if scenario:
        parts.append("场景配置（可在插件管理中应用）")
    miss = plugins_mod.missing_requires(plugin)
    note = f"\n⚠ 缺失依赖：{'、'.join(miss)}（pip install …，可在「依赖状态」查看）" if miss else ""
    return (
        f"✅ 插件「{name}」已生成并安装：{'；'.join(parts) or '空'}。\n"
        f"安装目录：{res.get('path')}{note}\n"
        "可在「工具 → 插件管理」查看、停用或卸载；插件可导出 .wtplugin 分享给他人。"
    )


TOOL_CALL_MAP = {
    "get_date": get_date,
    "ask_user": None,  # 特殊处理：chat() 中通过 on_ask 回调询问用户
    "request_permission": None,  # 特殊处理：chat() 中通过 on_request_permission 回调
    "write_memory": write_memory,
    "read_memory": read_memory,
    "query_memory_graph": query_memory_graph,
    "get_weather": get_weather,
    "run_python": run_python,
    "read_file": read_file,
    "fetch_url": fetch_url,
    "fetch_blocked": _run_fetch_blocked,
    "search_web": search_web,
    "search_github": search_github,
    "search_realtime": search_realtime,
    "call_api": call_api,
    "system_status": system_status,
    "database_query": database_query,
    "tts_save": tts_save,
    "image_process": image_process,
    "ocr_image": ocr_image,
    "read_csv": read_csv,
    "write_csv": write_csv,
    "read_excel": read_excel,
    "write_excel": write_excel,
    "chart_data": chart_data,
    "database_query_mysql": database_query_mysql,
    "database_query_postgres": database_query_postgres,
    "send_webhook": send_webhook,
    "im_send": im_send,
    "telegram_poll_updates": telegram_poll_updates,
    "download_file": download_file,
    "epub_read": epub_read,
    "mobi_read": mobi_read,
    "doc_read": doc_read,
    "msg_read": msg_read,
    "archive_list": archive_list,
    "subagent_run": subagent_run,
    "run_tests": run_tests,
    "verify_output": verify_output,
    "send_email": send_email,
    "pip_install": pip_install,
    "write_file": write_file,
    "edit_file": edit_file,
    "list_dir": list_dir,
    "run_command": run_command,
    "search_local": search_local,
    "create_doc": create_doc,
    "write_code_project": write_code_project,
    "browser_navigate": browser_navigate,
    "rpa_screen_size": rpa_screen_size,
    "rpa_click": rpa_click,
    "rpa_type": rpa_type,
    "rpa_hotkey": rpa_hotkey,
    "rpa_move": rpa_move,
    "rpa_scroll": rpa_scroll,
    "rpa_screenshot": rpa_screenshot,
    "web_screenshot": web_screenshot,
    "publish_draft": publish_draft,
    "start_process": start_process,
    "stop_process": stop_process,
    "list_processes": list_processes,
    "environment_info": environment_info,
    "project_info": project_info,
    "read_project_file": read_project_file,
    "create_evolution": create_evolution,
    "verify_files": verify_files,
    # ===== v2 能力层（全新分类） =====
    "schedule_task": schedule_task,
    "list_schedules": list_schedules,
    "cancel_schedule": cancel_schedule,
    "notify_desktop": notify_desktop,
    "clipboard_get": clipboard_get,
    "clipboard_set": clipboard_set,
    "delete_file": delete_file,
    "archive_files": archive_files,
    "extract_archive": extract_archive,
    "batch_rename": batch_rename,
    "image_understand": image_understand,
    "screen_see": screen_see,
    "chart_read": chart_read,
    "screenshot_to_html": screenshot_to_html,
    "debug_screenshot": debug_screenshot,
    "scan_read": scan_read,
    "image_batch": image_batch,
    "screen_capture": screen_capture,
    "speech_to_text": speech_to_text,
    "knowledge_index": knowledge_index,
    "knowledge_search": knowledge_search,
    "database_execute": database_execute,
    "read_email": read_email,
    "email_summary": email_summary,
    "agent_mail": agent_mail,
    "task_checkpoint_save": task_checkpoint_save,
    "task_checkpoint_load": task_checkpoint_load,
    "run_workflow": run_workflow,
    "image_generate": image_generate,
    "usage_report": usage_report,
    # ===== 文档 / 资讯 / 二维码 / KV / 媒体 / WebDAV =====
    "pdf_extract": pdf_extract,
    "pdf_create": pdf_create,
    "docx_read": docx_read,
    "pptx_read": pptx_read,
    "rss_fetch": rss_fetch,
    "qrcode": qrcode,
    "kv_store": kv_store,
    "secret_store": secret_store,
    "media_ffmpeg": media_ffmpeg,
    "webdav": webdav,
    "run_wechat_writer": run_wechat_writer,
    "daily_brief": daily_brief,
    "create_plugin": create_plugin,
}

MAX_TOOL_ROUNDS = 100
MAX_EMPTY_RETRIES = 1
MAX_SAME_TOOL_REPEATS = 3
MAX_PLAN_REJECTIONS = 3
_RESULT_INTO_CONTEXT_MAX = 40000  # 工具结果写入历史的字符上限（≈1 万 token）
# 停止后等待已提交工具结果的宽限期：副作用已发生的工具（发信/写文件/启进程）
# 要拿到真实结果写回历史，模型下轮才不会重试造成重复执行
_STOP_TOOL_GRACE_S = 1.5


class _StopRequested(Exception):
    """请求已停止（内部信号）：调用链内部转成干净 return，不向 UI 抛异常。"""

# 自我进化工具：不受 enabled_tools / tools_enabled 控制，始终对模型可用
SELF_EVOLUTION_TOOLS = {
    "project_info",
    "read_project_file",
    "create_evolution",
    "verify_files",
}

# ===== 智能工具调取（smart_tools）：索引 + 按需激活 =====
# 完全智能模式不再一次性注入全部工具 schema（≈15k token），
# 改为：常驻注入精简「工具索引」+ activate_tools 点菜工具；
# AI 按需激活后，下一轮注入激活工具的完整 schema。

ACTIVATE_TOOL = {
    "type": "function",
    "function": {
        "name": "activate_tools",
        "description": (
            "加载你拥有但尚未加载的能力定义。你的全部能力见系统消息中的能力地图"
            "（按分类分组，每组有组名如「数据与文档」「媒体与图像」）。"
            "传入工具名激活单个工具，传入组名一次激活整组。"
            "只需激活本次要用到的工具/组，不要全部激活。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "tools": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要激活的工具名或组名列表（来自能力地图）",
                }
            },
            "required": ["tools"],
        },
    },
}

_TOOL_INDEX_CACHE = None
_TOOL_INDEX_KEY = None

# 能力地图分类（精确感知：类别 + 完整工具名 + 核心动作）。全部内置工具全覆盖。
TOOL_GROUPS = [
    ("🌐 浏览器与网页", ["browser_navigate", "web_screenshot", "fetch_url", "fetch_blocked", "search_web", "search_realtime", "search_github", "webdav", "download_file", "rss_fetch"]),
    ("💻 编程与执行", ["run_python", "run_command", "pip_install", "run_tests", "write_code_project", "subagent_run", "verify_output", "start_process", "stop_process", "list_processes", "environment_info", "system_status"]),
    ("📁 文件与目录", ["read_file", "write_file", "edit_file", "list_dir", "search_local", "delete_file", "archive_files", "extract_archive", "batch_rename", "clipboard_get", "clipboard_set"]),
    ("📊 数据与文档", ["read_csv", "write_csv", "read_excel", "write_excel", "chart_data", "database_query", "database_query_mysql", "database_query_postgres", "database_execute", "create_doc", "docx_read", "pptx_read", "pdf_extract", "pdf_create", "epub_read", "mobi_read", "doc_read", "archive_list", "secret_store", "kv_store"]),
    ("📧 邮件与消息", ["send_email", "read_email", "email_summary", "agent_mail", "msg_read", "im_send", "telegram_poll_updates", "send_webhook", "publish_draft", "run_wechat_writer", "daily_brief"]),
    ("🎨 媒体与图像", ["image_process", "image_understand", "screen_see", "chart_read", "screenshot_to_html", "debug_screenshot", "scan_read", "image_batch", "image_generate", "ocr_image", "screen_capture", "speech_to_text", "tts_save", "media_ffmpeg", "qrcode"]),
    ("🖱 桌面自动化", ["rpa_screen_size", "rpa_click", "rpa_type", "rpa_hotkey", "rpa_move", "rpa_scroll", "rpa_screenshot", "notify_desktop"]),
    ("⏰ 定时与任务", ["schedule_task", "list_schedules", "cancel_schedule", "task_checkpoint_save", "task_checkpoint_load", "run_workflow"]),
    ("🧠 记忆与知识", ["write_memory", "read_memory", "query_memory_graph", "knowledge_index", "knowledge_search"]),
    ("🔧 系统与基础", ["get_date", "get_weather", "ask_user", "request_permission", "call_api", "project_info", "read_project_file", "create_evolution", "verify_files", "usage_report", "create_plugin"]),
]

# 组名 -> 成员工具名（activate_tools 支持按组激活：传组名一次激活整组）。
# 键含两种形式：原文（含 emoji）与去掉 emoji 的裸组名（如「数据与文档」）。
_TOOL_GROUP_NAME_MAP = {}
for _cat, _members in TOOL_GROUPS:
    _TOOL_GROUP_NAME_MAP[_cat] = list(_members)
    _bare = _cat.split(" ", 1)[-1] if " " in _cat else _cat
    if _bare != _cat:
        _TOOL_GROUP_NAME_MAP[_bare] = list(_members)


def _expand_activation(wanted, available_names, activated):
    """展开 activate_tools 的请求：支持工具名与组名（组名展开为整组工具）。"""
    for n in wanted:
        n = str(n).strip()
        if not n:
            continue
        group = _TOOL_GROUP_NAME_MAP.get(n)
        if group is not None:
            for m in group:
                if m in available_names:
                    activated.add(m)
        elif n in available_names:
            activated.add(n)
    return activated


# 关键词预激活（chat 层兜底）：扫描最近 user 消息，命中常见意图关键词时
# 预激活对应工具，让常见任务免点菜直接可用（仅提前加载定义，不改变权限）。
_PREACTIVATE_HINTS = [
    (("搜索", "搜一下", "查一下", "新闻", "资讯", "最新"), ["search_web", "search_realtime", "fetch_url"]),
    (("天气", "气温", "台风", "预报"), ["get_weather"]),
    (("下载",), ["download_file", "fetch_url"]),
    (("邮件", "发邮件", "收件箱"), ["send_email", "read_email", "email_summary"]),
    (("文件", "读取", "读一下", "打开"), ["read_file", "list_dir", "search_local"]),
    (("写", "保存", "创建", "生成"), ["write_file", "create_doc", "write_code_project"]),
    (("图片", "图像", "截图", "看图", "图表"), ["image_understand", "screen_see", "image_process", "chart_read", "chart_data", "ocr_image"]),
    (("表格", "excel", "csv", "报表"), ["read_excel", "write_excel", "read_csv", "chart_data"]),
    (("代码", "编程", "python", "bug", "脚本", "函数"), ["run_python", "read_file", "run_tests"]),
    (("定时", "提醒", "计划", "日程"), ["schedule_task"]),
    (("数据库", "sql", "mysql", "postgres"), ["database_query", "database_query_mysql", "database_query_postgres"]),
    (("网页", "url", "抓取", "爬"), ["fetch_url", "browser_navigate", "web_screenshot"]),
    (("搜索文件", "检索", "找文件"), ["search_local", "list_dir"]),
    (("记忆", "记住", "偏好"), ["write_memory", "read_memory", "query_memory_graph"]),
]


def _message_text(m):
    """提取消息的纯文本（兼容图片内联的内容块）。"""
    c = m.get("content")
    if isinstance(c, list):
        return " ".join(
            str(b.get("text", "")) for b in c
            if isinstance(b, dict) and b.get("type") == "text"
        )
    return str(c or "")


def _preactivate_from_messages(messages, activated):
    """chat 层关键词预激活：扫描最近的 user 消息，命中意图词即预激活对应工具。"""
    for m in reversed(messages):
        if not (isinstance(m, dict) and m.get("role") == "user"):
            continue
        text = _message_text(m)[:800]
        if not text.strip():
            continue
        for kws, tools in _PREACTIVATE_HINTS:
            if any(kw in text for kw in kws):
                activated.update(tools)
        return activated  # 只扫最近一条 user 消息
    return activated

# 核心动作短语（能力感知关键：一行说清「能做什么」）
_TOOL_ACTION_PHRASES = {
    "browser_navigate": "控制浏览器（打开网页/点击/输入/填表/提交/选择/取文本，共享登录态）",
    "web_screenshot": "网页截图保存",
    "fetch_url": "抓取网页/接口的文本或 JSON",
    "fetch_blocked": "抓取被墙站点（代理+指纹绕过封锁）",
    "search_web": "联网搜索（多引擎聚合）",
    "search_realtime": "实时热点/社区讨论搜索",
    "search_github": "搜索 GitHub 仓库",
    "webdav": "WebDAV 云盘（坚果云/Nextcloud）上传下载",
    "download_file": "下载文件到本地",
    "rss_fetch": "RSS 订阅/聚合阅读",
    "run_python": "执行 Python 代码",
    "run_command": "执行系统命令",
    "pip_install": "安装 Python 包",
    "run_tests": "运行项目测试",
    "write_code_project": "创建完整代码项目",
    "subagent_run": "派发子智能体并行处理",
    "verify_output": "核验产物/输出",
    "start_process": "启动后台进程（服务器/长驻任务）",
    "stop_process": "停止后台进程",
    "list_processes": "查看后台进程列表",
    "environment_info": "环境/依赖信息",
    "system_status": "系统资源自检",
    "read_file": "读取文件内容",
    "write_file": "写入文件",
    "edit_file": "编辑文件（局部修改）",
    "list_dir": "列出目录",
    "search_local": "在允许目录内全文检索文件",
    "delete_file": "删除文件/目录",
    "archive_files": "打包压缩",
    "extract_archive": "解压归档",
    "batch_rename": "批量重命名",
    "clipboard_get": "读取剪贴板",
    "clipboard_set": "写入剪贴板",
    "read_csv": "读 CSV",
    "write_csv": "写 CSV",
    "read_excel": "读 Excel",
    "write_excel": "写 Excel",
    "chart_data": "数据可视化图表（线/柱/饼/散点）",
    "database_query": "SQLite 只读查询",
    "database_query_mysql": "MySQL 只读查询",
    "database_query_postgres": "PostgreSQL 只读查询",
    "database_execute": "数据库写操作（SQLite/MySQL/PG，带审批）",
    "create_doc": "创建 Office 文档（docx/pptx/pdf）",
    "docx_read": "读取 Word 文档",
    "pptx_read": "读取 PPT",
    "pdf_extract": "提取 PDF 文本",
    "pdf_create": "生成 PDF",
    "epub_read": "读取 epub 电子书",
    "mobi_read": "读取 mobi 电子书",
    "doc_read": "读取 doc/rtf 等旧格式",
    "archive_list": "列出归档内容",
    "secret_store": "加密密钥存储",
    "kv_store": "轻量键值存储（缓存/状态）",
    "send_email": "发送邮件（SMTP）",
    "read_email": "读取邮件（IMAP）",
    "email_summary": "邮件摘要/统计",
    "agent_mail": "Agent 邮箱（查看/列表/搜索/回复/转发）",
    "msg_read": "读取邮件消息",
    "im_send": "IM 消息（Telegram/企业微信）",
    "telegram_poll_updates": "轮询 Telegram 更新",
    "send_webhook": "Webhook 推送（钉钉/ServerChan/Slack）",
    "publish_draft": "发布草稿",
    "run_wechat_writer": "公众号文章生成/排版",
    "daily_brief": "每日简报（采集当日资讯→提炼点评）",
    "image_process": "图像处理（缩放/裁剪/滤镜/格式转换）",
    "image_understand": "多模态看图理解",
    "screen_see": "截图并让视觉模型解读当前屏幕",
    "chart_read": "图表截图→结构化数据+解读",
    "screenshot_to_html": "UI截图→HTML/CSS前端还原",
    "debug_screenshot": "报错截图→诊断修复建议",
    "scan_read": "扫描件/文档图片读取（图表/公式/手写）",
    "image_batch": "批量视觉分析文件夹图片并汇总",
    "image_generate": "文生图",
    "ocr_image": "图片文字识别 OCR",
    "screen_capture": "屏幕截图",
    "speech_to_text": "语音转文字",
    "tts_save": "文字转语音",
    "media_ffmpeg": "音视频处理（ffmpeg）",
    "qrcode": "二维码生成/识别",
    "rpa_screen_size": "获取屏幕尺寸",
    "rpa_click": "模拟点击（屏幕坐标）",
    "rpa_type": "模拟键盘输入",
    "rpa_hotkey": "模拟快捷键",
    "rpa_move": "移动鼠标",
    "rpa_scroll": "滚动页面",
    "rpa_screenshot": "屏幕区域截图",
    "notify_desktop": "桌面通知",
    "schedule_task": "定时任务（cron/每日/周期）",
    "list_schedules": "查看定时任务",
    "cancel_schedule": "取消定时任务",
    "task_checkpoint_save": "保存任务断点",
    "task_checkpoint_load": "加载任务断点",
    "run_workflow": "执行工作流",
    "write_memory": "写入长期记忆",
    "read_memory": "检索长期记忆",
    "query_memory_graph": "记忆知识图谱查询",
    "knowledge_index": "建立知识库索引",
    "knowledge_search": "语义检索知识库",
    "get_date": "获取当前日期/时间",
    "get_weather": "查询天气",
    "ask_user": "向用户提问（澄清/确认）",
    "request_permission": "请求权限（白名单）",
    "call_api": "调用任意 HTTP API",
    "project_info": "项目信息/文件树",
    "read_project_file": "读取项目文件",
    "create_evolution": "创建自我进化提案",
    "verify_files": "核验项目文件完整性",
    "usage_report": "用量/费用统计",
    "create_plugin": "创建用户插件",
}


def build_tool_index(tools=None):
    """生成能力地图：分类 + 完整工具名 + 核心动作短语（AI 准确感知全部能力）。"""
    global _TOOL_INDEX_CACHE, _TOOL_INDEX_KEY
    tools = tools if tools is not None else TOOLS
    # 缓存键用内容指纹（工具名 + 描述），而非 id(tools)：传入深拷贝/重建列表时
    # id 会变导致缓存失效 → index_msg 内容漂移 → 前缀缓存不命中（成本翻几十倍）。
    key = tuple(sorted(
        (t["function"]["name"], t["function"].get("description", ""))
        for t in (tools or [])
    ))
    if _TOOL_INDEX_CACHE is not None and _TOOL_INDEX_KEY == key:
        return _TOOL_INDEX_CACHE
    by_name = {t["function"]["name"]: t for t in (tools or [])}
    lines = [
        "你拥有以下全部能力（工具），共 %d 项。需要某能力时，调用 activate_tools([\"工具名或组名\",...]) 激活，激活后立即可用；"
        "组名如「数据与文档」「媒体与图像」（见下方分类），传组名一次激活整组。"
        "未激活前也具备该能力，只是定义尚未加载。简单对话可以不激活任何工具。" % len(by_name)
    ]
    for cat, members in TOOL_GROUPS:
        rows = []
        for name in members:
            if name not in by_name:
                continue
            phrase = _TOOL_ACTION_PHRASES.get(name) or ""
            if not phrase:
                desc = by_name[name]["function"].get("description", "")
                phrase = re.sub(r"\s+", " ", desc).strip()[:60]
            rows.append(f"{name}({phrase})" if phrase else name)
        if rows:
            lines.append(f"{cat}: " + "、".join(rows))
    _TOOL_INDEX_CACHE = "\n".join(lines)
    _TOOL_INDEX_KEY = key
    return _TOOL_INDEX_CACHE


def compact_tool_schema(tool):
    """压缩工具 schema 描述（省 token）：去掉兜底废话、截断长描述。

    只压缩 description；name / type / required / enum / properties 结构一律保留，
    保证 strict 模式与工具解析不受影响。
    """
    t = json.loads(json.dumps(tool))
    fn = t["function"]
    desc = fn.get("description", "")
    # 兜底废话正则：更多冗余括号模式（重复短语/许可性/依赖提示等）
    for pat in (
        r"（[^）]*可能不严格[^）]*）", r"（[^）]*依赖[^）]*）",
        r"（[^）]*可选[^）]*）", r"（[^）]*保证生效[^）]*）",
        r"（[^）]*默认为[^）]*）", r"（[^）]*默认 [^）]*）",
        r"（[^）]*需审批[^）]*）", r"（[^）]*需用户[^）]*）",
        r"（[^）]*敏感[^）]*）", r"（[^）]*Beta[^）]*）",
        r"（[^）]*可选依赖[^）]*）", r"（[^）]*需安装[^）]*）",
    ):
        desc = re.sub(pat, "", desc)
    desc = re.sub(r"\s+", " ", desc).strip()
    if len(desc) > 130:
        desc = desc[:130].rstrip("，。；;:：, ") + "…"
    fn["description"] = desc
    for p in (fn.get("parameters", {}).get("properties") or {}).values():
        if isinstance(p, dict) and "description" in p:
            d = p["description"]
            d = re.sub(r"^可选[：:]\s*", "", d)
            d = re.sub(r"（[^）]*）", "", d)
            d = re.sub(r"\s+", " ", d).strip()
            if len(d) > 40:
                d = d[:40].rstrip("，。；;:：, ") + "…"
            p["description"] = d
    return t


def compact_tools_list(tools):
    """批量压缩工具 schema（保持顺序，安全返回原列表）。"""
    return [compact_tool_schema(t) for t in tools]


def check_balance(api_key, base_url=DEFAULT_BASE_URL, timeout=10.0):
    # balance 接口只在官方主端点，避免 base_url 带 /beta 等路径时拼接错误
    try:
        url = str(httpx.URL(str(base_url)).join("/user/balance"))
    except Exception:
        url = DEFAULT_BASE_URL + "/user/balance"
    response = _http_client().get(
        url,
        headers={"Authorization": f"Bearer {api_key}"},
        timeout=timeout,
    )
    response.raise_for_status()
    return response.json()


def _restore_text_content(m):
    """把发送时内联的图片内容块还原为纯文本 content（UI/存档只保留文本与图片路径）。"""
    if not (isinstance(m, dict) and isinstance(m.get("content"), list)):
        return m
    text_parts = []
    for b in m["content"]:
        if isinstance(b, dict) and b.get("type") == "text" and b.get("text"):
            text_parts.append(b["text"])
    m = dict(m)
    m["content"] = "\n".join(text_parts)
    return m


def _prune_reasoning_for_send(messages):
    """发送时剥离「无工具调用轮次」的思考内容（官方多轮拼接规则）。

    官方规则：两个 user 消息之间，若模型未进行工具调用，中间 assistant 的
    reasoning_content 在后续轮次传入 API 会被忽略——不传即可省下这部分
    输入 token（思考内容可占数千 token，是纯浪费）。
    带 tool_calls 的轮次必须完整回传 reasoning_content（缺失会 400），
    因此只剥离「无 tool_calls 的 assistant 消息」。
    仅对浅拷贝生效，不影响调用方内存中的会话历史（UI/存档仍需 reasoning）。
    """
    cleaned = []
    for m in messages:
        if (
            isinstance(m, dict)
            and m.get("role") == "assistant"
            and not m.get("tool_calls")
            and m.get("reasoning_content")
        ):
            m = dict(m)
            m.pop("reasoning_content", None)
        cleaned.append(m)
    return cleaned


def _strictify_schema(schema, is_root=False):
    """把 JSON Schema 转换为 strict 模式（Beta）兼容形式。

    服务端实测规则（2026-08 beta /chat/completions 校验）：
    - object 的 required 必须与 properties 键集完全一致（缺项/多出项均 400）
    - 每个 object 必须显式声明 properties，嵌套 object 的 properties 不能为空
      （根级可空：无参工具 200 OK，嵌套空 properties 400）
    - 每个 schema 节点必须携带 type / anyOf / $ref（宽松 `{}`、纯 enum、oneOf 均 400）
    - additionalProperties 仅允许 false 或缺省（true / map 形式 400）

    无法合规的结构（自由对象/宽松 items/纯 enum/oneOf 等）返回 None，
    调用方（_strictify_tools）必须放弃该工具的 strict 标记——
    服务端实测：strict 校验仅作用于标了 strict=true 的 function，
    混合列表正常通过，硬转则整轮请求直接 400。
    """
    if not isinstance(schema, dict):
        return None
    st = dict(schema)
    if "$ref" in st:
        return None  # 含 ref 的节点不再递归处理（保守放弃 strict）
    t = st.get("type")
    if t == "object":
        props = st.get("properties")
        if not isinstance(props, dict):
            return None  # 自由对象/任意键 map 无法在 strict 模式下表达
        if props:
            new_props = {}
            for k, v in props.items():
                nv = _strictify_schema(v)
                if nv is None:
                    return None
                new_props[k] = nv
            st["properties"] = new_props
            # 服务端要求 required 与 properties 一致：全量必填 + 不得多出
            st["required"] = list(new_props.keys())
        elif not is_root:
            return None  # 嵌套空 properties 被服务端拒绝
        else:
            st["required"] = []
        ap = st.get("additionalProperties")
        if ap is True or (ap is not None and not isinstance(ap, bool)):
            return None
        st["additionalProperties"] = False
        return st
    if t == "array":
        items = st.get("items")
        if not isinstance(items, dict):
            return None
        ni = _strictify_schema(items)
        if ni is None:
            return None
        st["items"] = ni
        return st
    if "anyOf" in st:
        new_branches = []
        for b in st["anyOf"]:
            nb = _strictify_schema(b)
            if nb is None:
                return None
            new_branches.append(nb)
        st["anyOf"] = new_branches
        return st
    if t in ("string", "integer", "number", "boolean", "null"):
        return st
    return None


_BASE_TOOLS_CACHE = None
_CUSTOM_TOOLS_CACHE = None
_CUSTOM_TOOLS_ID = None


def _cached_all_tools(custom_tools):
    """缓存内置工具/自定义工具的深拷贝，避免每轮 chat() 全量 deepcopy 100+ 工具。"""
    global _BASE_TOOLS_CACHE, _CUSTOM_TOOLS_CACHE, _CUSTOM_TOOLS_ID
    if _BASE_TOOLS_CACHE is None or _BASE_TOOLS_CACHE[0] != id(TOOLS):
        _BASE_TOOLS_CACHE = (id(TOOLS), copy.deepcopy(TOOLS))
    base = _BASE_TOOLS_CACHE[1]
    custom_tools = custom_tools or []
    if not custom_tools:
        return copy.deepcopy(base)
    cid = id(custom_tools)
    if _CUSTOM_TOOLS_CACHE is None or _CUSTOM_TOOLS_ID != cid:
        _CUSTOM_TOOLS_CACHE = copy.deepcopy(custom_tools)
        _CUSTOM_TOOLS_ID = cid
    return copy.deepcopy(base + _CUSTOM_TOOLS_CACHE)


def _strictify_tools(tools):
    """把 function 工具转换为 strict 模式（Beta）兼容形式。

    strict 校验仅作用于标了 strict=true 的 function；schema 无法合规化的工具
    （自由对象/宽松 items 等）不标 strict——服务端实测混合列表正常通过，
    而强行标 strict 会导致整轮请求 400。
    """
    out = []
    for t in tools or []:
        if not isinstance(t, dict):
            out.append(t)
            continue
        t = dict(t)
        fn = t.get("function")
        if isinstance(fn, dict):
            fn = dict(fn)
            params = fn.get("parameters")
            if isinstance(params, dict):
                clean = _strictify_schema(params, is_root=True)
                if clean is not None:
                    fn["parameters"] = clean
                    fn["strict"] = True
            t["function"] = fn
        out.append(t)
    return out


class DeepSeekClient:
    def __init__(self, api_key, base_url=DEFAULT_BASE_URL, model=DEFAULT_MODEL, timeout=120.0):
        self.api_key = api_key
        self.base_url = base_url
        self.model = model
        self.timeout = timeout
        self.client = OpenAI(
            api_key=api_key,
            base_url=base_url,
            timeout=(10.0, timeout),
        )

    @staticmethod
    def _sanitize_messages(messages):
        """过滤非法消息：assistant 消息必须包含 content 或 tool_calls。

        中断生成或模型仅返回思考时会产生 content 为空的 assistant 消息，
        DeepSeek API 会以 400 (invalid_request_error) 拒绝此类历史记录。
        同时清理「悬空 tool_calls」：assistant 声明了工具调用但历史中
        没有配对 tool 消息（stop/异常中断的残留），不清理会导致后续 400。
        """
        has_tool = set()
        for m in messages:
            if isinstance(m, dict) and m.get("role") == "tool" and m.get("tool_call_id"):
                has_tool.add(m["tool_call_id"])
        cleaned = []
        for m in messages:
            if not isinstance(m, dict) or m.get("role") != "assistant":
                cleaned.append(m)
                continue
            content = m.get("content")
            tcs = m.get("tool_calls")
            if tcs:
                valid = [
                    tc for tc in tcs
                    if isinstance(tc, dict) and tc.get("id") in has_tool
                ]
                if len(valid) != len(tcs):
                    m = dict(m)
                    if valid:
                        m["tool_calls"] = valid
                    else:
                        m.pop("tool_calls", None)
            if not content and not m.get("tool_calls"):
                continue
            m = dict(m)
            if content is None:
                m["content"] = ""
            cleaned.append(m)
        return cleaned

    def chat(
        self,
        messages,
        scenario="通用",
        thinking="high",
        max_tokens=16384,
        seed=None,
        tools_enabled=False,
        enabled_tools=None,
        on_reasoning=None,
        on_content=None,
        on_tool=None,
        on_tool_start=None,
        on_usage=None,
        stop_event=None,
        temperature=None,
        top_p=None,
        custom_tools=None,
        max_tool_rounds=None,
        on_loop_guard=None,
        on_tool_duration=None,
        json_output=False,
        continue_prefix=False,
        stop=None,
        logprobs=False,
        tool_choice=None,
        on_approval=None,
        on_plan=None,
        memory_text=None,
        pure_chat=False,
        on_ask=None,
        on_request_permission=None,
        on_truncated=None,
        trailing_text=None,
        strict_tools=False,
        smart_tools=False,
        preset_tools=None,
    ):
        cfg = SCENARIOS.get(scenario, SCENARIOS["通用"])
        thinking_key = thinking if thinking in THINKING_MODES else "high"

        messages[:] = self._sanitize_messages(messages)
        # 图片内联：仅替换受影响的 user 消息副本；原始消息对象（content 文本 +
        # images 路径）不受影响，chat() 结束时再同步回调用方（含新增 assistant/tool 消息）。
        work = embed_message_images(messages, self.model)
        json_hint = None
        memory_msg = None
        # 缓存友好消息布局（官方硬盘缓存按「前缀完整匹配」命中）：
        # - 恒定的 json_hint 保持在最前（前缀稳定 → 可命中）
        # - 系统提示词 messages[0] 紧随其后（稳定前缀主体）
        # - 可能变化的记忆注入追加在末尾（system 消息位置任意），
        #   记忆刷新只破坏尾部单元，稳定前缀继续命中
        if json_output:
            json_hint = {"role": "system", "content": JSON_HINT_MESSAGE}
            work = [json_hint] + work
        if memory_text:
            memory_msg = {"role": "system", "content": memory_text}
            work = work + [memory_msg]
        trailing = str(trailing_text or "")

        extra_body = {
            "thinking": {"type": "enabled" if thinking_key != "none" else "disabled"}
        }
        kwargs = {
            "model": self.model,
            "messages": work,
            "max_tokens": max_tokens,
            "stream": True,
            # 流式必须显式请求 usage（否则部分端点不返回末尾 usage chunk）
            "stream_options": {"include_usage": True},
            "extra_body": extra_body,
        }
        if json_output:
            kwargs["response_format"] = {"type": "json_object"}
        if stop:
            # 官方支持：string 或最多 16 个 string
            if isinstance(stop, str):
                kwargs["stop"] = [stop]
            elif isinstance(stop, (list, tuple)):
                kwargs["stop"] = [str(s) for s in stop][:16]
        if logprobs:
            kwargs["logprobs"] = True
        if tool_choice is not None:
            kwargs["tool_choice"] = tool_choice
        if continue_prefix and work and work[-1].get("role") == "assistant":
            last = dict(work[-1])
            last["prefix"] = True
            work[-1] = last
        if thinking_key == "none":
            kwargs["temperature"] = cfg["temperature"] if temperature is None else temperature
            kwargs["top_p"] = cfg["top_p"] if top_p is None else top_p
        else:
            effort = EFFORT_BY_THINKING.get(thinking_key)
            if thinking_key == "auto":
                effort = _auto_effort(work)
            if effort == "none":
                # auto 判定为简单任务：关闭思考（reasoning_effort 不支持 none，
                # 此前直接传 "none" 会被 API 拒绝或静默忽略）
                kwargs["temperature"] = cfg["temperature"] if temperature is None else temperature
                kwargs["top_p"] = cfg["top_p"] if top_p is None else top_p
            else:
                if effort is None:
                    effort = cfg["reasoning_effort"]
                kwargs["reasoning_effort"] = effort
        if seed is not None:
            kwargs["seed"] = seed
        all_tools = _cached_all_tools(custom_tools or [])
        # 兜底：array 参数缺 items 会导致 API 400（missing field 'items'），
        # 内置/自定义/插件工具一律补齐，防御用户自定义 schema 遗漏
        _patch_array_items(all_tools)
        # smart_tools（完全智能模式）：索引 + 按需激活，避免全量工具定义挤占上下文。
        # enabled_tools 为 None 或「覆盖全部工具」时视为全量 → 可启用点菜
        builtin_names = {t["function"]["name"] for t in TOOLS}
        if enabled_tools is not None:
            sel_names = set(enabled_tools)
            if not (sel_names >= builtin_names):
                enabled_tools = list(sel_names)  # 显式子集 → 非全量
                smart_avail = False
            else:
                enabled_tools = None  # 全量 → smart 模式
                smart_avail = True
        else:
            smart_avail = True
        smart_avail = bool(smart_tools and tools_enabled and smart_avail)
        activated = set(preset_tools or ())  # 预激活（关键词预筛）+ AI 点菜
        if smart_avail:
            # chat 层关键词预激活兜底：所有调用方（main/api_server/子代理等）受益，
            # 命中常见意图直接免点菜可用（仅提前加载定义，不改变权限）
            _preactivate_from_messages(work, activated)
        smart_round = smart_avail  # 索引阶段：点菜工具 + 预激活工具并注入
        index_msg = None
        _index_shown = False
        if smart_avail:
            index_msg = {
                "role": "system",
                "content": (
                    "你是一个拥有 100+ 项专业能力的桌面 AI 智能体，能力地图如下（你确实拥有这些能力，"
                    "不要拒绝用户请求）。能力定义未加载时，先调用 activate_tools 激活再使用。\n\n"
                    "视觉自检准则：执行截图/浏览器/RPA 操作、或生成图片/图表后，"
                    "用 screen_see / image_understand 查看结果并自查是否达到目标；未达标则继续修正"
                    "（点击/输入/重新生成），完成后才向用户汇报。\n\n"
                    + build_tool_index()
                ),
            }
        if pure_chat:
            # 纯对话模式：完全不传 tools schema（避免工具提示词污染对话能力）
            pass
        elif tools_enabled:
            tools = all_tools
            if enabled_tools is not None:
                tools = [
                    t for t in all_tools
                    if t["function"]["name"] in enabled_tools
                    or t["function"]["name"] in SELF_EVOLUTION_TOOLS
                ]
                smart_avail = False  # 显式子集时不启用点菜
            if tools:
                kwargs["tools"] = _strictify_tools(tools) if strict_tools else tools
        elif any(t["function"]["name"] in SELF_EVOLUTION_TOOLS for t in all_tools):
            # 工具总开关关闭时，自我进化工具仍可用（自我审查是安全只读+分支提案）
            kwargs["tools"] = [
                t for t in all_tools if t["function"]["name"] in SELF_EVOLUTION_TOOLS
            ]

        empty_retries = 0
        plan_rejections = 0
        json_retried = False  # JSON 输出自校验重试只允许一次
        rounds = max_tool_rounds if max_tool_rounds and max_tool_rounds > 0 else MAX_TOOL_ROUNDS
        last_tool_key = None
        same_repeats = 0
        try:
            for _ in range(rounds):
                if stop_event and stop_event.is_set():
                    return False
                # smart_tools 阶段切换：索引阶段注入 activate_tools 点菜；
                # 激活后注入激活工具的完整 schema（压缩版），并移除点菜工具
                if smart_avail:
                    if smart_round:
                        # 点菜阶段：点菜工具 + 已预激活工具（AI 可直接用，也可补充点菜）
                        preset_specs = [t for t in all_tools if t["function"]["name"] in activated]
                        kw_tools = [ACTIVATE_TOOL] + compact_tools_list(preset_specs)
                    else:
                        sel = [t for t in all_tools if t["function"]["name"] in activated]
                        kw_tools = compact_tools_list(sel) if sel else None
                    if kw_tools:
                        kwargs["tools"] = _strictify_tools(kw_tools) if strict_tools else kw_tools
                    else:
                        kwargs.pop("tools", None)
                    # 索引消息在「点菜/工具调用发生前」每轮注入（AI 可随时参考点菜），
                    # 首次点菜或工具调用后移除（AI 已上手，避免浪费）
                    if not _index_shown and index_msg is not None:
                        req_work = [index_msg] + list(work)
                    else:
                        req_work = work
                else:
                    req_work = work
                # 发送前构造请求消息：剥离无工具轮次的思考内容（省输入 token），
                # 动态上下文（trailing_text）追加到最近一条 user 消息尾部。
                # 仅作用于浅拷贝：不修改 work / 调用方内存历史（UI 与存档保留 reasoning）
                req_msgs = _prune_reasoning_for_send(req_work)
                if trailing:
                    req_msgs = list(req_msgs)
                    for i in range(len(req_msgs) - 1, -1, -1):
                        if req_msgs[i].get("role") == "user":
                            last = dict(req_msgs[i])
                            if isinstance(last.get("content"), list):
                                # 图片内联消息：动态上下文追加为 text 块
                                last["content"] = list(last["content"]) + [
                                    {"type": "text", "text": trailing}
                                ]
                            else:
                                last["content"] = str(last.get("content") or "") + "\n\n" + trailing
                            req_msgs[i] = last
                            break
                    else:
                        req_msgs.append({"role": "user", "content": trailing})
                kw = dict(kwargs) if req_msgs is not work else kwargs
                if req_msgs is not work:
                    kw["messages"] = req_msgs
                # 流中途断线（APIConnectionError）会抛在 _create_with_retry 之外：
                # 仅在「尚无任何增量送达 UI」时整体重试，避免已显示内容重复
                reasoning, content, tool_calls, finish_reason = "", "", {}, None
                stream_usage = None
                try:
                    for stream_attempt in range(2):
                        response = self._create_with_retry(kw, attempts=2, stop_event=stop_event)
                        try:
                            reasoning, content, tool_calls, finish_reason, stream_usage = self._consume_stream(
                                response, on_reasoning, on_content, stop_event
                            )
                            break
                        except (APIConnectionError, APITimeoutError) as e:
                            if reasoning or content or tool_calls or stream_attempt == 1:
                                raise
                            logger.warning("流式连接中途断开（尚未收到内容），重试: %s", e)
                except _StopRequested:
                    return False  # 停止请求：干净返回，不把半截内容当异常抛给 UI
                except (APIConnectionError, APITimeoutError) as e:
                    # 流中途断线且已有部分增量送达 UI：不再重试（避免已显示内容重复），
                    # 明确告知本轮未正常完成（finish_reason=aborted 语义）
                    logger.warning("流式连接中途断开，本轮生成未完成: %s", e)
                    if on_truncated:
                        on_truncated("网络中断：本轮回复不完整")
                    return False
                if stream_usage is not None:
                    if on_usage:
                        on_usage(self._usage_dict(stream_usage))
                elif getattr(response, "usage", None):
                    # 兼容旧版 openai SDK（Stream 对象自带聚合 usage）
                    if on_usage:
                        on_usage(self._usage_dict(response.usage))

                if not tool_calls and not content.strip() and not reasoning.strip():
                    if empty_retries < MAX_EMPTY_RETRIES:
                        empty_retries += 1
                        logger.warning("收到空响应（无内容/思考/工具调用），自动重试第 %s 次", empty_retries)
                        continue
                    logger.warning("空响应重试已达上限，按失败返回（不写入空历史）")
                    if on_truncated:
                        on_truncated("模型连续返回空响应，本轮生成失败")
                    return False

                if continue_prefix and work and work[-1].get("role") == "assistant":
                    prev = dict(work[-1])
                    prev["content"] = (
                        ((prev.get("content") or "") + content) if content else prev.get("content")
                    )
                    if reasoning:
                        prev["reasoning_content"] = (
                            (prev.get("reasoning_content") or "") + reasoning
                        )
                    if tool_calls:
                        prev["tool_calls"] = [
                            {
                                "id": tc["id"],
                                "type": "function",
                                "function": {"name": tc["name"], "arguments": tc["args"]},
                            }
                            for tc in tool_calls
                        ]
                    work[-1] = prev
                else:
                    assistant_msg = {
                        "role": "assistant",
                        "content": content if content else None,
                        "time": datetime.now().strftime("%H:%M:%S"),
                    }
                    if reasoning:
                        assistant_msg["reasoning_content"] = reasoning
                    if tool_calls:
                        assistant_msg["tool_calls"] = [
                            {
                                "id": tc["id"],
                                "type": "function",
                                "function": {"name": tc["name"], "arguments": tc["args"]},
                            }
                            for tc in tool_calls
                        ]
                    work.append(assistant_msg)

                if not tool_calls:
                    # finish_reason=length：输出达 max_tokens 上限被截断，
                    # 任务/回复未完成——此前静默返回 True 导致 UI 误报「任务完成」
                    if finish_reason == "length" and on_truncated:
                        on_truncated("输出已达上限（max_tokens）被截断，回复不完整")
                    if json_output and content.strip():
                        # JSON 输出自校验：解析失败自动修正重试一次
                        # （官方说明 JSON 输出有概率返回非法内容，这是应用层可救回的部分）
                        try:
                            json.loads(content)
                        except (TypeError, ValueError):
                            if not json_retried:
                                json_retried = True
                                if on_truncated:
                                    on_truncated("JSON 输出解析失败，正在自动修正重试")
                                if work and work[-1].get("role") == "assistant":
                                    work.pop()  # 移除刚追加的半截 assistant 消息
                                work.append(
                                    {
                                        "role": "user",
                                        "content": (
                                            "[系统] 你上次输出的内容无法解析为合法 JSON。"
                                            "请重新输出：仅返回一个合法的 JSON 对象，"
                                            "不要输出任何其他文字。"
                                        ),
                                    }
                                )
                                continue
                    return True

                # 停止/中断防护：用户已停止或 tool_calls 不完整（缺 id/name，流被中断）
                # 时不执行任何工具；移除半截 tool_calls，避免历史残留「悬空 tool_call」导致下次 400
                if (stop_event and stop_event.is_set()) or any(
                    not (tc.get("id") and tc.get("name")) for tc in tool_calls
                ):
                    work[-1].pop("tool_calls", None)
                    if on_truncated and not (stop_event and stop_event.is_set()):
                        on_truncated("工具调用流被截断，本轮工具未执行")
                    return False

                if smart_avail:
                    # 拦截 activate_tools 点菜调用：更新激活集，切换为完整工具注入
                    act_calls = [tc for tc in tool_calls if tc.get("name") == "activate_tools"]
                    if act_calls:
                        for tc in act_calls:
                            try:
                                args = json.loads(tc.get("args") or "{}")
                                wanted = args.get("tools") or []
                                # 支持工具名与组名（组名一次激活整组）
                                _expand_activation(
                                    wanted,
                                    {t["function"]["name"] for t in all_tools},
                                    activated,
                                )
                            except (TypeError, ValueError):
                                pass
                            work.append(
                                {
                                    "role": "tool",
                                    "tool_call_id": tc["id"],
                                    "content": f"已激活工具: {sorted(activated) or '（无）'}",
                                }
                            )
                        # 非点菜工具也一并执行；点菜轮若有内容则保留
                        rest = [tc for tc in tool_calls if tc.get("name") != "activate_tools"]
                        if not rest:
                            _index_shown = True
                            if smart_round:
                                smart_round = False  # 下一轮注入激活工具完整 schema
                            continue
                        tool_calls = rest
                    else:
                        # 直接调用工具（未点菜，如预激活场景）：索引已完成使命
                        _index_shown = True

                if on_plan is not None:
                    ok_plan, reason_plan = on_plan(
                        [(tc["name"], (tc["args"] or "")[:300]) for tc in tool_calls]
                    )
                    if not ok_plan:
                        plan_rejections += 1
                        if plan_rejections >= MAX_PLAN_REJECTIONS:
                            logger.warning("计划连续被拒绝 %s 次，终止工具循环", plan_rejections)
                            for tc in tool_calls:
                                work.append(
                                    {
                                        "role": "tool",
                                        "tool_call_id": tc["id"],
                                        "content": f"计划连续被拒绝 {plan_rejections} 次，已终止执行",
                                    }
                                )
                            return False
                        for tc in tool_calls:
                            work.append(
                                {
                                    "role": "tool",
                                    "tool_call_id": tc["id"],
                                    "content": reason_plan or "用户取消了本轮工具计划",
                                }
                            )
                        continue

                # 计划确认期间用户可能已点停止：执行前再查一次
                if stop_event and stop_event.is_set():
                    work[-1].pop("tool_calls", None)
                    return False

                # 循环防护：按 (name, args) 顺序预判连续重复，跨轮累计（每轮不重置，
                # 模型换参数时自然清零），命中则整轮终止
                guarded = False
                for tc in tool_calls:
                    key = (tc["name"], (tc["args"] or "").strip())
                    if key == last_tool_key:
                        same_repeats += 1
                    else:
                        same_repeats = 1
                        last_tool_key = key
                    if same_repeats >= MAX_SAME_TOOL_REPEATS:
                        guarded = True
                        break
                if guarded:
                    logger.warning(
                        "工具循环防护：%s 连续调用 %s 次相同参数，终止工具循环",
                        tc["name"], same_repeats,
                    )
                    if on_loop_guard:
                        on_loop_guard(tc["name"], same_repeats)
                    # 补齐本轮全部 tool 结果，避免历史残留「悬空 tool_call」
                    # （assistant 含 tool_calls 但无对应 tool 消息，下一次请求会被 API 以 400 拒绝）
                    for tc in tool_calls:
                        work.append(
                            {
                                "role": "tool",
                                "tool_call_id": tc["id"],
                                "content": (
                                    f"已终止：{tc['name']} 连续调用相同参数达到上限"
                                    f"（{MAX_SAME_TOOL_REPEATS} 次）"
                                ),
                            }
                        )
                    return True

                custom_map = {
                    t["function"]["name"]: t for t in (custom_tools or [])
                }

                def execute_tool(tc):
                    """执行单个工具调用，返回 (name, args, result, duration)。

                    交互工具（ask_user/request_permission）与普通工具分开执行，
                    普通工具可并行（同轮多工具场景提速，多轮串行仍是主链路）。
                    """
                    name = tc["name"]
                    raw_args = tc["args"] or ""
                    if on_tool_start is not None:
                        try:
                            on_tool_start(name, raw_args)
                        except Exception:
                            pass
                    fn = TOOL_CALL_MAP.get(name)
                    args = {}
                    t0 = time.monotonic()
                    if name == "ask_user":
                        # 询问用户：阻塞等待 UI 回答（on_ask 由 main 提供）
                        try:
                            qargs = json.loads(raw_args) if raw_args else {}
                            prompt = str(qargs.get("prompt") or "") if isinstance(qargs, dict) else ""
                        except json.JSONDecodeError:
                            prompt = raw_args[:200]
                        if not prompt:
                            prompt = "请提供需要用户回答的问题"
                        if on_ask is not None:
                            result = on_ask(prompt)
                            args = {"prompt": prompt}
                        else:
                            result = "错误：无法询问用户（当前环境不支持交互式询问）"
                    elif name == "request_permission":
                        # 请求用户将操作加入白名单（弹窗同意/拒绝）
                        try:
                            pargs = json.loads(raw_args) if raw_args else {}
                        except json.JSONDecodeError:
                            pargs = {}
                        atype = str(pargs.get("action_type") or "") if isinstance(pargs, dict) else ""
                        pvalue = str(pargs.get("value") or "") if isinstance(pargs, dict) else ""
                        args = {"action_type": atype, "value": pvalue}
                        if not atype:
                            result = "错误：action_type 必填（dir / command / write）"
                        elif on_request_permission is not None:
                            ok_rp, msg_rp = on_request_permission(atype, pvalue)
                            if ok_rp:
                                result = f"已加入白名单：{msg_rp}。可以重试刚才被拒绝的操作。"
                            else:
                                result = f"白名单请求被拒绝：{msg_rp}"
                        else:
                            result = "错误：当前环境不支持白名单请求"
                    elif not fn:
                        handler = custom_map.get(name)
                        if handler:
                            result = self._run_custom_tool(handler, raw_args)
                        else:
                            result = f"未知工具: {name}，可用工具: {list(TOOL_CALL_MAP)} + 自定义工具 {list(custom_map)}"
                    else:
                        approved = True
                        if on_approval is not None:
                            approved, reason_a = on_approval(name, raw_args)
                        if not approved:
                            result = reason_a or "权限拒绝：未获批准"
                        else:
                            try:
                                args = json.loads(raw_args)
                                if not isinstance(args, dict):
                                    raise ValueError("工具参数必须是 JSON 对象")
                                result = fn(**args)
                            except (json.JSONDecodeError, ValueError) as e:
                                result = (
                                    f"工具参数解析失败: {e}，原始参数: {raw_args!r}，请修正参数格式后重试"
                                )
                            except TypeError as e:
                                result = f"工具参数错误: {e}，收到的参数: {args}，请修正后重试"
                            except Exception as e:
                                result = f"工具执行失败: {e}"
                    duration = time.monotonic() - t0
                    return name, args, result, duration

                # 交互工具串行（弹窗不能并发），其余工具并行执行（同轮多工具提速）
                serial_tools = [
                    tc for tc in tool_calls if tc["name"] in ("ask_user", "request_permission")
                ]
                parallel_tools = [
                    tc for tc in tool_calls if tc["name"] not in ("ask_user", "request_permission")
                ]
                exec_results = {}
                if parallel_tools:
                    futs = {
                        tc["id"]: _tool_executor_for(tc["name"]).submit(execute_tool, tc)
                        for tc in parallel_tools
                    }
                    pending = set(futs.values())
                    while pending:
                        if stop_event and stop_event.is_set():
                            # 停止感知：给已提交工具短宽限期，副作用已发生的工具
                            # （发信/写文件/启进程）如实记录结果——历史写"已中断"会让
                            # 模型下轮重试同参数，造成重复执行
                            done2, pending = wait(pending, timeout=_STOP_TOOL_GRACE_S)
                            for f in done2:
                                tcid = next(k for k, v in futs.items() if v is f)
                                try:
                                    name, args, result, duration = f.result()
                                except Exception as e:
                                    name, args, result, duration = (
                                        next(
                                            (
                                                t["name"]
                                                for k2, t in futs.items()
                                                if futs[k2] is f
                                            ),
                                            "?",
                                        ),
                                        {},
                                        f"工具执行异常: {e}",
                                        None,
                                    )
                                exec_results[tcid] = (name, args, result, duration)
                                if on_tool:
                                    on_tool(name, args, result)
                                if on_tool_duration:
                                    on_tool_duration(name, duration)
                            break  # 仍未完成的工具后台继续跑（自带超时兜底）
                        done, pending = wait(
                            pending, timeout=0.25, return_when=FIRST_COMPLETED
                        )
                        for f in done:
                            tcid = next(k for k, v in futs.items() if v is f)
                            name, args, result, duration = f.result()
                            exec_results[tcid] = (name, args, result, duration)
                            # 完成即回调 UI：快的工具不再被慢的拖到最后一齐出现
                            if on_tool:
                                on_tool(name, args, result)
                            if on_tool_duration:
                                on_tool_duration(name, duration)
                for tc in serial_tools:
                    name, args, result, duration = execute_tool(tc)
                    exec_results[tc["id"]] = (name, args, result, duration)
                    if on_tool:
                        on_tool(name, args, result)
                    if on_tool_duration:
                        on_tool_duration(name, duration)
                # 按原始 tool_calls 顺序追加历史（保证历史消息顺序稳定）；
                # 已完成的上面已回调 UI，这里只为停止时未执行的补回调
                for tc in tool_calls:
                    entry = exec_results.get(tc["id"])
                    if entry is None:
                        entry = (tc["name"], {}, "工具执行已被中断（停止生成），未完成", None)
                        exec_results[tc["id"]] = entry
                        if on_tool:
                            on_tool(entry[0], entry[1], entry[2])
                    name, args, result, duration = entry
                    # 进上下文的工具结果截断：fetch_url 500KB/read_file 100KB 原样
                    # 重传给模型会白白消耗数万 token（费用 + 延迟 + 逼近 1M 上限）。
                    # 超长结果自动落盘到工作区，上下文只留路径 + 首尾摘要。
                    text = str(result)
                    if len(text) > _RESULT_INTO_CONTEXT_MAX:
                        text = _persist_long_result(name, text)
                    # 视觉自审（config vision_self_review 开启且为视觉模型）：工具产出图片时，
                    # 自动调用视觉模型审图，把审阅意见附在结果里，模型据此迭代（B 自我审图闭环）。
                    if (
                        VISION_SELF_REVIEW
                        and name in _IMAGE_PRODUCING_TOOLS
                        and is_vision_model(self.model)
                    ):
                        img_path = _extract_image_path(text)
                        if img_path:
                            try:
                                review = image_understand(
                                    img_path,
                                    question=(
                                        "（自动审图）请审阅这张图片：评估清晰度、构图、文字是否完整准确、"
                                        "是否完全满足创作/数据要求；若发现问题，给出具体、可执行的修改建议。"
                                    ),
                                )
                                if review and not review.startswith("错误"):
                                    text = text + "\n\n【AI 自审】\n" + review
                            except Exception:
                                logger.exception("视觉自审失败: %s", img_path)
                    work.append(
                        {"role": "tool", "tool_call_id": tc["id"], "content": text}
                    )
            # for 循环自然结束 = 工具轮数耗尽（任务可能未完成）：告知 UI
            if on_truncated:
                on_truncated(
                    f"工具调用轮数已达上限（{rounds} 轮），若任务未完成请追加指令继续"
                )
            return True
        finally:
            # work 恒为独立列表：把新增的 assistant/tool 消息同步回调用方，
            # 并将图片内联的消息还原为纯文本 content（UI/存档只保留文本 + images 路径）
            messages[:] = [
                _restore_text_content(m)
                for m in work
                if m is not json_hint and m is not memory_msg
            ]
            for m in messages:
                if isinstance(m, dict):
                    m.pop("prefix", None)

    def _create_with_retry(self, kwargs, attempts=3, stop_event=None):
        last_error = None
        for i in range(1, attempts + 1):
            try:
                return self.client.chat.completions.create(**kwargs)
            except RateLimitError as e:
                last_error = e
                logger.warning("限流(429)，第 %s 次重试", i)
            except (APIConnectionError, APITimeoutError) as e:
                last_error = e
                logger.warning("网络错误，第 %s 次重试", i)
            except APIError as e:
                if getattr(e, "status_code", None) in (500, 503):
                    last_error = e
                    logger.warning("服务端错误 %s，第 %s 次重试", getattr(e, "status_code", None), i)
                else:
                    raise
            if i < attempts:
                # 分段 sleep：用户点停止后立即感知，不等完整个退避窗口
                end = time.monotonic() + 2 ** (i - 1)
                while time.monotonic() < end:
                    if stop_event and stop_event.is_set():
                        raise _StopRequested()  # 干净信号：chat() 内部转 return False
                    time.sleep(min(0.1, end - time.monotonic()))
        raise RuntimeError(f"请求失败，已重试 {attempts} 次: {last_error}")

    def _consume_stream(self, response, on_reasoning, on_content, stop_event):
        reasoning = ""
        content = ""
        tool_calls = {}
        finish_reason = None
        usage = None
        try:
            for chunk in response:
                if stop_event and stop_event.is_set():
                    break
                # 流式 usage：服务端在末尾 chunk 返回（choices 为空数组，仅带 usage 字段）。
                # openai SDK 2.x 的 Stream 不再聚合 usage，必须在这里自行捕获。
                chunk_usage = getattr(chunk, "usage", None)
                if chunk_usage is not None:
                    usage = chunk_usage
                if not chunk.choices:
                    continue
                delta = chunk.choices[0].delta
                try:
                    fr = chunk.choices[0].finish_reason
                    if fr:
                        finish_reason = fr
                except Exception:
                    pass
                reasoning_part = getattr(delta, "reasoning_content", None)
                if reasoning_part:
                    reasoning += reasoning_part
                    if on_reasoning:
                        on_reasoning(reasoning_part)
                if delta.content:
                    content += delta.content
                    if on_content:
                        on_content(delta.content)
                if delta.tool_calls:
                    for tc in delta.tool_calls:
                        idx = tc.index if tc.index is not None else 0
                        entry = tool_calls.setdefault(idx, {"id": "", "name": "", "args": ""})
                        if tc.id:
                            entry["id"] = tc.id
                        if tc.function and tc.function.name:
                            entry["name"] = tc.function.name
                        if tc.function and tc.function.arguments:
                            entry["args"] += tc.function.arguments
        finally:
            # 显式关闭 SSE 流，连接回收不依赖 GC
            try:
                if hasattr(response, "close"):
                    response.close()
            except Exception:
                pass
        return reasoning, content, list(tool_calls.values()), finish_reason, usage

    @staticmethod
    def _usage_dict(usage):
        return {
            "prompt": getattr(usage, "prompt_tokens", 0) or 0,
            "completion": getattr(usage, "completion_tokens", 0) or 0,
            "cache_hit": getattr(usage, "prompt_cache_hit_tokens", 0) or 0,
            "cache_miss": getattr(usage, "prompt_cache_miss_tokens", 0) or 0,
        }

    @staticmethod
    def _run_custom_tool(handler, raw_args, timeout=15.0, max_chars=20000):
        """执行用户自定义工具：向 endpoint 发送 HTTP 请求。"""
        spec = handler.get("function", {})
        endpoint = spec.get("endpoint") or handler.get("endpoint")
        method = (spec.get("method") or handler.get("method") or "POST").upper()
        if not endpoint:
            return "错误：自定义工具未配置 endpoint"
        err = _safe_url(endpoint)
        if err:
            return f"错误：自定义工具 endpoint 不合法（{err}）"
        try:
            args = json.loads(raw_args) if raw_args else {}
            if not isinstance(args, dict):
                raise ValueError("工具参数必须是 JSON 对象")
        except (json.JSONDecodeError, ValueError) as e:
            return f"工具参数解析失败: {e}，原始参数: {raw_args!r}"
        try:
            if method in ("GET",):
                resp = _safe_request(
                    "GET", endpoint, params=args, timeout=timeout
                )
            else:
                resp = _safe_request(
                    "POST", endpoint, json=args, timeout=timeout
                )
            resp.raise_for_status()
            out = resp.text
            if len(out) > max_chars:
                out = out[:max_chars] + "\n[输出已截断]"
            return out
        except Exception as e:
            return f"错误：调用自定义工具失败: {e}"

    def fim_complete(self, prompt, suffix="", max_tokens=2048):
        """FIM 补全（Beta）：提供前缀与可选后缀，模型补全中间内容。

        需要 Beta 端点（https://api.deepseek.com/beta），最大补全长度 4K。
        """
        base = self.base_url.rstrip("/")
        if not base.endswith("/beta"):
            base += "/beta"
        # 按 base_url 缓存 SDK client：FIM 是菜单手动触发，避免每次新建连接池
        if (
            getattr(self, "_beta_client", None) is None
            or getattr(self, "_beta_base", None) != base
        ):
            old_beta = getattr(self, "_beta_client", None)
            if old_beta is not None and getattr(self, "_beta_base", None) != base:
                try:
                    old_beta.close()
                except Exception:
                    pass
            self._beta_client = OpenAI(
                api_key=self.api_key,
                base_url=base,
                timeout=(10.0, self.timeout),
            )
            self._beta_base = base
        resp = self._beta_client.completions.create(
            model=self.model,
            prompt=prompt,
            suffix=(suffix or None),
            max_tokens=max_tokens,
        )
        return (resp.choices[0].text or "").strip()
